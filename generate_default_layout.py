#!/usr/bin/env python3
"""
custom-template.pptx から default_layout.json を生成するスクリプト。

各スライドレイアウトのコンテンツ用プレースホルダー座標（EMU単位）を
JSON形式で抽出・出力する。フッター・日付・スライド番号プレースホルダー
（idx 10, 11, 12）は出力対象外とする。

使用方法:
    python generate_default_layout.py [テンプレートパス] [-o 出力パス]

出力形式は設計書4.9節のスキーマ定義に準拠する。
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

# フッター系プレースホルダーのidx（日付・フッター・スライド番号）—抽出対象外
_FOOTER_IDX: frozenset[int] = frozenset({10, 11, 12})

# 対象7レイアウトの idx → role マッピング（設計書4.9節 role定義準拠）
# role名はスライドレンダラーがJSONフォールバック時に参照するキーとなる
#
# 注意: Comparison レイアウトは設計書の簡略定義（left_content/right_content のみ）
# よりも実テンプレートのプレースホルダー構成が多い（左右各ヘッダー＋コンテンツの計5個）
# ため、実テンプレートの全構成をrolマッピングとして保持する。
# 注意: Content with Caption レイアウトは実テンプレートの配置に従い
# idx=1 を body（右側大エリア）、idx=2 を caption（左側テキストエリア）とする。
_ROLE_MAP: dict[str, dict[int, str]] = {
    "Title Slide": {
        0: "title",     # タイトルテキスト（CENTER_TITLE型）
        1: "subtitle",  # サブタイトル・author/date 表示エリア（SUBTITLE型）
    },
    "Title and Content": {
        0: "title",  # スライドタイトル（TITLE型）
        1: "body",   # メインコンテンツエリア（OBJECT型）
    },
    "Section Header": {
        0: "title",  # セクションタイトル（TITLE型）
        1: "body",   # セクション説明テキスト（BODY型）
    },
    "Two Content": {
        0: "title",          # スライドタイトル（TITLE型）
        1: "left_content",   # 左カラムコンテンツ（OBJECT型）
        2: "right_content",  # 右カラムコンテンツ（OBJECT型）
    },
    "Comparison": {
        0: "title",           # スライドタイトル（TITLE型）
        1: "left_header",     # 左カラム上部ヘッダーテキスト（BODY型）
        2: "left_content",    # 左カラムメインコンテンツ（OBJECT型）
        3: "right_header",    # 右カラム上部ヘッダーテキスト（BODY型）
        4: "right_content",   # 右カラムメインコンテンツ（OBJECT型）
    },
    "Content with Caption": {
        0: "title",    # スライドタイトル（TITLE型・左上）
        1: "body",     # メインコンテンツ（OBJECT型・右エリア）
        2: "caption",  # キャプションテキスト（BODY型・左下エリア）
    },
    "Blank": {},       # コンテンツなし（スピーカーノートのみ格納）
}


def _extract_placeholders(
    layout: Any,
    role_map: dict[int, str],
    layout_name: str,
) -> list[dict[str, Any]]:
    """
    1つのスライドレイアウトからコンテンツ用プレースホルダーを抽出する。

    フッター系（idx 10, 11, 12）はスキップする。
    role_map に未定義の idx が存在する場合は警告を出してスキップする。

    Parameters
    ----------
    layout : SlideLayout
        python-pptx のスライドレイアウトオブジェクト。
    role_map : dict[int, str]
        idx → role名 のマッピング辞書。
    layout_name : str
        警告メッセージ表示用のレイアウト名。

    Returns
    -------
    list[dict[str, Any]]
        プレースホルダー情報のリスト。idx 昇順でソート済み。
    """
    placeholders: list[dict[str, Any]] = []

    for ph in layout.placeholders:
        idx: int = ph.placeholder_format.idx

        # フッター・日付・スライド番号は除外
        if idx in _FOOTER_IDX:
            continue

        role = role_map.get(idx)
        if role is None:
            # ロールマッピングに未定義のidxは警告を出してスキップ
            print(
                f"  警告: '{layout_name}' の idx={idx} ('{ph.name}') は"
                f" ロールマッピング未定義のためスキップします",
                file=sys.stderr,
            )
            continue

        placeholders.append(
            {
                "idx": idx,
                "role": role,
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height,
            }
        )

    # idx 昇順でソートして返す
    placeholders.sort(key=lambda p: p["idx"])
    return placeholders


def generate_default_layout_json(
    template_path: Path,
    output_path: Path,
) -> None:
    """
    PPTXテンプレートを読み込み、default_layout.json を生成する。

    対象レイアウトが存在しない場合は警告を出して次のレイアウトへ進む。
    全対象レイアウトを処理後、output_path にJSON形式で書き出す。

    Parameters
    ----------
    template_path : Path
        入力PPTXテンプレートのファイルパス。
    output_path : Path
        出力JSONファイルのパス。
    """
    prs = Presentation(str(template_path))

    # スライド幅・高さをEMU単位で取得
    slide_width_emu: int = prs.slide_width
    slide_height_emu: int = prs.slide_height
    print(f"スライドサイズ: {slide_width_emu} x {slide_height_emu} EMU")

    # レイアウト名 → SlideLayout オブジェクト のマッピングを構築
    layout_map: dict[str, Any] = {
        layout.name: layout for layout in prs.slide_layouts
    }

    result: dict[str, Any] = {
        "slide_width_emu": slide_width_emu,
        "slide_height_emu": slide_height_emu,
        "layouts": {},
    }

    print("\nレイアウト抽出開始:")
    for layout_name, role_map in _ROLE_MAP.items():
        if layout_name not in layout_map:
            print(
                f"  警告: レイアウト '{layout_name}' がテンプレートに"
                f" 存在しないためスキップします",
                file=sys.stderr,
            )
            continue

        placeholders = _extract_placeholders(
            layout_map[layout_name], role_map, layout_name
        )
        result["layouts"][layout_name] = {"placeholders": placeholders}
        print(f"  '{layout_name}': {len(placeholders)} 個のプレースホルダーを抽出")

    # JSON書き出し（日本語文字はそのまま出力、インデント2スペース）
    with output_path.open("w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"\n出力完了: {output_path}")


def _build_arg_parser() -> argparse.ArgumentParser:
    """コマンドライン引数パーサーを構築して返す。"""
    parser = argparse.ArgumentParser(
        description=(
            "PPTXテンプレートからスライドレイアウトの座標情報を抽出し、"
            " default_layout.json を生成する。"
        )
    )
    parser.add_argument(
        "template",
        nargs="?",
        default="custom-template.pptx",
        help="入力PPTXテンプレートのパス（デフォルト: custom-template.pptx）",
    )
    parser.add_argument(
        "-o",
        "--output",
        default="default_layout.json",
        help="出力JSONファイルのパス（デフォルト: default_layout.json）",
    )
    return parser


def main() -> None:
    """エントリーポイント。引数を解析してJSON生成処理を実行する。"""
    parser = _build_arg_parser()
    args = parser.parse_args()

    template_path = Path(args.template)
    output_path = Path(args.output)

    if not template_path.exists():
        print(
            f"エラー: テンプレートファイル '{template_path}' が見つかりません",
            file=sys.stderr,
        )
        sys.exit(1)

    print(f"テンプレート読み込み: {template_path}")
    generate_default_layout_json(template_path, output_path)


if __name__ == "__main__":
    main()
