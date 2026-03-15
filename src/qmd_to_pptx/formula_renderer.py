"""
数式レンダラーモジュール。

arithmatexによりマーキングされたLaTeX数式テキストをOMML（Office Math Markup Language）
形式に変換し、python-pptxの数式オブジェクトとして現在のスライドに追加する。
"""

from __future__ import annotations

import xml.etree.ElementTree as ET
from lxml import etree

import latex2mathml.converter
import mathml2omml
from pptx.slide import Slide
from pptx.util import Emu, Pt


class FormulaRenderer:
    """
    数式レンダラークラス。

    arithmatexによりマーキングされたLaTeX数式テキストを
    OMML形式に変換してスライドに配置する。
    """

    def render_inline(
        self,
        run: object,
        element: ET.Element,
    ) -> None:
        """
        elementからLaTeXテキストを取り出し、OMMLに変換してrunのXMLに埋め込む。

        インライン数式（span.arithmatex）をテキストボックス内に埋め込む。

        Parameters
        ----------
        run : object
            python-pptxのRunオブジェクト。
        element : ET.Element
            インライン数式要素（span class="arithmatex"）。
        """
        latex_text = self._extract_latex(element)
        if not latex_text:
            return

        try:
            omml_str = self._latex_to_omml(latex_text)
            # runのXML要素にOMMLを追加する
            run_elem = run._r
            omml_elem = etree.fromstring(omml_str)
            run_elem.addnext(omml_elem)
        except Exception:
            # 変換失敗時はLaTeXテキストをそのままrunに設定する
            run.text = latex_text

    def render_block(
        self,
        slide: Slide,
        element: ET.Element,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        elementからLaTeXテキストを取り出し、OMMLに変換して指定座標の
        数式Shapeとしてスライドに配置する。

        ブロック数式（div.arithmatex）を独立した数式Shapeとして配置する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        element : ET.Element
            ブロック数式要素（div class="arithmatex"）。
        left : int
            左端座標（EMU）。
        top : int
            上端座標（EMU）。
        width : int
            幅（EMU）。
        height : int
            高さ（EMU）。
        """
        latex_text = self._extract_latex(element)
        if not latex_text:
            return

        try:
            omml_str = self._latex_to_omml(latex_text)
            # スライドのXML要素に直接OMMLを追加する
            spTree = slide.shapes._spTree
            omml_elem = etree.fromstring(omml_str)
            spTree.append(omml_elem)
        except Exception:
            # 変換失敗時はテキストボックスにLaTeXテキストを表示する
            shape = slide.shapes.add_textbox(
                Emu(left), Emu(top), Emu(width), Emu(height)
            )
            tf = shape.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = latex_text
            run.font.size = Pt(14)

    def _extract_latex(self, element: ET.Element) -> str:
        """
        arithmatex要素からLaTeXテキストを取り出す。

        arithmatexは ``\\(...\\)`` または ``\\[...\\]`` 形式でテキストを格納するため、
        デリミタを除去して純粋なLaTeXテキストを返す。

        Parameters
        ----------
        element : ET.Element
            arithmatex要素（span/div class="arithmatex"）。

        Returns
        -------
        str
            LaTeXテキスト（デリミタ除去済み）。
        """
        text = "".join(element.itertext()).strip()
        # インライン数式のデリミタ \(...\) を除去する
        if text.startswith("\\(") and text.endswith("\\)"):
            return text[2:-2].strip()
        # ブロック数式のデリミタ \[...\] を除去する
        if text.startswith("\\[") and text.endswith("\\]"):
            return text[2:-2].strip()
        # $$...$$ 形式も対応する
        if text.startswith("$$") and text.endswith("$$"):
            return text[2:-2].strip()
        return text

    def _latex_to_omml(self, latex_text: str) -> bytes:
        """
        LaTeXテキストをOMML要素（lxml Element）の bytes に変換する。

        latex2mathmlでMathMLに変換し、mathml2ommlでOMMLに変換する。
        mathml2omml が出力する OMML 文字列は名前空間宣言（xmlns:m=...）を
        持たないため、名前空間を付与してから lxml でパースできる bytes を返す。

        Parameters
        ----------
        latex_text : str
            LaTeXテキスト。

        Returns
        -------
        bytes
            名前空間付き OMML XML バイト列。
        """
        # LaTeXをMathMLに変換する
        mathml_str = latex2mathml.converter.convert(latex_text)
        # MathMLをOMMLに変換する（m: プレフィックスの名前空間宣言が欠落している）
        omml_str = mathml2omml.convert(mathml_str)
        # lxml で解析できるよう名前空間宣言を補完する
        OMML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        omml_with_ns = omml_str.replace(
            "<m:oMath>",
            f'<m:oMath xmlns:m="{OMML_NS}">',
            1,
        )
        return omml_with_ns.encode("utf-8")
