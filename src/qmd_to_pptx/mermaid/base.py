"""
Mermaid図レンダラーの基底クラスモジュール。

各ダイアグラム種別レンダラーが継承する共通ユーティリティメソッドを提供する。
"""

from __future__ import annotations

import logging
import math

from lxml import etree as lxml_etree
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

# モジュールロガーを取得する
logger = logging.getLogger(__name__)

# ノードのデフォルトサイズ（EMU）
NODE_WIDTH_EMU: int = 1200000
NODE_HEIGHT_EMU: int = 500000


class BaseDiagramRenderer:
    """
    全Mermaidダイアグラムレンダラーの基底クラス。

    正規化座標変換・ノード描画・エッジ描画・フォールバック描画の
    共通ユーティリティメソッドを提供する。
    各サブクラスはこのクラスを継承して専用のrender()を実装する。
    """

    def _pos_to_emu(
        self,
        x_norm: float,
        y_norm: float,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> tuple[int, int]:
        """
        spring_layoutで計算した正規化座標（-1.0〜1.0）をEMU座標に変換する。

        Parameters
        ----------
        x_norm : float
            正規化X座標（-1.0〜1.0）。
        y_norm : float
            正規化Y座標（-1.0〜1.0）。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。

        Returns
        -------
        tuple[int, int]
            (x_emu, y_emu) のタプル。
        """
        margin_x = NODE_WIDTH_EMU // 2
        margin_y = NODE_HEIGHT_EMU // 2
        usable_w = width - NODE_WIDTH_EMU
        usable_h = height - NODE_HEIGHT_EMU

        # -1.0〜1.0 を 0〜1 に正規化する
        x_ratio = (x_norm + 1.0) / 2.0
        y_ratio = (y_norm + 1.0) / 2.0

        x_emu = left + margin_x + int(x_ratio * usable_w)
        y_emu = top + margin_y + int(y_ratio * usable_h)
        return x_emu, y_emu

    def _connection_indices(
        self,
        dx: float,
        dy: float,
    ) -> tuple[int, int]:
        """
        方向ベクトル (dx, dy) から始点・終点の接続ポイントインデックスを決定する。

        python-pptx の begin_connect/end_connect に渡すインデックスの実態は
        ECMA-376 のプリセット定義と一部異なる。矩形・丸角矩形等の標準形状では：
        - 0: 上辺中点
        - 1: 左辺中点  ← ECMA-376 の idx=1(右) ではなく左
        - 2: 下辺中点
        - 3: 右辺中点  ← ECMA-376 の idx=3(左) ではなく右

        Parameters
        ----------
        dx : float
            X方向の差（dst.x - src.x）。正規化座標系での値。
        dy : float
            Y方向の差（dst.y - src.y）。正規化座標系での値。

        Returns
        -------
        tuple[int, int]
            (始点接続ポイントインデックス, 終点接続ポイントインデックス)
        """
        if abs(dx) >= abs(dy):
            # 水平方向が支配的な場合
            # python-pptx の接続ポイントインデックス実態（ECMA-376 と異なる）:
            #   idx=0: 上辺中点, idx=1: 左辺中点, idx=2: 下辺中点, idx=3: 右辺中点
            if dx >= 0:
                # 左から右: src右辺(idx=3) → dst左辺(idx=1)
                return (3, 1)
            else:
                # 右から左: src左辺(idx=1) → dst右辺(idx=3)
                return (1, 3)
        else:
            # 垂直方向が支配的な場合
            if dy >= 0:
                # 上から下: src下辺(idx=2) → dst上辺(idx=0)
                return (2, 0)
            else:
                # 下から上: src上辺(idx=0) → dst下辺(idx=2)
                return (0, 2)

    def _draw_nodes(
        self,
        slide: Slide,
        nodes: list[str],
        pos: dict[str, tuple[float, float]],
        left: int,
        top: int,
        width: int,
        height: int,
        label_map: dict[str, str] | None = None,
    ) -> dict[str, object]:
        """
        ノードを矩形Shapeとしてスライドに配置する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        nodes : list[str]
            ノードIDのリスト。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        label_map : dict[str, str] | None
            ノードIDをキー、表示ラベル文字列を値とする辞書。省略時はノードIDを表示する。

        Returns
        -------
        dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        """
        node_shapes: dict[str, object] = {}
        for node_id in nodes:
            if node_id not in pos:
                continue
            x_norm, y_norm = pos[node_id]
            cx, cy = self._pos_to_emu(x_norm, y_norm, left, top, width, height)
            shape_left = cx - NODE_WIDTH_EMU // 2
            shape_top = cy - NODE_HEIGHT_EMU // 2

            # 矩形Shapeを追加する
            shape = slide.shapes.add_shape(
                1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
                Emu(shape_left),
                Emu(shape_top),
                Emu(NODE_WIDTH_EMU),
                Emu(NODE_HEIGHT_EMU),
            )
            # 表示ラベルを決定する（label_mapがあればそれを優先する）
            label = label_map.get(node_id, node_id) if label_map else node_id
            shape.text = label
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)

            node_shapes[node_id] = shape

        return node_shapes

    def _draw_edges(
        self,
        slide: Slide,
        edges: list[tuple[str, str]],
        pos: dict[str, tuple[float, float]],
        node_shapes: dict[str, object],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        エッジをConnectorShapeとしてスライドに描画し、始点・終点のshapeに接続する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        edges : list[tuple[str, str]]
            (始点ノードID, 終点ノードID) のタプルリスト。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        node_shapes : dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        for src, dst in edges:
            if src not in pos or dst not in pos:
                continue
            sx_norm, sy_norm = pos[src]
            dx_norm, dy_norm = pos[dst]
            sx, sy = self._pos_to_emu(sx_norm, sy_norm, left, top, width, height)
            dx, dy = self._pos_to_emu(dx_norm, dy_norm, left, top, width, height)

            # 曲線コネクターを追加する（connector_type=3 = CURVE）
            connector = slide.shapes.add_connector(
                3,  # MSO_CONNECTOR_TYPE.CURVE
                Emu(sx),
                Emu(sy),
                Emu(dx),
                Emu(dy),
            )

            # 方向ベクトルから接続ポイントインデックスを決定する
            src_cp, dst_cp = self._connection_indices(
                dx_norm - sx_norm, dy_norm - sy_norm
            )

            # shapeに begin_connect/end_connect でブロックに接続する
            src_shape = node_shapes.get(src)
            dst_shape = node_shapes.get(dst)
            if src_shape is not None:
                connector.begin_connect(src_shape, src_cp)
            if dst_shape is not None:
                connector.end_connect(dst_shape, dst_cp)

    def _render_fallback(
        self,
        slide: Slide,
        mermaid_text: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        Mermaidパース失敗時のフォールバック処理。
        テキストボックスにMermaidテキストをそのまま表示する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        mermaid_text : str
            Mermaidテキスト。
        left : int
            左端座標（EMU）。
        top : int
            上端座標（EMU）。
        width : int
            幅（EMU）。
        height : int
            高さ（EMU）。
        """
        # 先頭行からダイアグラム種別を取得してログに記録する
        lines = mermaid_text.splitlines() if mermaid_text else []
        first_line = lines[0].strip() if lines else ""
        logger.warning(
            "Mermaidダイアグラムの描画にフォールバックしました（テキストボックス表示）: %s",
            first_line,
        )
        shape = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = mermaid_text
        run.font.name = "Courier New"
        run.font.size = Pt(10)

    def _add_edge_label_near_source(
        self,
        slide: Slide,
        text: str,
        sx: int,
        sy: int,
        dx: int,
        dy: int,
        font_size_pt: int = 10,
    ) -> object:
        """
        エッジラベルを始点ノードの枠線上に中心が来るよう配置する。

        始点から終点方向へのベクトルとノード矩形の交点（= 枠線上）を計算し、
        そこを中心とした背景・枠線なしのテキストボックスを追加する。
        コネクター中点固定方式と異なり、始点ブロックの枠に沿ってラベルを表示できる。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        text : str
            ラベルテキスト。
        sx, sy : int
            始点ノード中心の EMU 座標。
        dx, dy : int
            終点ノード中心の EMU 座標。
        font_size_pt : int
            フォントサイズ（ポイント）。デフォルト 10pt。

        Returns
        -------
        object
            追加したテキストボックス Shape。
        """
        vec_x = dx - sx
        vec_y = dy - sy
        length = math.sqrt(vec_x ** 2 + vec_y ** 2)

        if length < 1:
            # コネクターが極端に短い場合は始点に配置する
            label_cx = sx
            label_cy = sy
        else:
            ux = vec_x / length
            uy = vec_y / length
            # 方向ベクトルとノード矩形の交点を求める（= 始点ノードの枠線上の点）
            half_w = NODE_WIDTH_EMU / 2
            half_h = NODE_HEIGHT_EMU / 2
            if abs(ux) < 1e-9:
                t = half_h
            elif abs(uy) < 1e-9:
                t = half_w
            else:
                t = min(half_w / abs(ux), half_h / abs(uy))
            label_cx = sx + int(ux * t)
            label_cy = sy + int(uy * t)

        box_w = NODE_WIDTH_EMU
        box_h = NODE_HEIGHT_EMU // 2
        box_left = label_cx - box_w // 2
        box_top = label_cy - box_h // 2

        txBox = slide.shapes.add_textbox(
            Emu(box_left), Emu(box_top), Emu(box_w), Emu(box_h)
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.text = text
        for para in tf.paragraphs:
            para.alignment = None  # 左揃え
            for run in para.runs:
                run.font.size = Pt(font_size_pt)
                run.font.bold = False

        # 背景と枠線を透明にする
        txBox.fill.background()
        spPr_el = txBox._element.find(qn("p:spPr"))
        if spPr_el is not None:
            existing_ln = spPr_el.find(qn("a:ln"))
            if existing_ln is None:
                existing_ln = lxml_etree.SubElement(spPr_el, qn("a:ln"))
            lxml_etree.SubElement(existing_ln, qn("a:noFill"))

        return txBox

    @staticmethod
    def _next_shape_id(spTree: object) -> int:
        """
        spTree 内の全 cNvPr id の最大値 + 1 を返す。

        手動で OOXML 要素（grpSp 等）を追加する際の衝突しない ID 採番に使用する。

        Parameters
        ----------
        spTree : object
            スライドの spTree lxml 要素。

        Returns
        -------
        int
            既存 id の最大値 + 1。
        """
        max_id: int = 0
        for el in spTree.iter():
            if el.tag.endswith("}cNvPr"):
                try:
                    max_id = max(max_id, int(el.get("id", 0)))
                except ValueError:
                    pass
        return max_id + 1

    def _group_node_with_labels(
        self,
        slide: Slide,
        node_shape: object,
        txboxes: list,
    ) -> None:
        """
        始点ノードシェイプとラベルテキストボックス群をグループ化する。

        コネクターはトップレベルに残し、begin_connect/end_connect で設定された
        接続（シェイプIDによる参照）をそのまま維持する。
        ノードシェイプをグループ内に移動してもシェイプIDは保持されるため、
        PowerPoint はグループ内シェイプへの接続を正しく解決できる。
        1つのノードに複数のラベル付きエッジがある場合も1グループにまとめる。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        node_shape : object
            グループ化対象の始点ノード Shape オブジェクト。
        txboxes : list
            グループ化対象のラベルテキストボックス Shape オブジェクトのリスト。
        """
        spTree = slide.shapes._spTree
        node_el = node_shape._element
        txb_els = [tb._element for tb in txboxes]

        def _get_xywh(el: object) -> tuple[int, int, int, int]:
            """spPr/xfrm から left, top, cx, cy を取得する。"""
            spPr = el.find(qn("p:spPr"))
            if spPr is None:
                return 0, 0, 0, 0
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is None:
                return 0, 0, 0, 0
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            if off is None or ext is None:
                return 0, 0, 0, 0
            return (
                int(off.get("x", 0)),
                int(off.get("y", 0)),
                int(ext.get("cx", 0)),
                int(ext.get("cy", 0)),
            )

        # 全要素のバウンディングボックスを計算する
        all_els = [node_el] + txb_els
        lefts, tops, rights, bottoms = [], [], [], []
        for el in all_els:
            l, t, w, h = _get_xywh(el)
            lefts.append(l)
            tops.append(t)
            rights.append(l + w)
            bottoms.append(t + h)

        grp_left = min(lefts)
        grp_top = min(tops)
        grp_w = max(1, max(rights) - grp_left)
        grp_h = max(1, max(bottoms) - grp_top)

        # <p:grpSp> 要素を構築する（OOXML スキーマ: 第1子は nvGrpSpPr が必須）
        grp_id = self._next_shape_id(spTree)
        grpSp = lxml_etree.Element(qn("p:grpSp"))
        nvGrpSpPr = lxml_etree.SubElement(grpSp, qn("p:nvGrpSpPr"))
        cNvPr_el = lxml_etree.SubElement(nvGrpSpPr, qn("p:cNvPr"))
        cNvPr_el.set("id", str(grp_id))
        cNvPr_el.set("name", f"グループ {grp_id}")
        lxml_etree.SubElement(nvGrpSpPr, qn("p:cNvGrpSpPr"))
        lxml_etree.SubElement(nvGrpSpPr, qn("p:nvPr"))
        grpSpPr = lxml_etree.SubElement(grpSp, qn("p:grpSpPr"))
        xfrm = lxml_etree.SubElement(grpSpPr, qn("a:xfrm"))
        off_el = lxml_etree.SubElement(xfrm, qn("a:off"))
        off_el.set("x", str(grp_left))
        off_el.set("y", str(grp_top))
        ext_el = lxml_etree.SubElement(xfrm, qn("a:ext"))
        ext_el.set("cx", str(grp_w))
        ext_el.set("cy", str(grp_h))
        # 子座標系をスライド座標系と同一（恒等変換）にして絶対座標をそのまま保持する
        chOff = lxml_etree.SubElement(xfrm, qn("a:chOff"))
        chOff.set("x", str(grp_left))
        chOff.set("y", str(grp_top))
        chExt = lxml_etree.SubElement(xfrm, qn("a:chExt"))
        chExt.set("cx", str(grp_w))
        chExt.set("cy", str(grp_h))

        # spTree からシェイプを取り出して grpSp に移動する
        spTree.remove(node_el)
        for txb_el in txb_els:
            spTree.remove(txb_el)
        grpSp.append(node_el)
        for txb_el in txb_els:
            grpSp.append(txb_el)

        # grpSp を spTree に追加する
        spTree.append(grpSp)
