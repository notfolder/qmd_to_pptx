"""
クラス図レンダラーモジュール。

classDiagram を UMLクラス図形式（ヘッダー・属性・メソッドの3段構成）で描画する。
"""

from __future__ import annotations

import math
from lxml import etree as lxml_etree
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer


# クラス図ボックスのサイズ定数（EMU）
_CLASS_WIDTH_EMU: int = 1700000          # クラスボックスの幅
_CLASS_HEADER_HEIGHT_EMU: int = 440000   # クラス名ヘッダーの高さ
_CLASS_ROW_HEIGHT_EMU: int = 280000      # メンバー1行あたりの高さ
_CLASS_MIN_SECTION_HEIGHT_EMU: int = 320000  # 属性/メソッドセクションの最小高さ

# mermaid-parser-py の relationType 定数
_AGGREGATION = 0
_EXTENSION   = 1
_COMPOSITION = 2
_DEPENDENCY  = 3


class ClassDiagramRenderer(BaseDiagramRenderer):
    """
    classDiagram をUMLクラス図形式で描画するレンダラー。

    各クラスをヘッダー・属性・メソッドの3段構成の矩形グループとして描画し、
    関係（継承・コンポジション・集約・依存など）をUML矢印スタイルのコネクターで描画する。
    """

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        classDiagramをスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"classes"と"relations"キーを含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        classes = graph_data.get("classes", {})
        relations = graph_data.get("relations", [])
        if not classes:
            return

        nodes = list(classes.keys())

        # 格子配置でノードを均等に並べる（spring_layoutによる重なりを防ぐ）
        n = len(nodes)
        n_cols = max(1, math.ceil(math.sqrt(n)))
        n_rows = math.ceil(n / n_cols)

        # _draw_edges_classが参照するposを格子の正規化座標で構築する
        pos: dict[str, tuple[float, float]] = {}
        for i, node_id in enumerate(nodes):
            col = i % n_cols
            row = i // n_cols
            x_norm = -1.0 + 2.0 * col / max(n_cols - 1, 1) if n_cols > 1 else 0.0
            y_norm = -1.0 + 2.0 * row / max(n_rows - 1, 1) if n_rows > 1 else 0.0
            pos[node_id] = (x_norm, y_norm)

        # クラス図用レイアウト計算（クラス本体の高さに合わせる）
        max_class_h = max(
            (self._class_total_height(classes[nid]) for nid in nodes if nid in classes),
            default=800000,
        )
        margin_x = _CLASS_WIDTH_EMU // 2
        margin_y = max_class_h // 2
        usable_w = max(width - _CLASS_WIDTH_EMU, _CLASS_WIDTH_EMU)
        usable_h = max(height - max_class_h, max_class_h)

        node_shapes: dict[str, object] = {}
        for node_id in nodes:
            if node_id not in pos or node_id not in classes:
                continue
            x_norm, y_norm = pos[node_id]
            x_ratio = (x_norm + 1.0) / 2.0
            y_ratio = (y_norm + 1.0) / 2.0
            cx = left + margin_x + int(x_ratio * usable_w)
            cy = top + margin_y + int(y_ratio * usable_h)
            anchor = self._draw_class_box(slide, node_id, classes[node_id], cx, cy)
            node_shapes[node_id] = anchor

        # UML矢印スタイルでエッジを描画する
        self._draw_edges_class(slide, relations, pos, node_shapes, left, top, width, height)

    def _class_total_height(self, class_data: dict) -> int:
        """
        クラスボックスの合計高さ（EMU）を返す。

        Parameters
        ----------
        class_data : dict
            mermaid-parser-pyが返すクラス情報辞書。

        Returns
        -------
        int
            ヘッダー・属性セクション・メソッドセクションの合計高さ（EMU）。
        """
        n_members = len(class_data.get("members", []))
        n_methods = len(class_data.get("methods", []))
        attrs_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, n_members * _CLASS_ROW_HEIGHT_EMU)
        methods_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, n_methods * _CLASS_ROW_HEIGHT_EMU)
        return _CLASS_HEADER_HEIGHT_EMU + attrs_h + methods_h

    def _draw_class_box(
        self,
        slide: Slide,
        node_id: str,
        class_data: dict,
        cx: int,
        cy: int,
    ) -> object:
        """
        UMLクラス図の1クラス分のボックスを描画してヘッダーShapeを返す。

        ヘッダー（クラス名）・属性セクション・メソッドセクションの3段構成の矩形を
        縦に重ねて配置し、グループ化する。コネクターのアンカーとしてヘッダーShapeを返す。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        node_id : str
            クラスID（クラス名）。
        class_data : dict
            mermaid-parser-pyが返すクラス情報辞書。
        cx : int
            クラスボックス中心のX座標（EMU）。
        cy : int
            クラスボックス中心のY座標（EMU）。

        Returns
        -------
        object
            ヘッダー部のShapeオブジェクト（コネクターのアンカーとして使用）。
        """
        members = class_data.get("members", [])
        methods = class_data.get("methods", [])
        attrs_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, len(members) * _CLASS_ROW_HEIGHT_EMU)
        methods_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, len(methods) * _CLASS_ROW_HEIGHT_EMU)
        total_h = _CLASS_HEADER_HEIGHT_EMU + attrs_h + methods_h

        box_left = cx - _CLASS_WIDTH_EMU // 2
        header_top = cy - total_h // 2

        # ---- ヘッダーセクション（クラス名・太字・中央寄せ）----
        header_shape = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            Emu(box_left), Emu(header_top),
            Emu(_CLASS_WIDTH_EMU), Emu(_CLASS_HEADER_HEIGHT_EMU),
        )
        tf = header_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(50000)
        tf.margin_right = Emu(50000)
        tf.margin_top = Emu(60000)
        tf.margin_bottom = Emu(60000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = class_data.get("label", node_id)
        run.font.bold = True
        run.font.size = Pt(11)

        # ---- 属性セクション ----
        attrs_top = header_top + _CLASS_HEADER_HEIGHT_EMU
        attrs_shape = slide.shapes.add_shape(
            1,
            Emu(box_left), Emu(attrs_top),
            Emu(_CLASS_WIDTH_EMU), Emu(attrs_h),
        )
        tf_a = attrs_shape.text_frame
        tf_a.word_wrap = True
        tf_a.margin_left = Emu(80000)
        tf_a.margin_right = Emu(50000)
        tf_a.margin_top = Emu(50000)
        tf_a.margin_bottom = Emu(50000)
        if members:
            first = True
            for member in members:
                # バックスラッシュエスケープを除去する
                raw = member.get("text", "").lstrip("\\")
                if first:
                    p = tf_a.paragraphs[0]
                    first = False
                else:
                    p = tf_a.add_paragraph()
                run = p.add_run()
                run.text = raw
                run.font.size = Pt(9)

        # ---- メソッドセクション ----
        methods_top = attrs_top + attrs_h
        methods_shape = slide.shapes.add_shape(
            1,
            Emu(box_left), Emu(methods_top),
            Emu(_CLASS_WIDTH_EMU), Emu(methods_h),
        )
        tf_m = methods_shape.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = Emu(80000)
        tf_m.margin_right = Emu(50000)
        tf_m.margin_top = Emu(50000)
        tf_m.margin_bottom = Emu(50000)
        if methods:
            first = True
            for method in methods:
                raw = method.get("text", "").lstrip("\\")
                if first:
                    p = tf_m.paragraphs[0]
                    first = False
                else:
                    p = tf_m.add_paragraph()
                run = p.add_run()
                run.text = raw
                run.font.size = Pt(9)

        # 3つのShapeをグループ化する
        # グループ内でもshapeのIDは保持されるためコネクターの接続ポイントに使用可能
        slide.shapes.add_group_shape([header_shape, attrs_shape, methods_shape])

        return header_shape

    def _draw_edges_class(
        self,
        slide: Slide,
        relations: list[dict],
        pos: dict[str, tuple[float, float]],
        node_shapes: dict[str, object],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        classDiagram のエッジをUML矢印スタイルで描画する。

        mermaid-parser-py の relationType に応じて以下の矢印を設定する。
        - EXTENSION(1)  : 三角矢印（継承）
        - COMPOSITION(2): ひし形（コンポジション、大サイズ）
        - AGGREGATION(0): ひし形（集約、中サイズ）
        - DEPENDENCY(3) : 開放矢印（依存）
        lineType=1 の場合は点線コネクターとして描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        relations : list[dict]
            mermaid-parser-py の relations リスト。各要素に id1, id2, relation キーを持つ。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、spring_layout 正規化座標を値とする辞書。
        node_shapes : dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        for rel in relations:
            src = rel.get("id1")
            dst = rel.get("id2")
            if not src or not dst:
                continue
            if src not in pos or dst not in pos:
                continue

            relation = rel.get("relation", {})
            type1 = relation.get("type1", "none")    # id1(src)側の関係種別
            type2 = relation.get("type2", "none")    # id2(dst)側の関係種別
            line_type = relation.get("lineType", 0)  # 0=実線, 1=点線

            sx_norm, sy_norm = pos[src]
            dx_norm, dy_norm = pos[dst]
            sx, sy = self._pos_to_emu(sx_norm, sy_norm, left, top, width, height)
            dx, dy = self._pos_to_emu(dx_norm, dy_norm, left, top, width, height)

            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                Emu(sx), Emu(sy), Emu(dx), Emu(dy),
            )

            src_cp, dst_cp = self._connection_indices(dx_norm - sx_norm, dy_norm - sy_norm)
            src_shape = node_shapes.get(src)
            dst_shape = node_shapes.get(dst)
            if src_shape is not None:
                connector.begin_connect(src_shape, src_cp)
            if dst_shape is not None:
                connector.end_connect(dst_shape, dst_cp)

            # コネクター内の <p:spPr> から <a:ln> を取得または作成する
            cxn_el = connector._element
            spPr = cxn_el.find(qn("p:spPr"))
            if spPr is None:
                continue
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                ln = lxml_etree.SubElement(spPr, qn("a:ln"))

            # 点線スタイル（lineType=1）: <a:prstDash val="dash"/>
            if line_type == 1:
                prstDash = ln.find(qn("a:prstDash"))
                if prstDash is None:
                    prstDash = lxml_etree.SubElement(ln, qn("a:prstDash"))
                prstDash.set("val", "dash")

            # headEnd: コネクター終点（id2=dst）側の矢印を設定する
            if type2 == _EXTENSION:
                head = lxml_etree.SubElement(ln, qn("a:headEnd"))
                head.set("type", "triangle")
                head.set("w", "lg")
                head.set("len", "lg")
            elif type2 == _DEPENDENCY:
                head = lxml_etree.SubElement(ln, qn("a:headEnd"))
                head.set("type", "arrow")
                head.set("w", "med")
                head.set("len", "med")
            elif type2 == _COMPOSITION:
                head = lxml_etree.SubElement(ln, qn("a:headEnd"))
                head.set("type", "diamond")
                head.set("w", "lg")
                head.set("len", "lg")
            elif type2 == _AGGREGATION:
                head = lxml_etree.SubElement(ln, qn("a:headEnd"))
                head.set("type", "diamond")
                head.set("w", "med")
                head.set("len", "med")

            # tailEnd: コネクター始点（id1=src）側の矢印を設定する
            if type1 == _EXTENSION:
                tail = lxml_etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", "triangle")
                tail.set("w", "lg")
                tail.set("len", "lg")
            elif type1 == _COMPOSITION:
                tail = lxml_etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", "diamond")
                tail.set("w", "lg")
                tail.set("len", "lg")
            elif type1 == _AGGREGATION:
                tail = lxml_etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", "diamond")
                tail.set("w", "med")
                tail.set("len", "med")
            elif type1 == _DEPENDENCY:
                tail = lxml_etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", "arrow")
                tail.set("w", "med")
                tail.set("len", "med")
