"""
ER図レンダラーモジュール（拡張版）。

erDiagram のエンティティ（ヘッダー + アトリビュート行）とリレーションシップ
（カーディナリティ・ロールラベル・実線/破線コネクター）を描画する。

【グループ化の方針】
- エンティティ: full_rect（透明・境界線・コネクター接続用） + header + row_shapes
  + そのエンティティ境界に接するカーディナリティラベルを全部まとめて OOXML グループ化する
- リレーション: コネクター単体 + ロールラベル（中点）のみをグループ化する
- コネクターはエンティティグループ内の full_rect に begin/end_connect する
"""

from __future__ import annotations

import math
from collections import defaultdict

import networkx as nx
from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer


class ErDiagramRenderer(BaseDiagramRenderer):
    """
    erDiagram を描画するレンダラー。

    - エンティティ: ヘッダー（濃い青）+ アトリビュート行（型・名前・PK/FK/UK・コメント）
    - リレーション: カーディナリティテキスト付き直線コネクター
      - IDENTIFYING     → 実線
      - NON_IDENTIFYING → 破線
    - カーディナリティラベルはエンティティ側グループに含める
    - コネクターはロールラベルとのみグループ化する
    """

    # ---- エンティティボックス定数（EMU） ----
    _ENTITY_WIDTH: int = 2_400_000     # エンティティ幅
    _HEADER_HEIGHT: int = 420_000      # ヘッダー行高さ
    _ROW_HEIGHT: int = 300_000         # アトリビュート行高さ
    _MIN_BODY_HEIGHT: int = 320_000    # 属性なし時の最小 body 高さ

    # ---- 色定義（R, G, B タプル） ----
    _HEADER_BG: tuple[int, int, int] = (31, 73, 125)       # ヘッダー背景（濃い青）
    _HEADER_FG: tuple[int, int, int] = (255, 255, 255)     # ヘッダー文字（白）
    _ROW_BG_EVEN: tuple[int, int, int] = (255, 255, 255)   # 偶数行背景（白）
    _ROW_BG_ODD: tuple[int, int, int] = (235, 241, 252)    # 奇数行背景（薄青）
    _BORDER_COLOR: tuple[int, int, int] = (31, 73, 125)    # 枠線（濃い青）
    _PK_COLOR: tuple[int, int, int] = (196, 0, 0)          # PK テキスト（赤）
    _FK_COLOR: tuple[int, int, int] = (0, 70, 180)         # FK テキスト（青）
    _UK_COLOR: tuple[int, int, int] = (0, 128, 0)          # UK テキスト（緑）
    _COMMENT_COLOR: tuple[int, int, int] = (130, 130, 130) # コメント（薄灰）
    _TYPE_COLOR: tuple[int, int, int] = (100, 100, 100)    # 型名（グレー）
    _REL_LINE_COLOR: tuple[int, int, int] = (50, 50, 50)   # リレーション線

    # ---- カーディナリティラベル定数 ----
    _CARD_FONT_SIZE: int = 9     # カーディナリティフォントサイズ（pt）
    _ROLE_FONT_SIZE: int = 9     # ロールラベルフォントサイズ（pt）
    _CARD_BOX_W: int = 350_000  # カーディナリティテキストボックス幅（EMU）
    _CARD_BOX_H: int = 220_000  # カーディナリティテキストボックス高さ（EMU）
    _CARD_GAP: int = 60_000     # エンティティ境界からカーディナリティラベルまでのギャップ（EMU）

    # ---- パーサー値 → 表示テキスト変換マップ ----
    _CARD_MAP: dict[str, str] = {
        "ONLY_ONE": "1",
        "ZERO_OR_ONE": "0..1",
        "ZERO_OR_MORE": "0..*",
        "ONE_OR_MORE": "1..*",
    }

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        erDiagram をスライドに描画する。

        処理の流れ:
        1. NetworkX でエンティティ配置を計算する
        2. 各エンティティの Shape 群を生成する（full_rect + header + rows）
        3. 各リレーションのカーディナリティラベルを作成し、エンティティ別に振り分ける
        4. エンティティ Shape + 所属カーディナリティラベルを OOXML グループ化する
        5. コネクターをグループ内の full_rect に接続し、ロールラベルとグループ化する

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            "entities" と "relationships" キーを含む辞書。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        entities: dict = graph_data.get("entities", {})
        relationships: list[dict] = graph_data.get("relationships", [])
        if not entities:
            return

        # 内部IDからエンティティ名への逆引きマップを構築する
        id_to_label: dict[str, str] = {
            ent.get("id", ""): label
            for label, ent in entities.items()
            if ent.get("id")
        }

        nodes: list[str] = list(entities.keys())

        # エンティティごとの高さを算出する（ヘッダー + アトリビュート行数）
        entity_heights: dict[str, int] = {
            name: self._entity_height(len(entities[name].get("attributes", [])))
            for name in nodes
        }

        # NetworkX で配置計算する
        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for rel in relationships:
            a_label = id_to_label.get(rel.get("entityA", ""))
            b_label = id_to_label.get(rel.get("entityB", ""))
            if a_label and b_label:
                G.add_edge(a_label, b_label)
        try:
            pos: dict[str, tuple[float, float]] = nx.kamada_kawai_layout(G)
        except Exception:
            pos = nx.spring_layout(G, seed=42, k=2.0)

        # 正規化座標 → EMU 座標変換（エンティティの可変高さを考慮する）
        max_h = max(entity_heights.values()) if entity_heights else self._HEADER_HEIGHT
        emu_pos: dict[str, tuple[int, int]] = self._norm_to_emu_er(
            pos, nodes, max_h, left, top, width, height
        )

        # ---- Step 1: エンティティ Shape 群を生成する ----
        # (full_rect, [all_shapes]) の辞書を作る
        entity_shape_info: dict[str, tuple[object, list[object]]] = {}
        for name in nodes:
            entity = entities[name]
            cx, cy = emu_pos[name]
            eh = entity_heights[name]
            full_rect, all_shapes = self._draw_er_entity_shapes(
                slide, name, entity, cx, cy, eh
            )
            entity_shape_info[name] = (full_rect, all_shapes)

        # ---- Step 2: 全リレーションのカーディナリティラベルを生成し、エンティティ側に振り分ける ----
        # エンティティ名 → そのエンティティ側に接するカーディナリティラベルのリスト
        entity_card_labels: dict[str, list[object]] = defaultdict(list)
        # リレーションごとに (connector_info) を記録する
        connector_tasks: list[dict] = []

        for rel in relationships:
            a_label = id_to_label.get(rel.get("entityA", ""))
            b_label = id_to_label.get(rel.get("entityB", ""))
            if not a_label or not b_label:
                continue
            if a_label not in emu_pos or b_label not in emu_pos:
                continue

            relSpec: dict = rel.get("relSpec", {})
            card_b: str = relSpec.get("cardB", "ONLY_ONE")   # A 側ラベル
            card_a: str = relSpec.get("cardA", "ZERO_OR_MORE")  # B 側ラベル
            rel_type: str = relSpec.get("relType", "IDENTIFYING")
            role_a: str = rel.get("roleA", "")

            ax, ay = emu_pos[a_label]
            bx, by = emu_pos[b_label]
            dx = bx - ax
            dy = by - ay
            length = math.sqrt(dx * dx + dy * dy)
            if length < 1:
                continue
            nx_v = dx / length
            ny_v = dy / length

            hw_a = self._ENTITY_WIDTH // 2
            hh_a = entity_heights[a_label] // 2
            hw_b = self._ENTITY_WIDTH // 2
            hh_b = entity_heights[b_label] // 2

            # 各エンティティの境界点を計算する
            bpax, bpay = self._rect_border_point(ax, ay, hw_a, hh_a, nx_v, ny_v)
            bpbx, bpby = self._rect_border_point(bx, by, hw_b, hh_b, -nx_v, -ny_v)

            # A 側カーディナリティラベル（境界点から B 方向へ少しオフセット）
            offset = self._CARD_BOX_H // 2 + self._CARD_GAP
            ca_cx = bpax + int(nx_v * offset)
            ca_cy = bpay + int(ny_v * offset)
            tb_card_a = self._make_card_label(
                slide, self._CARD_MAP.get(card_b, card_b), ca_cx, ca_cy
            )
            entity_card_labels[a_label].append(tb_card_a)

            # B 側カーディナリティラベル（境界点から A 方向へ少しオフセット）
            cb_cx = bpbx - int(nx_v * offset)
            cb_cy = bpby - int(ny_v * offset)
            tb_card_b = self._make_card_label(
                slide, self._CARD_MAP.get(card_a, card_a), cb_cx, cb_cy
            )
            entity_card_labels[b_label].append(tb_card_b)

            # ロールラベルを A 側（始点側）境界点の近くに生成し、A のグループに含める
            if role_a:
                role_offset = offset + self._CARD_BOX_H + self._CARD_GAP
                role_cx = bpax + int(nx_v * role_offset)
                role_cy = bpay + int(ny_v * role_offset)
                tb_role = self._make_role_label(slide, role_a, role_cx, role_cy)
                entity_card_labels[a_label].append(tb_role)

            # コネクタータスクを記録する（グループ化後に実行）
            connector_tasks.append({
                "a_label": a_label,
                "b_label": b_label,
                "bpax": bpax, "bpay": bpay,
                "bpbx": bpbx, "bpby": bpby,
                "dx": dx, "dy": dy,
                "rel_type": rel_type,
            })

        # ---- Step 3: エンティティ Shape + カーディナリティラベルをグループ化する ----
        entity_group_shapes: dict[str, object] = {}
        for name in nodes:
            full_rect, all_shapes = entity_shape_info[name]
            card_labels = entity_card_labels.get(name, [])
            extra_shapes = all_shapes[1:] + card_labels  # full_rect 以外 + ラベル
            # full_rect をグループ代表として保持し、グループ化する
            grp_shape = self._group_entity_shapes(slide, full_rect, extra_shapes)
            entity_group_shapes[name] = grp_shape

        # ---- Step 4: コネクターを描画し、エンティティグループに接続する ----
        for task in connector_tasks:
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                Emu(task["bpax"]), Emu(task["bpay"]),
                Emu(task["bpbx"]), Emu(task["bpby"]),
            )
            connector.line.color.rgb = RGBColor(*self._REL_LINE_COLOR)
            if task["rel_type"] == "NON_IDENTIFYING":
                self._set_connector_dash(connector)

            # エンティティグループ（の代表 full_rect）にコネクターを接続する
            src_cp, dst_cp = self._connection_indices(task["dx"], task["dy"])
            src_grp = entity_group_shapes.get(task["a_label"])
            dst_grp = entity_group_shapes.get(task["b_label"])
            if src_grp is not None:
                connector.begin_connect(src_grp, src_cp)
            if dst_grp is not None:
                connector.end_connect(dst_grp, dst_cp)


    # ------------------------------------------------------------------
    # エンティティ描画
    # ------------------------------------------------------------------

    def _entity_height(self, attr_count: int) -> int:
        """
        アトリビュート数からエンティティボックスの高さ（EMU）を算出する。

        Parameters
        ----------
        attr_count : int
            アトリビュート数。

        Returns
        -------
        int
            エンティティ高さ（EMU）。
        """
        body_h = max(self._MIN_BODY_HEIGHT, attr_count * self._ROW_HEIGHT)
        return self._HEADER_HEIGHT + body_h

    def _norm_to_emu_er(
        self,
        pos: dict[str, tuple[float, float]],
        nodes: list[str],
        max_entity_height: int,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> dict[str, tuple[int, int]]:
        """
        NetworkX の正規化座標をEMU座標に変換する。

        エンティティの幅・最大高さをパディングとして考慮し、
        キャンバス端でエンティティが欠けないようにする。

        Parameters
        ----------
        pos : dict[str, tuple[float, float]]
            NetworkX 正規化座標辞書。
        nodes : list[str]
            ノード名リスト。
        max_entity_height : int
            全エンティティの最大高さ（パディング計算用）。
        left, top, width, height : int
            描画エリア EMU 座標。

        Returns
        -------
        dict[str, tuple[int, int]]
            エンティティ名 → (cx_emu, cy_emu) の辞書。
        """
        if not pos:
            return {}
        xs = [p[0] for p in pos.values()]
        ys = [p[1] for p in pos.values()]
        x_min, x_max = min(xs), max(xs)
        y_min, y_max = min(ys), max(ys)
        x_range = max(x_max - x_min, 1e-9)
        y_range = max(y_max - y_min, 1e-9)

        pad_x = self._ENTITY_WIDTH // 2 + 200_000
        pad_y = max_entity_height // 2 + 200_000
        usable_w = max(1, width - 2 * pad_x)
        usable_h = max(1, height - 2 * pad_y)

        result: dict[str, tuple[int, int]] = {}
        for name in nodes:
            if name not in pos:
                continue
            xn, yn = pos[name]
            x_ratio = (xn - x_min) / x_range
            y_ratio = (yn - y_min) / y_range
            cx = left + pad_x + int(x_ratio * usable_w)
            cy = top + pad_y + int(y_ratio * usable_h)
            result[name] = (cx, cy)
        return result

    def _draw_er_entity_shapes(
        self,
        slide: Slide,
        name: str,
        entity: dict,
        cx: int,
        cy: int,
        entity_height: int,
    ) -> tuple[object, list[object]]:
        """
        エンティティボックスの全 Shape を生成して返す（グループ化は呼び出し側で行う）。

        - full_rect: 全体バウンディング矩形（透明 fill・境界線・コネクター接続用代表）
        - header: ヘッダー矩形（エンティティ名・濃い青 fill）
        - row_shapes: アトリビュート行矩形リスト

        Parameters
        ----------
        slide : Slide
        name : str
            エンティティ名。
        entity : dict
            パーサーが返したエンティティ辞書。
        cx, cy : int
            エンティティボックス中心座標（EMU）。
        entity_height : int
            エンティティボックス全体高さ（EMU）。

        Returns
        -------
        tuple[object, list[object]]
            (full_rect, [full_rect, header, *row_shapes])
            ※ リストは full_rect を先頭に含む全 Shape のリスト。
        """
        ew = self._ENTITY_WIDTH
        entity_left = cx - ew // 2
        entity_top = cy - entity_height // 2
        all_shapes: list[object] = []

        # ① 全体バウンディング矩形（透明 fill + 境界線 + コネクター接続用）
        full_rect = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            Emu(entity_left),
            Emu(entity_top),
            Emu(ew),
            Emu(entity_height),
        )
        full_rect.fill.background()
        full_rect.line.color.rgb = RGBColor(*self._BORDER_COLOR)
        full_rect.text_frame.text = ""
        all_shapes.append(full_rect)

        # ② ヘッダー矩形（エンティティ名）
        header = slide.shapes.add_shape(
            1,
            Emu(entity_left),
            Emu(entity_top),
            Emu(ew),
            Emu(self._HEADER_HEIGHT),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*self._HEADER_BG)
        self._set_no_border(header)

        display_name = entity.get("alias") or name
        tf = header.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = display_name
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*self._HEADER_FG)
        all_shapes.append(header)

        # ③ アトリビュート行（各行の矩形 Shape のテキストとして描画する）
        attributes: list[dict] = entity.get("attributes", [])
        body_top = entity_top + self._HEADER_HEIGHT

        for i, attr in enumerate(attributes):
            row_top = body_top + i * self._ROW_HEIGHT
            bg_color = self._ROW_BG_EVEN if i % 2 == 0 else self._ROW_BG_ODD
            row_shape = slide.shapes.add_shape(
                1,
                Emu(entity_left),
                Emu(row_top),
                Emu(ew),
                Emu(self._ROW_HEIGHT),
            )
            row_shape.fill.solid()
            row_shape.fill.fore_color.rgb = RGBColor(*bg_color)
            self._set_no_border(row_shape)
            self._fill_attribute_row_text(row_shape, attr)
            all_shapes.append(row_shape)

        return full_rect, all_shapes

    def _group_entity_shapes(
        self,
        slide: Slide,
        full_rect: object,
        extra_shapes: list[object],
    ) -> object:
        """
        エンティティの全 Shape（full_rect + header + rows + card_labels）を
        OOXML グループ化し、グループ内の full_rect を返す。

        コネクターは後からグループ内の full_rect に begin/end_connect するため、
        full_rect の XML 要素はグループ内でも shape ID を保持し続ける。

        Parameters
        ----------
        slide : Slide
        full_rect : object
            グループ代表（コネクター接続用の透明矩形）。
        extra_shapes : list[object]
            header, row_shapes, card_labels 等の残り Shape。

        Returns
        -------
        object
            グループ化後の full_rect Shape（Shape ID は変わらず）。
        """
        spTree = slide.shapes._spTree
        all_els = [full_rect._element] + [s._element for s in extra_shapes]

        def _xywh(el: object) -> tuple[int, int, int, int]:
            spPr = el.find(qn("p:spPr"))
            if spPr is None:
                return 0, 0, 0, 0
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is None:
                return 0, 0, 0, 0
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            if off is None or ext is None:
                return 0, 0, 0, 0
            return (
                int(off.get("x", 0)),
                int(off.get("y", 0)),
                int(ext.get("cx", 0)),
                int(ext.get("cy", 0)),
            )

        lefts, tops, rights, bottoms = [], [], [], []
        for el in all_els:
            l, t, w, h = _xywh(el)
            lefts.append(l)
            tops.append(t)
            rights.append(l + w)
            bottoms.append(t + h)

        grp_left = min(lefts)
        grp_top = min(tops)
        grp_w = max(1, max(rights) - grp_left)
        grp_h = max(1, max(bottoms) - grp_top)

        grpSp = lxml_etree.Element(qn("p:grpSp"))
        grpSpPr = lxml_etree.SubElement(grpSp, qn("p:grpSpPr"))
        xfrm = lxml_etree.SubElement(grpSpPr, qn("a:xfrm"))
        off_el = lxml_etree.SubElement(xfrm, qn("a:off"))
        off_el.set("x", str(grp_left))
        off_el.set("y", str(grp_top))
        ext_el = lxml_etree.SubElement(xfrm, qn("a:ext"))
        ext_el.set("cx", str(grp_w))
        ext_el.set("cy", str(grp_h))
        # 子座標系をスライド座標系と同一にして絶対座標をそのまま保持する
        chOff = lxml_etree.SubElement(xfrm, qn("a:chOff"))
        chOff.set("x", str(grp_left))
        chOff.set("y", str(grp_top))
        chExt = lxml_etree.SubElement(xfrm, qn("a:chExt"))
        chExt.set("cx", str(grp_w))
        chExt.set("cy", str(grp_h))

        for el in all_els:
            try:
                spTree.remove(el)
            except ValueError:
                pass
            grpSp.append(el)

        spTree.append(grpSp)
        return full_rect

    def _fill_attribute_row_text(
        self,
        row_shape: object,
        attr: dict,
    ) -> None:
        """
        アトリビュート行 Shape のテキストフレームに型名・属性名・キー・コメントを設定する。

        独立したテキストボックスを作らず、行矩形 Shape 自体のテキストフレームを使う。
        内部パディングを lxml で調整して視覚的な余白を確保する。

        Parameters
        ----------
        row_shape : object
            アトリビュート行矩形 Shape。
        attr : dict
            パーサーのアトリビュート辞書（"type"・"name"・"keys"・"comment"）。
        """
        tf = row_shape.text_frame
        tf.word_wrap = False

        # lxml で内部パディング（insFocus, lIns 等）を設定する
        txBody = tf._txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("lIns", "80000")   # 左余白（EMU）
            bodyPr.set("tIns", "30000")   # 上余白（EMU）
            bodyPr.set("rIns", "80000")   # 右余白（EMU）
            bodyPr.set("bIns", "30000")   # 下余白（EMU）
            bodyPr.set("anchor", "ctr")   # 縦方向中央揃え

        para = tf.paragraphs[0]

        # 型名（グレー・等幅で幅を揃える）
        r_type = para.add_run()
        r_type.text = f"{attr.get('type', ''):<10}"
        r_type.font.size = Pt(9)
        r_type.font.color.rgb = RGBColor(*self._TYPE_COLOR)

        # 属性名（キーがあれば太字）
        keys: list[str] = attr.get("keys", [])
        r_name = para.add_run()
        r_name.text = f"  {attr.get('name', ''):<18}"
        r_name.font.size = Pt(9)
        r_name.font.bold = len(keys) > 0

        # キー（PK / FK / UK）— 色分けして表示する
        if keys:
            r_key = para.add_run()
            r_key.text = f"  {','.join(keys)}"
            r_key.font.size = Pt(8)
            r_key.font.bold = True
            if "PK" in keys and "FK" in keys:
                r_key.font.color.rgb = RGBColor(*self._PK_COLOR)
            elif "PK" in keys:
                r_key.font.color.rgb = RGBColor(*self._PK_COLOR)
            elif "FK" in keys:
                r_key.font.color.rgb = RGBColor(*self._FK_COLOR)
            else:
                r_key.font.color.rgb = RGBColor(*self._UK_COLOR)

        # コメント（イタリック・薄グレー）
        comment: str = attr.get("comment", "")
        if comment:
            r_cmt = para.add_run()
            r_cmt.text = f'  "{comment}"'
            r_cmt.font.size = Pt(8)
            r_cmt.font.italic = True
            r_cmt.font.color.rgb = RGBColor(*self._COMMENT_COLOR)

    # ------------------------------------------------------------------
    # ジオメトリユーティリティ
    # ------------------------------------------------------------------

    def _rect_border_point(
        self,
        cx: int,
        cy: int,
        hw: int,
        hh: int,
        nx_v: float,
        ny_v: float,
    ) -> tuple[int, int]:
        """
        矩形の中心から方向ベクトル (nx_v, ny_v) 方向の境界点を計算する。

        Parameters
        ----------
        cx, cy : int
            矩形中心 EMU。
        hw, hh : int
            矩形の半幅・半高さ EMU。
        nx_v, ny_v : float
            方向単位ベクトル。

        Returns
        -------
        tuple[int, int]
            境界点 (x, y) EMU。
        """
        if abs(nx_v) < 1e-9 and abs(ny_v) < 1e-9:
            return cx, cy
        if abs(nx_v) < 1e-9:
            t = hh / abs(ny_v)
        elif abs(ny_v) < 1e-9:
            t = hw / abs(nx_v)
        else:
            t = min(hw / abs(nx_v), hh / abs(ny_v))
        return int(cx + nx_v * t), int(cy + ny_v * t)

    # ------------------------------------------------------------------
    # ラベル作成ユーティリティ
    # ------------------------------------------------------------------

    def _make_card_label(
        self,
        slide: Slide,
        text: str,
        center_x: int,
        center_y: int,
    ) -> object:
        """
        カーディナリティラベルのテキストボックスを作成する。

        Parameters
        ----------
        slide : Slide
        text : str
            表示テキスト（"1", "0..*" 等）。
        center_x, center_y : int
            テキストボックス中心 EMU 座標。

        Returns
        -------
        object
            追加したテキストボックス Shape。
        """
        bw = self._CARD_BOX_W
        bh = self._CARD_BOX_H
        txbox = slide.shapes.add_textbox(
            Emu(center_x - bw // 2),
            Emu(center_y - bh // 2),
            Emu(bw),
            Emu(bh),
        )
        tf = txbox.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(self._CARD_FONT_SIZE)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*self._REL_LINE_COLOR)
        txbox.fill.background()
        self._set_no_textbox_border(txbox)
        return txbox

    def _make_role_label(
        self,
        slide: Slide,
        text: str,
        center_x: int,
        center_y: int,
    ) -> object:
        """
        ロールラベルのテキストボックスを作成する（コネクター中点）。

        Parameters
        ----------
        slide : Slide
        text : str
            ロールテキスト（関係名）。
        center_x, center_y : int
            テキストボックス中心 EMU 座標。

        Returns
        -------
        object
            追加したテキストボックス Shape。
        """
        bw = 1_200_000
        bh = 240_000
        txbox = slide.shapes.add_textbox(
            Emu(center_x - bw // 2),
            Emu(center_y - bh // 2),
            Emu(bw),
            Emu(bh),
        )
        tf = txbox.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(self._ROLE_FONT_SIZE)
        run.font.italic = True
        run.font.color.rgb = RGBColor(80, 80, 80)
        txbox.fill.solid()
        txbox.fill.fore_color.rgb = RGBColor(245, 245, 255)
        self._set_no_textbox_border(txbox)
        return txbox

    # ------------------------------------------------------------------
    # XML ユーティリティ
    # ------------------------------------------------------------------

    def _set_no_border(self, shape: object) -> None:
        """Shape の枠線を noFill（非表示）に設定する。"""
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            return
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr, qn("a:ln"))
        if ln.find(qn("a:noFill")) is None:
            lxml_etree.SubElement(ln, qn("a:noFill"))

    def _set_no_textbox_border(self, txbox: object) -> None:
        """テキストボックスの枠線を noFill（非表示）に設定する。"""
        spPr_el = txbox._element.find(qn("p:spPr"))
        if spPr_el is None:
            return
        ln = spPr_el.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr_el, qn("a:ln"))
        if ln.find(qn("a:noFill")) is None:
            lxml_etree.SubElement(ln, qn("a:noFill"))

    def _set_connector_dash(self, connector: object) -> None:
        """
        コネクターを破線スタイルに設定する（NON_IDENTIFYING リレーション用）。

        python-pptx では connector.line.dash_style を直接設定できないため、
        lxml で p:spPr/a:ln/a:prstDash を注入する。
        """
        cxnSp_el = connector._element
        spPr_el = cxnSp_el.find(qn("p:spPr"))
        if spPr_el is None:
            return
        ln = spPr_el.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr_el, qn("a:ln"))
        # 既存の prstDash を除去してから再設定する
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        prstDash = lxml_etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")

    def _group_connector_with_labels(
        self,
        slide: Slide,
        connector: object,
        txboxes: list,
    ) -> None:
        """
        コネクターとラベルテキストボックス群を OOXML グループ化する。

        connector（<p:cxnSp>）と txboxes（<p:sp>）を一つの <p:grpSp> にまとめる。
        グループ内でも connector の begin/end connect（shape ID 参照）は保持される。

        Parameters
        ----------
        slide : Slide
        connector : object
            python-pptx のコネクター Shape。
        txboxes : list
            カーディナリティ・ロールラベル等のテキストボックス Shape リスト。
        """
        spTree = slide.shapes._spTree
        conn_el = connector._element
        txb_els = [tb._element for tb in txboxes]

        def _xywh_of(el: object) -> tuple[int, int, int, int]:
            """Shape / Connector の xfrm から (left, top, width, height) EMU を返す。"""
            # <p:sp> は p:spPr/<a:xfrm>、<p:cxnSp> も同様に p:spPr/<a:xfrm>
            for pr_tag in (qn("p:spPr"), qn("p:cxnSpPr")):
                pr_el = el.find(pr_tag)
                if pr_el is not None:
                    xfrm = pr_el.find(qn("a:xfrm"))
                    if xfrm is not None:
                        off = xfrm.find(qn("a:off"))
                        ext = xfrm.find(qn("a:ext"))
                        if off is not None and ext is not None:
                            return (
                                int(off.get("x", 0)),
                                int(off.get("y", 0)),
                                int(ext.get("cx", 0)),
                                int(ext.get("cy", 0)),
                            )
            return 0, 0, 0, 0

        # グループのバウンディングボックスを計算する
        all_els = [conn_el] + txb_els
        lefts, tops, rights, bottoms = [], [], [], []
        for el in all_els:
            l, t, w, h = _xywh_of(el)
            lefts.append(l)
            tops.append(t)
            rights.append(l + w)
            bottoms.append(t + h)

        grp_left = min(lefts)
        grp_top = min(tops)
        grp_w = max(1, max(rights) - grp_left)
        grp_h = max(1, max(bottoms) - grp_top)

        # <p:grpSp> を構築する（子座標系 = スライド座標系に設定して絶対座標をそのまま保持する）
        grpSp = lxml_etree.Element(qn("p:grpSp"))
        grpSpPr = lxml_etree.SubElement(grpSp, qn("p:grpSpPr"))
        xfrm = lxml_etree.SubElement(grpSpPr, qn("a:xfrm"))
        off_el = lxml_etree.SubElement(xfrm, qn("a:off"))
        off_el.set("x", str(grp_left))
        off_el.set("y", str(grp_top))
        ext_el = lxml_etree.SubElement(xfrm, qn("a:ext"))
        ext_el.set("cx", str(grp_w))
        ext_el.set("cy", str(grp_h))
        # 子座標系をスライド座標系と同一（恒等変換）にする
        chOff = lxml_etree.SubElement(xfrm, qn("a:chOff"))
        chOff.set("x", str(grp_left))
        chOff.set("y", str(grp_top))
        chExt = lxml_etree.SubElement(xfrm, qn("a:chExt"))
        chExt.set("cx", str(grp_w))
        chExt.set("cy", str(grp_h))

        # spTree から各要素を取り出してグループに移動する
        for el in [conn_el] + txb_els:
            try:
                spTree.remove(el)
            except ValueError:
                pass
            grpSp.append(el)

        # グループを spTree に追加する
        spTree.append(grpSp)
