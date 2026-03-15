"""
フローチャートレンダラーモジュール。

flowchart / graph 系のMermaid記法を解析し、14種類のノード形状・
7種類のエッジ矢印・4種類の線種・エッジラベルに対応してスライドに描画する。
描画方向（TD/TB/BT/LR/RL）を考慮したDAG階層レイアウトをサポートする。
"""

from __future__ import annotations

import re

import networkx as nx
from lxml import etree as lxml_etree
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer, NODE_WIDTH_EMU, NODE_HEIGHT_EMU


# ---------------------- ノード形状マップ ----------------------
# Mermaid頂点タイプ → MSO_AUTO_SHAPE_TYPE整数値
# python-pptx add_shape() に渡す値（MSO_AUTO_SHAPE_TYPE enum）
_SHAPE_MAP: dict[str, int] = {
    "square":        1,   # RECTANGLE（デフォルト矩形 []）
    "round":         5,   # ROUNDED_RECTANGLE（丸角矩形 ()）
    "stadium":       5,   # ROUNDED_RECTANGLE（スタジアム型 ([])）
    "subroutine":    1,   # RECTANGLE（サブルーチン [[]] ※二重枠は省略）
    "cylinder":      13,  # CAN（シリンダー [()]）
    "circle":        9,   # OVAL（円 (())）
    "odd":           51,  # PENTAGON（非対称 >]）
    "diamond":       4,   # DIAMOND（ひし形 {}）
    "hexagon":       10,  # HEXAGON（六角形 {{}}）
    "lean_right":    2,   # PARALLELOGRAM（右傾斜 [/]）
    "lean_left":     2,   # PARALLELOGRAM + flipH（左傾斜 [\]）
    "trapezoid":     3,   # TRAPEZOID（台形 [/\]）
    "inv_trapezoid": 3,   # TRAPEZOID + flipV（逆台形 [\/]）
    "doublecircle":  9,   # OVAL（二重円 ((())) ※二重枠は省略）
}

# ---------------------- エッジ矢印マップ ----------------------
# mermaid-parser-py の edge["type"] → headEnd/tailEnd 設定辞書
# headEnd: コネクター終点（dst側）の矢印種別
# tailEnd: コネクター始点（src側）の矢印種別（両方向エッジのみ）
_EDGE_ARROW_MAP: dict[str, dict[str, str]] = {
    "arrow_open":          {},                                      # 矢印なし（開放線）
    "arrow_point":         {"headEnd": "arrow"},                    # →（通常矢印）
    "double_arrow_point":  {"headEnd": "arrow", "tailEnd": "arrow"}, # ←→（両方向）
    "arrow_circle":        {"headEnd": "oval"},                     # 丸端矢印
    "double_arrow_circle": {"headEnd": "oval", "tailEnd": "oval"},  # 両方向丸端
    "arrow_cross":         {"headEnd": "arrow"},                    # ×端（arrow で代替）
    "double_arrow_cross":  {"headEnd": "arrow", "tailEnd": "arrow"},# 両方向×端
}


class FlowchartRenderer(BaseDiagramRenderer):
    """
    flowchart / graph 系のMermaid記法を描画するレンダラー。

    mermaid-parser-py が返す graph_data（vertices/edges）を受け取り、
    ノード形状・エッジスタイル・エッジラベルに対応してスライドに描画する。
    """

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        mermaid_text: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        フローチャートをスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            mermaid-parser-py が返す graph_data 辞書（vertices/edges を含む）。
        mermaid_text : str
            元のMermaidテキスト（フォールバック用）。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        vertices: dict = graph_data.get("vertices", {})
        raw_edges: list = graph_data.get("edges", [])

        if not vertices:
            self._render_fallback(slide, mermaid_text, left, top, width, height)
            return

        nodes = list(vertices.keys())

        # NetworkXグラフを構築してレイアウトを計算する
        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for edge in raw_edges:
            if isinstance(edge, dict):
                src = edge.get("start")
                dst = edge.get("end")
                if src is not None and dst is not None:
                    G.add_edge(str(src), str(dst))

        # mermaid_textから描画方向を取得して階層レイアウトを計算する
        # canvasサイズを渡して同一世代にノードが入りきらない時の自動折り返しを有効にする
        direction = self._extract_direction(mermaid_text)
        pos: dict[str, tuple[float, float]] = self._hierarchical_layout(
            G, direction, canvas_w=width, canvas_h=height
        )

        # 頂点をシェイプ種別に応じて描画する
        node_shapes = self._draw_nodes_flowchart(
            slide, vertices, pos, left, top, width, height
        )

        # エッジを矢印・線種・ラベルに応じて描画する
        self._draw_edges_flowchart(
            slide, raw_edges, pos, node_shapes, left, top, width, height
        )

    def _draw_nodes_flowchart(
        self,
        slide: Slide,
        vertices: dict,
        pos: dict[str, tuple[float, float]],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> dict[str, object]:
        """
        フローチャートのノードをシェイプ種別に応じて描画する。

        vertices の各エントリから "type" フィールドを取得し、
        _SHAPE_MAP に基づいてAutoShapeTypeを選択する。
        lean_left は水平フリップ、inv_trapezoid は垂直フリップを追加で適用する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        vertices : dict
            ノードID → {"text": str, "type": str, ...} の辞書。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        left, top, width, height : int
            描画エリアのEMU座標。

        Returns
        -------
        dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        """
        node_shapes: dict[str, object] = {}

        for node_id, vertex in vertices.items():
            if node_id not in pos:
                continue

            x_norm, y_norm = pos[node_id]
            cx, cy = self._pos_to_emu(x_norm, y_norm, left, top, width, height)
            shape_left = cx - NODE_WIDTH_EMU // 2
            shape_top = cy - NODE_HEIGHT_EMU // 2

            # ノード種別に対応するAutoShapeTypeを取得する（未知の場合は矩形）
            node_type = (
                vertex.get("type", "square")
                if isinstance(vertex, dict)
                else "square"
            )
            shape_type = _SHAPE_MAP.get(node_type, 1)

            shape = slide.shapes.add_shape(
                shape_type,
                Emu(shape_left),
                Emu(shape_top),
                Emu(NODE_WIDTH_EMU),
                Emu(NODE_HEIGHT_EMU),
            )

            # lean_left（左傾き平行四辺形）は水平フリップを適用する
            if node_type == "lean_left":
                self._apply_flip(shape, flip_h=True, flip_v=False)

            # inv_trapezoid（逆台形）は垂直フリップを適用する。
            # invertedTrapezoid はOOXML正規プリセット名に存在しないため、
            # lean_left と同様に xfrm の flipV="1" で実現する。
            if node_type == "inv_trapezoid":
                self._apply_flip(shape, flip_h=False, flip_v=True)

            # 表示ラベルを設定する（text フィールドを優先、なければノードID）
            label = (
                vertex.get("text", node_id)
                if isinstance(vertex, dict)
                else str(node_id)
            )
            shape.text = label
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)

            node_shapes[node_id] = shape

        return node_shapes

    def _apply_flip(self, shape: object, flip_h: bool, flip_v: bool) -> None:
        """
        shapeのDrawingML XML要素にフリップ属性を設定する。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト。
        flip_h : bool
            水平フリップを適用する場合True。
        flip_v : bool
            垂直フリップを適用する場合True。
        """
        sp_el = shape._element
        spPr = sp_el.find(qn("p:spPr"))
        if spPr is None:
            return
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            return
        if flip_h:
            xfrm.set("flipH", "1")
        if flip_v:
            xfrm.set("flipV", "1")

    def _draw_edges_flowchart(
        self,
        slide: Slide,
        raw_edges: list,
        pos: dict[str, tuple[float, float]],
        node_shapes: dict[str, object],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        フローチャートのエッジを矢印種別・線種・ラベルに応じて描画する。

        stroke="invisible" のエッジはコネクター自体を生成しない。
        stroke="dotted" は破線、stroke="thick" は3pt実線を適用する。
        edge_type で headEnd/tailEnd の矢印形状を決定する。
        text が設定されている場合はコネクターXMLにラベルテキストを直接設定する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        raw_edges : list
            mermaid-parser-py の edges リスト。各要素に start/end/stroke/type/text を持つ。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        node_shapes : dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        # ノードID → [ラベルtxBox] マップ（ループ後にノードとグループ化する）
        node_label_map: dict[str, list] = {}

        for edge in raw_edges:
            if not isinstance(edge, dict):
                continue

            src = str(edge.get("start", ""))
            dst = str(edge.get("end", ""))
            if not src or not dst:
                continue
            if src not in pos or dst not in pos:
                continue

            stroke = edge.get("stroke", "normal")
            edge_type = edge.get("type", "arrow_point")
            edge_label: str = edge.get("text", "") or ""

            # invisible エッジはコネクターを生成しない
            if stroke == "invisible":
                continue

            sx_norm, sy_norm = pos[src]
            dx_norm, dy_norm = pos[dst]
            sx, sy = self._pos_to_emu(sx_norm, sy_norm, left, top, width, height)
            dx, dy = self._pos_to_emu(dx_norm, dy_norm, left, top, width, height)

            # 曲線コネクターを追加する
            connector = slide.shapes.add_connector(
                3,  # MSO_CONNECTOR_TYPE.CURVE
                Emu(sx), Emu(sy), Emu(dx), Emu(dy),
            )

            # 接続ポイントを決定してshapeに接続する
            src_cp, dst_cp = self._connection_indices(
                dx_norm - sx_norm, dy_norm - sy_norm
            )
            src_shape = node_shapes.get(src)
            dst_shape = node_shapes.get(dst)
            if src_shape is not None:
                connector.begin_connect(src_shape, src_cp)
            if dst_shape is not None:
                connector.end_connect(dst_shape, dst_cp)

            # コネクター内の <a:ln> XML要素を取得または作成する
            cxn_el = connector._element
            spPr = cxn_el.find(qn("p:spPr"))
            if spPr is None:
                continue
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                ln = lxml_etree.SubElement(spPr, qn("a:ln"))

            # 線種を適用する
            if stroke == "dotted":
                # 破線スタイル
                prstDash = lxml_etree.SubElement(ln, qn("a:prstDash"))
                prstDash.set("val", "dash")
            elif stroke == "thick":
                # 太線（3pt = 38100 EMU）
                ln.set("w", "38100")

            # 矢印種別を適用する（未知の場合は通常矢印）
            arrow_conf = _EDGE_ARROW_MAP.get(edge_type, {"headEnd": "arrow"})
            if "headEnd" in arrow_conf:
                head = lxml_etree.SubElement(ln, qn("a:headEnd"))
                head.set("type", arrow_conf["headEnd"])
                head.set("w", "med")
                head.set("len", "med")
            if "tailEnd" in arrow_conf:
                tail = lxml_etree.SubElement(ln, qn("a:tailEnd"))
                tail.set("type", arrow_conf["tailEnd"])
                tail.set("w", "med")
                tail.set("len", "med")

            # ラベルがある場合は始点近くにテキストボックスを配置し、
            # ノードIDをキーにしてマップに登録する（ループ後にノードとグループ化する）
            if edge_label:
                txBox = self._add_edge_label_near_source(
                    slide, edge_label, sx, sy, dx, dy
                )
                if src not in node_label_map:
                    node_label_map[src] = []
                node_label_map[src].append(txBox)

        # ラベルテキストボックスを始点ノードとグループ化する。
        # コネクターはトップレベルに残すことで begin_connect/end_connect の接続が維持される。
        for node_id, txboxes in node_label_map.items():
            src_shape = node_shapes.get(node_id)
            if src_shape is not None:
                self._group_node_with_labels(slide, src_shape, txboxes)

    def _extract_direction(self, mermaid_text: str) -> str:
        """
        Mermaidテキストの宣言行から描画方向を抽出する。

        "flowchart LR" や "graph TD" のような宣言行から
        方向キーワードを取得する。見つからない場合は "TD"（上→下）を返す。

        Parameters
        ----------
        mermaid_text : str
            Mermaidフローチャートのテキスト。

        Returns
        -------
        str
            描画方向（"TD" / "TB" / "BT" / "LR" / "RL"）。デフォルトは "TD"。
        """
        match = re.search(
            r'(?:flowchart|graph)\s+(TD|TB|LR|RL|BT)\b',
            mermaid_text,
            re.IGNORECASE,
        )
        return match.group(1).upper() if match else "TD"

    def _hierarchical_layout(
        self,
        G: nx.DiGraph,
        direction: str,
        canvas_w: int | None = None,
        canvas_h: int | None = None,
    ) -> dict[str, tuple[float, float]]:
        """
        描画方向に基づいたDAG階層レイアウト座標を計算する。

        グラフがDAGである場合はトポロジカル世代でレベルを決定し、
        各レベル内のノードを均等配置した正規化座標（-1.0〜1.0）を返す。
        canvasに入りきらない場合は以下の2軸で自動折り返しを行う。
          ・同一世代のノード数超過 → スパン軸方向にサブレベルへ分割
          ・世代数（レベル数）超過 → レベル軸方向を複数レーンに折り返し
        サイクルを含むグラフはkamada_kawai → spring_layoutにフォールバックする。

        Parameters
        ----------
        G : nx.DiGraph
            レイアウト計算対象の有向グラフ。
        direction : str
            描画方向（"TD" / "TB" / "BT" / "LR" / "RL"）。
        canvas_w : int | None
            描画エリアの幅（EMU）。
        canvas_h : int | None
            描画エリアの高さ（EMU）。

        Returns
        -------
        dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標(-1.0〜1.0)のタプルを値とする辞書。
        """
        # サイクルを含む場合はkamada_kawai → spring_layoutにフォールバック
        if not nx.is_directed_acyclic_graph(G):
            try:
                return nx.kamada_kawai_layout(G)
            except Exception:
                return nx.spring_layout(G, seed=42, k=2.0)

        generations = list(nx.topological_generations(G))

        if direction in ("TD", "TB", "BT"):
            # レベル軸=Y（世代が縦に並ぶ）、スパン軸=X（同一世代が横に並ぶ）
            # スパン軸の上限: 1世代に横方向に入るノード数
            max_span = (
                max(1, canvas_w // NODE_WIDTH_EMU) if canvas_w is not None else None
            )
            # レベル軸の上限: ノード高さの1.5倍を1レベル分の占有高さとして計算する。
            # 実際の間隔がノード高さ約1.5倍になるよう余裕を持たせた世代数にする。
            max_levels = (
                max(1, canvas_h // int(NODE_HEIGHT_EMU * 1.5)) if canvas_h is not None else None
            )
        else:  # LR, RL
            # レベル軸=X（世代が横に並ぶ）、スパン軸=Y（同一世代が縦に並ぶ）
            # スパン軸の上限: 1世代に縦方向に入るノード数
            max_span = (
                max(1, canvas_h // NODE_HEIGHT_EMU) if canvas_h is not None else None
            )
            # レベル軸の上限: 横方向に入る世代数（1レーン分）
            max_levels = (
                max(1, canvas_w // NODE_WIDTH_EMU) if canvas_w is not None else None
            )

        # ① 同一世代のノード数がmax_spanを超える場合、サブレベルに分割する
        expanded_levels: list[list[str]] = []
        for level_nodes in generations:
            sorted_nodes = sorted(level_nodes)
            if max_span is None or len(sorted_nodes) <= max_span:
                expanded_levels.append(sorted_nodes)
            else:
                for i in range(0, len(sorted_nodes), max_span):
                    expanded_levels.append(sorted_nodes[i : i + max_span])

        # ② 世代数（expanded_levels数）がmax_levelsを超える場合、複数レーンに折り返す
        # 各レーンはmax_levels個の世代を持ち、レーン同士はスパン軸方向に積み上げる
        if max_levels is not None and len(expanded_levels) > max_levels:
            lanes: list[list[list[str]]] = []
            for i in range(0, len(expanded_levels), max_levels):
                lanes.append(expanded_levels[i : i + max_levels])
        else:
            lanes = [expanded_levels]

        n_lanes = len(lanes)
        pos: dict[str, tuple[float, float]] = {}

        # レーン間のギャップをノード幅/高さの半分に設定する（正規化座標）。
        # スパン軸（TD=X, LR=Y）の usable 幅から算出し、canvasがない場合は0とする。
        if n_lanes > 1 and canvas_w is not None and canvas_h is not None:
            if direction in ("TD", "TB", "BT"):
                # スパン軸=X: usable_span = canvas_w - NODE_WIDTH_EMU
                # ギャップ = NODE_WIDTH_EMU/2 を正規化座標に換算
                usable_span = max(1, canvas_w - NODE_WIDTH_EMU)
                gap_norm = NODE_WIDTH_EMU / usable_span
            else:  # LR, RL
                # スパン軸=Y: usable_span = canvas_h - NODE_HEIGHT_EMU
                # ギャップ = NODE_HEIGHT_EMU/2 を正規化座標に換算
                usable_span = max(1, canvas_h - NODE_HEIGHT_EMU)
                gap_norm = NODE_HEIGHT_EMU / usable_span
            # 総ギャップがスパン全体を超えないよう上限を設ける
            total_gap = (n_lanes - 1) * gap_norm
            if total_gap >= 1.8:  # 2.0 のうち 1.8 以上はギャップに使わない
                gap_norm = 1.8 / max(1, n_lanes - 1)
                total_gap = (n_lanes - 1) * gap_norm
            slot_size = (2.0 - total_gap) / n_lanes
        else:
            gap_norm = 0.0
            slot_size = 2.0 / n_lanes

        for lane_idx, lane_levels in enumerate(lanes):
            n_lvl = len(lane_levels)
            # レベル軸の範囲: 各レーンはレベル軸全体（-1.0〜+1.0）を使う。
            # 折り返し後もCR（レベル軸リセット）して幅いっぱいに広がる。
            level_axis_start = -1.0
            level_axis_end = 1.0
            # スパン軸の範囲: ギャップを考慮してレーンを均等配置する。
            span_range_start = -1.0 + lane_idx * (slot_size + gap_norm)
            span_range_end = span_range_start + slot_size

            for level_idx, level_nodes in enumerate(lane_levels):
                n_in_level = len(level_nodes)
                # レベル軸の正規化座標（各レーン共通で -1.0〜+1.0 全体を使う）
                if n_lvl > 1:
                    level_t = level_idx / (n_lvl - 1)  # 0.0〜1.0
                else:
                    level_t = 0.5
                level_norm = level_axis_start + (level_axis_end - level_axis_start) * level_t

                for node_idx, node_id in enumerate(level_nodes):
                    # スパン軸の正規化座標（このレーン内で均等配置）
                    if n_in_level > 1:
                        span_t = node_idx / (n_in_level - 1)  # 0.0〜1.0
                    else:
                        span_t = 0.5
                    span_norm = span_range_start + (span_range_end - span_range_start) * span_t

                    if direction in ("TD", "TB"):
                        # 上→下: level_normがY軸、span_normがX軸
                        pos[node_id] = (span_norm, level_norm)
                    elif direction == "BT":
                        # 下→上: Y軸を反転
                        pos[node_id] = (span_norm, -level_norm)
                    elif direction == "LR":
                        # 左→右: level_normがX軸、span_normがY軸
                        pos[node_id] = (level_norm, span_norm)
                    elif direction == "RL":
                        # 右→左: X軸を反転
                        pos[node_id] = (-level_norm, span_norm)
                    else:
                        # 未知の方向はTD（上→下）として扱う
                        pos[node_id] = (span_norm, level_norm)

        return pos
