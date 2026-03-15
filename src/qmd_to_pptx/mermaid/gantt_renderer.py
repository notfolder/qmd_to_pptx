"""
Mermaidガントチャートレンダラーモジュール。

GanttChart オブジェクトを入力として PPTXスライドにガントチャートを描画する。
レイアウト:
  - OOXMLテーブル: 1列目=タスク名, 2列目以降=時間軸グリッド
  - タスクバー: rightArrow図形 (MSO_AUTO_SHAPE_TYPE=33) をテーブル上にオーバーレイ
  - マイルストーン: diamond図形 (MSO_AUTO_SHAPE_TYPE=4) を配置

時間軸粒度の自動選択:
  - 全期間 < 30日: 日単位
  - 全期間 < 90日: 週単位（月曜始まり）
  - それ以上: 月単位
"""

from __future__ import annotations

import math
from datetime import date, timedelta
from typing import Optional

from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .gantt_parser import GanttChart, GanttTask

# ---------------------------------------------------------------------------
# 時間軸粒度定数
# ---------------------------------------------------------------------------
_GRANULARITY_DAY = "day"
_GRANULARITY_WEEK = "week"
_GRANULARITY_MONTH = "month"

# ---------------------------------------------------------------------------
# カラー定義 (RGB タプル)
# ---------------------------------------------------------------------------
_COLOR_DONE: tuple[int, int, int] = (170, 170, 170)       # 完了タスク: グレー
_COLOR_ACTIVE: tuple[int, int, int] = (70, 130, 180)      # 進行中タスク: スチールブルー
_COLOR_CRIT: tuple[int, int, int] = (220, 70, 70)         # クリティカルタスク: 赤
_COLOR_MILESTONE: tuple[int, int, int] = (220, 170, 0)    # マイルストーン: 黄
_COLOR_NORMAL: tuple[int, int, int] = (100, 180, 230)     # 通常タスク: 水色
_COLOR_SECTION_BG: tuple[int, int, int] = (60, 60, 110)   # セクション行背景: 濃紺
_COLOR_SECTION_FG: tuple[int, int, int] = (255, 255, 255) # セクション行文字: 白
_COLOR_HEADER_BG: tuple[int, int, int] = (40, 40, 80)     # ヘッダー行背景: 最濃紺
_COLOR_HEADER_FG: tuple[int, int, int] = (255, 255, 255)  # ヘッダー行文字: 白
_COLOR_TASK_ROW_BG: tuple[int, int, int] = (245, 245, 250) # タスク行背景: ごく薄い青白

# ---------------------------------------------------------------------------
# MSO_AUTO_SHAPE_TYPE 整数値
# ---------------------------------------------------------------------------
_SHAPE_RIGHT_ARROW: int = 33  # MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW
_SHAPE_DIAMOND: int = 4       # MSO_AUTO_SHAPE_TYPE.DIAMOND

# ---------------------------------------------------------------------------
# レイアウト定数
# ---------------------------------------------------------------------------
_TASK_COL_RATIO: float = 0.22    # タスク名列の幅比率（全体の22%）
_MIN_BAR_WIDTH_EMU: int = 60_000  # タスクバーの最小幅 (EMU)


class GanttRenderer:
    """
    Mermaidガントチャートレンダラークラス。

    GanttChart オブジェクトを入力として PPTXスライドに
    テーブルと図形オーバーレイでガントチャートを描画する。
    """

    def render(
        self,
        slide: Slide,
        chart: GanttChart,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        ガントチャートをスライドに描画する。

        Parameters
        ----------
        slide : Slide
            描画対象のスライド。
        chart : GanttChart
            パース済みのガントチャートデータ。
        left : int
            描画エリアの左端座標 (EMU)。
        top : int
            描画エリアの上端座標 (EMU)。
        width : int
            描画エリアの幅 (EMU)。
        height : int
            描画エリアの高さ (EMU)。
        """
        all_tasks = chart.all_tasks()
        if not all_tasks:
            return

        # 時間軸を計算する
        axis_start, axis_end, granularity, col_dates = self._compute_time_axis(all_tasks)
        total_days = max(1, (axis_end - axis_start).days)

        # テーブルの行構造を組み立てる
        # rows_info: (row_type, section_name_or_None, task_or_None)
        rows_info: list[tuple[str, Optional[str], Optional[GanttTask]]] = []
        rows_info.append(("header", None, None))
        for section in chart.sections:
            if section.name:
                rows_info.append(("section", section.name, None))
            for task in section.tasks:
                rows_info.append(("task", section.name, task))

        n_rows = len(rows_info)
        n_axis_cols = len(col_dates)
        n_cols = 1 + n_axis_cols  # タスク名列 + 時間軸列

        # 列幅を計算する: タスク名列 = 全体の _TASK_COL_RATIO, 時間軸列 = 残りを等分
        task_col_w = int(width * _TASK_COL_RATIO)
        axis_total_w = width - task_col_w
        axis_col_w = axis_total_w // n_axis_cols if n_axis_cols > 0 else axis_total_w

        # 行高を計算する（均等割り）
        row_height = height // n_rows if n_rows > 0 else height

        # テーブルをスライドに追加する
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        # 列幅を設定する
        table.columns[0].width = task_col_w
        for i in range(n_axis_cols):
            table.columns[1 + i].width = axis_col_w

        # 行高を設定する
        for i in range(n_rows):
            table.rows[i].height = row_height

        # ヘッダー行を記入する
        self._fill_header_row(table, col_dates, granularity)

        # セクション行・タスク行を記入する
        for row_idx, (row_type, section_name, task) in enumerate(rows_info):
            if row_type == "header":
                continue
            if row_type == "section" and section_name is not None:
                self._fill_section_row(table, row_idx, n_cols, section_name)
            elif row_type == "task" and task is not None:
                self._fill_task_name_cell(table, row_idx, task.title)

        # タスクバーをテーブル上にオーバーレイで描画する
        axis_emu_left = left + task_col_w
        for row_idx, (row_type, _section_name, task) in enumerate(rows_info):
            if row_type != "task" or task is None:
                continue

            # この行のtop位置を計算する
            bar_top = top + row_idx * row_height
            bar_height = row_height

            # タスクバーの横位置を開始/終了日の比率から計算する
            start_days = (task.start_date - axis_start).days
            end_days = (task.end_date - axis_start).days
            start_ratio = start_days / total_days
            end_ratio = end_days / total_days

            # EMU座標に変換する（描画エリア外にはみ出さないようクリップ）
            bar_emu_left = axis_emu_left + int(max(0.0, start_ratio) * axis_total_w)
            bar_emu_right = axis_emu_left + int(min(1.0, end_ratio) * axis_total_w)
            bar_emu_width = max(_MIN_BAR_WIDTH_EMU, bar_emu_right - bar_emu_left)

            # タスクの色を決定する
            color = self._get_task_color(task)

            # マイルストーン: ひし形図形を行中央に配置する
            if task.is_milestone:
                diamond_size = int(bar_height * 0.75)
                diamond_left = bar_emu_left - diamond_size // 2
                diamond_top = bar_top + (bar_height - diamond_size) // 2
                self._draw_shape(
                    slide, _SHAPE_DIAMOND,
                    diamond_left, diamond_top,
                    diamond_size, diamond_size,
                    color,
                )
            else:
                # 通常タスク: 行高の70%の高さで右矢印図形を配置する
                inner_top = bar_top + int(bar_height * 0.15)
                inner_height = int(bar_height * 0.70)
                self._draw_shape(
                    slide, _SHAPE_RIGHT_ARROW,
                    bar_emu_left, inner_top,
                    bar_emu_width, inner_height,
                    color,
                )

    # ------------------------------------------------------------------
    # 時間軸計算
    # ------------------------------------------------------------------

    def _compute_time_axis(
        self,
        tasks: list[GanttTask],
    ) -> tuple[date, date, str, list[date]]:
        """
        全タスクの開始/終了日から時間軸列の日付リストを計算する。

        Returns
        -------
        tuple[date, date, str, list[date]]
            (axis_start, axis_end, granularity, col_start_dates)
        """
        raw_start = min(t.start_date for t in tasks)
        raw_end = max(t.end_date for t in tasks)
        total_days = (raw_end - raw_start).days

        # 全期間に応じて粒度を選択する
        if total_days < 30:
            granularity = _GRANULARITY_DAY
        elif total_days < 90:
            granularity = _GRANULARITY_WEEK
        else:
            granularity = _GRANULARITY_MONTH

        col_dates = self._build_col_dates(raw_start, raw_end, granularity)

        # 軸の開始/終了を列日付の範囲に合わせる
        axis_start = col_dates[0] if col_dates else raw_start
        # 軸終端: 最終列の次の期間まで伸ばして全タスクが収まるようにする
        axis_end = max(raw_end, axis_start)

        return axis_start, axis_end, granularity, col_dates

    def _build_col_dates(
        self,
        raw_start: date,
        raw_end: date,
        granularity: str,
    ) -> list[date]:
        """
        粒度に応じた列日付リストを生成する（最大12〜24列）。

        Parameters
        ----------
        raw_start : date
            全タスクの最早開始日。
        raw_end : date
            全タスクの最遅終了日。
        granularity : str
            粒度定数 (_GRANULARITY_DAY / _GRANULARITY_WEEK / _GRANULARITY_MONTH)。

        Returns
        -------
        list[date]
            時間軸の各列の開始日リスト。
        """
        if granularity == _GRANULARITY_DAY:
            # 日単位: 日ごとに1列（最大12列）
            col_start = raw_start
            col_end = raw_end + timedelta(days=1)
            all_dates: list[date] = []
            d = col_start
            while d < col_end:
                all_dates.append(d)
                d += timedelta(days=1)
            max_cols = 12
            if len(all_dates) > max_cols:
                step = math.ceil(len(all_dates) / max_cols)
                all_dates = all_dates[::step]
            return all_dates

        if granularity == _GRANULARITY_WEEK:
            # 週単位: 月曜始まりの週（最大12列）
            col_start = raw_start - timedelta(days=raw_start.weekday())
            col_end = raw_end + timedelta(days=7)
            all_dates = []
            d = col_start
            while d < col_end:
                all_dates.append(d)
                d += timedelta(weeks=1)
            return all_dates[:12]

        # 月単位: 月初（最大24列）
        col_start = date(raw_start.year, raw_start.month, 1)
        col_end_year = raw_end.year + (1 if raw_end.month == 12 else 0)
        col_end_month = 1 if raw_end.month == 12 else raw_end.month + 1
        col_end = date(col_end_year, col_end_month, 1)
        all_dates = []
        d = col_start
        while d < col_end:
            all_dates.append(d)
            # 翌月に進む
            m = d.month + 1
            y = d.year
            if m > 12:
                m = 1
                y += 1
            d = date(y, m, 1)
        return all_dates[:24]

    # ------------------------------------------------------------------
    # テーブル記入ヘルパー
    # ------------------------------------------------------------------

    def _fill_header_row(
        self,
        table: object,
        col_dates: list[date],
        granularity: str,
    ) -> None:
        """ヘッダー行にタスク列ラベルと時間軸ラベルを記入する。"""
        # 列0: "タスク" ラベル
        cell0 = table.cell(0, 0)
        self._set_cell_text(cell0, "タスク", bold=True, align=PP_ALIGN.CENTER, font_size=9)
        self._set_cell_bg(cell0, _COLOR_HEADER_BG)
        self._set_cell_text_color(cell0, _COLOR_HEADER_FG)

        # 時間軸列: 粒度に応じたラベル
        for i, d in enumerate(col_dates):
            cell = table.cell(0, 1 + i)
            label = self._format_col_label(d, granularity)
            self._set_cell_text(cell, label, bold=True, align=PP_ALIGN.CENTER, font_size=7)
            self._set_cell_bg(cell, _COLOR_HEADER_BG)
            self._set_cell_text_color(cell, _COLOR_HEADER_FG)

    def _fill_section_row(
        self,
        table: object,
        row_idx: int,
        n_cols: int,
        section_name: str,
    ) -> None:
        """セクション行を記入する（全列結合してセクション名を配置）。"""
        # 全列を結合する
        table.cell(row_idx, 0).merge(table.cell(row_idx, n_cols - 1))
        cell = table.cell(row_idx, 0)
        self._set_cell_text(cell, f"  {section_name}", bold=True, align=PP_ALIGN.LEFT, font_size=10)
        self._set_cell_bg(cell, _COLOR_SECTION_BG)
        self._set_cell_text_color(cell, _COLOR_SECTION_FG)

    def _fill_task_name_cell(
        self,
        table: object,
        row_idx: int,
        task_name: str,
    ) -> None:
        """タスク行の列0（タスク名セル）を記入する。"""
        cell = table.cell(row_idx, 0)
        self._set_cell_text(cell, f" {task_name}", bold=False, align=PP_ALIGN.LEFT, font_size=9)
        self._set_cell_bg(cell, _COLOR_TASK_ROW_BG)

    # ------------------------------------------------------------------
    # 色・図形ヘルパー
    # ------------------------------------------------------------------

    @staticmethod
    def _get_task_color(task: GanttTask) -> tuple[int, int, int]:
        """タスクの状態から塗りつぶし色を決定する。"""
        if task.is_done:
            return _COLOR_DONE
        if task.is_crit:
            return _COLOR_CRIT
        if task.is_active:
            return _COLOR_ACTIVE
        if task.is_milestone:
            return _COLOR_MILESTONE
        return _COLOR_NORMAL

    @staticmethod
    def _draw_shape(
        slide: Slide,
        shape_type: int,
        left: int,
        top: int,
        width: int,
        height: int,
        rgb: tuple[int, int, int],
    ) -> None:
        """図形をスライドに追加し、塗りつぶし色を設定して枠線を除去する。"""
        if width <= 0 or height <= 0:
            return
        shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        shape.line.fill.background()  # 枠線なし
        # テキストフレームを空にする
        if shape.has_text_frame:
            shape.text_frame.text = ""

    @staticmethod
    def _format_col_label(d: date, granularity: str) -> str:
        """粒度に応じた列ヘッダーラベル文字列を生成する。"""
        if granularity == _GRANULARITY_DAY:
            return str(d.day)       # 日付の日のみ（例: "15"）
        if granularity == _GRANULARITY_WEEK:
            return d.strftime("%m/%d")  # 例: "01/15"
        # MONTH
        return d.strftime("%m月")   # 例: "01月"

    # ------------------------------------------------------------------
    # OOXMLセル操作ヘルパー
    # ------------------------------------------------------------------

    @staticmethod
    def _set_cell_bg(cell: object, rgb: tuple[int, int, int]) -> None:
        """セルの背景色をOOXML直接操作で設定する。"""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        # 既存のsolidFillを削除する
        for elem in tcPr.findall(qn("a:solidFill")):
            tcPr.remove(elem)
        sf = lxml_etree.SubElement(tcPr, qn("a:solidFill"))
        srgb = lxml_etree.SubElement(sf, qn("a:srgbClr"))
        srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

    @staticmethod
    def _set_cell_text(
        cell: object,
        text: str,
        bold: bool = False,
        align: int = PP_ALIGN.LEFT,
        font_size: int = 9,
    ) -> None:
        """セルのテキスト・フォントサイズ・太字・整列を設定する。"""
        tf = cell.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.text = text
        para.alignment = align
        if para.runs:
            run = para.runs[0]
            run.font.size = Pt(font_size)
            run.font.bold = bold

    @staticmethod
    def _set_cell_text_color(cell: object, rgb: tuple[int, int, int]) -> None:
        """セルのテキスト文字色を設定する。"""
        tf = cell.text_frame
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].font.color.rgb = RGBColor(*rgb)
