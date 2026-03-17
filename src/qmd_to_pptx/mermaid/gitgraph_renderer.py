"""
Mermaid Gitグラフ（gitGraph）レンダラーモジュール。

GitGraph データクラスを入力として PowerPoint スライドに Git グラフを描画する。

レイアウト方針:
  LR（デフォルト）: ブランチを水平レーン（上から下）に積み重ね、
                     コミットを時系列順に左から右へ配置する。
  TB:              ブランチを垂直レーン（左から右）に並べ、
                     コミットを時系列順に上から下へ配置する。
  BT:              TB と同じレーン配置だが、コミットを下から上へ配置する。

ノード描画仕様（OOXML 制約への代替案含む）:
  NORMAL       → 塗りつぶし楕円（ブランチ色）
  REVERSE      → 楕円（ブランチ色）+ 白テキスト "✕" オーバーレイ
  HIGHLIGHT    → 塗りつぶし矩形（ブランチ色）
  MERGE        → 楕円の上に白い小楕円を重ねた二重円
  CHERRY_PICK  → 明るい紫色の楕円

タグ描画:
  コミット円の上部に小矩形テキストボックス（薄黄背景・細枠線）を配置する。

ブランチレーン:
  ブランチの最初のコミット〜最後のコミット間を水平/垂直の直線で結ぶ。
  ブランチ名は始端に角丸矩形ラベルを配置する。

マージ接続線:
  マージコミットと親コミット（マージ元）を結ぶ直線コネクターを描画する。

コミットIDラベル:
  コミット円の下（LR）または右（TB/BT）にテキストボックスを配置する。

ブランチカラーパレット（最大8色、循環使用）:
  index 0: #E8A838（オレンジ/main 系）
  index 1: #58D68D（緑/develop 系）
  index 2: #EC7063（赤/hotfix 系）
  index 3: #5DADE2（青/feature 系）
  index 4: #AF7AC5（紫/release 系）
  index 5: #F0B27A（薄橙）
  index 6: #45B39D（シアン）
  index 7: #DC7633（茶）
"""

from __future__ import annotations

import logging

from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .gitgraph_parser import GitBranch, GitCommit, GitGraph

# モジュールロガーを取得する
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# ブランチカラーパレット（RGB タプル）
# ---------------------------------------------------------------------------
_BRANCH_COLORS: list[tuple[int, int, int]] = [
    (232, 168,  56),   # 0: オレンジ（main）
    ( 88, 214, 141),   # 1: 緑（develop）
    (236, 112,  99),   # 2: 赤（hotfix/release）
    ( 93, 173, 226),   # 3: 青（feature）
    (175, 122, 197),   # 4: 紫
    (240, 178, 122),   # 5: 薄橙
    ( 69, 179, 157),   # 6: シアン
    (220, 118,  51),   # 7: 茶
]

# CHERRY_PICK コミット色（紫）
_CHERRY_COLOR: tuple[int, int, int] = (167, 105, 200)

# タグ背景色
_TAG_BG_COLOR: tuple[int, int, int] = (255, 253, 206)

# タグ枠線色
_TAG_BORDER_COLOR: tuple[int, int, int] = (150, 130, 60)

# ブランチラベル文字色
_LABEL_TEXT_COLOR: tuple[int, int, int] = (255, 255, 255)

# コミットIDラベル色（灰色）
_COMMIT_LABEL_COLOR: tuple[int, int, int] = (80, 80, 80)

# レーン線色（灰色調）
_LANE_LINE_ALPHA = (160, 160, 160)

# ---------------------------------------------------------------------------
# サイズ定数（EMU）
# ---------------------------------------------------------------------------
_COMMIT_R: int = 180_000          # コミット円の半径（直径 = 2 * R）
_COMMIT_D: int = _COMMIT_R * 2   # コミット円の直径

_INNER_R: int = 105_000           # MERGE 二重円の内側円半径
_INNER_D: int = _INNER_R * 2

_TAG_H: int = 200_000             # タグ矩形の高さ
_TAG_W: int = 500_000             # タグ矩形の幅（最大、テキストに合わせる）
_TAG_GAP: int = 50_000            # タグとコミット円の間隔

_LABEL_H: int = 200_000           # ブランチラベル矩形の高さ
_LABEL_W: int = 700_000           # ブランチラベル矩形の幅

_COMMIT_LABEL_H: int = 180_000    # コミットIDラベルの高さ
_COMMIT_LABEL_W: int = 800_000    # コミットIDラベルの幅

# ---------------------------------------------------------------------------
# 間隔定数（EMU）
# ---------------------------------------------------------------------------
_LANE_GAP: int = 750_000          # ブランチレーン間の距離（中央線間）
_COMMIT_GAP: int = 800_000        # コミット間の距離（中央点間）
_MARGIN: int = 500_000            # 端余白


def _branch_color(idx: int) -> tuple[int, int, int]:
    """ブランチインデックスに対応するブランチ色を循環して返す。"""
    return _BRANCH_COLORS[idx % len(_BRANCH_COLORS)]


def _rgb(r: int, g: int, b: int) -> RGBColor:
    """(r, g, b) タプルから RGBColor を生成するユーティリティ。"""
    return RGBColor(r, g, b)


def _darker(rgb: tuple[int, int, int], factor: float = 0.7) -> tuple[int, int, int]:
    """色を暗くする（枠線色の計算に使用）。"""
    return (int(rgb[0] * factor), int(rgb[1] * factor), int(rgb[2] * factor))


def _set_no_fill(shape: object) -> None:
    """シェープの塗りつぶしを「なし」に設定する。"""
    shape.fill.background()  # type: ignore[attr-defined]


def _set_solid_fill(shape: object, rgb: tuple[int, int, int]) -> None:
    """シェープを単色塗りつぶしに設定する。"""
    shape.fill.solid()  # type: ignore[attr-defined]
    shape.fill.fore_color.rgb = _rgb(*rgb)  # type: ignore[attr-defined]


def _set_line_color(shape: object, rgb: tuple[int, int, int], pt: float = 1.0) -> None:
    """シェープの枠線色と幅を設定する。"""
    shape.line.color.rgb = _rgb(*rgb)  # type: ignore[attr-defined]
    shape.line.width = Pt(pt)


def _no_line(shape: object) -> None:
    """シェープの枠線を非表示にする。"""
    shape.line.fill.background()  # type: ignore[attr-defined]


def _add_ellipse(
    slide: Slide,
    cx: int,
    cy: int,
    w: int,
    h: int,
) -> object:
    """
    中心座標 (cx, cy) に幅 w 高さ h の楕円を追加する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    cx, cy : int
        楕円中心のEMU座標。
    w, h : int
        楕円の幅・高さ（EMU）。

    Returns
    -------
    Shape
        追加した楕円シェープ。
    """
    from pptx.util import Emu as _Emu  # ローカルインポートで循環回避
    shape = slide.shapes.add_shape(
        9,  # MSO_AUTO_SHAPE_TYPE.OVAL
        _Emu(cx - w // 2),
        _Emu(cy - h // 2),
        _Emu(w),
        _Emu(h),
    )
    return shape


def _add_rectangle(
    slide: Slide,
    left: int,
    top: int,
    w: int,
    h: int,
    rounded: bool = False,
) -> object:
    """
    矩形（または角丸矩形）を追加する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    left, top : int
        矩形左上のEMU座標。
    w, h : int
        幅・高さ（EMU）。
    rounded : bool
        True の場合は角丸矩形を追加する。

    Returns
    -------
    Shape
        追加したシェープ。
    """
    shape_type = 5 if rounded else 1  # 5=ROUNDED_RECTANGLE, 1=RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Emu(left),
        Emu(top),
        Emu(w),
        Emu(h),
    )
    return shape


def _set_text(
    shape: object,
    text: str,
    font_size_pt: float = 8.0,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    align_center: bool = True,
) -> None:
    """
    シェープのテキストフレームにテキストを設定する。

    Parameters
    ----------
    shape : Shape
        テキストフレームを持つシェープ。
    text : str
        設定するテキスト文字列。
    font_size_pt : float
        フォントサイズ（ポイント）。
    bold : bool
        太字かどうか。
    color : tuple[int, int, int]
        テキスト色 (R, G, B)。
    align_center : bool
        True の場合は中央揃え。
    """
    from pptx.enum.text import PP_ALIGN
    tf = shape.text_frame  # type: ignore[attr-defined]
    tf.word_wrap = False
    para = tf.paragraphs[0]
    if align_center:
        para.alignment = PP_ALIGN.CENTER
    else:
        para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(*color)


def _add_connector(
    slide: Slide,
    x1: int, y1: int,
    x2: int, y2: int,
    rgb: tuple[int, int, int] = (150, 150, 150),
    pt: float = 1.5,
    dashed: bool = False,
) -> object:
    """
    2点間を結ぶ直線コネクターを追加する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    x1, y1 : int
        始点のEMU座標。
    x2, y2 : int
        終点のEMU座標。
    rgb : tuple[int, int, int]
        線の色。
    pt : float
        線の太さ（ポイント）。
    dashed : bool
        True の場合は破線にする（OOXML prstDash="dash" 設定）。

    Returns
    -------
    Connector
        追加したコネクターシェープ。
    """
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(x1), Emu(y1),
        Emu(x2), Emu(y2),
    )
    connector.line.color.rgb = _rgb(*rgb)
    connector.line.width = Pt(pt)
    if dashed:
        # OOXML 直接操作で破線設定
        ln_elem = connector.line._ln  # type: ignore[attr-defined]
        if ln_elem is not None:
            prstDash = lxml_etree.SubElement(ln_elem, qn("a:prstDash"))
            prstDash.set("val", "dash")
    return connector


# ---------------------------------------------------------------------------
# メインレンダラークラス
# ---------------------------------------------------------------------------

class GitGraphRenderer:
    """
    GitGraph データクラスを受け取り、PowerPoint スライドに Git グラフを描画するクラス。

    対応する描画機能:
    - LR / TB / BT の 3 方向レイアウト
    - NORMAL / REVERSE / HIGHLIGHT / MERGE / CHERRY_PICK の 5 コミットタイプ
    - タグ矩形ラベル
    - ブランチレーン線とブランチ名ラベル
    - マージコミット、cherry-pick コミットの接続線
    - コミット ID ラベル（コミット円の下または右）
    """

    def render(
        self,
        slide: Slide,
        graph: GitGraph,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        GitGraph をスライドに描画する。

        コミットが 0 件の場合は何も描画しない。

        Parameters
        ----------
        slide : Slide
            描画先の python-pptx Slide オブジェクト。
        graph : GitGraph
            parse_gitgraph() で生成した GitGraph データクラス。
        left, top, width, height : int
            描画エリアのEMU座標・サイズ。
        """
        if not graph.commits:
            return

        direction = graph.direction.upper()

        # ブランチの表示順序を決定する（order 指定 → 定義順）
        ordered_branches = _sort_branches(graph.branches)

        # ブランチインデックスマップ（名前 → 表示レーンインデックス）
        branch_lane: dict[str, int] = {b.name: i for i, b in enumerate(ordered_branches)}
        n_lanes = len(ordered_branches)

        # コミット列インデックス（コミット発生順の通し番号）
        commit_col: dict[str, int] = {c.commit_id: i for i, c in enumerate(graph.commits)}
        n_cols = len(graph.commits)

        # 各ブランチ色の割り当て
        branch_color_map: dict[str, tuple[int, int, int]] = {
            b.name: _branch_color(i) for i, b in enumerate(ordered_branches)
        }

        if direction == "LR":
            self._render_lr(
                slide, graph, left, top, width, height,
                branch_lane, branch_color_map, commit_col, n_lanes, n_cols,
            )
        elif direction in ("TB", "BT"):
            self._render_tb(
                slide, graph, left, top, width, height,
                branch_lane, branch_color_map, commit_col, n_lanes, n_cols,
                bottom_to_top=(direction == "BT"),
            )
        else:
            # 不明な方向は LR にフォールバックする
            logger.warning(
                "gitグラフで不明な方向 %r が指定されました。LR にフォールバックします。",
                direction,
            )
            self._render_lr(
                slide, graph, left, top, width, height,
                branch_lane, branch_color_map, commit_col, n_lanes, n_cols,
            )

    # ------------------------------------------------------------------
    # LR レイアウト
    # ------------------------------------------------------------------

    def _render_lr(
        self,
        slide: Slide,
        graph: GitGraph,
        left: int,
        top: int,
        width: int,
        height: int,
        branch_lane: dict[str, int],
        branch_color_map: dict[str, tuple[int, int, int]],
        commit_col: dict[str, int],
        n_lanes: int,
        n_cols: int,
    ) -> None:
        """
        LR（左→右）方向で Git グラフを描画する。

        ブランチは水平レーン（上から下）として配置する。
        コミットは時系列順に左から右へ配置する。
        """
        # --- 座標計算 ---
        # ブランチレーン Y 座標（レーン中央線）
        # ラベル幅分の余白を確保する
        label_area_w = _LABEL_W + _MARGIN // 2
        usable_w = width - label_area_w - _MARGIN
        usable_h = height - _MARGIN * 2

        lane_y = _calc_positions_1d(n_lanes, usable_h, _LANE_GAP, top + _MARGIN)
        commit_x = _calc_positions_1d(n_cols, usable_w, _COMMIT_GAP, left + label_area_w)

        # コミットID → (cx, cy) のマップ
        commit_pos: dict[str, tuple[int, int]] = {}
        for commit in graph.commits:
            col = commit_col[commit.commit_id]
            lane = branch_lane.get(commit.branch, 0)
            cx = commit_x[col]
            cy = lane_y[lane]
            commit_pos[commit.commit_id] = (cx, cy)

        # --- ブランチレーン線の描画 ---
        for branch in graph.branches:
            lane = branch_lane.get(branch.name, 0)
            cy = lane_y[lane]
            color = branch_color_map.get(branch.name, _BRANCH_COLORS[0])

            # このブランチに属するコミットのX座標範囲
            branch_commits = [c for c in graph.commits if c.branch == branch.name]
            if not branch_commits:
                # コミットがないブランチは先祖コミットとの接続用に点だけ描く
                continue

            # ブランチ線の開始X: ブランチ作成時の親コミット位置を取得する
            first_commit = branch_commits[0]
            if first_commit.parents:
                parent_pos = commit_pos.get(first_commit.parents[0])
                x_start = parent_pos[0] if parent_pos else commit_x[commit_col[first_commit.commit_id]]
            else:
                x_start = commit_x[commit_col[first_commit.commit_id]]

            last_commit = branch_commits[-1]
            x_end = commit_x[commit_col[last_commit.commit_id]]

            # ブランチ線（コミット円の上を通る太めの線）
            _add_connector(
                slide,
                x_start, cy,
                x_end, cy,
                rgb=color,
                pt=3.0,
            )

        # --- ブランチラベルの描画 ---
        for branch in graph.branches:
            lane = branch_lane.get(branch.name, 0)
            cy = lane_y[lane]
            color = branch_color_map.get(branch.name, _BRANCH_COLORS[0])

            lbl = _add_rectangle(
                slide,
                left + _MARGIN // 2,
                cy - _LABEL_H // 2,
                _LABEL_W,
                _LABEL_H,
                rounded=True,
            )
            _set_solid_fill(lbl, color)
            _set_line_color(lbl, _darker(color), pt=0.75)
            _set_text(lbl, branch.name, font_size_pt=8.0, bold=True,
                      color=_LABEL_TEXT_COLOR, align_center=True)

        # --- マージ接続線の描画（コミット円より先に描いてコミットで上書きする）---
        for commit in graph.commits:
            if commit.commit_type in ("MERGE", "CHERRY_PICK") and len(commit.parents) >= 2:
                cx, cy = commit_pos[commit.commit_id]
                # 2番目の親（マージ元）への接続線
                src_id = commit.parents[1]
                if src_id in commit_pos:
                    sx, sy = commit_pos[src_id]
                    src_color = branch_color_map.get(
                        next((c.branch for c in graph.commits if c.commit_id == src_id), ""),
                        _BRANCH_COLORS[0],
                    )
                    _add_connector(slide, sx, sy, cx, cy, rgb=src_color, pt=2.0, dashed=True)

        # --- コミット円の描画 ---
        for commit in graph.commits:
            cx, cy = commit_pos[commit.commit_id]
            color = branch_color_map.get(commit.branch, _BRANCH_COLORS[0])
            self._draw_commit(slide, commit, cx, cy, color)

        # --- タグの描画 ---
        for commit in graph.commits:
            if commit.tag:
                cx, cy = commit_pos[commit.commit_id]
                self._draw_tag(slide, commit.tag, cx, cy - _COMMIT_R - _TAG_GAP, direction="LR")

        # --- コミットIDラベルの描画 ---
        for commit in graph.commits:
            cx, cy = commit_pos[commit.commit_id]
            label = commit.commit_id
            # ラベルテキストボックス（コミット円下部）
            tb = slide.shapes.add_textbox(
                Emu(cx - _COMMIT_LABEL_W // 2),
                Emu(cy + _COMMIT_R + 30_000),
                Emu(_COMMIT_LABEL_W),
                Emu(_COMMIT_LABEL_H),
            )
            tf = tb.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            from pptx.enum.text import PP_ALIGN
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = label
            run.font.size = Pt(6.0)
            run.font.color.rgb = _rgb(*_COMMIT_LABEL_COLOR)

    # ------------------------------------------------------------------
    # TB / BT レイアウト
    # ------------------------------------------------------------------

    def _render_tb(
        self,
        slide: Slide,
        graph: GitGraph,
        left: int,
        top: int,
        width: int,
        height: int,
        branch_lane: dict[str, int],
        branch_color_map: dict[str, tuple[int, int, int]],
        commit_col: dict[str, int],
        n_lanes: int,
        n_cols: int,
        bottom_to_top: bool = False,
    ) -> None:
        """
        TB（上→下）または BT（下→上）方向で Git グラフを描画する。

        ブランチは垂直レーン（左から右）として配置する。
        コミットは時系列順に上→下（TB）または下→上（BT）へ配置する。
        """
        # ラベルエリアをレーン上部に取る
        label_area_h = _LABEL_H + _MARGIN // 2
        usable_w = width - _MARGIN * 2
        usable_h = height - label_area_h - _MARGIN

        lane_x = _calc_positions_1d(n_lanes, usable_w, _LANE_GAP, left + _MARGIN)

        # BT の場合はコミット位置を反転する
        if bottom_to_top:
            commit_y_base = top + label_area_h + usable_h
            step = -_COMMIT_GAP
        else:
            commit_y_base = top + label_area_h
            step = _COMMIT_GAP

        commit_y_list = [commit_y_base + i * step for i in range(n_cols)]

        # コミットID → (cx, cy) のマップ
        commit_pos: dict[str, tuple[int, int]] = {}
        for commit in graph.commits:
            col = commit_col[commit.commit_id]
            lane = branch_lane.get(commit.branch, 0)
            cx = lane_x[lane]
            cy = commit_y_list[col]
            commit_pos[commit.commit_id] = (cx, cy)

        # --- ブランチレーン線の描画 ---
        for branch in graph.branches:
            lane = branch_lane.get(branch.name, 0)
            cx = lane_x[lane]
            color = branch_color_map.get(branch.name, _BRANCH_COLORS[0])

            branch_commits = [c for c in graph.commits if c.branch == branch.name]
            if not branch_commits:
                continue

            first_commit = branch_commits[0]
            if first_commit.parents:
                parent_pos = commit_pos.get(first_commit.parents[0])
                y_start = parent_pos[1] if parent_pos else commit_y_list[commit_col[first_commit.commit_id]]
            else:
                y_start = commit_y_list[commit_col[first_commit.commit_id]]

            last_commit = branch_commits[-1]
            y_end = commit_y_list[commit_col[last_commit.commit_id]]

            _add_connector(
                slide,
                cx, y_start,
                cx, y_end,
                rgb=color,
                pt=3.0,
            )

        # --- ブランチラベルの描画（上部に横長矩形として配置）---
        for branch in graph.branches:
            lane = branch_lane.get(branch.name, 0)
            cx = lane_x[lane]
            color = branch_color_map.get(branch.name, _BRANCH_COLORS[0])

            lbl = _add_rectangle(
                slide,
                cx - _LABEL_W // 2,
                top + _MARGIN // 4,
                _LABEL_W,
                _LABEL_H,
                rounded=True,
            )
            _set_solid_fill(lbl, color)
            _set_line_color(lbl, _darker(color), pt=0.75)
            _set_text(lbl, branch.name, font_size_pt=8.0, bold=True,
                      color=_LABEL_TEXT_COLOR, align_center=True)

        # --- マージ接続線の描画 ---
        for commit in graph.commits:
            if commit.commit_type in ("MERGE", "CHERRY_PICK") and len(commit.parents) >= 2:
                cx, cy = commit_pos[commit.commit_id]
                src_id = commit.parents[1]
                if src_id in commit_pos:
                    sx, sy = commit_pos[src_id]
                    src_color = branch_color_map.get(
                        next((c.branch for c in graph.commits if c.commit_id == src_id), ""),
                        _BRANCH_COLORS[0],
                    )
                    _add_connector(slide, sx, sy, cx, cy, rgb=src_color, pt=2.0, dashed=True)

        # --- コミット円の描画 ---
        for commit in graph.commits:
            cx, cy = commit_pos[commit.commit_id]
            color = branch_color_map.get(commit.branch, _BRANCH_COLORS[0])
            self._draw_commit(slide, commit, cx, cy, color)

        # --- タグの描画 ---
        for commit in graph.commits:
            if commit.tag:
                cx, cy = commit_pos[commit.commit_id]
                self._draw_tag(slide, commit.tag, cx - _COMMIT_R - _TAG_GAP, cy, direction="TB")

        # --- コミットIDラベルの描画（コミット円右側）---
        for commit in graph.commits:
            cx, cy = commit_pos[commit.commit_id]
            label = commit.commit_id
            tb = slide.shapes.add_textbox(
                Emu(cx + _COMMIT_R + 30_000),
                Emu(cy - _COMMIT_LABEL_H // 2),
                Emu(_COMMIT_LABEL_W),
                Emu(_COMMIT_LABEL_H),
            )
            tf = tb.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            from pptx.enum.text import PP_ALIGN
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = label
            run.font.size = Pt(6.0)
            run.font.color.rgb = _rgb(*_COMMIT_LABEL_COLOR)

    # ------------------------------------------------------------------
    # コミット円の描画
    # ------------------------------------------------------------------

    def _draw_commit(
        self,
        slide: Slide,
        commit: GitCommit,
        cx: int,
        cy: int,
        color: tuple[int, int, int],
    ) -> None:
        """
        コミットタイプに応じてコミットシェープを描画する。

        コミットタイプと描画方法の対応:
          NORMAL       → 塗りつぶし楕円
          REVERSE      → 楕円 + 白テキスト "✕"
          HIGHLIGHT    → 塗りつぶし矩形
          MERGE        → 楕円 + 上に白い小楕円（二重円）
          CHERRY_PICK  → 紫色の楕円

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        commit : GitCommit
            描画するコミット情報。
        cx, cy : int
            コミット中心のEMU座標。
        color : tuple[int, int, int]
            ブランチ色 (R, G, B)。
        """
        ctype = commit.commit_type

        if ctype == "HIGHLIGHT":
            # 矩形で表現（枠線はブランチ色を暗くする）
            shape = _add_rectangle(
                slide,
                cx - _COMMIT_R,
                cy - _COMMIT_R,
                _COMMIT_D,
                _COMMIT_D,
            )
            _set_solid_fill(shape, color)
            _set_line_color(shape, _darker(color), pt=1.5)

        elif ctype == "MERGE":
            # 外側円
            outer = _add_ellipse(slide, cx, cy, _COMMIT_D, _COMMIT_D)
            _set_solid_fill(outer, color)
            _set_line_color(outer, _darker(color), pt=1.5)
            # 内側白円（二重円の内側）
            inner = _add_ellipse(slide, cx, cy, _INNER_D, _INNER_D)
            _set_solid_fill(inner, (255, 255, 255))
            _set_line_color(inner, color, pt=1.5)

        elif ctype == "REVERSE":
            # 楕円 + "✕" テキスト
            shape = _add_ellipse(slide, cx, cy, _COMMIT_D, _COMMIT_D)
            _set_solid_fill(shape, color)
            _set_line_color(shape, _darker(color), pt=1.5)
            # テキストとして "✕" を設定する
            _set_text(shape, "✕", font_size_pt=9.0, bold=True,
                      color=(255, 255, 255), align_center=True)

        elif ctype == "CHERRY_PICK":
            # 紫色の楕円
            shape = _add_ellipse(slide, cx, cy, _COMMIT_D, _COMMIT_D)
            _set_solid_fill(shape, _CHERRY_COLOR)
            _set_line_color(shape, _darker(_CHERRY_COLOR), pt=1.5)
            # "cp" テキスト（チェリーピックを示す）
            _set_text(shape, "cp", font_size_pt=7.0, bold=True,
                      color=(255, 255, 255), align_center=True)

        else:
            # NORMAL （デフォルト）
            shape = _add_ellipse(slide, cx, cy, _COMMIT_D, _COMMIT_D)
            _set_solid_fill(shape, color)
            _set_line_color(shape, _darker(color), pt=1.5)

    # ------------------------------------------------------------------
    # タグ矩形の描画
    # ------------------------------------------------------------------

    def _draw_tag(
        self,
        slide: Slide,
        tag_text: str,
        x: int,
        y: int,
        direction: str = "LR",
    ) -> None:
        """
        タグ文字列を小矩形テキストボックスとして描画する。

        OOXML では旗型シェープが使えないため、角丸矩形で代替する。

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        tag_text : str
            タグのテキスト文字列。
        x, y : int
            タグ矩形の中心X座標, 上端Y座標（LR方向の場合）。
        direction : str
            "LR" の場合はコミット上部、"TB" の場合はコミット左側に配置する。
        """
        if direction == "LR":
            # タグをコミット上部に配置する
            tag = _add_rectangle(
                slide,
                x - _TAG_W // 2,
                y - _TAG_H,
                _TAG_W,
                _TAG_H,
                rounded=True,
            )
        else:
            # タグをコミット左側に配置する
            tag = _add_rectangle(
                slide,
                x - _TAG_W,
                y - _TAG_H // 2,
                _TAG_W,
                _TAG_H,
                rounded=True,
            )

        _set_solid_fill(tag, _TAG_BG_COLOR)
        _set_line_color(tag, _TAG_BORDER_COLOR, pt=0.75)
        _set_text(tag, tag_text, font_size_pt=6.5, bold=False,
                  color=(80, 70, 0), align_center=True)


# ---------------------------------------------------------------------------
# レイアウトユーティリティ
# ---------------------------------------------------------------------------

def _sort_branches(branches: list[GitBranch]) -> list[GitBranch]:
    """
    ブランチリストを表示順にソートする。

    ソート優先度:
    1. order=None のブランチを定義順で並べる（元のリスト順を維持）
    2. order 指定のブランチをその値の昇順で並べる

    この結果、order なし → order 小→大 の順でレーンが割り当てられる。

    Parameters
    ----------
    branches : list[GitBranch]
        元の定義順ブランチリスト。

    Returns
    -------
    list[GitBranch]
        表示順にソートされたブランチリスト。
    """
    no_order = [b for b in branches if b.order is None]
    with_order = sorted([b for b in branches if b.order is not None], key=lambda b: b.order)  # type: ignore[arg-type]
    return no_order + with_order


def _calc_positions_1d(
    n: int,
    total_size: int,
    gap: int,
    start: int,
) -> list[int]:
    """
    n 個の要素を 1 次元に等間隔配置した中心座標リストを計算する。

    描画エリアに収まるよう間隔を自動調整する。

    Parameters
    ----------
    n : int
        要素数。
    total_size : int
        配置可能な総幅/総高さ（EMU）。
    gap : int
        要素間の理想間隔（EMU）。
    start : int
        先頭方向の開始座標（EMU）。

    Returns
    -------
    list[int]
        各要素の中心座標リスト（EMU）。
    """
    if n == 0:
        return []
    if n == 1:
        return [start + total_size // 2]

    # 実際の間隔 = min(理想間隔, 総サイズ / (n-1)) で調整する
    actual_gap = min(gap, total_size // (n - 1)) if n > 1 else gap
    # 全体の幅 = (n - 1) * gap
    total_span = (n - 1) * actual_gap
    # 中央に寄せる
    offset = start + (total_size - total_span) // 2

    return [offset + i * actual_gap for i in range(n)]
