"""
Mermaid ユーザージャーニー図レンダラーモジュール。

JourneyChart オブジェクトを入力として PPTXスライドにユーザージャーニー図を描画する。

レイアウト構成（上から下）:
  1. タイトル行（全幅テキストボックス）
  2. グリッドエリア（左端にアクター凡例、右側にタスク列）:
       - セクションヘッダー帯（色付き丸角矩形）
       - 感情アイコン行（😊 / 😐 / 😢 表示の円シェープ）
       - タスクカード行（タスク名入り丸角矩形）
       - アクタードット行（参加アクターを色付き小円で表示）

感情マッピング（Mermaid 公式仕様準拠）:
  score > 3  → 😊 （笑顔: 高満足）
  score == 3 → 😐 （中立）
  score < 3  → 😢 （悲しい: 低満足）

セクション色は 7 色サイクルで自動割り当てする。
アクター色は 6 色サイクルで自動割り当てする。
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .journey_parser import JourneyChart, JourneyTask

# ---------------------------------------------------------------------------
# 感情スコア → 絵文字マッピング
# ---------------------------------------------------------------------------

def _score_to_emoji(score: int) -> str:
    """
    Mermaid 仕様に準拠して感情スコアを絵文字に変換する。

    Parameters
    ----------
    score : int
        感情スコア (1〜5)。

    Returns
    -------
    str
        感情を表す絵文字 (😊 / 😐 / 😢)。
    """
    if score > 3:
        return "😊"
    elif score < 3:
        return "😢"
    else:
        return "😐"


# ---------------------------------------------------------------------------
# カラー定義 (RGB タプル)
# ---------------------------------------------------------------------------

# セクションヘッダー背景色（7 色サイクル、Mermaid 公式 sectionFills に近い暗めの色）
_SECTION_COLORS: list[tuple[int, int, int]] = [
    (91, 74, 160),    # 紫
    (46, 100, 160),   # 青
    (46, 140, 90),    # 緑
    (160, 70, 40),    # 赤橙
    (90, 90, 30),     # 黄緑
    (40, 120, 140),   # 水色
    (120, 40, 90),    # ピンク
]

# セクションヘッダー文字色（白）
_SECTION_TEXT_RGB: tuple[int, int, int] = (255, 255, 255)

# アクタードット色（6 色サイクル、Mermaid 公式 actorColours に対応）
_ACTOR_COLORS: list[tuple[int, int, int]] = [
    (143, 188, 143),  # DarkSeaGreen (#8FBC8F)
    (70, 130, 180),   # SteelBlue
    (255, 165,   0),  # Orange
    (220,  20,  60),  # Crimson
    (138,  43, 226),  # BlueViolet
    (32, 178, 170),   # LightSeaTeal
]

# タスクカード背景色（薄いグレー）
_TASK_BG_RGB: tuple[int, int, int] = (245, 245, 250)

# タスクカード枠色（明るいグレー）
_TASK_BORDER_RGB: tuple[int, int, int] = (200, 200, 210)

# タスクカード文字色（濃紺）
_TASK_TEXT_RGB: tuple[int, int, int] = (40, 40, 60)

# タイトル文字色（ほぼ黒）
_TITLE_TEXT_RGB: tuple[int, int, int] = (30, 30, 60)

# 感情アイコン円の背景色（薄黄色）
_EMOTION_BG_RGB: tuple[int, int, int] = (255, 253, 200)

# 感情アイコン円の枠色（薄い金色）
_EMOTION_BORDER_RGB: tuple[int, int, int] = (200, 180, 80)

# アクタードット枠色（暗いグレー）
_DOT_STROKE_RGB: tuple[int, int, int] = (80, 80, 80)

# アクター凡例テキスト色
_LEGEND_TEXT_RGB: tuple[int, int, int] = (40, 40, 60)


# ---------------------------------------------------------------------------
# レイアウト定数
# ---------------------------------------------------------------------------

# タイトル行の高さ比率（全体高さに対する割合）
_TITLE_RATIO: float = 0.12

# グリッド内の各行高さ比率（タイトルを除いた残り高さに対する割合）
_SECTION_ROW_RATIO: float = 0.20   # セクションヘッダー行
_EMOTION_ROW_RATIO: float = 0.28   # 感情アイコン行
_TASK_ROW_RATIO: float = 0.34      # タスクカード行
_DOT_ROW_RATIO: float = 0.18       # アクタードット行

# アクター凡例エリアの幅（EMU。1.2 inch 相当）
_LEGEND_WIDTH_EMU: int = 1_100_000

# タスク列の最小幅（EMU。0.76 inch 相当）
_MIN_TASK_WIDTH_EMU: int = 700_000

# タスク列間の隙間（EMU）
_TASK_GAP_EMU: int = 50_000

# アクタードットの直径（EMU）
_DOT_DIAMETER_EMU: int = 190_000

# アクター凡例ドットの直径（EMU）
_LEGEND_DOT_DIAMETER_EMU: int = 160_000

# タスクカード内余白（EMU）
_CARD_PAD_EMU: int = 40_000

# シェープ型番（MSO_AUTO_SHAPE_TYPE）
_SHAPE_RECT: int = 1          # 矩形
_SHAPE_ROUNDED_RECT: int = 5  # 丸角矩形
_SHAPE_OVAL: int = 9          # 楕円（円）


class JourneyRenderer:
    """
    JourneyChart データクラスを受け取り、PowerPoint スライドに
    ユーザージャーニー図を描画するクラス。
    """

    def render(
        self,
        slide: Slide,
        chart: JourneyChart,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        ユーザージャーニー図をスライドに描画する。

        タスクが 0 件の場合は何も描画しない。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        chart : JourneyChart
            parse_journey() で生成されたジャーニーデータ。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        if not chart.tasks:
            return

        n_tasks = len(chart.tasks)
        has_actors = bool(chart.actors)

        # ------------------------------------------------------------------
        # 座標・サイズを計算する
        # ------------------------------------------------------------------

        # アクター凡例エリアの幅（アクターがいない場合は 0）
        legend_w = _LEGEND_WIDTH_EMU if has_actors else 0

        # グリッドエリアの左端・幅
        grid_left = left + legend_w
        grid_w = width - legend_w

        # タスク列の幅: グリッド幅をタスク数 + ギャップで等分する
        task_total_gap = _TASK_GAP_EMU * max(0, n_tasks - 1)
        task_w = max(
            _MIN_TASK_WIDTH_EMU,
            (grid_w - task_total_gap) // n_tasks,
        )

        # タイトル行の高さ
        title_h = max(300_000, int(height * _TITLE_RATIO))

        # グリッドエリアの上端・高さ（タイトル行の下）
        grid_top = top + title_h
        grid_h = height - title_h

        # グリッド内の各行高さを比率から計算する
        section_h = max(250_000, int(grid_h * _SECTION_ROW_RATIO))
        emotion_h = max(350_000, int(grid_h * _EMOTION_ROW_RATIO))
        task_h    = max(350_000, int(grid_h * _TASK_ROW_RATIO))
        dot_h     = max(180_000, int(grid_h * _DOT_ROW_RATIO))

        # 各行の top 座標
        section_top = grid_top
        emotion_top = section_top + section_h
        task_top    = emotion_top + emotion_h
        dot_top     = task_top + task_h

        # ------------------------------------------------------------------
        # アクター → カラーマップを構築する
        # ------------------------------------------------------------------
        actor_color_map: dict[str, tuple[int, int, int]] = {
            actor: _ACTOR_COLORS[idx % len(_ACTOR_COLORS)]
            for idx, actor in enumerate(chart.actors)
        }

        # ------------------------------------------------------------------
        # タイトルを描画する
        # ------------------------------------------------------------------
        if chart.title:
            self._draw_title(slide, chart.title, left, top, width, title_h)

        # ------------------------------------------------------------------
        # アクター凡例を描画する（グリッド左端）
        # ------------------------------------------------------------------
        if has_actors:
            self._draw_actor_legend(
                slide,
                chart.actors,
                actor_color_map,
                left,
                grid_top,
                legend_w,
                grid_h,
            )

        # ------------------------------------------------------------------
        # セクション色インデックスマップを作成する
        # ------------------------------------------------------------------
        section_color_idx: dict[str, int] = {
            sec: idx % len(_SECTION_COLORS)
            for idx, sec in enumerate(chart.sections)
        }

        # ------------------------------------------------------------------
        # タスク列を左から順に描画する
        # ------------------------------------------------------------------
        task_x = grid_left
        for task_idx, task in enumerate(chart.tasks):
            sec_color = _SECTION_COLORS[section_color_idx.get(task.section, 0)]

            # セクションヘッダー: 同一セクションの先頭タスクでのみ描画する
            is_section_start = (
                task_idx == 0
                or chart.tasks[task_idx - 1].section != task.section
            )
            if is_section_start and task.section:
                # このセクションに含まれる連続タスク数を計算する
                span = self._count_section_span(chart.tasks, task_idx)
                # セクションヘッダーの幅: span 個のタスク列 + (span-1) 個のギャップ
                header_w = span * task_w + _TASK_GAP_EMU * max(0, span - 1)
                self._draw_section_header(
                    slide,
                    task.section,
                    sec_color,
                    task_x,
                    section_top,
                    header_w,
                    section_h,
                )

            # 感情アイコンを描画する
            self._draw_emotion_icon(
                slide,
                task.score,
                task_x,
                emotion_top,
                task_w,
                emotion_h,
            )

            # タスクカードを描画する
            self._draw_task_card(
                slide,
                task.task,
                sec_color,
                task_x,
                task_top,
                task_w,
                task_h,
            )

            # アクタードットを描画する（アクターがいる場合のみ）
            if task.people:
                self._draw_actor_dots(
                    slide,
                    task.people,
                    actor_color_map,
                    task_x,
                    dot_top,
                    task_w,
                    dot_h,
                )

            # 次のタスク列の左端へ進む
            task_x += task_w + _TASK_GAP_EMU

    # -------------------------------------------------------------------------
    # 描画ヘルパーメソッド
    # -------------------------------------------------------------------------

    def _draw_title(
        self,
        slide: Slide,
        title: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """タイトルテキストボックスを描画する（中央揃え・太字）。"""
        shape = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        tf = shape.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*_TITLE_TEXT_RGB)

    def _draw_actor_legend(
        self,
        slide: Slide,
        actors: list[str],
        color_map: dict[str, tuple[int, int, int]],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        アクター凡例を左サイドに縦並びで描画する。

        各アクターを「色付き小円 + アクター名テキスト」として表示する。
        """
        n = len(actors)
        if n == 0:
            return

        # 1アクターあたりの行高（最大 400,000 EMU）
        row_h = min(height // n, 400_000)
        dot_d = _LEGEND_DOT_DIAMETER_EMU
        left_pad = 80_000   # 円の左マージン
        text_gap = 60_000   # 円と名前の隙間

        for idx, actor_name in enumerate(actors):
            rgb = color_map.get(actor_name, _ACTOR_COLORS[0])
            row_top = top + idx * row_h
            row_cy = row_top + row_h // 2

            # 色付き小円を描画する
            dot_left = left + left_pad
            dot_top_pos = row_cy - dot_d // 2
            circle = slide.shapes.add_shape(
                _SHAPE_OVAL,
                Emu(dot_left),
                Emu(dot_top_pos),
                Emu(dot_d),
                Emu(dot_d),
            )
            self._fill_shape(circle, rgb)
            circle.line.color.rgb = RGBColor(*_DOT_STROKE_RGB)
            circle.line.width = Emu(9525)

            # アクター名テキストボックスを描画する
            text_left = dot_left + dot_d + text_gap
            text_w = width - (text_left - left) - left_pad
            if text_w > 0:
                tb = slide.shapes.add_textbox(
                    Emu(text_left),
                    Emu(row_top),
                    Emu(text_w),
                    Emu(row_h),
                )
                tf = tb.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = actor_name
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*_LEGEND_TEXT_RGB)

    def _draw_section_header(
        self,
        slide: Slide,
        section_name: str,
        color: tuple[int, int, int],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        セクションヘッダー帯を描画する（色付き丸角矩形 + セクション名テキスト）。

        同一セクションの全タスク列にまたがる幅で描画する。
        """
        shape = slide.shapes.add_shape(
            _SHAPE_ROUNDED_RECT,
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
        )
        self._fill_shape(shape, color)
        shape.line.fill.background()  # 枠線なし

        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = section_name
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*_SECTION_TEXT_RGB)

    def _draw_emotion_icon(
        self,
        slide: Slide,
        score: int,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        感情絵文字アイコンを円シェープで描画する。

        Mermaid の drawFace 仕様に準拠して score > 3 → 😊、
        score == 3 → 😐、score < 3 → 😢 の絵文字を使用する。
        円の下部にスコア数値を小さく表示する。
        """
        emoji = _score_to_emoji(score)

        # 円サイズ: 行高・列幅の小さい方から余白を引いた値
        icon_size = min(width - 80_000, height - 60_000)
        icon_size = max(icon_size, 200_000)

        # 行内中央に配置する
        cx = left + (width - icon_size) // 2
        cy = top + (height - icon_size) // 2

        # 薄黄色の円シェープを追加する
        circle = slide.shapes.add_shape(
            _SHAPE_OVAL,
            Emu(cx),
            Emu(cy),
            Emu(icon_size),
            Emu(icon_size),
        )
        self._fill_shape(circle, _EMOTION_BG_RGB)
        circle.line.color.rgb = RGBColor(*_EMOTION_BORDER_RGB)
        circle.line.width = Emu(9525)

        tf = circle.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = emoji
        run.font.size = Pt(16)

    def _draw_task_card(
        self,
        slide: Slide,
        task_name: str,
        section_color: tuple[int, int, int],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        タスクカードを描画する（丸角矩形 + タスク名テキスト）。

        枠線にセクション色を適用してセクションとの関連を視覚化する。
        """
        card_left = left + _CARD_PAD_EMU
        card_top = top + _CARD_PAD_EMU
        card_w = max(10_000, width - _CARD_PAD_EMU * 2)
        card_h = max(10_000, height - _CARD_PAD_EMU * 2)

        shape = slide.shapes.add_shape(
            _SHAPE_ROUNDED_RECT,
            Emu(card_left),
            Emu(card_top),
            Emu(card_w),
            Emu(card_h),
        )
        self._fill_shape(shape, _TASK_BG_RGB)
        shape.line.color.rgb = RGBColor(*section_color)
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = task_name
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(*_TASK_TEXT_RGB)

    def _draw_actor_dots(
        self,
        slide: Slide,
        people: list[str],
        color_map: dict[str, tuple[int, int, int]],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        タスクに参加するアクターを色付き小円として横並びで描画する。

        ドット群はタスク列内で水平中央に配置する。
        """
        n = len(people)
        dot_d = min(height - 20_000, _DOT_DIAMETER_EMU)
        dot_d = max(dot_d, 80_000)
        dot_gap = 40_000  # ドット間隙

        # 全ドットの幅を計算してタスク列内で中央配置する
        total_w = dot_d * n + dot_gap * max(0, n - 1)
        start_x = left + max(0, (width - total_w) // 2)
        dot_top_pos = top + (height - dot_d) // 2

        for idx, person in enumerate(people):
            rgb = color_map.get(person, _ACTOR_COLORS[0])
            dot_left = start_x + idx * (dot_d + dot_gap)

            circle = slide.shapes.add_shape(
                _SHAPE_OVAL,
                Emu(dot_left),
                Emu(dot_top_pos),
                Emu(dot_d),
                Emu(dot_d),
            )
            self._fill_shape(circle, rgb)
            circle.line.color.rgb = RGBColor(*_DOT_STROKE_RGB)
            circle.line.width = Emu(9525)

    # -------------------------------------------------------------------------
    # 静的ユーティリティ
    # -------------------------------------------------------------------------

    @staticmethod
    def _fill_shape(shape: object, rgb: tuple[int, int, int]) -> None:
        """
        シェープの塗りつぶし色を設定する。

        Parameters
        ----------
        shape : object
            python-pptx の Shape オブジェクト。
        rgb : tuple[int, int, int]
            塗りつぶし色の (R, G, B) タプル。
        """
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)

    @staticmethod
    def _count_section_span(tasks: list[JourneyTask], start_idx: int) -> int:
        """
        指定タスクインデックスから始まる同一セクションの連続タスク数を返す。

        Parameters
        ----------
        tasks : list[JourneyTask]
            全タスクリスト。
        start_idx : int
            カウント開始インデックス。

        Returns
        -------
        int
            同一セクションに連続するタスクの数（最小 1）。
        """
        target_section = tasks[start_idx].section
        count = 0
        for ti in range(start_idx, len(tasks)):
            if tasks[ti].section == target_section:
                count += 1
            else:
                break
        return max(1, count)
