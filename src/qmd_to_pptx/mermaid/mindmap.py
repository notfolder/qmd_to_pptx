"""
マインドマップレンダラーモジュール。

mindmap を放射状レイアウト（ルート中心・全方位ブランチ展開）で描画する。
葉ノード数に比例した角度配分により各ブランチが均等な空間を占有する。
"""

from __future__ import annotations

import math

from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer

# ---- ノードサイズ定義（深さ別、EMU）----
# 深い階層ほど小さくしてキャンバス内に収まりやすくする
_NODE_SIZES: list[tuple[int, int]] = [
    (1_200_000, 480_000),   # depth 0: ルートノード
    (900_000, 370_000),     # depth 1: L1ノード
    (680_000, 280_000),     # depth 2以降: 末端ノード
]

# ---- ブランチカラーパレット（L1ブランチ用、最大8色） ----
_BRANCH_PALETTE: list[tuple[int, int, int]] = [
    (231, 76, 60),    # 赤
    (230, 126, 34),   # オレンジ
    (241, 196, 15),   # 黄
    (39, 174, 96),    # 緑
    (26, 188, 156),   # ティール
    (52, 152, 219),   # 青
    (155, 89, 182),   # 紫
    (233, 30, 99),    # ピンク
]

# ---- ルートノードの色定義 ----
_ROOT_FILL_RGB: tuple[int, int, int] = (70, 70, 180)
_ROOT_TEXT_RGB: tuple[int, int, int] = (255, 255, 255)

# ---- type値 → prstGeom名 マッピング ----
_TYPE_TO_PRST: dict[int, str] = {
    0: "roundRect",       # DEFAULT（デフォルト - 角丸矩形）
    1: "roundRect",       # ROUNDED_RECT（角丸矩形）
    2: "rect",            # RECT（矩形）
    3: "ellipse",         # CIRCLE（楕円）
    4: "cloud",           # CLOUD（雲形）
    5: "irregularSeal1",  # BANG（爆発形）
    6: "hexagon",         # HEXAGON（六角形）
}


class MindmapRenderer(BaseDiagramRenderer):
    """
    mindmap を放射状（ラジアル）レイアウトで描画するレンダラー。

    ルートノードをスライド描画エリア中央に配置し、各ブランチを
    葉ノード数比例の角度セクターに割り当てて全方位に展開する。
    ノード形状はMermaidのtype値に応じて変化し、ブランチごとに
    異なるカラーを適用する。コネクターはcurvedConnector3を使用する。
    """

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        mindmap をスライドに放射状に描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"nodes"キーにルートから始まるツリー構造を含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        nodes_data = graph_data.get("nodes", [])
        if not nodes_data:
            return

        # ルートノードはnodes_dataの先頭要素（完全なツリー構造を持つ）
        root = nodes_data[0]
        root_id = root.get("nodeId", "")
        if not root_id:
            return

        # ツリーを再帰走査してnodeId→nodeデータのマップを構築する
        node_info: dict[str, dict] = {}
        self._collect_node_info(root, node_info)
        if not node_info:
            return

        # 最大深さを算出する（depth = level // 2 - 1）
        max_depth = max(self._node_depth(n) for n in node_info.values())
        max_depth = max(max_depth, 1)  # ゼロ除算防止

        # キャンバス中心（EMU）を計算する
        cx_emu = left + width // 2
        cy_emu = top + height // 2

        # 外側楕円の半径：末端ノードが枠内に収まるようマージンを引く
        # アスペクト比に合わせてX・Y半径を独立に設定することで歪みを防ぐ
        outer_w, outer_h = _NODE_SIZES[2]
        padding_x = outer_w // 2 + 80_000
        padding_y = outer_h // 2 + 60_000
        rx_outer = max(1, width // 2 - padding_x)
        ry_outer = max(1, height // 2 - padding_y)

        # 放射状座標（EMU）とブランチカラーを計算する
        emu_pos: dict[str, tuple[int, int]] = {root_id: (cx_emu, cy_emu)}
        branch_color_map: dict[str, tuple[int, int, int]] = {
            root_id: _ROOT_FILL_RGB
        }
        self._assign_radial_emu_positions(
            root, 0.0, 2 * math.pi, 1, max_depth,
            emu_pos, branch_color_map,
            cx_emu, cy_emu, rx_outer, ry_outer,
        )

        # 親子エッジを収集する
        edges: list[tuple[str, str]] = []
        self._collect_edges(root, edges)

        # ノードShapeを描画する（EMU中心座標を直接渡す）
        node_shapes: dict[str, object] = {}
        for node_id, node in node_info.items():
            if node_id not in emu_pos:
                continue
            nx, ny = emu_pos[node_id]
            depth = self._node_depth(node)
            fill_rgb = branch_color_map.get(node_id, (200, 200, 200))
            shape = self._draw_mindmap_node(slide, node, nx, ny, depth, fill_rgb)
            node_shapes[node_id] = shape

        # 曲線コネクターを描画する
        for src_id, dst_id in edges:
            if src_id not in emu_pos or dst_id not in emu_pos:
                continue
            color = branch_color_map.get(dst_id, (150, 150, 150))
            self._draw_curved_mindmap_edge(
                slide,
                emu_pos[src_id], emu_pos[dst_id],
                node_shapes.get(src_id), node_shapes.get(dst_id),
                color,
            )

    # ------------------------------------------------------------------
    # レイアウト計算ヘルパー
    # ------------------------------------------------------------------

    def _node_depth(self, node: dict) -> int:
        """
        ノードの深さを返す（root=0, L1=1, L2=2, ...）。

        パーサーのlevelフィールドは2の倍数（root=2, L1=4, L2=6）で表されるため、
        depth = level // 2 - 1 に変換する。

        Parameters
        ----------
        node : dict
            パーサーが返したノード辞書。

        Returns
        -------
        int
            ノードの深さ（0以上）。
        """
        return max(0, node.get("level", 2) // 2 - 1)

    def _count_leaves(self, node: dict) -> int:
        """
        サブツリーの葉ノード数を再帰的に数える。

        Parameters
        ----------
        node : dict
            基点ノード。

        Returns
        -------
        int
            葉ノードの数（自身が葉であれば1）。
        """
        children = node.get("children", [])
        if not children:
            return 1
        return sum(self._count_leaves(c) for c in children)

    def _collect_node_info(
        self, node: dict, node_info: dict[str, dict]
    ) -> None:
        """
        ツリーを再帰走査してnodeId→nodeデータのマップを構築する。

        Parameters
        ----------
        node : dict
            現在のノード。
        node_info : dict[str, dict]
            構築対象のマップ（更新対象）。
        """
        node_id = node.get("nodeId", "")
        if node_id:
            node_info[node_id] = node
        for child in node.get("children", []):
            self._collect_node_info(child, node_info)

    def _collect_edges(
        self, node: dict, edges: list[tuple[str, str]]
    ) -> None:
        """
        ツリーを再帰走査して親子エッジを収集する。

        Parameters
        ----------
        node : dict
            現在のノード。
        edges : list[tuple[str, str]]
            (親nodeId, 子nodeId) を格納するリスト（更新対象）。
        """
        node_id = node.get("nodeId", "")
        for child in node.get("children", []):
            child_id = child.get("nodeId", "")
            if node_id and child_id:
                edges.append((node_id, child_id))
            self._collect_edges(child, edges)

    def _assign_radial_emu_positions(
        self,
        node: dict,
        angle_start: float,
        angle_end: float,
        depth: int,
        max_depth: int,
        emu_pos: dict[str, tuple[int, int]],
        branch_color_map: dict[str, tuple[int, int, int]],
        cx: int,
        cy: int,
        rx_outer: int,
        ry_outer: int,
        parent_color: tuple[int, int, int] | None = None,
        l1_index: int = 0,
    ) -> int:
        """
        ノードの子を楕円放射状に配置し、emu_pos と branch_color_map を更新する。

        キャンバスの X 半径 rx_outer と Y 半径 ry_outer を独立に管理することで
        アスペクト比の歪みを防ぎ、全方向に均等なノード間隔を実現する。
        各子の角度セクターは配下葉ノード数に比例して配分する。

        Parameters
        ----------
        node : dict
            子を配置する親ノード。
        angle_start, angle_end : float
            子に割り当てる角度セクター（ラジアン）。
        depth : int
            配置対象の子ノードの深さ（1=L1, 2=L2, ...）。
        max_depth : int
            ツリー全体の最大深さ（半径計算に使用）。
        emu_pos : dict
            ノードID→EMU中心座標のマップ（更新対象）。
        branch_color_map : dict
            ノードID→RGB色のマップ（更新対象）。
        cx, cy : int
            キャンバス中心のEMU座標。
        rx_outer, ry_outer : int
            外側楕円のX・Y半径（EMU）。アスペクト比に合わせて個別設定する。
        parent_color : tuple[int, int, int] | None
            親ノードの色（L2以降の色継承に使用）。
        l1_index : int
            L1ブランチのパレットインデックス（呼び出し元で管理）。

        Returns
        -------
        int
            次のL1インデックス値（l1_indexの更新値）。
        """
        children = node.get("children", [])
        if not children:
            return l1_index

        # 現ノードの全葉ノード数（角度比率の分母として使用）
        node_leaves = self._count_leaves(node)

        # この深さでのX・Y楕円半径（外辺半径の92%まで使用）
        r_frac = (depth / max_depth) * 0.92
        rx = int(r_frac * rx_outer)
        ry = int(r_frac * ry_outer)

        current_angle = angle_start
        next_l1_index = l1_index

        for child in children:
            child_id = child.get("nodeId", "")
            child_leaves = self._count_leaves(child)

            # 子の角度幅 = 現ノードのセクター × (子の葉数 / 現ノードの葉数)
            angle_span = (child_leaves / node_leaves) * (angle_end - angle_start)
            angle_mid = current_angle + angle_span / 2.0

            # 楕円極座標をEMU座標に変換する（X・Yで独立した半径を使用）
            x_emu = cx + int(rx * math.cos(angle_mid))
            y_emu = cy + int(ry * math.sin(angle_mid))
            emu_pos[child_id] = (x_emu, y_emu)

            # ブランチカラーを決定する
            if depth == 1:
                # L1: パレットから順番に選択する
                color: tuple[int, int, int] = _BRANCH_PALETTE[
                    next_l1_index % len(_BRANCH_PALETTE)
                ]
                next_l1_index += 1
            else:
                # L2以降: 親色を白方向に40%明るくする
                pc = parent_color if parent_color else (200, 200, 200)
                color = (
                    pc[0] + int((255 - pc[0]) * 0.4),
                    pc[1] + int((255 - pc[1]) * 0.4),
                    pc[2] + int((255 - pc[2]) * 0.4),
                )
            branch_color_map[child_id] = color

            # 子サブツリーを再帰的に配置する
            next_l1_index = self._assign_radial_emu_positions(
                child,
                current_angle, current_angle + angle_span,
                depth + 1, max_depth,
                emu_pos, branch_color_map,
                cx, cy, rx_outer, ry_outer,
                parent_color=color,
                l1_index=next_l1_index,
            )
            current_angle += angle_span

        return next_l1_index

    # ------------------------------------------------------------------
    # ノード描画ヘルパー
    # ------------------------------------------------------------------

    def _node_size(self, depth: int) -> tuple[int, int]:
        """
        深さに応じたノードサイズ（幅, 高さ）をEMUで返す。

        Parameters
        ----------
        depth : int
            ノードの深さ（0=ルート, 1=L1, 2以降=末端）。

        Returns
        -------
        tuple[int, int]
            (width_emu, height_emu) のタプル。
        """
        return _NODE_SIZES[min(depth, 2)]

    def _draw_mindmap_node(
        self,
        slide: Slide,
        node: dict,
        cx_emu: int,
        cy_emu: int,
        depth: int,
        fill_rgb: tuple[int, int, int],
    ) -> object:
        """
        ノードをtype別の形状でスライドに描画し、カラーを適用する。

        type値とprstGeomの対応はモジュールレベルの _TYPE_TO_PRST を参照。
        テキストフォントサイズは深さに応じて変化する（ルート14pt, L1 11pt, L2以降 10pt）。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        node : dict
            ノードデータ（nodeId, descr, type等を持つ）。
        cx_emu, cy_emu : int
            ノード中心のEMU座標（楕円ラジアルレイアウトで計算済み）。
        depth : int
            ノードの深さ（0=ルート, 1=L1, 2以降）。
        fill_rgb : tuple[int, int, int]
            塗り色の(R, G, B)。

        Returns
        -------
        object
            追加したShapeオブジェクト。
        """
        node_type = node.get("type", 0)
        label = node.get("descr", node.get("nodeId", ""))
        w_emu, h_emu = self._node_size(depth)

        sl = cx_emu - w_emu // 2
        st = cy_emu - h_emu // 2

        # いったんRECTANGLEで追加してからXMLでprstGeomを書き換える
        shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Emu(sl), Emu(st), Emu(w_emu), Emu(h_emu),
        )

        # テキストとフォントを設定する
        shape.text = label
        tf = shape.text_frame
        tf.word_wrap = True
        font_pt = Pt(14 if depth == 0 else 11 if depth == 1 else 10)
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = font_pt
                if depth == 0:
                    # ルートは白文字
                    run.font.color.rgb = RGBColor(*_ROOT_TEXT_RGB)
                else:
                    # その他は濃いグレー文字
                    run.font.color.rgb = RGBColor(30, 30, 30)

        # 塗り色と枠線色を設定する
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
        line_rgb = (
            max(0, fill_rgb[0] - 40),
            max(0, fill_rgb[1] - 40),
            max(0, fill_rgb[2] - 40),
        )
        shape.line.color.rgb = RGBColor(*line_rgb)

        # prstGeomをtype値に応じて書き換える
        self._apply_node_prst_geom(shape, node_type)

        return shape

    def _apply_node_prst_geom(self, shape: object, node_type: int) -> None:
        """
        ShapeのprstGeom属性をノードtype値に応じて書き換える。

        Parameters
        ----------
        shape : Shape
            書き換え対象のShapeオブジェクト。
        node_type : int
            Mermaidのtype値（0〜6）。
        """
        prst_name = _TYPE_TO_PRST.get(node_type, "roundRect")
        sp_el = shape._element
        spPr = sp_el.find(qn("p:spPr"))
        if spPr is None:
            return
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            prstGeom.set("prst", prst_name)

    # ------------------------------------------------------------------
    # エッジ描画ヘルパー
    # ------------------------------------------------------------------

    def _draw_curved_mindmap_edge(
        self,
        slide: Slide,
        src_emu: tuple[int, int],
        dst_emu: tuple[int, int],
        src_shape: object | None,
        dst_shape: object | None,
        color_rgb: tuple[int, int, int],
    ) -> None:
        """
        2ノード間にcurvedConnector3を描画し、ブランチカラーを適用する。

        始点・終点の方向に基づいて接続ポイントを選択し、
        コネクター線色をXML操作でブランチカラーに設定する。
        EMU座標を直接受け取ることでアスペクト比の計算を省く。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        src_emu, dst_emu : tuple[int, int]
            始点・終点の中心EMU座標。
        src_shape, dst_shape : Shape | None
            始点・終点のShapeオブジェクト（Noneの場合は接続スキップ）。
        color_rgb : tuple[int, int, int]
            コネクター線のRGB色。
        """
        sx, sy = src_emu
        dx, dy = dst_emu

        # curvedConnector3 を追加する（MSO_CONNECTOR_TYPE.CURVE = 3）
        connector = slide.shapes.add_connector(
            3, Emu(sx), Emu(sy), Emu(dx), Emu(dy)
        )

        # 方向ベクトルから接続ポイントインデックスを決定する（EMU差分を使用）
        src_cp, dst_cp = self._connection_indices(dx - sx, dy - sy)
        if src_shape is not None:
            connector.begin_connect(src_shape, src_cp)
        if dst_shape is not None:
            connector.end_connect(dst_shape, dst_cp)

        # コネクター線色をXMLで設定する
        sp_el = connector._element
        spPr = sp_el.find(qn("p:spPr"))
        if spPr is None:
            spPr = lxml_etree.SubElement(sp_el, qn("p:spPr"))

        ln_el = spPr.find(qn("a:ln"))
        if ln_el is None:
            ln_el = lxml_etree.SubElement(spPr, qn("a:ln"))

        # 既存のfill要素を除去してから新たに設定する
        for tag in (qn("a:solidFill"), qn("a:noFill")):
            existing = ln_el.find(tag)
            if existing is not None:
                ln_el.remove(existing)

        solidFill = lxml_etree.SubElement(ln_el, qn("a:solidFill"))
        srgbClr = lxml_etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", "{:02X}{:02X}{:02X}".format(*color_rgb))

        # 線幅を2.0pt相当に設定する（1pt = 12700 EMU）
        ln_el.set("w", "25400")
