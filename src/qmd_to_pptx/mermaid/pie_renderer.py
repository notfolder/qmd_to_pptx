"""
Mermaid 円グラフ（pie chart）レンダラーモジュール。

python-pptx のネイティブ円グラフ（XL_CHART_TYPE.PIE）を使い、
PieChart データクラスの情報を PowerPoint スライドに描画する。

描画仕様:
    - タイトル: chart.chart_title （タイトルが空でも非表示にする）
    - データ: ChartData.categories + add_series でカテゴリ名・数値を渡す
    - 凡例（判例）: chart.legend（位置 BOTTOM・常時表示）
    - データラベル:
        - パーセント表示は常時オン
        - showData=True の場合は実数値も表示する
        - カテゴリ名もデータラベルに含める（スライス上で判読しやすくする）
    - データラベル位置: PieChart.text_position を OOXML の dLblPos にマッピングする

textPosition → dLblPos マッピング:
    0.0 〜 0.39  → "ctr"     （スライス中央）
    0.4 〜 0.69  → "inEnd"   （スライス内縁）
    0.7 〜 0.99  → "bestFit" （自動最適、デフォルト）
    1.0          → "outEnd"  （スライス外側）
"""

from __future__ import annotations

import logging

from lxml import etree as lxml_etree
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .pie_parser import PieChart

# モジュールロガーを取得する
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# データラベル位置マッピング定数
# ---------------------------------------------------------------------------
_TEXT_POSITION_MAP: list[tuple[float, str]] = [
    # (閾値以下なら → 対応する OOXML dLblPos 値)
    (0.39, "ctr"),
    (0.69, "inEnd"),
    (0.99, "bestFit"),
]
_TEXT_POSITION_DEFAULT = "bestFit"  # 0.75 デフォルト値が対応する位置

# ---------------------------------------------------------------------------
# タイトルフォントサイズ
# ---------------------------------------------------------------------------
_TITLE_FONT_SIZE_PT: int = 18


class PieChartRenderer:
    """
    PieChart データクラスを受け取り、PowerPoint スライドに円グラフを描画するクラス。

    python-pptx のネイティブ円グラフ機能（ChartData / add_chart）を使用する。
    """

    def render(
        self,
        slide: Slide,
        pie_chart: PieChart,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        PieChart をスライドに描画する。

        セクションが 0 件の場合は何も描画しない。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        pie_chart : PieChart
            parse_pie() で生成された円グラフデータ。
        left, top, width, height : int
            描画エリアの EMU 座標。
        """
        if not pie_chart.sections:
            return

        # ChartData にカテゴリ名と数値を設定する
        chart_data = ChartData()
        chart_data.categories = [s.label for s in pie_chart.sections]
        chart_data.add_series("", tuple(s.value for s in pie_chart.sections))

        # スライドにネイティブ円グラフフレームを追加する
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
            chart_data,
        )
        chart = chart_frame.chart

        # ---- タイトル設定 ----
        if pie_chart.title:
            chart.has_title = True
            tf = chart.chart_title.text_frame
            tf.text = pie_chart.title
            tf.paragraphs[0].runs[0].font.size = Pt(_TITLE_FONT_SIZE_PT)
            tf.paragraphs[0].runs[0].font.bold = True
        else:
            chart.has_title = False

        # ---- 凡例（判例）設定 ----
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        # ---- データラベル設定 ----
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        # パーセント表示は常時オンにする
        data_labels.show_percentage = True

        # showData=True の場合は実数値も合わせて表示する
        # showData=False の場合はデフォルト OOXML が showVal="1" になりうるため明示的に無効化する
        data_labels.show_value = pie_chart.show_data

        # カテゴリ名（ラベル）をデータラベルに含める
        data_labels.show_category_name = True

        # データラベルの位置を text_position から OOXML の dLblPos に変換して設定する
        self._set_data_label_position(chart, pie_chart.text_position)

        # スライスごとに色を変える（デフォルト動作を明示的に有効化する）
        plot.vary_by_categories = True

    # ------------------------------------------------------------------
    # プライベートユーティリティ
    # ------------------------------------------------------------------

    def _text_position_to_dLblPos(self, text_position: float) -> str:
        """
        Mermaid の textPosition 値（0.0〜1.0）を OOXML の dLblPos 文字列に変換する。

        Parameters
        ----------
        text_position : float
            0.0（中心）〜 1.0（外縁）の実数値。

        Returns
        -------
        str
            OOXML の dLblPos 属性値。
        """
        if text_position >= 1.0:
            return "outEnd"
        for threshold, val in _TEXT_POSITION_MAP:
            if text_position <= threshold:
                return val
        return _TEXT_POSITION_DEFAULT

    def _set_data_label_position(self, chart: object, text_position: float) -> None:
        """
        グラフの全データラベルの位置を OOXML を直接操作して設定する。

        python-pptx の DataLabels クラスは dLblPos を直接設定する API を持たないため、
        lxml で OOXML を直接編集する。

        OOXML スキーマ（CT_DLbls）の要素順序にしたがって c:dLblPos を
        c:showLegendKey の直前に挿入する。順序違反があると PowerPoint が
        ファイル修復を実行するため、必ず正しい位置に挿入する必要がある。

        Parameters
        ----------
        chart : object
            python-pptx の Chart オブジェクト。
        text_position : float
            textPosition 値。
        """
        dLblPos_val = self._text_position_to_dLblPos(text_position)

        # chart.plots[0] の OOXML 要素（<c:pieChart>）から <c:dLbls> を探す
        plot_el = chart.plots[0]._element
        dLbls = plot_el.find(qn("c:dLbls"))
        if dLbls is None:
            return

        # 既存の <c:dLblPos> 要素があれば削除する（正しい位置に再挿入するため）
        existing = dLbls.find(qn("c:dLblPos"))
        if existing is not None:
            dLbls.remove(existing)

        # OOXML スキーマの定義順: ...dLblPos?, showLegendKey, showVal, ...
        # c:showLegendKey の直前に c:dLblPos を挿入する
        show_legend_key_el = dLbls.find(qn("c:showLegendKey"))
        dLblPos_el = lxml_etree.Element(qn("c:dLblPos"))
        dLblPos_el.set("val", dLblPos_val)

        if show_legend_key_el is not None:
            # showLegendKey の index を取得して、その直前に挿入する
            idx = list(dLbls).index(show_legend_key_el)
            dLbls.insert(idx, dLblPos_el)
        else:
            # showLegendKey が見つからない場合は先頭に追加する（フォールバック）
            logger.warning(
                "円グラフのOOXML要素 showLegendKey が見つかりませんでした。dLblPos を先頭に挿入します。"
            )
            dLbls.insert(0, dLblPos_el)
