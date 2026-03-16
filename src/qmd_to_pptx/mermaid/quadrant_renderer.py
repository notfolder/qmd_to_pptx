"""
Mermaid 象限チャート（quadrantChart）レンダラーモジュール。

QuadrantChart データクラスを入力として PowerPoint スライドに象限チャートを描画する。

レイアウト構成（上から下・左から右）:
  ┌─────────────────────────────────────────┐
  │  タイトル帯（全幅 × 高さ 8%）           │
  ├──────┬──────────────────────────────────┤
  │      │  チャートエリア（4象限）           │
  │ Y軸  │  ┌────────────┬───────────────┐  │
  │ ラ   │  │ Quadrant-2 │  Quadrant-1   │  │
  │ ベ   │  │（左上・薄青）│（右上・薄紫） │  │
  │ ル   ├  ├────────────┼───────────────┤  │
  │ 列   │  │ Quadrant-3 │  Quadrant-4   │  │
  │ 7%   │  │（左下・薄黄）│（右下・薄緑） │  │
  │      │  └────────────┴───────────────┘  │
  ├──────┴──────────────────────────────────┤
  │  X軸ラベル帯（チャート幅 × 高さ 9%）   │
  └─────────────────────────────────────────┘

象限番号とMermaid公式の対応:
  quadrant-1 = 右上、quadrant-2 = 左上、
  quadrant-3 = 左下、quadrant-4 = 右下

ポイントスタイル適用優先順位（高い → 低い）:
  1. インライン直接指定
  2. classDef クラス定義
  3. デフォルト値

OOXML 制約と代替実装:
  - Y軸ラベル縦書き → テキストボックスを rotation=270° で水平回転
  - 軸線 → add_connector(MSO_CONNECTOR.STRAIGHT, ...) で実線コネクタ
  - ポイント → OVAL シェープ（サイズは radius × 2 から EMU 変換）
  - 象限背景 → 単色矩形（グラデーションは未対応）
  - ポイントラベル重なり → 回避なし（重なりはそのまま表示）
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .quadrant_parser import PointStyle, QuadrantChart, QuadrantPoint

# ---------------------------------------------------------------------------
# デフォルトカラーパレット（Mermaid 公式テーマ変数に準拠した固定色）
# ---------------------------------------------------------------------------

# 象限背景色 (Mermaid quadrant1Fill 〜 quadrant4Fill に対応)
_QUADRANT_FILL: dict[int, tuple[int, int, int]] = {
    1: (236, 234, 255),  # 右上: 薄紫 (#ECEAFF)
    2: (208, 228, 255),  # 左上: 薄青 (#D0E4FF)
    3: (255, 244, 208),  # 左下: 薄黄 (#FFF4D0)
    4: (208, 255, 228),  # 右下: 薄緑 (#D0FFE4)
}

# 象限ラベルテキスト色 (quadrant1TextFill 〜 quadrant4TextFill)
_QUADRANT_TEXT: dict[int, tuple[int, int, int]] = {
    1: (80, 60, 160),    # 右上: 紫系
    2: (40, 80, 160),    # 左上: 青系
    3: (120, 100, 30),   # 左下: 黄茶系
    4: (30, 120, 70),    # 右下: 緑系
}

# ポイントデフォルト塗りつぶし色 (quadrantPointFill)
_POINT_FILL_DEFAULT: tuple[int, int, int] = (128, 128, 128)

# ポイントデフォルト枠線色
_POINT_STROKE_DEFAULT: tuple[int, int, int] = (80, 80, 80)

# 軸線（内側）色 (quadrantInternalBorderStrokeFill)
_AXIS_LINE_RGB: tuple[int, int, int] = (160, 160, 160)

# 外枠線色 (quadrantExternalBorderStrokeFill)
_OUTER_BORDER_RGB: tuple[int, int, int] = (100, 100, 100)

# タイトルテキスト色 (quadrantTitleFill)
_TITLE_TEXT_RGB: tuple[int, int, int] = (30, 30, 60)

# 軸ラベルテキスト色 (quadrantXAxisTextFill / quadrantYAxisTextFill)
_AXIS_LABEL_RGB: tuple[int, int, int] = (80, 80, 100)

# ポイントラベルテキスト色 (quadrantPointTextFill)
_POINT_LABEL_RGB: tuple[int, int, int] = (50, 50, 50)

# ---------------------------------------------------------------------------
# レイアウト比率定数
# ---------------------------------------------------------------------------

# タイトル帯の高さ比率（全体高さに対する割合）
_TITLE_H_RATIO: float = 0.08

# X軸ラベル帯の高さ比率（全体高さに対する割合）
_X_AXIS_H_RATIO: float = 0.09

# Y軸ラベル列の幅比率（全体幅に対する割合）
_Y_AXIS_W_RATIO: float = 0.07

# ---------------------------------------------------------------------------
# ポイント描画定数
# ---------------------------------------------------------------------------

# ポイントのデフォルト半径 (px 相当)
_POINT_RADIUS_DEFAULT_PX: int = 5

# 1px を EMU に換算した値（9525）
_PX_TO_EMU: int = 9525

# ポイントラベルとポイント下端のオフセット（EMU）
_LABEL_OFFSET_EMU: int = 25_000

# ポイントラベルの高さ（EMU）
_LABEL_HEIGHT_EMU: int = 200_000

# ポイントラベルの幅（EMU）
_LABEL_WIDTH_EMU: int = 800_000

# ---------------------------------------------------------------------------
# シェープ型番（MSO_AUTO_SHAPE_TYPE 値）
# ---------------------------------------------------------------------------
_SHAPE_RECT: int = 1      # 矩形
_SHAPE_OVAL: int = 9      # 楕円（円）


# ---------------------------------------------------------------------------
# ユーティリティ関数
# ---------------------------------------------------------------------------


def _parse_color_rgb(hex_color: str | None, default: tuple[int, int, int]) -> tuple[int, int, int]:
    """
    "#rrggbb" 形式の色文字列を (R, G, B) タプルに変換する。

    Parameters
    ----------
    hex_color : str | None
        "#rrggbb" 形式の文字列。None の場合は default を返す。
    default : tuple[int, int, int]
        変換失敗時のデフォルト色。

    Returns
    -------
    tuple[int, int, int]
        (R, G, B) タプル（各要素 0〜255）。
    """
    if not hex_color:
        return default
    raw = hex_color.lstrip("#")
    try:
        r = int(raw[0:2], 16)
        g = int(raw[2:4], 16)
        b = int(raw[4:6], 16)
        return (r, g, b)
    except (ValueError, IndexError):
        return default


def _resolve_style(
    point: QuadrantPoint,
    class_defs: dict[str, PointStyle],
) -> tuple[tuple[int, int, int], int, int, tuple[int, int, int]]:
    """
    ポイントのスタイルをインライン → classDef → デフォルトの優先順で解決する。

    Parameters
    ----------
    point : QuadrantPoint
        解決対象のポイントデータ。
    class_defs : dict[str, PointStyle]
        classDef 定義辞書。

    Returns
    -------
    tuple[tuple[int, int, int], int, int, tuple[int, int, int]]
        (fill_rgb, radius_px, stroke_width_px, stroke_rgb) のタプル。
    """
    # classDef が適用されている場合は取得する（なければデフォルトの空スタイル）
    cls_style = PointStyle()
    if point.class_name and point.class_name in class_defs:
        cls_style = class_defs[point.class_name]

    inline = point.inline_style

    # 塗りつぶし色（インライン → classDef → デフォルト）
    if inline.color is not None:
        fill_rgb = _parse_color_rgb(inline.color, _POINT_FILL_DEFAULT)
    elif cls_style.color is not None:
        fill_rgb = _parse_color_rgb(cls_style.color, _POINT_FILL_DEFAULT)
    else:
        fill_rgb = _POINT_FILL_DEFAULT

    # 半径（インライン → classDef → デフォルト）
    if inline.radius is not None:
        radius_px = inline.radius
    elif cls_style.radius is not None:
        radius_px = cls_style.radius
    else:
        radius_px = _POINT_RADIUS_DEFAULT_PX

    # 枠線幅（インライン → classDef → デフォルト）
    if inline.stroke_width is not None:
        stroke_w_px = inline.stroke_width
    elif cls_style.stroke_width is not None:
        stroke_w_px = cls_style.stroke_width
    else:
        stroke_w_px = 1

    # 枠線色（インライン → classDef → デフォルト）
    if inline.stroke_color is not None:
        stroke_rgb = _parse_color_rgb(inline.stroke_color, _POINT_STROKE_DEFAULT)
    elif cls_style.stroke_color is not None:
        stroke_rgb = _parse_color_rgb(cls_style.stroke_color, _POINT_STROKE_DEFAULT)
    else:
        stroke_rgb = _POINT_STROKE_DEFAULT

    return fill_rgb, radius_px, stroke_w_px, stroke_rgb


# ---------------------------------------------------------------------------
# レンダラークラス
# ---------------------------------------------------------------------------


class QuadrantRenderer:
    """
    QuadrantChart データクラスを受け取り、PowerPoint スライドに
    象限チャートを描画するクラス。

    4象限の背景矩形・軸線・象限ラベル・データポイント・各種軸ラベル・
    タイトルを python-pptx のシェープとして配置する。
    """

    def render(
        self,
        slide: Slide,
        chart: QuadrantChart,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        象限チャートをスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        chart : QuadrantChart
            parse_quadrant() で生成された象限チャートデータ。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        # ------------------------------------------------------------------
        # ゾーンサイズを計算する
        # ------------------------------------------------------------------

        # タイトル帯の高さ
        title_h = max(200_000, int(height * _TITLE_H_RATIO))

        # X軸ラベル帯の高さ
        x_axis_h = max(150_000, int(height * _X_AXIS_H_RATIO))

        # Y軸ラベル列の幅
        y_axis_w = max(150_000, int(width * _Y_AXIS_W_RATIO))

        # チャートエリアの位置とサイズ
        chart_left = left + y_axis_w
        chart_top = top + title_h
        chart_width = width - y_axis_w
        chart_height = height - title_h - x_axis_h

        # 中央軸の EMU 座標（チャートエリア内の 50% 位置）
        mid_x = chart_left + chart_width // 2
        mid_y = chart_top + chart_height // 2

        # ------------------------------------------------------------------
        # タイトルを描画する
        # ------------------------------------------------------------------
        if chart.title:
            self._draw_title(
                slide, chart.title,
                left, top, width, title_h,
            )

        # ------------------------------------------------------------------
        # 4象限の背景矩形を描画する
        # ------------------------------------------------------------------
        # 象限2（左上）
        self._draw_quadrant_bg(
            slide, 2,
            chart_left, chart_top,
            chart_width // 2, chart_height // 2,
        )
        # 象限1（右上）
        self._draw_quadrant_bg(
            slide, 1,
            mid_x, chart_top,
            chart_width - chart_width // 2, chart_height // 2,
        )
        # 象限3（左下）
        self._draw_quadrant_bg(
            slide, 3,
            chart_left, mid_y,
            chart_width // 2, chart_height - chart_height // 2,
        )
        # 象限4（右下）
        self._draw_quadrant_bg(
            slide, 4,
            mid_x, mid_y,
            chart_width - chart_width // 2, chart_height - chart_height // 2,
        )

        # ------------------------------------------------------------------
        # 外枠線を描画する（チャートエリア外周）
        # ------------------------------------------------------------------
        self._draw_border(
            slide,
            chart_left, chart_top, chart_width, chart_height,
        )

        # ------------------------------------------------------------------
        # 軸線（中央交差線）を描画する
        # ------------------------------------------------------------------
        # 水平軸線（X軸境界: y = mid_y）
        self._draw_axis_line(
            slide,
            chart_left, mid_y,
            chart_left + chart_width, mid_y,
        )
        # 垂直軸線（Y軸境界: x = mid_x）
        self._draw_axis_line(
            slide,
            mid_x, chart_top,
            mid_x, chart_top + chart_height,
        )

        # ------------------------------------------------------------------
        # 象限ラベルを描画する（各象限の内部上端）
        # ------------------------------------------------------------------
        # 象限2（左上）
        if 2 in chart.quadrant_labels:
            self._draw_quadrant_label(
                slide, 2, chart.quadrant_labels[2],
                chart_left, chart_top,
                chart_width // 2, chart_height // 2,
            )
        # 象限1（右上）
        if 1 in chart.quadrant_labels:
            self._draw_quadrant_label(
                slide, 1, chart.quadrant_labels[1],
                mid_x, chart_top,
                chart_width - chart_width // 2, chart_height // 2,
            )
        # 象限3（左下）
        if 3 in chart.quadrant_labels:
            self._draw_quadrant_label(
                slide, 3, chart.quadrant_labels[3],
                chart_left, mid_y,
                chart_width // 2, chart_height - chart_height // 2,
            )
        # 象限4（右下）
        if 4 in chart.quadrant_labels:
            self._draw_quadrant_label(
                slide, 4, chart.quadrant_labels[4],
                mid_x, mid_y,
                chart_width - chart_width // 2, chart_height - chart_height // 2,
            )

        # ------------------------------------------------------------------
        # X軸ラベルを描画する（チャートエリア下部）
        # ------------------------------------------------------------------
        x_label_top = chart_top + chart_height
        self._draw_x_axis_labels(
            slide,
            chart.x_label_left, chart.x_label_right,
            chart_left, x_label_top, chart_width, x_axis_h,
        )

        # ------------------------------------------------------------------
        # Y軸ラベルを描画する（チャートエリア左部、縦書き）
        # ------------------------------------------------------------------
        self._draw_y_axis_labels(
            slide,
            chart.y_label_bottom, chart.y_label_top,
            left, chart_top, y_axis_w, chart_height,
        )

        # ------------------------------------------------------------------
        # データポイントを描画する
        # ------------------------------------------------------------------
        for point in chart.points:
            fill_rgb, radius_px, stroke_w_px, stroke_rgb = _resolve_style(
                point, chart.class_defs
            )
            # ポイント直径（EMU）
            diameter_emu = radius_px * _PX_TO_EMU * 2

            # チャートエリア内の EMU 座標に変換する
            # x: 0.0 = 左端、1.0 = 右端
            # y: 0.0 = 下端、1.0 = 上端（画面座標は上が 0 なので反転）
            px_center = chart_left + int(point.x * chart_width)
            py_center = chart_top + int((1.0 - point.y) * chart_height)

            self._draw_point(
                slide,
                px_center, py_center,
                diameter_emu,
                fill_rgb, stroke_w_px, stroke_rgb,
            )
            self._draw_point_label(
                slide,
                point.name,
                px_center, py_center,
                diameter_emu,
            )

    # -------------------------------------------------------------------------
    # 描画ヘルパーメソッド
    # -------------------------------------------------------------------------

    def _draw_title(
        self,
        slide: Slide,
        title: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """タイトルテキストボックスを描画する（中央揃え・太字）。"""
        tb = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*_TITLE_TEXT_RGB)

    def _draw_quadrant_bg(
        self,
        slide: Slide,
        quadrant_num: int,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        象限の背景矩形を描画する。

        Parameters
        ----------
        quadrant_num : int
            象限番号（1〜4）。
        """
        fill_color = _QUADRANT_FILL.get(quadrant_num, (240, 240, 240))
        shape = slide.shapes.add_shape(
            _SHAPE_RECT,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        self._fill_shape(shape, fill_color)
        # 枠線なし（軸線を別途描画するため）
        shape.line.fill.background()

    def _draw_border(
        self,
        slide: Slide,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        チャートエリアの外枠線を矩形シェープで描画する。

        塗りつぶしなしの矩形に枠線のみ設定することで外枠を表現する。
        """
        shape = slide.shapes.add_shape(
            _SHAPE_RECT,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        # 塗りつぶしなし（透明）
        shape.fill.background()
        shape.line.color.rgb = RGBColor(*_OUTER_BORDER_RGB)
        shape.line.width = Pt(1.5)

    def _draw_axis_line(
        self,
        slide: Slide,
        x1: int,
        y1: int,
        x2: int,
        y2: int,
    ) -> None:
        """
        直線コネクタで軸線を描画する。

        Parameters
        ----------
        x1, y1 : int
            始点の EMU 座標。
        x2, y2 : int
            終点の EMU 座標。
        """
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(x1), Emu(y1), Emu(x2), Emu(y2),
        )
        connector.line.color.rgb = RGBColor(*_AXIS_LINE_RGB)
        connector.line.width = Pt(1.0)

    def _draw_quadrant_label(
        self,
        slide: Slide,
        quadrant_num: int,
        label_text: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        象限ラベルを各象限の内部上端に描画する。

        Parameters
        ----------
        quadrant_num : int
            象限番号（1〜4）。スタイルの選択に使用する。
        label_text : str
            表示するラベルテキスト。
        """
        text_rgb = _QUADRANT_TEXT.get(quadrant_num, (60, 60, 60))
        # 内側余白（上・左）
        pad = 50_000
        label_top = top + pad
        label_left = left + pad
        label_width = max(10_000, width - pad * 2)
        label_height = max(10_000, min(300_000, height // 4))

        tb = slide.shapes.add_textbox(
            Emu(label_left), Emu(label_top),
            Emu(label_width), Emu(label_height),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        # 奇数象限（1,3）は右上・右下なので中央揃え、偶数（2,4）は左上・右下なので左端
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label_text
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*text_rgb)

    def _draw_x_axis_labels(
        self,
        slide: Slide,
        left_text: str,
        right_text: str,
        chart_left: int,
        label_top: int,
        chart_width: int,
        label_height: int,
    ) -> None:
        """
        X軸ラベルを描画する（左端・右端テキストボックス）。

        左ラベルは左揃え、右ラベルは右揃えで配置する。
        """
        label_w = chart_width // 2

        # 左ラベル（低い側）
        if left_text:
            tb = slide.shapes.add_textbox(
                Emu(chart_left), Emu(label_top),
                Emu(label_w), Emu(label_height),
            )
            tf = tb.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = left_text
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*_AXIS_LABEL_RGB)

        # 右ラベル（高い側）
        if right_text:
            tb = slide.shapes.add_textbox(
                Emu(chart_left + chart_width - label_w), Emu(label_top),
                Emu(label_w), Emu(label_height),
            )
            tf = tb.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = right_text
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*_AXIS_LABEL_RGB)

    def _draw_y_axis_labels(
        self,
        slide: Slide,
        bottom_text: str,
        top_text: str,
        axis_left: int,
        chart_top: int,
        axis_width: int,
        chart_height: int,
    ) -> None:
        """
        Y軸ラベルを描画する（上端・下端テキストボックス、rotation=270° で縦向き）。

        OOXML 代替: テキストボックスを 270° 回転させて縦向きに表現する。
        上方向が「高い側」（top_text）、下方向が「低い側」（bottom_text）。
        """
        # テキストボックスのサイズ（回転前の寸法）
        label_w = max(10_000, chart_height // 2)   # 回転後に高さ方向の長さになる
        label_h = max(10_000, axis_width - 20_000)  # 回転後に幅方向の長さになる

        def _add_rotated_label(text: str, cx: int, cy: int, align: PP_ALIGN) -> None:
            """回転テキストボックスを中心座標 (cx, cy) に配置する。"""
            # テキストボックスの左上座標は中心から label_w/2, label_h/2 を引いた位置
            tb_left = cx - label_w // 2
            tb_top = cy - label_h // 2
            tb = slide.shapes.add_textbox(
                Emu(tb_left), Emu(tb_top),
                Emu(label_w), Emu(label_h),
            )
            tb.rotation = 270.0  # 反時計回り270° = 時計回り90°で縦向き

            tf = tb.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = align
            run = para.add_run()
            run.text = text
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*_AXIS_LABEL_RGB)

        # 下ラベル（低い側）: チャート下側 1/4 の中心に配置する
        if bottom_text:
            cy = chart_top + chart_height * 3 // 4  # 下寄り
            cx = axis_left + axis_width // 2
            _add_rotated_label(bottom_text, cx, cy, PP_ALIGN.LEFT)

        # 上ラベル（高い側）: チャート上側 1/4 の中心に配置する
        if top_text:
            cy = chart_top + chart_height // 4  # 上寄り
            cx = axis_left + axis_width // 2
            _add_rotated_label(top_text, cx, cy, PP_ALIGN.LEFT)

    def _draw_point(
        self,
        slide: Slide,
        center_x: int,
        center_y: int,
        diameter_emu: int,
        fill_rgb: tuple[int, int, int],
        stroke_w_px: int,
        stroke_rgb: tuple[int, int, int],
    ) -> None:
        """
        データポイントを OVAL シェープとして描画する。

        Parameters
        ----------
        center_x, center_y : int
            ポイント中心の EMU 座標。
        diameter_emu : int
            ポイント直径（EMU）。
        fill_rgb : tuple[int, int, int]
            塗りつぶし色。
        stroke_w_px : int
            枠線幅（px 相当）。
        stroke_rgb : tuple[int, int, int]
            枠線色。
        """
        pt_left = center_x - diameter_emu // 2
        pt_top = center_y - diameter_emu // 2

        shape = slide.shapes.add_shape(
            _SHAPE_OVAL,
            Emu(pt_left), Emu(pt_top),
            Emu(diameter_emu), Emu(diameter_emu),
        )
        self._fill_shape(shape, fill_rgb)
        shape.line.color.rgb = RGBColor(*stroke_rgb)
        shape.line.width = Emu(stroke_w_px * _PX_TO_EMU)

    def _draw_point_label(
        self,
        slide: Slide,
        name: str,
        center_x: int,
        center_y: int,
        diameter_emu: int,
    ) -> None:
        """
        ポイント名ラベルをポイントの真下に描画する。

        Parameters
        ----------
        name : str
            ポイント名テキスト。
        center_x, center_y : int
            対応するポイント中心の EMU 座標。
        diameter_emu : int
            対応するポイント直径（EMU）。ラベル上端オフセットの計算に使用する。
        """
        # ラベルの左端をポイント中心から左にずらして中央揃えに見せる
        label_left = center_x - _LABEL_WIDTH_EMU // 2
        # ラベルの上端はポイント下端 + 余白
        label_top = center_y + diameter_emu // 2 + _LABEL_OFFSET_EMU

        tb = slide.shapes.add_textbox(
            Emu(label_left), Emu(label_top),
            Emu(_LABEL_WIDTH_EMU), Emu(_LABEL_HEIGHT_EMU),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = name
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(*_POINT_LABEL_RGB)

    # -------------------------------------------------------------------------
    # 静的ユーティリティ
    # -------------------------------------------------------------------------

    @staticmethod
    def _fill_shape(shape: object, rgb: tuple[int, int, int]) -> None:
        """
        シェープの塗りつぶし色を設定する。

        Parameters
        ----------
        shape : object
            python-pptx の Shape オブジェクト。
        rgb : tuple[int, int, int]
            塗りつぶし色の (R, G, B) タプル。
        """
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
