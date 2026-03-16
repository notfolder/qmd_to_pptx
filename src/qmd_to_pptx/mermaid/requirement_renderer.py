"""
Mermaid 要件図（requirementDiagram）レンダラーモジュール。

RequirementDiagram データクラスを入力として PowerPoint スライドに要件図を描画する。

レイアウト:
  - 各ノード（要件・エレメント）をグリッド配置（direction に応じて列数を変える）
  - 各ノードを 2 段ヘッダー（ステレオタイプ・名前）＋ボディの矩形で表現
  - リレーションを直線コネクター＋中点ラベルテキストボックスで表現

ノード構造:
  ┌──────────────────────────────┐
  │ <<Stereotype>>               │ ← ヘッダー上段（ステレオタイプ・italic・小）
  │ NodeName                     │ ← ヘッダー下段（名前・太字）
  ├──────────────────────────────┤
  │ ID: value                    │ ← ボディ行（左寄せ）
  │ Text: 説明文テキスト          │
  │ Risk: High                   │
  │ Verification: Test           │ ← 要件のみ
  │ Type: simulation             │ ← エレメントのみ
  │ Doc Ref: reqs/...            │ ← エレメントのみ（docref がある場合）
  └──────────────────────────────┘

OOXML 制約と代替案:
  - リレーションのラベル ('«traces»' 等): コネクター中点テキストボックス
  - 点線コネクター: OOXML a:ln/a:prstDash で "dash" 設定（直接 XML 操作）
  - ノードタイプ別ヘッダー色: 矩形の fill.solid() で実装
  - Markdown 書式 (**bold** / *italic*): run ごとに書式指定
  - direction LR/RL: グリッドの列優先 vs 行優先で対応

ノードタイプ別ヘッダー色:
  requirement            : #D0E8FF（薄青）
  functionalrequirement  : #E8D0FF（薄紫）
  interfacerequirement   : #D0FFE8（薄緑）
  performancerequirement : #FFE8D0（薄橙）
  physicalrequirement    : #FFFFE0（薄黄）
  designconstraint       : #FFD0D0（薄赤）
  element                : #E0E0E0（薄グレー）
"""

from __future__ import annotations

import math

from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer
from .requirement_parser import (
    ElementNode,
    NodeStyle,
    RequirementDiagram,
    RequirementNode,
    parse_inline_markdown,
    resolve_node_style,
)

# ---------------------------------------------------------------------------
# ノードサイズ定数（EMU）
# ---------------------------------------------------------------------------

_NODE_WIDTH_EMU: int = 2_000_000        # ノード幅
_HEADER_STEREO_H: int = 320_000         # ヘッダー上段（ステレオタイプ）の高さ
_HEADER_NAME_H: int = 380_000           # ヘッダー下段（名前）の高さ
_BODY_ROW_H: int = 290_000              # ボディ行 1 行あたりの高さ
_BODY_MIN_H: int = 290_000              # ボディ最小高さ

# ノード間余白（EMU）
_H_GAP: int = 400_000                   # 水平余白
_V_GAP: int = 300_000                   # 垂直余白

# ---------------------------------------------------------------------------
# ヘッダー色パレット（ノードタイプ → (ヘッダー上段 RGB, ヘッダー下段 RGB, ボディ RGB)）
# ---------------------------------------------------------------------------

_HEADER_COLORS: dict[str, tuple[tuple[int, int, int], tuple[int, int, int], tuple[int, int, int]]] = {
    "requirement": (
        (160, 200, 240),   # ヘッダー上段（濃い青）
        (208, 232, 255),   # ヘッダー下段（薄青）
        (238, 247, 255),   # ボディ（最薄青）
    ),
    "functionalrequirement": (
        (200, 160, 240),
        (232, 208, 255),
        (247, 238, 255),
    ),
    "interfacerequirement": (
        (160, 240, 200),
        (208, 255, 232),
        (238, 255, 247),
    ),
    "performancerequirement": (
        (240, 200, 160),
        (255, 232, 208),
        (255, 247, 238),
    ),
    "physicalrequirement": (
        (230, 230, 160),
        (255, 255, 208),
        (255, 255, 238),
    ),
    "designconstraint": (
        (240, 160, 160),
        (255, 208, 208),
        (255, 238, 238),
    ),
    "element": (
        (180, 180, 180),
        (220, 220, 220),
        (245, 245, 245),
    ),
}

_DEFAULT_HEADER_COLORS = _HEADER_COLORS["requirement"]

# リレーションラベルテキスト色
_REL_LABEL_RGB = (60, 60, 80)

# ヘッダーテキスト色（共通）
_HEADER_TEXT_RGB = (30, 30, 60)

# ボディテキスト色（共通）
_BODY_TEXT_RGB = (40, 40, 40)

# 枠線色
_BORDER_RGB = (100, 110, 130)

# コネクター色
_CONNECTOR_RGB = (80, 90, 110)

# ---------------------------------------------------------------------------
# ユーティリティ関数
# ---------------------------------------------------------------------------


def _parse_hex_color(hex_str: str | None) -> tuple[int, int, int] | None:
    """
    "#rrggbb" または "rrggbb" 形式の色文字列を (R, G, B) タプルに変換する。

    Parameters
    ----------
    hex_str : str | None
        色文字列。None または無効形式の場合は None を返す。

    Returns
    -------
    tuple[int, int, int] | None
        (R, G, B) タプル。変換失敗時は None。
    """
    if not hex_str:
        return None
    raw = hex_str.strip().lstrip("#")
    if len(raw) == 6:
        try:
            r = int(raw[0:2], 16)
            g = int(raw[2:4], 16)
            b = int(raw[4:6], 16)
            return (r, g, b)
        except ValueError:
            pass
    return None


def _apply_fill(shape: object, rgb: tuple[int, int, int]) -> None:
    """シェープの塗りつぶし色を設定する。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)


def _apply_border(shape: object, rgb: tuple[int, int, int], pt: float = 1.0) -> None:
    """シェープの枠線色・幅を設定する。"""
    shape.line.color.rgb = RGBColor(*rgb)
    shape.line.width = Pt(pt)


def _node_body_rows(node: RequirementNode | ElementNode) -> list[str]:
    """
    ノードのボディ行テキストリストを返す。

    Parameters
    ----------
    node : RequirementNode | ElementNode
        描画対象ノード。

    Returns
    -------
    list[str]
        表示する行リスト（フィールド名: 値 の形式）。
    """
    rows: list[str] = []
    if isinstance(node, RequirementNode):
        if node.req_id:
            rows.append(f"ID: {node.req_id}")
        if node.text:
            rows.append(f"Text: {node.text}")
        if node.risk:
            rows.append(f"Risk: {node.risk}")
        if node.verify_method:
            rows.append(f"Verification: {node.verify_method}")
    else:
        if node.elem_type:
            rows.append(f"Type: {node.elem_type}")
        if node.docref:
            rows.append(f"Doc Ref: {node.docref}")
    return rows


def _node_total_height(node: RequirementNode | ElementNode) -> int:
    """ノード全体の EMU 高さを計算する。"""
    body_rows = _node_body_rows(node)
    body_h = max(_BODY_MIN_H, len(body_rows) * _BODY_ROW_H)
    return _HEADER_STEREO_H + _HEADER_NAME_H + body_h


# ---------------------------------------------------------------------------
# レンダラークラス
# ---------------------------------------------------------------------------


class RequirementRenderer(BaseDiagramRenderer):
    """
    RequirementDiagram データクラスを受け取り、PowerPoint スライドに
    要件図を描画するクラス。

    ノードをグリッド配置し、各ノードを 2 段ヘッダー＋ボディの矩形で描画する。
    リレーションを STRAIGHT コネクター＋ラベルテキストボックスで描画する。
    """

    def render(
        self,
        slide: Slide,
        diagram: RequirementDiagram,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        要件図をスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        diagram : RequirementDiagram
            parse_requirement() で生成された要件図データ。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        # 全ノード（要件 + エレメント）の辞書を構築する
        all_nodes: dict[str, RequirementNode | ElementNode] = {}
        all_nodes.update(diagram.requirements)
        all_nodes.update(diagram.elements)

        if not all_nodes:
            return

        node_list = list(all_nodes.keys())
        n = len(node_list)

        # グリッド配置の列数を direction に応じて決定する
        # LR / RL は横方向に広がるため行数を少なく（列数を多く）する
        direction = diagram.direction.upper()
        if direction in ("LR", "RL"):
            n_cols = n  # 全ノードを 1 行に並べる（多い場合は折り返す）
            n_cols = max(1, min(n, math.ceil(math.sqrt(n * 2))))
        else:
            # TB / BT: 縦方向に広がるため列数を少なく
            n_cols = max(1, math.ceil(math.sqrt(n)))
        n_rows = math.ceil(n / n_cols)

        # 最大ノード高さを計算してマージンに使う
        max_node_h = max(_node_total_height(all_nodes[nid]) for nid in node_list)
        col_w = _NODE_WIDTH_EMU + _H_GAP
        row_h = max_node_h + _V_GAP

        # 描画に必要な全幅・全高
        total_w = n_cols * col_w - _H_GAP
        total_h = n_rows * row_h - _V_GAP

        # 描画エリア内でセンタリングする
        x_offset = left + max(0, (width - total_w) // 2)
        y_offset = top + max(0, (height - total_h) // 2)

        # BT は下から上に並べる（インデックスを逆順で列に割り当てる）
        if direction == "BT":
            node_list = node_list[::-1]

        # 各ノードの中心座標（EMU）を計算する
        node_centers: dict[str, tuple[int, int]] = {}
        for i, nid in enumerate(node_list):
            col = i % n_cols
            row = i // n_cols
            cx = x_offset + col * col_w + _NODE_WIDTH_EMU // 2
            cy = y_offset + row * row_h + _node_total_height(all_nodes[nid]) // 2
            node_centers[nid] = (cx, cy)

        # ノードシェープを描画する（コネクター接続用にヘッダー下段シェープを記録）
        node_anchor_shapes: dict[str, object] = {}
        for nid in all_nodes:
            cx, cy = node_centers[nid]
            node = all_nodes[nid]
            eff_style = resolve_node_style(
                node.classes, node.style, diagram.class_defs
            )
            anchor = self._draw_node(slide, node, cx, cy, eff_style)
            node_anchor_shapes[nid] = anchor

        # リレーションを描画する
        for rel in diagram.relations:
            if rel.src not in node_centers or rel.dst not in node_centers:
                continue
            src_cx, src_cy = node_centers[rel.src]
            dst_cx, dst_cy = node_centers[rel.dst]
            src_shape = node_anchor_shapes.get(rel.src)
            dst_shape = node_anchor_shapes.get(rel.dst)
            self._draw_relation(
                slide,
                rel.rel_type, rel.src, rel.dst,
                src_cx, src_cy, dst_cx, dst_cy,
                src_shape, dst_shape,
                all_nodes,
            )

    # -------------------------------------------------------------------------
    # ノード描画メソッド
    # -------------------------------------------------------------------------

    def _draw_node(
        self,
        slide: Slide,
        node: RequirementNode | ElementNode,
        cx: int,
        cy: int,
        eff_style: NodeStyle,
    ) -> object:
        """
        ノードを 2 段ヘッダー + ボディの矩形 3 枚として描画し、
        コネクタ接続用にヘッダー下段シェープを返す。

        Parameters
        ----------
        node : RequirementNode | ElementNode
            描画対象ノード。
        cx, cy : int
            ノード中心の EMU 座標。
        eff_style : NodeStyle
            解決済み有効スタイル。

        Returns
        -------
        object
            ヘッダー下段シェープ（コネクター接続のアンカーとして使用）。
        """
        body_rows = _node_body_rows(node)
        body_h = max(_BODY_MIN_H, len(body_rows) * _BODY_ROW_H)
        total_h = _HEADER_STEREO_H + _HEADER_NAME_H + body_h

        box_left = cx - _NODE_WIDTH_EMU // 2
        box_top = cy - total_h // 2

        # ノードタイプを決定する（エレメントは "element" キーを使う）
        if isinstance(node, RequirementNode):
            type_key = node.req_type
            stereotype_text = f"<<{node.stereotype}>>"
            name_text = node.name
        else:
            type_key = "element"
            stereotype_text = "<<Element>>"
            name_text = node.name

        # カラーパレットを決定する（eff_style の fill が指定されていれば上書き）
        colors = _HEADER_COLORS.get(type_key, _DEFAULT_HEADER_COLORS)
        hdr_upper_rgb = _parse_hex_color(eff_style.fill) or colors[0]
        hdr_lower_rgb = colors[1]
        body_rgb = colors[2]
        text_rgb = _parse_hex_color(eff_style.color) or _BODY_TEXT_RGB
        border_rgb = _parse_hex_color(eff_style.stroke) or _BORDER_RGB

        # ----------------------------------------------------------------
        # ヘッダー上段（ステレオタイプ）
        # ----------------------------------------------------------------
        stereo_top = box_top
        stereo_shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Emu(box_left), Emu(stereo_top),
            Emu(_NODE_WIDTH_EMU), Emu(_HEADER_STEREO_H),
        )
        _apply_fill(stereo_shape, hdr_upper_rgb)
        _apply_border(stereo_shape, border_rgb)
        tf = stereo_shape.text_frame
        tf.margin_left = Emu(60_000)
        tf.margin_right = Emu(60_000)
        tf.margin_top = Emu(30_000)
        tf.margin_bottom = Emu(30_000)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = stereotype_text
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*_HEADER_TEXT_RGB)

        # ----------------------------------------------------------------
        # ヘッダー下段（名前）← コネクターのアンカーとして使用
        # ----------------------------------------------------------------
        name_top = stereo_top + _HEADER_STEREO_H
        name_shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Emu(box_left), Emu(name_top),
            Emu(_NODE_WIDTH_EMU), Emu(_HEADER_NAME_H),
        )
        _apply_fill(name_shape, hdr_lower_rgb)
        _apply_border(name_shape, border_rgb)
        tf = name_shape.text_frame
        tf.margin_left = Emu(60_000)
        tf.margin_right = Emu(60_000)
        tf.margin_top = Emu(40_000)
        tf.margin_bottom = Emu(40_000)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        bold = (eff_style.font_weight == "bold") if eff_style.font_weight else True
        # 名前の Markdown 書式を適用する
        for seg_text, seg_bold, seg_italic in parse_inline_markdown(name_text):
            run = para.add_run()
            run.text = seg_text
            run.font.size = Pt(11)
            run.font.bold = seg_bold or bold
            run.font.italic = seg_italic
            run.font.color.rgb = RGBColor(*_HEADER_TEXT_RGB)

        # ----------------------------------------------------------------
        # ボディ（フィールド行）
        # ----------------------------------------------------------------
        body_top = name_top + _HEADER_NAME_H
        body_shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Emu(box_left), Emu(body_top),
            Emu(_NODE_WIDTH_EMU), Emu(body_h),
        )
        _apply_fill(body_shape, body_rgb)
        _apply_border(body_shape, border_rgb)
        tf = body_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(80_000)
        tf.margin_right = Emu(60_000)
        tf.margin_top = Emu(40_000)
        tf.margin_bottom = Emu(40_000)

        first = True
        for row_text in body_rows:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            # "Text: " フィールドのみ Markdown 書式を適用する
            if row_text.startswith("Text: "):
                prefix = "Text: "
                content = row_text[len(prefix):]
                prefix_run = para.add_run()
                prefix_run.text = prefix
                prefix_run.font.size = Pt(9)
                prefix_run.font.bold = False
                prefix_run.font.color.rgb = RGBColor(*text_rgb)
                for seg_text, seg_bold, seg_italic in parse_inline_markdown(content):
                    seg_run = para.add_run()
                    seg_run.text = seg_text
                    seg_run.font.size = Pt(9)
                    seg_run.font.bold = seg_bold
                    seg_run.font.italic = seg_italic
                    seg_run.font.color.rgb = RGBColor(*text_rgb)
            else:
                run = para.add_run()
                run.text = row_text
                run.font.size = Pt(9)
                run.font.bold = False
                run.font.color.rgb = RGBColor(*text_rgb)

        # 3 つのシェープをグループ化する
        slide.shapes.add_group_shape([stereo_shape, name_shape, body_shape])

        # ヘッダー下段シェープをコネクターのアンカーとして返す
        return name_shape

    # -------------------------------------------------------------------------
    # リレーション描画メソッド
    # -------------------------------------------------------------------------

    def _draw_relation(
        self,
        slide: Slide,
        rel_type: str,
        src_name: str,
        dst_name: str,
        src_cx: int,
        src_cy: int,
        dst_cx: int,
        dst_cy: int,
        src_shape: object | None,
        dst_shape: object | None,
        all_nodes: dict[str, RequirementNode | ElementNode],
    ) -> None:
        """
        リレーションを直線コネクター＋ラベルテキストボックスで描画する。

        コネクターの始点・終点はソース・デスティネーションノードの中心座標とし、
        begin_connect / end_connect でヘッダー下段シェープに接続する。
        ラベル「«relType»」はコネクター中点に配置するテキストボックスで表現する。

        Parameters
        ----------
        rel_type : str
            リレーションタイプ（小文字: contains / copies / ...）。
        src_cx, src_cy : int
            始点ノード中心の EMU 座標。
        dst_cx, dst_cy : int
            終点ノード中心の EMU 座標。
        src_shape, dst_shape : object | None
            コネクター接続先のシェープ（ヘッダー下段）。
        all_nodes : dict[str, RequirementNode | ElementNode]
            ノード辞書（高さ計算に使用）。
        """
        # 接続ポイントをノードのヘッダー下段中心で計算する
        src_node = all_nodes.get(src_name)
        dst_node = all_nodes.get(dst_name)
        src_h = _node_total_height(src_node) if src_node else 0
        dst_h = _node_total_height(dst_node) if dst_node else 0

        # コネクターの始終点として各ノード中心を使う
        # （高さオフセットでヘッダー下段の Y 座標を近似する）
        sx = src_cx
        sy = src_cy - src_h // 2 + _HEADER_STEREO_H + _HEADER_NAME_H // 2
        dx = dst_cx
        dy = dst_cy - dst_h // 2 + _HEADER_STEREO_H + _HEADER_NAME_H // 2

        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(sx), Emu(sy),
            Emu(dx), Emu(dy),
        )
        connector.line.color.rgb = RGBColor(*_CONNECTOR_RGB)
        connector.line.width = Pt(1.0)

        # 矢印先端を設定する（OOXML 直接操作）
        self._apply_arrow_end(connector)

        # 始点・終点シェープに接続する
        if src_shape is not None:
            src_cp, dst_cp = self._connection_indices(dx - sx, dy - sy)
            connector.begin_connect(src_shape, src_cp)
        if dst_shape is not None:
            _, dst_cp = self._connection_indices(dx - sx, dy - sy)
            connector.end_connect(dst_shape, dst_cp)

        # ラベルテキストボックスを中点に配置する
        mid_x = (sx + dx) // 2
        mid_y = (sy + dy) // 2
        label = f"\u00ab{rel_type}\u00bb"  # 《labelType》
        label_w = 800_000
        label_h = 250_000
        tb = slide.shapes.add_textbox(
            Emu(mid_x - label_w // 2), Emu(mid_y - label_h // 2),
            Emu(label_w), Emu(label_h),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*_REL_LABEL_RGB)

        # テキストボックスの背景・枠線を透明にする
        tb.fill.background()
        spPr_el = tb._element.find(qn("p:spPr"))
        if spPr_el is not None:
            existing_ln = spPr_el.find(qn("a:ln"))
            if existing_ln is None:
                existing_ln = lxml_etree.SubElement(spPr_el, qn("a:ln"))
            if existing_ln.find(qn("a:noFill")) is None:
                lxml_etree.SubElement(existing_ln, qn("a:noFill"))

    # -------------------------------------------------------------------------
    # OOXML ユーティリティ
    # -------------------------------------------------------------------------

    @staticmethod
    def _apply_arrow_end(connector: object) -> None:
        """
        コネクターの終点に矢印を設定する（OOXML a:tailEnd / a:headEnd 直接操作）。

        python-pptx の公開 API には矢印設定がないため、
        OOXML を直接操作して <a:headEnd type="arrow"> を設定する。

        Parameters
        ----------
        connector : object
            python-pptx の Connector オブジェクト。
        """
        cxn_el = connector._element
        spPr = cxn_el.find(qn("p:spPr"))
        if spPr is None:
            return
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr, qn("a:ln"))

        # 既存の headEnd を除去してから再設定する
        for tag in [qn("a:headEnd"), qn("a:tailEnd")]:
            existing = ln.find(tag)
            if existing is not None:
                ln.remove(existing)

        # 終点（dst 側）に矢印を設定する
        head_end = lxml_etree.SubElement(ln, qn("a:headEnd"))
        head_end.set("type", "arrow")
        head_end.set("w", "med")
        head_end.set("len", "med")

        # 始点は矢印なし（flat）
        tail_end = lxml_etree.SubElement(ln, qn("a:tailEnd"))
        tail_end.set("type", "none")
