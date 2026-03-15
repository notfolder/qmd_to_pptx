"""
Mermaidシーケンス図レンダラーモジュール。

sequenceDiagram 構文を mermaid-parser-py でパースした結果（graph_type="sequence"）を入力とし、
参加者ボックス・ライフライン・メッセージ矢印・アクティベーションバー・ノート・
フレームボックス（loop/alt/opt/par/critical/break）・背景（rect）を
PPTXスライドに描画する。

OOXML制限による代替仕様:
- actor人型シンボル: OOXMLプリセットなし → 角丸矩形で代替
- "-x"/"-)" クロス・非同期矢印: OOXML headEnd に cross なし → ラベルに「✕」付加 / "open" 矢印
- Noteの折り目（dog-ear）: 複雑なパス不要 → 単純黄色矩形
- 自己メッセージ: ELBOWコネクターでL字型に近似
- フレームタブ: 五角形困難 → 左上に小矩形タブを別途配置
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import Optional

from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer

# ---------------------------------------------------------------------------
# レイアウト定数（EMU単位）
# ---------------------------------------------------------------------------
_PARTICIPANT_H: int = 500_000       # 参加者ボックス高さ
_COL_MIN_STEP: int = 1_800_000      # 列間最小ピッチ
_MSG_STEP: int = 600_000            # メッセージ間垂直ピッチ
_LIFELINE_TOP_MARGIN: int = 80_000  # 参加者ボックス下端〜ライフライン開始
_NOTE_H: int = 380_000              # ノートボックス高さ
_FRAME_TAB_H: int = 280_000         # フレームタブ高さ（ヘッダー用スペース）
_ACTIVATION_W: int = 140_000        # アクティベーションバー幅
_SELF_MSG_INDENT: int = 800_000     # 自己メッセージの右折れ幅

# ---------------------------------------------------------------------------
# LINETYPEマッピング（mermaid-parser-pyのtype値）
# ---------------------------------------------------------------------------
_LT_SOLID: int = 0           # 実線・塗り矢印 (->>)
_LT_DOTTED: int = 1          # 破線・塗り矢印 (-->>)
_LT_NOTE: int = 2             # ノート
_LT_SOLID_CROSS: int = 3     # 実線・クロス端 (-x)
_LT_DOTTED_CROSS: int = 4    # 破線・クロス端 (--x)
_LT_SOLID_OPEN: int = 5      # 実線・開き矢印 (-) async
_LT_DOTTED_OPEN: int = 6     # 破線・開き矢印 (--) async
_LT_LOOP_START: int = 10     # loop 開始
_LT_LOOP_END: int = 11       # loop 終了
_LT_ALT_START: int = 12      # alt 開始
_LT_ALT_ELSE: int = 13       # else 区切り
_LT_ALT_END: int = 14        # alt 終了
_LT_OPT_START: int = 15      # opt 開始
_LT_OPT_END: int = 16        # opt 終了
_LT_ACTIVE_START: int = 17   # activate 開始
_LT_ACTIVE_END: int = 18     # activate 終了
_LT_PAR_START: int = 19      # par 開始
_LT_PAR_AND: int = 20        # and 区切り
_LT_PAR_END: int = 21        # par 終了
_LT_RECT_START: int = 22     # rect 背景開始
_LT_RECT_END: int = 23       # rect 背景終了
_LT_SOLID_POINT: int = 24    # 実線・ポイント矢印 (->) ←矢印なし寄り
_LT_DOTTED_POINT: int = 25   # 破線・ポイント矢印 (-->)
_LT_AUTONUMBER: int = 26     # autonumber 有効化
_LT_CRITICAL_START: int = 27 # critical 開始
_LT_CRITICAL_OPT: int = 28   # option 区切り
_LT_CRITICAL_END: int = 29   # critical 終了
_LT_BREAK_START: int = 30    # break 開始
_LT_BREAK_END: int = 31      # break 終了
_LT_PAR_OVER_START: int = 32 # par over 開始
_LT_BIDIR_SOLID: int = 33    # 双方向実線 (<<->>)
_LT_BIDIR_DOTTED: int = 34   # 双方向破線 (<<-->>)

# 実際のメッセージ矢印型セット
_MSG_TYPES: frozenset[int] = frozenset({
    _LT_SOLID, _LT_DOTTED, _LT_SOLID_CROSS, _LT_DOTTED_CROSS,
    _LT_SOLID_OPEN, _LT_DOTTED_OPEN, _LT_SOLID_POINT, _LT_DOTTED_POINT,
    _LT_BIDIR_SOLID, _LT_BIDIR_DOTTED,
})

# フレーム開始型セット
_FRAME_START_TYPES: frozenset[int] = frozenset({
    _LT_LOOP_START, _LT_ALT_START, _LT_OPT_START,
    _LT_PAR_START, _LT_CRITICAL_START, _LT_BREAK_START,
})

# フレーム終了型セット
_FRAME_END_TYPES: frozenset[int] = frozenset({
    _LT_LOOP_END, _LT_ALT_END, _LT_OPT_END,
    _LT_PAR_END, _LT_CRITICAL_END, _LT_BREAK_END,
})

# フレームセクション区切り型セット（else/and/option）
_FRAME_SECTION_TYPES: frozenset[int] = frozenset({
    _LT_ALT_ELSE, _LT_PAR_AND, _LT_CRITICAL_OPT,
})

# フレーム開始型 → キーワードラベル
_FRAME_KEYWORD: dict[int, str] = {
    _LT_LOOP_START: "loop",
    _LT_ALT_START: "alt",
    _LT_OPT_START: "opt",
    _LT_PAR_START: "par",
    _LT_CRITICAL_START: "critical",
    _LT_BREAK_START: "break",
}


# ---------------------------------------------------------------------------
# 内部データクラス
# ---------------------------------------------------------------------------

@dataclass
class _ActorInfo:
    """参加者（participant/actor）の情報。"""

    actor_id: str   # Mermaid上のID
    label: str      # 表示ラベル（description または actor_id）
    kind: str       # "participant" または "actor"
    col_idx: int    # 列インデックス（0始まり）


@dataclass
class _ActivationSpan:
    """アクティベーションバーのスパン情報。"""

    actor_id: str   # 対象参加者ID
    start_y: int    # 開始Y座標（EMU）
    end_y: int      # 終了Y座標（EMU）
    depth: int      # スタック深度（0始まり）


@dataclass
class _FrameSpan:
    """フレームボックス（loop/alt/opt/par/critical/break）のスパン情報。"""

    kind: int                           # LINETYPEの開始型
    label: str                          # フレームラベル
    start_y: int                        # 開始Y座標（EMU）
    end_y: int                          # 終了Y座標（EMU）
    sections: list[tuple[int, str]] = field(default_factory=list)  # (Y座標, セクションラベル)のリスト


# ---------------------------------------------------------------------------
# レンダラー本体
# ---------------------------------------------------------------------------

class SequenceDiagramRenderer(BaseDiagramRenderer):
    """
    Mermaidシーケンス図をPPTXスライドに描画するレンダラー。

    mermaid-parser-pyが返すgraph_data（graph_type="sequence"）を入力とし、
    以下の要素を描画する:
    - 参加者ボックス（participant=矩形, actor=角丸矩形）
    - ライフライン（垂直破線）
    - メッセージ矢印（10種類 + ラベルテキストボックス）
    - アクティベーションバー（activate/deactivate）
    - Noteボックス（left of / right of / over）
    - フレームボックス（loop / alt / opt / par / critical / break）
    - 背景矩形（rect rgb(...)）
    - 参加者グループBox
    """

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        シーケンス図をスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            mermaid-parser-pyのgraph_data辞書（graph_type="sequence"時の値）。
        left : int
            描画エリア左端X座標（EMU）。
        top : int
            描画エリア上端Y座標（EMU）。
        width : int
            描画エリア幅（EMU）。
        height : int
            描画エリア高さ（EMU）。
        """
        records: dict = graph_data.get("state", {}).get("records", {})
        actors_raw: dict = records.get("actors", {})
        messages: list[dict] = records.get("messages", [])
        boxes_raw: list = records.get("boxes", [])

        if not actors_raw:
            self._render_fallback(slide, "sequenceDiagram", left, top, width, height)
            return

        # 参加者を表示順に並べる
        actors = self._build_actor_order(actors_raw)
        n = len(actors)
        if n == 0:
            return

        # 各列の中央X座標を計算する
        col_step = max(_COL_MIN_STEP, width // n)
        participant_w = min(col_step - 200_000, 1_200_000)
        cx_map: dict[str, int] = {
            a.actor_id: left + col_step * a.col_idx + col_step // 2
            for a in actors
        }

        # autonumber フラグを確認する
        autonumber: bool = any(m.get("type") == _LT_AUTONUMBER for m in messages)

        # ライフライン開始Y座標を計算する
        lifeline_top_y = top + _PARTICIPANT_H + _LIFELINE_TOP_MARGIN

        # メッセージを走査してイベントY座標・アクティベーション・フレーム・背景を解析する
        event_ys, activation_spans, frame_spans, rect_backgrounds = self._analyze_events(
            messages, actors, lifeline_top_y, height, left, width,
        )

        # ライフライン終端Y座標を決定する
        lifeline_bot_y = (
            max(event_ys.values()) + _MSG_STEP
            if event_ys
            else lifeline_top_y + _MSG_STEP
        )

        # ---- 描画（背景から前面へ順に描く） ----

        # 1. rect背景（最下層）
        self._draw_rect_backgrounds(slide, rect_backgrounds, left, width)

        # 2. 参加者グループBox背景
        self._draw_participant_boxes(
            slide, boxes_raw, actors, cx_map, left, top, participant_w, col_step,
        )

        # 3. 参加者ボックス
        self._draw_participants(slide, actors, cx_map, top, participant_w)

        # 4. ライフライン
        self._draw_lifelines(slide, actors, cx_map, lifeline_top_y, lifeline_bot_y)

        # 5. フレームボックス
        self._draw_frames(slide, frame_spans, left, width)

        # 6. アクティベーションバー
        self._draw_activations(slide, activation_spans, cx_map)

        # 7. メッセージ矢印とラベル
        self._draw_messages(slide, messages, event_ys, cx_map, autonumber)

        # 8. ノートボックス
        self._draw_notes(slide, messages, event_ys, cx_map, participant_w)

    # ------------------------------------------------------------------
    # 参加者順序の構築
    # ------------------------------------------------------------------

    def _build_actor_order(self, actors_raw: dict) -> list[_ActorInfo]:
        """
        actors辞書からnextActorチェーンをたどって表示順のリストを構築する。

        チェーンが見つからない場合はdict挿入順を使用する。

        Parameters
        ----------
        actors_raw : dict
            mermaid-parser-pyが返すactors辞書。

        Returns
        -------
        list[_ActorInfo]
            表示順に並んだ参加者情報リスト。
        """
        all_ids = set(actors_raw.keys())
        # prevActorフィールドを持つ参加者IDの集合（= チェーン先頭ではないID）
        # v["prevActor"] の「値」ではなく、prevActorを持つ「ID自身」を集める
        has_prev = {
            aid
            for aid, v in actors_raw.items()
            if v.get("prevActor") and v["prevActor"] in all_ids
        }
        heads = all_ids - has_prev

        if not heads:
            # チェーンが見つからない場合はdictの挿入順を使う
            ordered_ids = list(actors_raw.keys())
        else:
            # 先頭からnextActorをたどって順序を決める
            head = sorted(heads)[0]
            ordered_ids: list[str] = []
            cur: Optional[str] = head
            visited: set[str] = set()
            while cur and cur not in visited:
                ordered_ids.append(cur)
                visited.add(cur)
                next_id = actors_raw[cur].get("nextActor")
                cur = next_id if next_id and next_id in actors_raw else None

        result: list[_ActorInfo] = []
        for idx, aid in enumerate(ordered_ids):
            a = actors_raw[aid]
            result.append(_ActorInfo(
                actor_id=aid,
                label=a.get("description") or aid,
                kind=a.get("type", "participant"),
                col_idx=idx,
            ))
        return result

    # ------------------------------------------------------------------
    # イベントY座標・アクティベーション・フレーム解析
    # ------------------------------------------------------------------

    def _analyze_events(
        self,
        messages: list[dict],
        actors: list[_ActorInfo],
        lifeline_top_y: int,
        height: int,
        left: int,
        width: int,
    ) -> tuple[dict[int, int], list[_ActivationSpan], list[_FrameSpan], list[tuple]]:
        """
        messagesリストを走査してイベントY座標・各スパン・背景情報を収集する。

        Parameters
        ----------
        messages : list[dict]
            mermaid-parser-pyのmessagesリスト。
        actors : list[_ActorInfo]
            参加者リスト（アクティベーション初期化用）。
        lifeline_top_y : int
            ライフライン開始Y座標（EMU）。
        height : int
            描画エリア高さ（EMU）。
        left : int
            描画エリア左端（EMU）。
        width : int
            描画エリア幅（EMU）。

        Returns
        -------
        tuple
            (event_ys, activation_spans, frame_spans, rect_backgrounds)
            - event_ys: {メッセージインデックス: Y座標（EMU）}
            - activation_spans: _ActivationSpanのリスト
            - frame_spans: _FrameSpanのリスト
            - rect_backgrounds: (y_start, y_end, color_str) のリスト
        """
        # 描画可能な高さからMSG_STEPを算出する（最小300,000 EMU）
        n_slots = sum(
            1 for m in messages
            if m.get("type") in _MSG_TYPES
            or m.get("type") == _LT_NOTE
            or m.get("type") in _FRAME_START_TYPES
            or m.get("type") in _FRAME_SECTION_TYPES
        )
        available_h = height - _PARTICIPANT_H - _LIFELINE_TOP_MARGIN
        if n_slots > 0:
            msg_step = max(300_000, min(_MSG_STEP, available_h // (n_slots + 1)))
        else:
            msg_step = _MSG_STEP

        event_ys: dict[int, int] = {}
        current_y: int = lifeline_top_y + msg_step // 2
        # 最後に確定したイベントY（アクティベーション開始Y決定に使用）
        last_event_y: int = lifeline_top_y

        # アクティベーション管理: {actor_id: [start_y, ...]} のスタック
        activation_stacks: dict[str, list[int]] = {
            a.actor_id: [] for a in actors
        }
        activation_spans: list[_ActivationSpan] = []

        # フレームスタック
        frame_stack: list[dict] = []
        frame_spans: list[_FrameSpan] = []

        # rect背景スタック
        rect_stack: list[dict] = []
        rect_backgrounds: list[tuple] = []

        for i, msg in enumerate(messages):
            t = msg.get("type")

            if t == _LT_AUTONUMBER:
                # シーケンス番号有効化マーカー。Y割り当て不要。
                continue

            elif t in _MSG_TYPES:
                # 実際のメッセージ矢印にY座標を割り当てる
                event_ys[i] = current_y
                last_event_y = current_y
                current_y += msg_step

            elif t == _LT_NOTE:
                # ノートボックスにY座標を割り当てる
                event_ys[i] = current_y
                last_event_y = current_y
                current_y += _NOTE_H + msg_step // 4

            elif t in _FRAME_START_TYPES:
                # フレーム開始: タブ分のスペースを確保する
                event_ys[i] = current_y
                frame_stack.append({
                    "type": t,
                    "label": msg.get("message", "") or "",
                    "start_y": current_y,
                    "sections": [],
                })
                current_y += _FRAME_TAB_H

            elif t in _FRAME_SECTION_TYPES:
                # セクション区切り: 現在のYをセクション境界として記録する
                if frame_stack:
                    frame_stack[-1]["sections"].append(
                        (current_y, msg.get("message", "") or "")
                    )
                # Y座標は進めない（区切り線は空間を消費しない）

            elif t in _FRAME_END_TYPES:
                # フレーム終了: スタックからポップしてFrameSpanを作成する
                if frame_stack:
                    f = frame_stack.pop()
                    frame_spans.append(_FrameSpan(
                        kind=f["type"],
                        label=f["label"],
                        start_y=f["start_y"],
                        end_y=current_y,
                        sections=f["sections"],
                    ))

            elif t == _LT_ACTIVE_START:
                # アクティベーション開始: 直前イベントのYをstart_yとする
                actor_id = msg.get("from", "")
                if actor_id in activation_stacks:
                    activation_stacks[actor_id].append(last_event_y)

            elif t == _LT_ACTIVE_END:
                # アクティベーション終了: スパンを確定する
                actor_id = msg.get("from", "")
                if actor_id in activation_stacks and activation_stacks[actor_id]:
                    depth = len(activation_stacks[actor_id]) - 1
                    start_y = activation_stacks[actor_id].pop()
                    activation_spans.append(_ActivationSpan(
                        actor_id=actor_id,
                        start_y=start_y,
                        end_y=last_event_y + msg_step // 4,
                        depth=depth,
                    ))

            elif t == _LT_RECT_START:
                # rect背景開始
                rect_stack.append({
                    "color": msg.get("message", "") or "",
                    "start_y": current_y,
                })

            elif t == _LT_RECT_END:
                # rect背景終了: スパンを確定する
                if rect_stack:
                    r = rect_stack.pop()
                    rect_backgrounds.append((r["start_y"], current_y, r["color"]))

        return event_ys, activation_spans, frame_spans, rect_backgrounds

    # ------------------------------------------------------------------
    # rect背景描画
    # ------------------------------------------------------------------

    def _draw_rect_backgrounds(
        self,
        slide: Slide,
        rect_backgrounds: list[tuple],
        left: int,
        width: int,
    ) -> None:
        """
        rect rgb(...) で指定された背景矩形を最下層に描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        rect_backgrounds : list[tuple]
            (y_start, y_end, color_str) のリスト。
        left : int
            描画エリア左端（EMU）。
        width : int
            描画エリア幅（EMU）。
        """
        for y_start, y_end, color_str in rect_backgrounds:
            rect_h = y_end - y_start
            if rect_h <= 0:
                continue

            r, g, b = self._parse_rgb(color_str)
            sp = slide.shapes.add_shape(
                1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
                Emu(left),
                Emu(y_start),
                Emu(width),
                Emu(rect_h),
            )
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(r, g, b)

            # 透明度を50%に設定する（OOXMLのa:alpha経由）
            spPr = sp._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(".//" + qn("a:solidFill"))
                if solidFill is not None:
                    srgb_el = solidFill.find(qn("a:srgbClr"))
                    if srgb_el is not None:
                        alpha_el = lxml_etree.SubElement(srgb_el, qn("a:alpha"))
                        alpha_el.set("val", "50000")  # 50% 透明度

            sp.line.fill.background()

            # zオーダーを最下層に移動する（spTree の先頭近くに挿入）
            spTree = slide.shapes._spTree
            sp_el = sp._element
            spTree.remove(sp_el)
            spTree.insert(2, sp_el)

    # ------------------------------------------------------------------
    # 参加者グループBox描画
    # ------------------------------------------------------------------

    def _draw_participant_boxes(
        self,
        slide: Slide,
        boxes_raw: list,
        actors: list[_ActorInfo],
        cx_map: dict[str, int],
        left: int,
        top: int,
        participant_w: int,
        col_step: int,
    ) -> None:
        """
        box ... end で定義されたグループBoxの背景を参加者ヘッダー領域に描画する。

        Parameters
        ----------
        slide : Slide
        boxes_raw : list
            mermaid-parser-pyのboxesリスト。
        actors : list[_ActorInfo]
            表示順の参加者リスト。
        cx_map : dict
            参加者IDをキー、列中央X座標を値とする辞書。
        left : int
            描画エリア左端（EMU）。
        top : int
            描画エリア上端（EMU）。
        participant_w : int
            参加者ボックス幅（EMU）。
        col_step : int
            列ピッチ（EMU）。
        """
        if not boxes_raw:
            return

        actor_id_to_col = {a.actor_id: a.col_idx for a in actors}
        for box in boxes_raw:
            box_actors: list = box.get("actors", [])
            color_str: str = box.get("fill", "") or ""
            label_text: str = box.get("name", "") or ""

            if not box_actors:
                continue

            col_indices = [
                actor_id_to_col[aid]
                for aid in box_actors
                if aid in actor_id_to_col
            ]
            if not col_indices:
                continue

            min_col = min(col_indices)
            max_col = max(col_indices)
            box_left = left + min_col * col_step
            box_w = (max_col - min_col + 1) * col_step

            sp = slide.shapes.add_shape(
                1,  # RECTANGLE
                Emu(box_left),
                Emu(top),
                Emu(box_w),
                Emu(_PARTICIPANT_H + 60_000),
            )
            r, g, b = self._parse_color_str(color_str, default=(230, 230, 255))
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(r, g, b)
            sp.line.color.rgb = RGBColor(180, 180, 220)

            if label_text:
                sp.text = label_text
                tf = sp.text_frame
                for para in tf.paragraphs:
                    from pptx.enum.text import PP_ALIGN
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(80, 80, 120)

            # zオーダーを最下層に移動する
            spTree = slide.shapes._spTree
            sp_el = sp._element
            spTree.remove(sp_el)
            spTree.insert(2, sp_el)

    # ------------------------------------------------------------------
    # 参加者ボックス描画
    # ------------------------------------------------------------------

    def _draw_participants(
        self,
        slide: Slide,
        actors: list[_ActorInfo],
        cx_map: dict[str, int],
        top: int,
        participant_w: int,
    ) -> dict[str, object]:
        """
        参加者ボックスを描画する。

        - participant → 矩形（MSO RECTANGLE）
        - actor → 角丸矩形（MSO ROUNDED_RECTANGLE）で区別する

        Parameters
        ----------
        slide : Slide
        actors : list[_ActorInfo]
            参加者リスト。
        cx_map : dict
            参加者IDをキー、列中央X座標を値とする辞書。
        top : int
            描画エリア上端Y座標（EMU）。
        participant_w : int
            参加者ボックス幅（EMU）。

        Returns
        -------
        dict[str, object]
            参加者IDをキー、Shapeオブジェクトを値とする辞書。
        """
        shapes: dict[str, object] = {}
        for actor in actors:
            cx = cx_map[actor.actor_id]
            box_left = cx - participant_w // 2

            if actor.kind == "actor":
                # 角丸矩形（5 = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE）
                sp = slide.shapes.add_shape(
                    5,
                    Emu(box_left), Emu(top),
                    Emu(participant_w), Emu(_PARTICIPANT_H),
                )
            else:
                # 通常矩形（1 = MSO_AUTO_SHAPE_TYPE.RECTANGLE）
                sp = slide.shapes.add_shape(
                    1,
                    Emu(box_left), Emu(top),
                    Emu(participant_w), Emu(_PARTICIPANT_H),
                )

            sp.text = actor.label
            tf = sp.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True

            shapes[actor.actor_id] = sp

        return shapes

    # ------------------------------------------------------------------
    # ライフライン描画
    # ------------------------------------------------------------------

    def _draw_lifelines(
        self,
        slide: Slide,
        actors: list[_ActorInfo],
        cx_map: dict[str, int],
        lifeline_top_y: int,
        lifeline_bot_y: int,
    ) -> None:
        """
        各参加者のライフライン（垂直破線コネクター）を描画する。

        Parameters
        ----------
        slide : Slide
        actors : list[_ActorInfo]
        cx_map : dict
        lifeline_top_y : int
            ライフライン開始Y座標（EMU）。
        lifeline_bot_y : int
            ライフライン終了Y座標（EMU）。
        """
        for actor in actors:
            cx = cx_map[actor.actor_id]
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                Emu(cx), Emu(lifeline_top_y),
                Emu(cx), Emu(lifeline_bot_y),
            )
            # 破線・細線スタイルをOOXML直接設定する
            cxn_el = connector._element
            spPr = cxn_el.find(qn("p:spPr"))
            if spPr is not None:
                ln_el = spPr.find(qn("a:ln"))
                if ln_el is None:
                    ln_el = lxml_etree.SubElement(spPr, qn("a:ln"))
                ln_el.set("w", "12700")  # 1pt
                prstDash = lxml_etree.SubElement(ln_el, qn("a:prstDash"))
                prstDash.set("val", "dash")

    # ------------------------------------------------------------------
    # フレームボックス描画
    # ------------------------------------------------------------------

    def _draw_frames(
        self,
        slide: Slide,
        frame_spans: list[_FrameSpan],
        left: int,
        width: int,
    ) -> None:
        """
        フレームボックス（loop/alt/opt/par/critical/break）を描画する。

        各フレームは破線矩形で囲み、左上に小矩形タブでキーワードを表示する。
        セクション区切り（else/and/option）は水平破線で描画する。

        Parameters
        ----------
        slide : Slide
        frame_spans : list[_FrameSpan]
            描画するフレームリスト。
        left : int
            描画エリア左端（EMU）。
        width : int
            描画エリア幅（EMU）。
        """
        FRAME_MARGIN = 60_000   # フレーム枠の水平マージン
        TAB_W = 500_000         # キーワードタブ幅
        TAB_H = 240_000         # キーワードタブ高さ

        for frame in frame_spans:
            frame_left = left + FRAME_MARGIN
            frame_top = frame.start_y
            frame_w = width - 2 * FRAME_MARGIN
            frame_h = frame.end_y - frame.start_y

            if frame_h <= 0:
                continue

            # 外枠矩形（破線・透明塗り）
            sp = slide.shapes.add_shape(
                1,  # RECTANGLE
                Emu(frame_left), Emu(frame_top),
                Emu(frame_w), Emu(frame_h),
            )
            sp.fill.background()
            sp.line.color.rgb = RGBColor(100, 100, 180)
            sp.line.width = Emu(19050)  # 1.5pt

            # 破線指定（OOXML直接設定）
            spPr = sp._element.find(qn("p:spPr"))
            if spPr is not None:
                ln_el = spPr.find(qn("a:ln"))
                if ln_el is not None:
                    prstDash = lxml_etree.SubElement(ln_el, qn("a:prstDash"))
                    prstDash.set("val", "dash")

            # キーワードタブ（左上の小矩形）
            keyword = _FRAME_KEYWORD.get(frame.kind, "")
            tab = slide.shapes.add_shape(
                1,  # RECTANGLE
                Emu(frame_left), Emu(frame_top),
                Emu(TAB_W), Emu(TAB_H),
            )
            tab.fill.solid()
            tab.fill.fore_color.rgb = RGBColor(200, 200, 240)
            tab.line.color.rgb = RGBColor(100, 100, 180)
            tab.text = keyword
            tf = tab.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True

            # フレームラベルテキスト（タブ右隣りに[ label ]形式で表示）
            if frame.label:
                label_sp = slide.shapes.add_textbox(
                    Emu(frame_left + TAB_W + 40_000),
                    Emu(frame_top),
                    Emu(frame_w - TAB_W - 40_000),
                    Emu(TAB_H),
                )
                tf = label_sp.text_frame
                tf.text = f"[{frame.label}]"
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.italic = True
                        run.font.color.rgb = RGBColor(60, 60, 140)

            # セクション区切り線（alt の else、par の and など）
            for section_y, section_label in frame.sections:
                # 水平コネクターで区切り線を描画する
                divider = slide.shapes.add_connector(
                    1,  # STRAIGHT
                    Emu(frame_left), Emu(section_y),
                    Emu(frame_left + frame_w), Emu(section_y),
                )
                cxn_el = divider._element
                spPr = cxn_el.find(qn("p:spPr"))
                if spPr is not None:
                    ln_el = spPr.find(qn("a:ln"))
                    if ln_el is None:
                        ln_el = lxml_etree.SubElement(spPr, qn("a:ln"))
                    ln_el.set("w", "12700")
                    # 色を設定する
                    solidFill = lxml_etree.SubElement(ln_el, qn("a:solidFill"))
                    srgb = lxml_etree.SubElement(solidFill, qn("a:srgbClr"))
                    srgb.set("val", "6464B4")
                    # 破線
                    prstDash = lxml_etree.SubElement(ln_el, qn("a:prstDash"))
                    prstDash.set("val", "dash")

                # セクションラベル（区切り線のすぐ下）
                if section_label:
                    lbl = slide.shapes.add_textbox(
                        Emu(frame_left + TAB_W + 40_000),
                        Emu(section_y),
                        Emu(frame_w - TAB_W - 40_000),
                        Emu(200_000),
                    )
                    tf = lbl.text_frame
                    tf.text = f"[{section_label}]"
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(8)
                            run.font.italic = True
                            run.font.color.rgb = RGBColor(60, 60, 140)

    # ------------------------------------------------------------------
    # アクティベーションバー描画
    # ------------------------------------------------------------------

    def _draw_activations(
        self,
        slide: Slide,
        activation_spans: list[_ActivationSpan],
        cx_map: dict[str, int],
    ) -> None:
        """
        アクティベーションバー（薄い青の細矩形）を描画する。

        深度に応じてX方向にオフセットしてスタック表現を実現する。

        Parameters
        ----------
        slide : Slide
        activation_spans : list[_ActivationSpan]
            描画するアクティベーションスパンリスト。
        cx_map : dict
            参加者IDをキー、列中央X座標を値とする辞書。
        """
        ACT_FILL = RGBColor(200, 220, 255)
        ACT_BORDER = RGBColor(100, 130, 200)

        for span in activation_spans:
            cx = cx_map.get(span.actor_id)
            if cx is None:
                continue

            bar_h = span.end_y - span.start_y
            if bar_h <= 0:
                continue

            # 深度に応じてX座標をずらす（スタックバーの表現）
            bar_left = cx - _ACTIVATION_W // 2 + span.depth * _ACTIVATION_W

            sp = slide.shapes.add_shape(
                1,  # RECTANGLE
                Emu(bar_left), Emu(span.start_y),
                Emu(_ACTIVATION_W), Emu(bar_h),
            )
            sp.fill.solid()
            sp.fill.fore_color.rgb = ACT_FILL
            sp.line.color.rgb = ACT_BORDER
            sp.line.width = Emu(9525)  # 0.75pt

    # ------------------------------------------------------------------
    # メッセージ描画
    # ------------------------------------------------------------------

    def _draw_messages(
        self,
        slide: Slide,
        messages: list[dict],
        event_ys: dict[int, int],
        cx_map: dict[str, int],
        autonumber: bool,
    ) -> None:
        """
        メッセージ矢印とラベルテキストボックスを描画する。

        autonumber が有効の場合はラベルの先頭にシーケンス番号を付与する。
        クロス矢印（-x/--x）はラベル末尾に「✕」を付加する（OOXML代替）。

        Parameters
        ----------
        slide : Slide
        messages : list[dict]
            mermaid-parser-pyのmessagesリスト。
        event_ys : dict[int, int]
            メッセージインデックス→Y座標の辞書。
        cx_map : dict
            参加者IDをキー、列中央X座標を値とする辞書。
        autonumber : bool
            シーケンス番号を付与するか。
        """
        seq_num = 0

        for i, msg in enumerate(messages):
            t = msg.get("type")
            if t not in _MSG_TYPES:
                continue
            if i not in event_ys:
                continue

            y = event_ys[i]
            src_id = msg.get("from", "")
            dst_id = msg.get("to", src_id)
            text: str = msg.get("message", "") or ""

            src_cx = cx_map.get(src_id)
            dst_cx = cx_map.get(dst_id, src_cx)
            if src_cx is None:
                continue

            seq_num += 1
            label = f"{seq_num}. {text}" if autonumber else text

            # クロス矢印はラベルに「✕」を付加する（OOXML headEnd に cross 種別なし）
            if t in (_LT_SOLID_CROSS, _LT_DOTTED_CROSS):
                label = f"{label} ✕" if label else "✕"

            # 自己メッセージ（src == dst）はELBOWコネクターで描画する
            if src_id == dst_id:
                self._draw_self_message(slide, src_cx, y, label, t)
            else:
                self._draw_normal_message(slide, src_cx, dst_cx, y, label, t)

    def _draw_normal_message(
        self,
        slide: Slide,
        src_cx: int,
        dst_cx: int,
        y: int,
        label: str,
        msg_type: int,
    ) -> None:
        """
        通常メッセージ（異なる参加者間）の矢印とラベルを描画する。

        コネクターとラベルテキストボックスは別シェイプとして配置する。

        Parameters
        ----------
        slide : Slide
        src_cx : int
            送信者の列中央X座標（EMU）。
        dst_cx : int
            受信者の列中央X座標（EMU）。
        y : int
            メッセージY座標（EMU）。
        label : str
            ラベルテキスト。
        msg_type : int
            LINETYPEの値。
        """
        connector = slide.shapes.add_connector(
            1,  # STRAIGHT
            Emu(src_cx), Emu(y),
            Emu(dst_cx), Emu(y),
        )
        self._apply_line_style(connector, msg_type)

        # ラベルをコネクター上方に独立テキストボックスとして配置する
        if label:
            tb_left = min(src_cx, dst_cx)
            tb_w = max(abs(dst_cx - src_cx), 500_000)
            tb = slide.shapes.add_textbox(
                Emu(tb_left),
                Emu(y - 220_000),
                Emu(tb_w),
                Emu(200_000),
            )
            tf = tb.text_frame
            tf.text = label
            for para in tf.paragraphs:
                from pptx.enum.text import PP_ALIGN
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(9)

    def _draw_self_message(
        self,
        slide: Slide,
        cx: int,
        y: int,
        label: str,
        msg_type: int,
    ) -> None:
        """
        自己メッセージ（同一参加者への矢印）をELBOWコネクターで描画する。

        ELBOWコネクターで右方向に折れた「L字型」として近似する。

        Parameters
        ----------
        slide : Slide
        cx : int
            参加者の列中央X座標（EMU）。
        y : int
            メッセージ開始Y座標（EMU）。
        label : str
            ラベルテキスト。
        msg_type : int
            LINETYPEの値。
        """
        connector = slide.shapes.add_connector(
            2,  # MSO_CONNECTOR_TYPE.ELBOW
            Emu(cx), Emu(y),
            Emu(cx), Emu(y + _MSG_STEP // 2),
        )
        self._apply_line_style(connector, msg_type)

        # ラベルを折れ部分の右上に配置する
        if label:
            tb = slide.shapes.add_textbox(
                Emu(cx + 60_000),
                Emu(y - 200_000),
                Emu(_SELF_MSG_INDENT - 60_000),
                Emu(200_000),
            )
            tf = tb.text_frame
            tf.text = label
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)

    def _apply_line_style(self, connector: object, msg_type: int) -> None:
        """
        コネクターのOOXMLに矢印・線種スタイルを適用する。

        OOXML制限:
        - クロス矢印（-x）は headEnd に "cross" 型がないため "stealth" で代用
          （ラベルへの「✕」付加で補完する）
        - 半矢印・その他非標準型は "stealth" にフォールバックする

        Parameters
        ----------
        connector : object
            python-pptxのConnectorオブジェクト。
        msg_type : int
            LINETYPEの値。
        """
        cxn_el = connector._element
        spPr = cxn_el.find(qn("p:spPr"))
        if spPr is None:
            return

        ln_el = spPr.find(qn("a:ln"))
        if ln_el is None:
            ln_el = lxml_etree.SubElement(spPr, qn("a:ln"))

        # 破線種別の設定
        is_dotted = msg_type in (
            _LT_DOTTED, _LT_DOTTED_CROSS, _LT_DOTTED_OPEN,
            _LT_DOTTED_POINT, _LT_BIDIR_DOTTED,
        )
        if is_dotted:
            prstDash = lxml_etree.SubElement(ln_el, qn("a:prstDash"))
            prstDash.set("val", "dash")

        # 矢印の先端形状を設定する
        if msg_type in (_LT_SOLID_OPEN, _LT_DOTTED_OPEN):
            # 非同期メッセージ: 開き矢印
            head = lxml_etree.SubElement(ln_el, qn("a:headEnd"))
            head.set("type", "open")
            head.set("w", "med")
            head.set("len", "med")

        elif msg_type in (_LT_BIDIR_SOLID, _LT_BIDIR_DOTTED):
            # 双方向矢印: 両端に矢印を設定する
            tail = lxml_etree.SubElement(ln_el, qn("a:tailEnd"))
            tail.set("type", "stealth")
            tail.set("w", "med")
            tail.set("len", "med")
            head = lxml_etree.SubElement(ln_el, qn("a:headEnd"))
            head.set("type", "stealth")
            head.set("w", "med")
            head.set("len", "med")

        elif msg_type in (
            _LT_SOLID, _LT_DOTTED,
            _LT_SOLID_CROSS, _LT_DOTTED_CROSS,
            _LT_SOLID_POINT, _LT_DOTTED_POINT,
        ):
            # 標準矢印: 塗り矢印（stealth）
            head = lxml_etree.SubElement(ln_el, qn("a:headEnd"))
            head.set("type", "stealth")
            head.set("w", "med")
            head.set("len", "med")

        # _LT_SOLID_POINT (24) / _LT_DOTTED_POINT (25) は矢印なし（->）相当だが
        # OOXML上はtype0と区別できないため同じ stealthを設定する

    # ------------------------------------------------------------------
    # ノート描画
    # ------------------------------------------------------------------

    def _draw_notes(
        self,
        slide: Slide,
        messages: list[dict],
        event_ys: dict[int, int],
        cx_map: dict[str, int],
        participant_w: int,
    ) -> None:
        """
        Noteボックス（黄色矩形）を描画する。

        PLACEMENT:
        - 0 (LEFTOF): 参加者の左側に配置
        - 1 (RIGHTOF): 参加者の右側に配置
        - 2 (OVER): 参加者の上に配置（複数参加者にまたがる場合はspan）

        Parameters
        ----------
        slide : Slide
        messages : list[dict]
            mermaid-parser-pyのmessagesリスト。
        event_ys : dict[int, int]
            メッセージインデックス→Y座標の辞書。
        cx_map : dict
            参加者IDをキー、列中央X座標を値とする辞書。
        participant_w : int
            参加者ボックス幅（EMU）。
        """
        NOTE_FILL = RGBColor(255, 253, 173)
        NOTE_BORDER = RGBColor(200, 190, 100)
        NOTE_MARGIN = 80_000  # ライフラインからの余白

        for i, msg in enumerate(messages):
            if msg.get("type") != _LT_NOTE:
                continue
            if i not in event_ys:
                continue

            y = event_ys[i]
            actor_id: str = msg.get("from", "")
            to_id: str = msg.get("to", actor_id) or actor_id
            text: str = msg.get("message", "") or ""
            placement: int = msg.get("placement", 2)

            cx_from = cx_map.get(actor_id)
            if cx_from is None:
                continue

            cx_to = cx_map.get(to_id) if to_id != actor_id else cx_from
            note_w = participant_w + 100_000

            if placement == 0:
                # LEFTOF: 参加者の左側
                note_left = cx_from - NOTE_MARGIN - note_w
            elif placement == 1:
                # RIGHTOF: 参加者の右側
                note_left = cx_from + NOTE_MARGIN
            else:
                # OVER: 参加者の上（複数参加者にまたがる場合はspan）
                if cx_to and cx_to != cx_from:
                    span_left = min(cx_from, cx_to)
                    span_right = max(cx_from, cx_to)
                    note_left = span_left - participant_w // 2
                    note_w = span_right - span_left + participant_w
                else:
                    note_left = cx_from - participant_w // 2 - 50_000
                    note_w = participant_w + 100_000

            sp = slide.shapes.add_shape(
                1,  # RECTANGLE
                Emu(note_left), Emu(y),
                Emu(note_w), Emu(_NOTE_H),
            )
            sp.fill.solid()
            sp.fill.fore_color.rgb = NOTE_FILL
            sp.line.color.rgb = NOTE_BORDER
            sp.text = text
            tf = sp.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)

    # ------------------------------------------------------------------
    # ユーティリティ
    # ------------------------------------------------------------------

    def _parse_rgb(self, color_str: str) -> tuple[int, int, int]:
        """
        "rgb(r, g, b)" 形式の文字列からRGB値を解析する。

        Parameters
        ----------
        color_str : str
            "rgb(200, 150, 255)" などのCSS rgb()形式の文字列。

        Returns
        -------
        tuple[int, int, int]
            (r, g, b) のタプル（各値0〜255）。解析失敗時は (240, 240, 255) を返す。
        """
        m = re.search(
            r"rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)",
            color_str,
        )
        if m:
            return int(m.group(1)), int(m.group(2)), int(m.group(3))
        return 240, 240, 255

    def _parse_color_str(
        self,
        color_str: str,
        default: tuple[int, int, int] = (220, 220, 240),
    ) -> tuple[int, int, int]:
        """
        カラー文字列からRGB値を解析する。

        rgb()形式を優先し、失敗時はCSSカラー名の簡易テーブルを参照する。

        Parameters
        ----------
        color_str : str
            カラー文字列（"rgb(200,150,255)" 形式またはCSSカラー名）。
        default : tuple[int, int, int]
            解析失敗時のデフォルト色。

        Returns
        -------
        tuple[int, int, int]
            (r, g, b) のタプル。
        """
        if not color_str:
            return default

        # rgb()形式を試みる
        m = re.search(
            r"rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)",
            color_str,
        )
        if m:
            return int(m.group(1)), int(m.group(2)), int(m.group(3))

        # CSSカラー名の簡易対応テーブル
        _CSS_COLORS: dict[str, tuple[int, int, int]] = {
            "aqua": (0, 255, 255), "blue": (0, 0, 255),
            "cyan": (0, 255, 255), "green": (0, 128, 0),
            "lime": (0, 255, 0), "magenta": (255, 0, 255),
            "orange": (255, 165, 0), "pink": (255, 192, 203),
            "purple": (128, 0, 128), "red": (255, 0, 0),
            "violet": (238, 130, 238), "white": (255, 255, 255),
            "yellow": (255, 255, 0),
        }
        return _CSS_COLORS.get(color_str.lower().strip(), default)
