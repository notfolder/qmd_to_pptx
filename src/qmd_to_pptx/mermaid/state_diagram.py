"""
ステートダイアグラムレンダラーモジュール（改版）。

stateDiagram-v2 の各種要素を正確に描画するレンダラー。

対応要素:
  - 通常状態 (shape="rect"): 角丸矩形（白背景・黒枠）
  - 開始擬似状態 (shape="stateStart"): 黒塗り小OVAL
  - 終了擬似状態 (shape="stateEnd"): bull's-eye（黒OVAL + 内側白OVAL の 2枚重ね）
  - 複合状態 (shape="roundedWithTitle"): 外枠角丸矩形 + タイトルテキスト + 子ノード内部配置
  - 並行セクション (shape="divider"): 親複合状態内の破線枠区画
  - 選択 (shape="choice"): DIAMOND（薄黄背景）
  - フォーク/ジョイン (shape="fork"/"join"): 扁平黒RECTANGLE バー
  - ノート (shape="noteGroup"): 黄色矩形のテキストメモ
  - 遷移ラベル (edges[].label): コネクター中点付近にテキストボックスで配置
  - 曲線コネクター: curvedConnector3 (connector_type=3)
  - 向き指定 (direction TB/LR/BT/RL): デフォルト TB（上→下）。
      renderer.py が graph_data["_direction"] に注入した値を使用する。

OOXML制限と代替案:
  - 終了状態 bull's-eye: OOXMLプリセットなし →
      黒OVAL + 内側小白OVAL の 2図形重ね合わせで表現
  - 複合状態グループ: python-pptxにGroupShape APIなし →
      外枠矩形を先に描き（背景）、子シェイプを独立配置（z-order: 外枠→子）
  - fork/join バー: OOXMLに横バー専用プリセットなし →
      扁平RECTANGLE（高さ小・幅大、黒塗り）で代替
  - divider（並行区画）: 破線枠 ROUNDED_RECTANGLE で区画を表現
"""

from __future__ import annotations

import math
import re
from typing import Optional

import networkx as nx
from lxml import etree as lxml_etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .base import BaseDiagramRenderer, NODE_WIDTH_EMU, NODE_HEIGHT_EMU

# ---------------------------------------------------------------------------
# 形状サイズ定数（EMU）
# ---------------------------------------------------------------------------
_START_SIZE: int = 380_000        # 開始擬似状態: 直径
_END_OUTER: int = 450_000         # 終了擬似状態: 外円直径
_END_INNER: int = 250_000         # 終了擬似状態: 内側白円直径
_CHOICE_SIZE: int = 500_000       # <<choice>> 菱形: 幅/高さ
_FORK_W: int = 1_200_000          # fork/join バー: 幅
_FORK_H: int = 120_000            # fork/join バー: 高さ
_STATE_W: int = NODE_WIDTH_EMU    # 通常状態: 幅
_STATE_H: int = NODE_HEIGHT_EMU   # 通常状態: 高さ
_COMPOSITE_PADDING: int = 500_000 # 複合状態コンテナの内側パディング（EMU）
_COMPOSITE_TITLE_H: int = 420_000 # 複合状態コンテナのタイトル行高さ（EMU）

# ---------------------------------------------------------------------------
# MSO_AUTO_SHAPE_TYPE 整数値
# ---------------------------------------------------------------------------
_MSO_RECT: int = 1         # RECTANGLE
_MSO_ROUNDED: int = 5      # ROUNDED_RECTANGLE
_MSO_OVAL: int = 9         # OVAL
_MSO_DIAMOND: int = 4      # DIAMOND

# ---------------------------------------------------------------------------
# カラー定義
# ---------------------------------------------------------------------------
_BLACK = RGBColor(0, 0, 0)
_WHITE = RGBColor(255, 255, 255)
_STATE_BG = RGBColor(255, 255, 255)          # 通常状態: 白背景
_STATE_BORDER = RGBColor(50, 50, 80)          # 通常状態: 濃い紺枠
_COMPOSITE_BG = RGBColor(238, 244, 255)       # 複合状態背景: 薄い青白
_COMPOSITE_BORDER = RGBColor(70, 100, 170)    # 複合状態枠: 青
_DIVIDER_BG = RGBColor(230, 240, 255)         # 並行区画背景: 複合状態より少し暗い青
_CHOICE_BG = RGBColor(255, 250, 210)          # choice 背景: 薄い黄
_CHOICE_BORDER = RGBColor(180, 140, 0)        # choice 枠: 黄茶
_FORK_COLOR = RGBColor(30, 30, 30)            # fork/join バー: ほぼ黒
_LABEL_FG = RGBColor(50, 50, 50)              # 遷移ラベル文字色
_NOTE_BG = RGBColor(255, 253, 200)            # ノート背景: 薄い黄色
_NOTE_BORDER = RGBColor(160, 140, 60)         # ノート枠: 黄茶


class StateDiagramRenderer(BaseDiagramRenderer):
    """
    stateDiagram-v2 を描画するレンダラー（改版）。

    mermaid-parser-py が返す graph_data（nodes/edges）を解析し、
    各ノードの shape フィールドに応じた形状をスライドに描画する。
    複合状態（composite state）は外枠矩形 + 子ノードの内部配置で表現し、
    遷移ラベルはコネクター中点のテキストボックスで、
    コネクターはすべて curvedConnector3 で描画する。
    """

    def render(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        stateDiagram-v2 をスライドに描画する。

        renderer.py が graph_data["_direction"] に方向を注入した上でこのメソッドを
        呼び出す。注入された値がない場合はデフォルト "TB"（上→下）で描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"nodes"と"edges"キーを含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        # renderer.py が注入した描画方向を取得する（取得後は削除する）
        direction: str = graph_data.pop("_direction", "TB")

        raw_nodes: list[dict] = graph_data.get("nodes", [])
        raw_edges: list[dict] = graph_data.get("edges", [])
        if not raw_nodes:
            return

        # ノードを ID でインデックスする
        node_by_id: dict[str, dict] = {
            n["id"]: n for n in raw_nodes if "id" in n
        }

        # ノート系（shape="note"/"noteGroup"）を通常描画から除外する
        _NOTE_SHAPES: frozenset[str] = frozenset(("note", "noteGroup"))
        note_ids: set[str] = {
            n["id"] for n in raw_nodes if n.get("shape") in _NOTE_SHAPES
        }

        # 複合状態コンテナ（isGroup=True、ノートでない）を収集する
        composite_ids: set[str] = {
            n["id"] for n in raw_nodes
            if n.get("isGroup") and n["id"] not in note_ids
        }

        # 親子関係マップ（ノート系ノードを除く）
        child_parent: dict[str, str] = {
            n["id"]: n["parentId"]
            for n in raw_nodes
            if "parentId" in n and n["id"] not in note_ids
        }

        # トップレベルノード（parentIdなし、ノートでない）
        top_level_ids: list[str] = [
            n["id"] for n in raw_nodes
            if "id" in n and "parentId" not in n and n["id"] not in note_ids
        ]

        # トップレベルのエッジ（両端がともにトップレベルのエッジのみ）
        top_level_set = set(top_level_ids)
        top_edges: list[tuple[str, str]] = [
            (e.get("start", ""), e.get("end", ""))
            for e in raw_edges
            if e.get("start", "") in top_level_set and e.get("end", "") in top_level_set
        ]

        # トップレベルグラフのレイアウトを計算する
        top_pos = self._calc_layout(top_level_ids, top_edges, direction)

        # 全ノードの EMU 中心座標マップと複合状態の外枠マップを構築する
        node_center_emu: dict[str, tuple[int, int]] = {}
        composite_emu_map: dict[str, tuple[int, int, int, int]] = {}

        # トップレベルの中心座標を登録する
        for nid in top_level_ids:
            if nid in top_pos:
                nx_n, ny_n = top_pos[nid]
                node_center_emu[nid] = self._pos_to_emu(
                    nx_n, ny_n, left, top, width, height
                )

        # トップレベル複合状態の子レイアウトを再帰的に確定する
        for comp_id in top_level_ids:
            if comp_id not in composite_ids:
                continue
            if comp_id not in node_center_emu:
                continue
            comp_cx, comp_cy = node_center_emu[comp_id]
            self._layout_composite_children(
                comp_id, comp_cx, comp_cy,
                child_parent, composite_ids, raw_edges,
                node_center_emu, composite_emu_map, direction,
            )

        # -----------------------------------------------------------------
        # 描画フェーズ（z-orderのため複合状態背景→通常ノード→ノート→エッジの順）
        # -----------------------------------------------------------------
        shape_map: dict[str, object] = {}

        # ① 複合状態の背景矩形＋タイトル（divider は破線枠スタイル）
        for comp_id in composite_ids:
            if comp_id not in composite_emu_map:
                continue
            cl, ct, cw, ch = composite_emu_map[comp_id]
            comp_node = node_by_id.get(comp_id, {})
            comp_label = comp_node.get("label", "")
            is_divider = comp_node.get("shape") == "divider"
            shape = self._draw_composite_bg(
                slide, comp_label, cl, ct, cw, ch, is_divider=is_divider
            )
            shape_map[comp_id] = shape

        # ② 個別ノード（複合状態コンテナとノートを除く）
        all_node_ids = list(top_level_ids) + list(child_parent.keys())
        for nid in all_node_ids:
            if nid in composite_ids or nid in note_ids:
                continue
            if nid not in node_center_emu:
                continue
            cx, cy = node_center_emu[nid]
            node = node_by_id.get(nid, {})
            shape_type = node.get("shape", "rect")
            label = node.get("label", nid)
            shape = self._draw_state_node(slide, shape_type, label, cx, cy)
            if shape is not None:
                shape_map[nid] = shape

        # ③ ノート（ターゲット状態の右または左に配置する）
        self._draw_note_nodes(slide, raw_nodes, node_center_emu, left, top, width, height)

        # ④ 遷移エッジ（コネクター + ラベル）
        self._draw_state_edges(slide, raw_edges, node_center_emu, shape_map)

    # ------------------------------------------------------------------
    # 複合状態の子レイアウト（再帰）
    # ------------------------------------------------------------------

    def _layout_composite_children(
        self,
        comp_id: str,
        comp_cx: int,
        comp_cy: int,
        child_parent: dict[str, str],
        composite_ids: set[str],
        raw_edges: list[dict],
        node_center_emu: dict[str, tuple[int, int]],
        composite_emu_map: dict[str, tuple[int, int, int, int]],
        direction: str,
    ) -> None:
        """
        複合状態（comp_id）の子ノードを再帰的にレイアウトする。

        子ノードに対して _calc_layout を実行し各子の EMU 中心座標を
        node_center_emu・composite_emu_map に書き込む。
        子が複合状態である場合はさらに再帰する。

        Parameters
        ----------
        comp_id : str
            親の複合状態ノードID。
        comp_cx, comp_cy : int
            親複合状態の中心座標（EMU）。
        child_parent : dict[str, str]
            ノードID → 親ノードID の辞書。
        composite_ids : set[str]
            複合状態（isGroup）のノードIDセット。
        raw_edges : list[dict]
            全エッジリスト。
        node_center_emu : dict
            ノードID → 中心EMU座標 の辞書（更新対象）。
        composite_emu_map : dict
            複合状態ID → (left, top, width, height)EMU の辞書（更新対象）。
        direction : str
            レイアウト方向（"TB"/"LR" など）。
        """
        children = [nid for nid, pid in child_parent.items() if pid == comp_id]
        if not children:
            return

        n_children = len(children)

        # 外枠サイズを子ノード数から決定する
        inner_w = _COMPOSITE_PADDING * 2 + _STATE_W * max(2, n_children)
        inner_h = _COMPOSITE_TITLE_H + _COMPOSITE_PADDING * 2 + _STATE_H * 2

        comp_left = comp_cx - inner_w // 2
        comp_top = comp_cy - inner_h // 2
        composite_emu_map[comp_id] = (comp_left, comp_top, inner_w, inner_h)

        # 子ノード間のエッジを収集してサブレイアウトを計算する
        child_set = set(children)
        child_edges = [
            (e.get("start", ""), e.get("end", ""))
            for e in raw_edges
            if e.get("start", "") in child_set and e.get("end", "") in child_set
        ]
        child_pos = self._calc_layout(children, child_edges, direction)

        # 有効キャンバス内に各子ノードの絶対EMUを確定する
        canvas_l = comp_left + _COMPOSITE_PADDING
        canvas_t = comp_top + _COMPOSITE_TITLE_H + _COMPOSITE_PADDING // 2
        canvas_w = inner_w - _COMPOSITE_PADDING * 2
        canvas_h = inner_h - _COMPOSITE_TITLE_H - _COMPOSITE_PADDING

        for cid, (nx_n, ny_n) in child_pos.items():
            x_ratio = max(0.0, min(1.0, (nx_n + 1.0) / 2.0))
            y_ratio = max(0.0, min(1.0, (ny_n + 1.0) / 2.0))
            cx_abs = canvas_l + int(x_ratio * canvas_w)
            cy_abs = canvas_t + int(y_ratio * canvas_h)
            node_center_emu[cid] = (cx_abs, cy_abs)

            # 子が複合状態（divider 含む）なら再帰的にレイアウトする
            if cid in composite_ids:
                self._layout_composite_children(
                    cid, cx_abs, cy_abs,
                    child_parent, composite_ids, raw_edges,
                    node_center_emu, composite_emu_map, direction,
                )

    # ------------------------------------------------------------------
    # レイアウト計算
    # ------------------------------------------------------------------

    def _calc_layout(
        self,
        node_ids: list[str],
        edges: list[tuple[str, str]],
        direction: str = "TB",
    ) -> dict[str, tuple[float, float]]:
        """
        ノードリストとエッジから正規化レイアウト座標（-1.0〜1.0）を計算する。

        DAG（有向非巡回グラフ）の場合はトポロジカル世代を使った階層レイアウトを行う。
        サイクルを含む場合は kamada_kawai → spring_layout にフォールバックする。

        Parameters
        ----------
        node_ids : list[str]
            レイアウト対象ノードIDのリスト。
        edges : list[tuple[str, str]]
            エッジのリスト（始点ID, 終点ID）。
        direction : str
            レイアウト方向（"TB"/"TD"/"BT"/"LR"/"RL"）。デフォルト "TB"。

        Returns
        -------
        dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標のタプルを値とする辞書。
        """
        if not node_ids:
            return {}
        if len(node_ids) == 1:
            return {node_ids[0]: (0.0, 0.0)}

        G = nx.DiGraph()
        G.add_nodes_from(node_ids)
        for src, dst in edges:
            if src in G and dst in G:
                G.add_edge(src, dst)

        # DAG の場合はトポロジカル階層レイアウトを使用する
        if nx.is_directed_acyclic_graph(G):
            try:
                return self._topological_layout(G, direction)
            except Exception:
                pass

        # サイクルあり: kamada_kawai → spring にフォールバックする
        try:
            return nx.kamada_kawai_layout(G)
        except Exception:
            return nx.spring_layout(G, seed=42, k=1.5)

    def _topological_layout(
        self,
        G: nx.DiGraph,
        direction: str,
    ) -> dict[str, tuple[float, float]]:
        """
        トポロジカル世代に基づく階層レイアウト座標（-1.0〜1.0）を返す。

        "TB"/"TD": 上→下（世代が Y 軸で連続）
        "BT": 下→上（Y 軸を反転）
        "LR": 左→右（世代が X 軸で連続）
        "RL": 右→左（X 軸を反転）

        Parameters
        ----------
        G : nx.DiGraph
            DAG グラフ（サイクルなし前提）。
        direction : str
            描画方向。

        Returns
        -------
        dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標(-1.0〜1.0)のタプルを値とする辞書。
        """
        generations = list(nx.topological_generations(G))
        n_levels = len(generations)
        pos: dict[str, tuple[float, float]] = {}

        for level_idx, gen_nodes in enumerate(generations):
            sorted_nodes = sorted(gen_nodes)
            n_nodes = len(sorted_nodes)

            # レベル軸（主軸）: -0.9 〜 +0.9 で等間隔配置
            if n_levels == 1:
                level_pos = 0.0
            else:
                level_pos = (level_idx / (n_levels - 1)) * 1.8 - 0.9

            # スパン軸（副軸）: 同世代ノードを等間隔配置
            for span_idx, nid in enumerate(sorted_nodes):
                if n_nodes == 1:
                    span_pos = 0.0
                else:
                    span_pos = (span_idx / (n_nodes - 1)) * 1.8 - 0.9

                if direction in ("TB", "TD"):
                    pos[nid] = (span_pos, level_pos)
                elif direction == "BT":
                    pos[nid] = (span_pos, -level_pos)
                elif direction == "LR":
                    pos[nid] = (level_pos, span_pos)
                elif direction == "RL":
                    pos[nid] = (-level_pos, span_pos)
                else:
                    # 不明な方向は TB にフォールバックする
                    pos[nid] = (span_pos, level_pos)

        return pos

    # ------------------------------------------------------------------
    # ノード描画
    # ------------------------------------------------------------------

    def _draw_state_node(
        self,
        slide: Slide,
        shape_type: str,
        label: str,
        cx: int,
        cy: int,
    ) -> Optional[object]:
        """
        ノードを種別に応じた形状でスライドに描画する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        shape_type : str
            mermaid-parser-py の node["shape"] フィールド値。
        label : str
            表示ラベル。
        cx, cy : int
            形状の中心座標（EMU）。

        Returns
        -------
        Shape | None
            描画した主シェイプ（コネクター接続用）。
        """
        if shape_type == "stateStart":
            return self._draw_start_state(slide, cx, cy)
        if shape_type == "stateEnd":
            return self._draw_end_state(slide, cx, cy)
        if shape_type == "choice":
            return self._draw_choice_state(slide, cx, cy)
        if shape_type in ("fork", "join"):
            return self._draw_fork_join(slide, cx, cy)
        # デフォルト: 通常状態（角丸矩形）
        return self._draw_normal_state(slide, label, cx, cy)

    def _draw_start_state(self, slide: Slide, cx: int, cy: int) -> object:
        """開始擬似状態（[*] から出発）を黒塗り小OVALで描画する。"""
        sl = cx - _START_SIZE // 2
        st = cy - _START_SIZE // 2
        shape = slide.shapes.add_shape(
            _MSO_OVAL, Emu(sl), Emu(st), Emu(_START_SIZE), Emu(_START_SIZE)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _BLACK
        shape.line.fill.background()
        if shape.has_text_frame:
            shape.text_frame.text = ""
        return shape

    def _draw_end_state(self, slide: Slide, cx: int, cy: int) -> object:
        """
        終了擬似状態（[*] に到着）を bull's-eye で描画する。

        OOXMLに bull's-eye プリセットが存在しないため、
        外側を黒 OVAL、内側を小さい白 OVAL の 2図形を重ねて表現する。
        コネクター接続には外側 OVAL のシェイプを返す。
        """
        # 外側の黒 OVAL
        sl = cx - _END_OUTER // 2
        st = cy - _END_OUTER // 2
        outer = slide.shapes.add_shape(
            _MSO_OVAL, Emu(sl), Emu(st), Emu(_END_OUTER), Emu(_END_OUTER)
        )
        outer.fill.solid()
        outer.fill.fore_color.rgb = _BLACK
        outer.line.fill.background()
        if outer.has_text_frame:
            outer.text_frame.text = ""

        # 内側の白 OVAL（オーバーレイ）
        il = cx - _END_INNER // 2
        it = cy - _END_INNER // 2
        inner = slide.shapes.add_shape(
            _MSO_OVAL, Emu(il), Emu(it), Emu(_END_INNER), Emu(_END_INNER)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = _WHITE
        inner.line.fill.background()
        if inner.has_text_frame:
            inner.text_frame.text = ""

        # コネクター接続には外側 OVAL を返す
        return outer

    def _draw_choice_state(self, slide: Slide, cx: int, cy: int) -> object:
        """<<choice>> 擬似状態（菱形）を描画する。"""
        sl = cx - _CHOICE_SIZE // 2
        st = cy - _CHOICE_SIZE // 2
        shape = slide.shapes.add_shape(
            _MSO_DIAMOND, Emu(sl), Emu(st), Emu(_CHOICE_SIZE), Emu(_CHOICE_SIZE)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _CHOICE_BG
        shape.line.color.rgb = _CHOICE_BORDER
        if shape.has_text_frame:
            shape.text_frame.text = ""
        return shape

    def _draw_fork_join(self, slide: Slide, cx: int, cy: int) -> object:
        """
        <<fork>> / <<join>> 擬似状態を扁平黒矩形バーで描画する。

        OOXMLに横バー専用プリセットがないため、
        幅広・高さ小の黒塗り RECTANGLE で代替する。
        """
        sl = cx - _FORK_W // 2
        st = cy - _FORK_H // 2
        shape = slide.shapes.add_shape(
            _MSO_RECT, Emu(sl), Emu(st), Emu(_FORK_W), Emu(_FORK_H)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _FORK_COLOR
        shape.line.fill.background()
        if shape.has_text_frame:
            shape.text_frame.text = ""
        return shape

    def _draw_normal_state(
        self, slide: Slide, label: str, cx: int, cy: int
    ) -> object:
        """通常状態（角丸矩形）を白背景・黒枠・ラベルありで描画する。"""
        sl = cx - _STATE_W // 2
        st = cy - _STATE_H // 2
        shape = slide.shapes.add_shape(
            _MSO_ROUNDED, Emu(sl), Emu(st), Emu(_STATE_W), Emu(_STATE_H)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _STATE_BG
        shape.line.color.rgb = _STATE_BORDER
        shape.text = label
        tf = shape.text_frame
        tf.word_wrap = False
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = _BLACK
        return shape

    def _draw_composite_bg(
        self,
        slide: Slide,
        label: str,
        cl: int,
        ct: int,
        cw: int,
        ch: int,
        is_divider: bool = False,
    ) -> object:
        """
        複合状態の外枠（背景角丸矩形）とタイトルテキストを描画する。

        is_divider=True の場合は並行区画として破線枠で描画し、タイトルは省略する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        label : str
            表示ラベル（複合状態名）。
        cl, ct, cw, ch : int
            外枠の左/上/幅/高さ（EMU）。
        is_divider : bool
            True のとき divider（並行区画）として破線枠スタイルで描画する。

        Returns
        -------
        Shape
            外枠シェイプ（コネクター接続用）。
        """
        shape = slide.shapes.add_shape(
            _MSO_ROUNDED, Emu(cl), Emu(ct), Emu(cw), Emu(ch)
        )
        if is_divider:
            # 並行区画: 破線枠・やや異なる背景色・タイトルなし
            shape.fill.solid()
            shape.fill.fore_color.rgb = _DIVIDER_BG
            # OOXML で破線枠を設定する
            self._apply_dashed_border(shape, _COMPOSITE_BORDER)
            if shape.has_text_frame:
                shape.text_frame.text = ""
        else:
            # 通常複合状態: 実線枠・タイトルあり
            shape.fill.solid()
            shape.fill.fore_color.rgb = _COMPOSITE_BG
            shape.line.color.rgb = _COMPOSITE_BORDER
            if shape.has_text_frame:
                shape.text_frame.text = ""
            # タイトルを上部テキストボックスとして配置する
            if label:
                tb = slide.shapes.add_textbox(
                    Emu(cl), Emu(ct), Emu(cw), Emu(_COMPOSITE_TITLE_H)
                )
                tf = tb.text_frame
                tf.word_wrap = False
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = label
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = _COMPOSITE_BORDER

        return shape

    def _apply_dashed_border(self, shape: object, color: RGBColor) -> None:
        """
        OOXML の <a:ln> に破線スタイルと色を設定する。

        Parameters
        ----------
        shape : Shape
            境界線を変更するシェイプ。
        color : RGBColor
            破線の色。
        """
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            return
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr, qn("a:ln"))
        # 既存の塗りを削除して solidFill を再設定する
        for old in ln.findall(qn("a:solidFill")):
            ln.remove(old)
        solid = lxml_etree.SubElement(ln, qn("a:solidFill"))
        srgb = lxml_etree.SubElement(solid, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02x}{color[1]:02x}{color[2]:02x}")
        # 破線スタイル
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        dash = lxml_etree.SubElement(ln, qn("a:prstDash"))
        dash.set("val", "dash")

    # ------------------------------------------------------------------
    # ノート描画
    # ------------------------------------------------------------------

    def _draw_note_nodes(
        self,
        slide: Slide,
        raw_nodes: list[dict],
        node_center_emu: dict[str, tuple[int, int]],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        noteGroup ノードをターゲット状態の近くに黄色矩形として描画する。

        ノートの ID パターン "{targetStateId}----parent" からターゲット状態を特定し、
        position="right of" / "left of" に応じてオフセット配置する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        raw_nodes : list[dict]
            全ノードリスト。
        node_center_emu : dict
            ノードID → 中心EMU座標 の辞書。
        left, top, width, height : int
            描画エリアのEMU座標（ターゲットが見つからない場合のフォールバック用）。
        """
        seen: set[str] = set()
        for node in raw_nodes:
            if node.get("shape") != "noteGroup":
                continue
            note_id = node["id"]
            if note_id in seen:
                continue
            seen.add(note_id)

            label = node.get("label", "")
            position = node.get("position", "right of")

            # ID パターン "A----parent" からターゲット状態 A を特定する
            target_id = note_id.split("----")[0] if "----" in note_id else ""
            if target_id and target_id in node_center_emu:
                tcx, tcy = node_center_emu[target_id]
            elif note_id in node_center_emu:
                tcx, tcy = node_center_emu[note_id]
            else:
                # フォールバック: 描画エリア右端中央に配置する
                tcx = left + width - _STATE_W
                tcy = top + height // 2

            self._draw_note(slide, label, tcx, tcy, position)

    def _draw_note(
        self,
        slide: Slide,
        label: str,
        target_cx: int,
        target_cy: int,
        position: str = "right of",
    ) -> object:
        """
        ノートを黄色矩形でターゲット状態の右または左に描画する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        label : str
            ノートテキスト。
        target_cx, target_cy : int
            ターゲット状態の中心座標（EMU）。
        position : str
            "left of" の場合は左側、それ以外は右側に配置する。

        Returns
        -------
        Shape
            描画したシェイプ。
        """
        note_w = _STATE_W
        note_h = _STATE_H
        gap = 250_000  # ターゲット状態との水平ギャップ（EMU）

        if "left" in position:
            nl = target_cx - _STATE_W // 2 - note_w - gap
        else:
            nl = target_cx + _STATE_W // 2 + gap
        nt = target_cy - note_h // 2

        shape = slide.shapes.add_shape(
            _MSO_RECT, Emu(nl), Emu(nt), Emu(note_w), Emu(note_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _NOTE_BG
        shape.line.color.rgb = _NOTE_BORDER
        shape.text = label
        tf = shape.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = _BLACK
        return shape

    # ------------------------------------------------------------------
    # エッジ描画
    # ------------------------------------------------------------------

    def _draw_state_edges(
        self,
        slide: Slide,
        raw_edges: list[dict],
        node_center_emu: dict[str, tuple[int, int]],
        shape_map: dict[str, object],
    ) -> None:
        """
        遷移（エッジ）を曲線コネクター（curvedConnector3）で描画する。

        ラベルが設定されている場合はコネクター中点近くにテキストボックスを追加する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        raw_edges : list[dict]
            edges リスト（start/end/label を持つ辞書）。
        node_center_emu : dict[str, tuple[int, int]]
            ノードIDをキー、中心座標（EMU）を値とする辞書。
        shape_map : dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        """
        for edge in raw_edges:
            if not isinstance(edge, dict):
                continue
            src = edge.get("start", "")
            dst = edge.get("end", "")
            if not src or not dst:
                continue
            if src not in node_center_emu or dst not in node_center_emu:
                continue

            sx, sy = node_center_emu[src]
            dx, dy = node_center_emu[dst]

            # 曲線コネクター（curvedConnector3）を追加する
            connector = slide.shapes.add_connector(
                3,  # MSO_CONNECTOR_TYPE.CURVE
                Emu(sx), Emu(sy), Emu(dx), Emu(dy),
            )

            # 接続ポイントをノードシェイプに接続する
            vec_x = dx - sx
            vec_y = dy - sy
            mag = math.sqrt(vec_x ** 2 + vec_y ** 2) or 1.0
            src_cp, dst_cp = self._connection_indices(vec_x / mag, vec_y / mag)
            src_shape = shape_map.get(src)
            dst_shape = shape_map.get(dst)
            if src_shape is not None:
                connector.begin_connect(src_shape, src_cp)
            if dst_shape is not None:
                connector.end_connect(dst_shape, dst_cp)

            # 矢印スタイル（headEnd: arrow）を設定する
            self._apply_arrow_to_connector(connector)

            # 遷移ラベルをコネクター中点付近に配置する
            label = edge.get("label", "")
            if label:
                self._add_transition_label(slide, label, sx, sy, dx, dy)

    def _apply_arrow_to_connector(self, connector: object) -> None:
        """
        コネクターの OOXML <a:ln> に headEnd 矢印スタイルを設定する。

        Parameters
        ----------
        connector : object
            python-pptx のコネクターシェイプ。
        """
        cxn_el = connector._element
        spPr = cxn_el.find(qn("p:spPr"))
        if spPr is None:
            return
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = lxml_etree.SubElement(spPr, qn("a:ln"))
        head = lxml_etree.SubElement(ln, qn("a:headEnd"))
        head.set("type", "arrow")
        head.set("w", "med")
        head.set("len", "med")

    def _add_transition_label(
        self,
        slide: Slide,
        label: str,
        sx: int,
        sy: int,
        dx: int,
        dy: int,
    ) -> None:
        """
        遷移ラベルをコネクター中点近くにテキストボックスで追加する。

        Parameters
        ----------
        slide : Slide
            描画対象スライド。
        label : str
            ラベルテキスト。
        sx, sy : int
            コネクター始点の EMU 座標。
        dx, dy : int
            コネクター終点の EMU 座標。
        """
        mx = (sx + dx) // 2
        my = (sy + dy) // 2
        tb_w = 1_000_000
        tb_h = 300_000
        tb_left = mx - tb_w // 2
        tb_top = my - tb_h // 2

        txbox = slide.shapes.add_textbox(
            Emu(tb_left), Emu(tb_top), Emu(tb_w), Emu(tb_h)
        )
        tf = txbox.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = _LABEL_FG
