"""
Mermaid タイムライン（timeline）レンダラーモジュール。

TimelineData オブジェクトを入力として PowerPoint スライドにタイムライン図を描画する。

レイアウト構成（上から下）:
  1. タイトル行（全幅テキストボックス、タイトルが空の場合は省略）
  2. セクションラベル帯（section が存在する場合のみ）
  3. イベントカードエリア（各 period のイベントを縦積み）
  4. 軸ライン + period 円ノード
  5. period ラベルエリア（軸の下）

主な仕様:

  水平軸:
    スライド幅を periods 数で等分した列幅（PERIOD_W）で各 period を配置する。
    軸は細い矩形（全幅 × 高さ 3px 相当）で表現する。

  period 円ノード:
    軸の中央に円を配置する。色は section ベース色（section なしの場合は period インデックス色）。
    section あり: 同一 section 内の period は同じベース色を使用する。
    section なし: period ごとに 12 色パレットを循環させる。

  イベントカード:
    period 円から上方向に event を縦積みする。
    各カードは角丸矩形（ベース色を白と 55% 混合した薄色、枠線はベース色）。
    テキストは自動折り返し対応（TextFrame の word_wrap=True）。
    カード内の \\n は段落区切りとして描画する。

  セクション帯:
    section が存在する場合、section 内の最初〜最後の period 列を跨ぐ帯を描画する。
    背景色はベース色を白と 75% 混合した薄色。
    帯の上部に section ラベルを配置する。

  period ラベル:
    period 円の下にテキストボックスを配置する。テキストは自動折り返し。

カラーパレット（Mermaid cScale0〜11 に対応): 12 色循環。
枠線・軸ライン色: ベース色を 20% 暗くした色を使用する。

OOXML 制約への代替案:
  - CSS テーマ変数（cScale0〜11）→ 固定 12 色パレットで近似
  - disableMulticolor オプション → 常にマルチカラー（PPTX では実用上問題なし）
  - 動的コンテンツ依存列幅 → 均等列幅
  - SVG 曲線コネクタ → 直線コネクタ
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .timeline_parser import TimelineData, TimelinePeriod

# ---------------------------------------------------------------------------
# カラーパレット（Mermaid cScale0〜11 に近似した 12 色）
# ---------------------------------------------------------------------------
_SCALE_COLORS: list[tuple[int, int, int]] = [
    (255, 179,  71),   #  0: オレンジ
    ( 88, 196, 150),   #  1: 緑
    (102, 178, 255),   #  2: 青
    (255, 111, 105),   #  3: 赤
    (180, 136, 209),   #  4: 紫
    (255, 207, 140),   #  5: 黄
    (100, 192, 192),   #  6: シアン
    (255, 160, 180),   #  7: ピンク
    (150, 200, 220),   #  8: 水色
    (200, 230, 150),   #  9: 黄緑
    (220, 150, 100),   # 10: テラコッタ
    (160, 180, 240),   # 11: 青紫
]

# ---------------------------------------------------------------------------
# レイアウト比率定数（全体高さに対する割合）
# ---------------------------------------------------------------------------
_TITLE_H_RATIO: float = 0.10       # タイトル行の高さ割合
_SECTION_H_RATIO: float = 0.13     # セクション帯の高さ割合（section が存在する場合）
_EVENT_AREA_RATIO: float = 0.42    # イベントカードエリアの高さ割合
_AXIS_AREA_RATIO: float = 0.14     # 軸＋円ノードエリアの高さ割合
_PERIOD_LABEL_RATIO: float = 0.31  # period ラベルエリアの高さ割合

# ---------------------------------------------------------------------------
# 各パーツのサイズ定数（EMU）
# ---------------------------------------------------------------------------
_CARD_GAP_H: int = 50_000          # イベントカード間の縦スペース
_CARD_GAP_W: int = 60_000          # カード左右マージン（列幅からの縮小量）
_CIRCLE_RATIO: float = 0.70        # 軸エリア内における period 円の高さ占有率
_MIN_CARD_H: int = 270_000         # イベントカードの最低高さ（EMU）
_MIN_CIRCLE_D: int = 200_000       # period 円の最低直径（EMU）
_SECTION_TOP_PAD: int = 30_000     # セクション帯内の上余白
_PERIOD_LABEL_PAD: int = 60_000    # period ラベルの上余白
_AXIS_LINE_H: int = 30_000         # 軸ラインの高さ（EMU）
_EVENT_BOTTOM_PAD: int = 80_000    # イベントエリア下端から軸中心までの余白

# テキストフォントサイズ
_TITLE_FONT_PT: float = 20.0
_SECTION_FONT_PT: float = 11.0
_EVENT_FONT_PT: float = 9.0
_PERIOD_FONT_PT: float = 10.0


# ---------------------------------------------------------------------------
# カラーユーティリティ
# ---------------------------------------------------------------------------

def _lighten(rgb: tuple[int, int, int], factor: float) -> tuple[int, int, int]:
    """
    ベース色を白（255,255,255）と混合して明るくする。

    Parameters
    ----------
    rgb : tuple[int, int, int]
        元の RGB 値。
    factor : float
        白の混合率（0.0=元の色、1.0=完全な白）。

    Returns
    -------
    tuple[int, int, int]
        明るくなった RGB 値。
    """
    r, g, b = rgb
    return (
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor),
    )


def _darken(rgb: tuple[int, int, int], factor: float = 0.8) -> tuple[int, int, int]:
    """
    ベース色を暗くする（枠線色の計算に使用）。

    Parameters
    ----------
    rgb : tuple[int, int, int]
        元の RGB 値。
    factor : float
        暗さ係数（1.0=変化なし、0.0=黒）。

    Returns
    -------
    tuple[int, int, int]
        暗くなった RGB 値。
    """
    r, g, b = rgb
    return (int(r * factor), int(g * factor), int(b * factor))


def _to_rgb(rgb: tuple[int, int, int]) -> RGBColor:
    """(r, g, b) タプルから RGBColor を生成する。"""
    return RGBColor(rgb[0], rgb[1], rgb[2])


# ---------------------------------------------------------------------------
# シェープ描画ユーティリティ
# ---------------------------------------------------------------------------

_SHAPE_RECT: int = 1           # MSO_AUTO_SHAPE_TYPE.RECTANGLE
_SHAPE_ROUNDED_RECT: int = 5   # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
_SHAPE_OVAL: int = 9           # MSO_AUTO_SHAPE_TYPE.OVAL


def _add_shape(
    slide: Slide,
    shape_type: int,
    left: int,
    top: int,
    width: int,
    height: int,
) -> object:
    """
    指定した型のシェープをスライドに追加する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    shape_type : int
        MSO_AUTO_SHAPE_TYPE 番号。
    left, top, width, height : int
        EMU 座標・サイズ。

    Returns
    -------
    Shape
        追加したシェープオブジェクト。
    """
    return slide.shapes.add_shape(
        shape_type,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )


def _add_textbox(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
) -> object:
    """
    テキストボックスをスライドに追加する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    left, top, width, height : int
        EMU 座標・サイズ。

    Returns
    -------
    Shape
        追加したテキストボックスシェープ。
    """
    return slide.shapes.add_textbox(
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )


def _set_fill_solid(shape: object, rgb: tuple[int, int, int]) -> None:
    """シェープを単色塗りつぶしにする。"""
    shape.fill.solid()  # type: ignore[attr-defined]
    shape.fill.fore_color.rgb = _to_rgb(rgb)  # type: ignore[attr-defined]


def _set_fill_none(shape: object) -> None:
    """シェープの塗りつぶしを透明（なし）にする。"""
    shape.fill.background()  # type: ignore[attr-defined]


def _set_line(shape: object, rgb: tuple[int, int, int], pt: float = 1.5) -> None:
    """シェープの枠線色と幅を設定する。"""
    shape.line.color.rgb = _to_rgb(rgb)  # type: ignore[attr-defined]
    shape.line.width = Pt(pt)


def _no_line(shape: object) -> None:
    """シェープの枠線を非表示にする。"""
    shape.line.fill.background()  # type: ignore[attr-defined]


def _set_text_simple(
    shape: object,
    text: str,
    font_size: float,
    bold: bool = False,
    color: tuple[int, int, int] = (30, 30, 60),
    align: PP_ALIGN = PP_ALIGN.CENTER,
    wrap: bool = True,
) -> None:
    """
    シェープのテキストフレームにシンプルなテキストを設定する。

    \\n を段落区切りとして扱う。

    Parameters
    ----------
    shape : object
        テキストフレームを持つシェープ。
    text : str
        設定するテキスト（\\n で段落区切り）。
    font_size : float
        フォントサイズ（ポイント）。
    bold : bool
        太字フラグ。
    color : tuple[int, int, int]
        文字色（RGB）。
    align : PP_ALIGN
        段落の水平アライメント。
    wrap : bool
        テキスト自動折り返し。
    """
    tf = shape.text_frame  # type: ignore[attr-defined]
    tf.word_wrap = wrap
    tf.auto_size = None  # type: ignore[assignment]

    lines = text.split("\n")
    # 最初の段落を設定（既存の paragraph[0] を使う）
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _to_rgb(color)


def _add_event_card(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    base_color: tuple[int, int, int],
) -> None:
    """
    イベントカード（角丸矩形）をスライドに描画する。

    背景はベース色を白と 55% 混合した薄色、
    枠線はベース色を 20% 暗くした色で描画する。

    Parameters
    ----------
    slide : Slide
        描画先スライド。
    left, top, width, height : int
        カードの EMU 座標・サイズ。
    text : str
        カード内テキスト（\\n を段落区切りとして扱う）。
    base_color : tuple[int, int, int]
        section/period のベース色。
    """
    bg_color = _lighten(base_color, 0.55)
    border_color = _darken(base_color, 0.80)

    card = _add_shape(slide, _SHAPE_ROUNDED_RECT, left, top, width, height)
    _set_fill_solid(card, bg_color)
    _set_line(card, border_color, pt=1.0)

    # 角丸の半径を小さめに設定（adjustments[0] = 0 〜 50000）
    card.adjustments[0] = 20000  # type: ignore[index]

    _set_text_simple(
        card, text,
        font_size=_EVENT_FONT_PT,
        color=(30, 30, 60),
        align=PP_ALIGN.CENTER,
    )
    card.text_frame.margin_top = Emu(30_000)  # type: ignore[attr-defined]
    card.text_frame.margin_bottom = Emu(30_000)  # type: ignore[attr-defined]
    card.text_frame.margin_left = Emu(40_000)  # type: ignore[attr-defined]
    card.text_frame.margin_right = Emu(40_000)  # type: ignore[attr-defined]


# ---------------------------------------------------------------------------
# カラー割り当てロジック
# ---------------------------------------------------------------------------

def _build_color_map(data: TimelineData) -> dict[int, tuple[int, int, int]]:
    """
    各 period インデックスにベース色を割り当てる辞書を構築する。

    section が存在する場合: section 単位で色を付ける（同 section = 同色）。
    section が存在しない場合: period ごとに 12 色パレットを循環させる。

    Parameters
    ----------
    data : TimelineData
        パース済みタイムラインデータ。

    Returns
    -------
    dict[int, tuple[int, int, int]]
        period インデックス → RGB タプルの辞書。
    """
    color_map: dict[int, tuple[int, int, int]] = {}

    if data.sections:
        # section ごとに色インデックスを割り当てる
        section_color_idx: dict[str, int] = {}
        color_counter = 0
        for i, period in enumerate(data.periods):
            sec = period.section
            if sec is None:
                # section 未定義の period は個別に色を付ける
                color_map[i] = _SCALE_COLORS[color_counter % len(_SCALE_COLORS)]
                color_counter += 1
            else:
                if sec not in section_color_idx:
                    section_color_idx[sec] = color_counter % len(_SCALE_COLORS)
                    color_counter += 1
                color_map[i] = _SCALE_COLORS[section_color_idx[sec]]
    else:
        # section なし: period ごとに循環
        for i in range(len(data.periods)):
            color_map[i] = _SCALE_COLORS[i % len(_SCALE_COLORS)]

    return color_map


# ---------------------------------------------------------------------------
# メインレンダラークラス
# ---------------------------------------------------------------------------

class TimelineRenderer:
    """
    TimelineData データクラスを受け取り、PowerPoint スライドにタイムラインを描画するクラス。

    描画は render() メソッドで行う。
    """

    def render(
        self,
        slide: Slide,
        data: TimelineData,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        TimelineData をスライドに描画する。

        period が 0 件の場合は何も描画しない。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        data : TimelineData
            parse_timeline() で生成されたタイムラインデータ。
        left, top, width, height : int
            描画エリアの EMU 座標。
        """
        if not data.periods:
            return

        # --- エリア分割 ---
        title_h = int(height * _TITLE_H_RATIO) if data.title else 0
        body_top = top + title_h
        body_h = height - title_h

        has_sections = bool(data.sections)
        section_h = int(body_h * _SECTION_H_RATIO) if has_sections else 0
        event_area_h = int(body_h * _EVENT_AREA_RATIO)
        axis_area_h = int(body_h * _AXIS_AREA_RATIO)
        period_label_h = body_h - section_h - event_area_h - axis_area_h

        section_top = body_top
        event_area_top = body_top + section_h
        axis_area_top = event_area_top + event_area_h
        period_label_top = axis_area_top + axis_area_h

        # 軸ラインの Y 中心
        axis_cy = axis_area_top + axis_area_h // 2

        # --- 列幅計算 ---
        n = len(data.periods)
        col_w = width // n if n > 0 else width

        # period 円のサイズ
        circle_d = max(_MIN_CIRCLE_D, int(axis_area_h * _CIRCLE_RATIO))
        circle_r = circle_d // 2

        # イベントカードのサイズ計算
        max_events = max((len(p.events) for p in data.periods), default=0)
        if max_events > 0:
            available_card_h = event_area_h - _EVENT_BOTTOM_PAD
            card_h = max(
                _MIN_CARD_H,
                (available_card_h - _CARD_GAP_H * (max_events - 1)) // max_events,
            )
        else:
            card_h = _MIN_CARD_H
        card_w = col_w - _CARD_GAP_W * 2

        # --- カラーマップ構築 ---
        color_map = _build_color_map(data)

        # --- タイトル描画 ---
        if data.title:
            self._render_title(slide, data.title, left, top, width, title_h)

        # --- 軸ライン描画 ---
        self._render_axis_line(slide, left, axis_cy, width)

        # --- セクション帯描画 ---
        if has_sections:
            self._render_sections(
                slide, data,
                color_map,
                left, section_top, col_w, section_h,
            )

        # --- 各 period の描画 ---
        for i, period in enumerate(data.periods):
            cx = left + i * col_w + col_w // 2  # period 列の中央 X

            base_color = color_map[i]

            # period 円ノード
            self._render_period_circle(
                slide, cx, axis_cy, circle_r, base_color,
            )

            # イベントカード（軸の上方向に積み上げる）
            self._render_events(
                slide, period, cx, col_w, card_w, card_h,
                event_area_top, event_area_h,
                _EVENT_BOTTOM_PAD, base_color,
            )

            # period ラベル
            self._render_period_label(
                slide, period.label, cx, col_w,
                period_label_top, period_label_h,
            )

    # -----------------------------------------------------------------------
    # 各要素の描画メソッド
    # -----------------------------------------------------------------------

    def _render_title(
        self,
        slide: Slide,
        title: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """タイトルテキストボックスを描画する。"""
        tb = _add_textbox(slide, left, top, width, height)
        _set_fill_none(tb)
        _no_line(tb)
        _set_text_simple(
            tb, title,
            font_size=_TITLE_FONT_PT,
            bold=True,
            color=(30, 30, 60),
            align=PP_ALIGN.CENTER,
        )
        tb.text_frame.margin_top = Emu(20_000)  # type: ignore[attr-defined]

    def _render_axis_line(
        self,
        slide: Slide,
        left: int,
        axis_cy: int,
        width: int,
    ) -> None:
        """水平軸ラインを細い矩形で描画する。"""
        ax = _add_shape(
            slide, _SHAPE_RECT,
            left,
            axis_cy - _AXIS_LINE_H // 2,
            width,
            _AXIS_LINE_H,
        )
        _set_fill_solid(ax, (180, 180, 200))
        _no_line(ax)

    def _render_sections(
        self,
        slide: Slide,
        data: TimelineData,
        color_map: dict[int, tuple[int, int, int]],
        left: int,
        section_top: int,
        col_w: int,
        section_h: int,
    ) -> None:
        """
        セクション帯をスライドに描画する。

        各セクションが跨ぐ列範囲を特定し、背景帯と section ラベルを描画する。

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        data : TimelineData
            タイムラインデータ。
        color_map : dict[int, tuple]
            period インデックス → ベース色。
        left : int
            描画エリアの左端 EMU。
        section_top : int
            セクション帯の上端 EMU。
        col_w : int
            1列あたりの幅（EMU）。
        section_h : int
            セクション帯の高さ（EMU）。
        """
        # section 名 → 最初・最後の period インデックスを収集
        section_indices: dict[str, list[int]] = {}
        for i, period in enumerate(data.periods):
            if period.section is not None:
                sec = period.section
                if sec not in section_indices:
                    section_indices[sec] = []
                section_indices[sec].append(i)

        for sec_name, indices in section_indices.items():
            first_idx = indices[0]
            last_idx = indices[-1]

            # 帯の色はこのセクションの最初の period の色を使う
            base_color = color_map[first_idx]
            band_color = _lighten(base_color, 0.70)
            border_color = _darken(base_color, 0.85)

            band_left = left + first_idx * col_w
            band_width = (last_idx - first_idx + 1) * col_w

            band = _add_shape(
                slide, _SHAPE_ROUNDED_RECT,
                band_left, section_top,
                band_width, section_h,
            )
            _set_fill_solid(band, band_color)
            _set_line(band, border_color, pt=1.0)
            # 角丸を小さめに
            band.adjustments[0] = 15000  # type: ignore[index]

            # section ラベル
            _set_text_simple(
                band,
                _process_br_for_display(sec_name),
                font_size=_SECTION_FONT_PT,
                bold=True,
                color=_darken(base_color, 0.60),
                align=PP_ALIGN.CENTER,
            )
            band.text_frame.margin_top = Emu(_SECTION_TOP_PAD)  # type: ignore[attr-defined]

    def _render_period_circle(
        self,
        slide: Slide,
        cx: int,
        cy: int,
        radius: int,
        base_color: tuple[int, int, int],
    ) -> None:
        """
        period 円ノードを軸上に描画する。

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        cx : int
            円の中心 X（EMU）。
        cy : int
            円の中心 Y（EMU）。
        radius : int
            円の半径（EMU）。
        base_color : tuple[int, int, int]
            section/period のベース色。
        """
        d = radius * 2
        circle = _add_shape(
            slide, _SHAPE_OVAL,
            cx - radius, cy - radius,
            d, d,
        )
        _set_fill_solid(circle, base_color)
        border_color = _darken(base_color, 0.80)
        _set_line(circle, border_color, pt=1.5)

    def _render_events(
        self,
        slide: Slide,
        period: TimelinePeriod,
        cx: int,
        col_w: int,
        card_w: int,
        card_h: int,
        event_area_top: int,
        event_area_h: int,
        bottom_pad: int,
        base_color: tuple[int, int, int],
    ) -> None:
        """
        period のイベントカードを軸上方に縦積みで描画する。

        イベントは軸に近い順（events[0] が軸に最も近い）で下から上へ積み上げる。

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        period : TimelinePeriod
            イベントの親 period。
        cx : int
            列の中央 X（EMU）。
        col_w : int
            列幅（EMU）。
        card_w : int
            カードの幅（EMU）。
        card_h : int
            カードの高さ（EMU）。
        event_area_top : int
            イベントエリアの上端 Y（EMU）。
        event_area_h : int
            イベントエリアの高さ（EMU）。
        bottom_pad : int
            イベントエリア下端から軸中心までの余白（EMU）。
        base_color : tuple[int, int, int]
            section/period のベース色。
        """
        # カード左端 X
        card_left = cx - card_w // 2

        # 軸側（下側）から積み上げるため、最も下のカードの top を計算する
        # event_area の下端 = event_area_top + event_area_h - bottom_pad
        area_bottom = event_area_top + event_area_h - bottom_pad

        for j, event in enumerate(period.events):
            # 下から j 番目のカード top
            card_top = area_bottom - (j + 1) * card_h - j * _CARD_GAP_H
            if card_top < event_area_top:
                # イベントエリアを超えてしまう場合はクリップ
                card_top = event_area_top

            _add_event_card(
                slide,
                card_left, card_top,
                card_w, card_h,
                event.text,
                base_color,
            )

    def _render_period_label(
        self,
        slide: Slide,
        label: str,
        cx: int,
        col_w: int,
        period_label_top: int,
        period_label_h: int,
    ) -> None:
        """
        period ラベルテキストボックスを軸の下に描画する。

        Parameters
        ----------
        slide : Slide
            描画先スライド。
        label : str
            period ラベルテキスト（\\n を含む場合は段落区切り）。
        cx : int
            列の中央 X（EMU）。
        col_w : int
            列幅（EMU）。
        period_label_top : int
            period ラベルエリアの上端 Y（EMU）。
        period_label_h : int
            period ラベルエリアの高さ（EMU）。
        """
        label_w = col_w - _CARD_GAP_W
        label_left = cx - label_w // 2
        tb = _add_textbox(
            slide,
            label_left,
            period_label_top + _PERIOD_LABEL_PAD,
            label_w,
            period_label_h - _PERIOD_LABEL_PAD,
        )
        _set_fill_none(tb)
        _no_line(tb)
        _set_text_simple(
            tb, label,
            font_size=_PERIOD_FONT_PT,
            bold=True,
            color=(30, 30, 60),
            align=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# 内部ユーティリティ（セクション名のテキスト処理）
# ---------------------------------------------------------------------------

def _process_br_for_display(text: str) -> str:
    """
    改行（\\n）をそのまま返す（_set_text_simple が \\n を段落区切りとして扱う）。

    既に parse_timeline 内で <br> → \\n 変換済みのため、ここでは何もしない。
    """
    return text
