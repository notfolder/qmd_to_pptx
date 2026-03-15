"""
Mermaidレンダラーモジュール。

Mermaid記法のテキストをパースしてグラフ構造を取得し、
python-pptxのShapeとして現在のスライドに図を描画する。
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

import networkx as nx
from mermaid_parser import MermaidParser
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt


# ノードのデフォルトサイズ（EMU）
_NODE_WIDTH_EMU: int = 1200000
_NODE_HEIGHT_EMU: int = 500000

# クラス図ボックスのサイズ定数（EMU）
_CLASS_WIDTH_EMU: int = 1700000          # クラスボックスの幅
_CLASS_HEADER_HEIGHT_EMU: int = 440000   # クラス名ヘッダーの高さ
_CLASS_ROW_HEIGHT_EMU: int = 280000      # メンバー1行あたりの高さ
_CLASS_MIN_SECTION_HEIGHT_EMU: int = 320000  # 属性/メソッドセクションの最小高さ


class MermaidRenderer:
    """
    Mermaidレンダラークラス。

    Mermaidテキストをmermaid-parser-pyで解析してNetworkXグラフに変換し、
    spring_layoutで座標計算後、python-pptxのShape/Connectorとして描画する。
    """

    def render(
        self,
        slide: Slide,
        element: ET.Element,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        elementからMermaidテキストを取り出し、スライドにグラフを描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        element : ET.Element
            Mermaidコード要素（code class="language-mermaid"）。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        mermaid_text = "".join(element.itertext()).strip()

        # JSエンジンが対応していないダイアグラムタイプはパーサーを呼ばずフォールバックする
        # これにより mermaid-parser-py の JS エンジンが stderr に出力するエラーを抑止する
        _UNSUPPORTED_PREFIXES = (
            "zenuml",
        )
        first_line = mermaid_text.splitlines()[0].strip().lower() if mermaid_text else ""
        if any(first_line.startswith(p) for p in _UNSUPPORTED_PREFIXES):
            self._render_fallback(slide, mermaid_text, left, top, width, height)
            return

        try:
            # mermaid-parser-pyでノードとエッジを取得する
            mp = MermaidParser()
            result = mp.parse(mermaid_text)
            graph_data = result.get("graph_data", {})
            graph_type = result.get("graph_type", "")
        except Exception:
            # パース失敗時はテキストボックスにそのまま表示する
            self._render_fallback(slide, mermaid_text, left, top, width, height)
            return

        # graph_typeに応じて専用レンダラーへ分岐する
        if graph_type == "stateDiagram":
            self._render_state_diagram(slide, graph_data, left, top, width, height)
            return
        if graph_type == "class":
            self._render_class_diagram(slide, graph_data, left, top, width, height)
            return
        if graph_type == "er":
            self._render_er_diagram(slide, graph_data, left, top, width, height)
            return
        if graph_type == "mindmap":
            self._render_mindmap(slide, graph_data, left, top, width, height)
            return

        # flowchart / graph 系: 頂点(vertices)とエッジから描画する
        nodes = self._extract_nodes(graph_data)
        edges = self._extract_edges(graph_data)

        if not nodes:
            self._render_fallback(slide, mermaid_text, left, top, width, height)
            return

        # NetworkXグラフを構築してspring_layoutで座標計算する
        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for src, dst in edges:
            G.add_edge(src, dst)

        # spring_layoutで正規化された座標（-1.0〜1.0）を計算する
        pos: dict[str, tuple[float, float]] = nx.spring_layout(G, seed=42)

        # 座標をEMUに変換してノードとエッジを描画する
        node_shapes = self._draw_nodes(slide, nodes, pos, left, top, width, height)
        self._draw_edges(slide, edges, pos, node_shapes, left, top, width, height)

    def _extract_nodes(self, graph_data: dict) -> list[str]:
        """
        mermaid-parser-pyの解析結果からノードIDのリストを取得する。

        Parameters
        ----------
        graph_data : dict
            graph_data辞書（"vertices"キーにノード情報を含む）。

        Returns
        -------
        list[str]
            ノードIDのリスト。
        """
        vertices = graph_data.get("vertices", {})
        if isinstance(vertices, dict):
            return list(vertices.keys())
        return []

    def _extract_edges(self, graph_data: dict) -> list[tuple[str, str]]:
        """
        mermaid-parser-pyの解析結果からエッジのリストを取得する。

        Parameters
        ----------
        graph_data : dict
            graph_data辞書（"edges"キーにエッジ情報を含む）。

        Returns
        -------
        list[tuple[str, str]]
            (始点ノードID, 終点ノードID) のタプルリスト。
        """
        raw_edges = graph_data.get("edges", [])
        edges: list[tuple[str, str]] = []
        for edge in raw_edges:
            if isinstance(edge, dict):
                src = edge.get("start")
                dst = edge.get("end")
                if src is not None and dst is not None:
                    edges.append((str(src), str(dst)))
        return edges

    def _pos_to_emu(
        self,
        x_norm: float,
        y_norm: float,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> tuple[int, int]:
        """
        spring_layoutで計算した正規化座標（-1.0〜1.0）をEMU座標に変換する。

        Parameters
        ----------
        x_norm : float
            正規化X座標（-1.0〜1.0）。
        y_norm : float
            正規化Y座標（-1.0〜1.0）。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。

        Returns
        -------
        tuple[int, int]
            (x_emu, y_emu) のタプル。
        """
        margin_x = _NODE_WIDTH_EMU // 2
        margin_y = _NODE_HEIGHT_EMU // 2
        usable_w = width - _NODE_WIDTH_EMU
        usable_h = height - _NODE_HEIGHT_EMU

        # -1.0〜1.0 を 0〜1 に正規化する
        x_ratio = (x_norm + 1.0) / 2.0
        y_ratio = (y_norm + 1.0) / 2.0

        x_emu = left + margin_x + int(x_ratio * usable_w)
        y_emu = top + margin_y + int(y_ratio * usable_h)
        return x_emu, y_emu

    def _draw_nodes(
        self,
        slide: Slide,
        nodes: list[str],
        pos: dict[str, tuple[float, float]],
        left: int,
        top: int,
        width: int,
        height: int,
        label_map: dict[str, str] | None = None,
    ) -> dict[str, object]:
        """
        ノードを矩形Shapeとしてスライドに配置する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        nodes : list[str]
            ノードIDのリスト。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        label_map : dict[str, str] | None
            ノードIDをキー、表示ラベル文字列を値とする辞書。省略時はノードIDを表示する。

        Returns
        -------
        dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        """
        node_shapes: dict[str, object] = {}
        for node_id in nodes:
            if node_id not in pos:
                continue
            x_norm, y_norm = pos[node_id]
            cx, cy = self._pos_to_emu(x_norm, y_norm, left, top, width, height)
            shape_left = cx - _NODE_WIDTH_EMU // 2
            shape_top = cy - _NODE_HEIGHT_EMU // 2

            # add_shape で矩形を追加する
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Emu(shape_left),
                Emu(shape_top),
                Emu(_NODE_WIDTH_EMU),
                Emu(_NODE_HEIGHT_EMU),
            )
            # ノードの表示ラベルを決定する（label_mapがあればそれを優先する）
            label = label_map.get(node_id, node_id) if label_map else node_id
            shape.text = label
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)

            node_shapes[node_id] = shape

        return node_shapes

    def _connection_indices(
        self,
        dx: float,
        dy: float,
    ) -> tuple[int, int]:
        """
        方向ベクトル (dx, dy) から始点・終点の接続ポイントインデックスを決定する。

        DrawingML 標準矩形の接続ポイントは以下のとおり。
        - 0: 上辺中点
        - 1: 右辺中点
        - 2: 下辺中点
        - 3: 左辺中点

        Parameters
        ----------
        dx : float
            X方向の差（dst.x - src.x）。正規化座標系での値。
        dy : float
            Y方向の差（dst.y - src.y）。正規化座標系での値。

        Returns
        -------
        tuple[int, int]
            (始点接続ポイントインデックス, 終点接続ポイントインデックス)
        """
        if abs(dx) >= abs(dy):
            # 水平方向が支配的な場合
            if dx >= 0:
                # 左から右: src右辺 → dst左辺
                return (1, 3)
            else:
                # 右から左: src左辺 → dst右辺
                return (3, 1)
        else:
            # 垂直方向が支配的な場合
            if dy >= 0:
                # 上から下: src下辺 → dst上辺
                return (2, 0)
            else:
                # 下から上: src上辺 → dst下辺
                return (0, 2)

    def _draw_edges(
        self,
        slide: Slide,
        edges: list[tuple[str, str]],
        pos: dict[str, tuple[float, float]],
        node_shapes: dict[str, object],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        エッジをConnectorShapeとしてスライドに描画し、始点・終点のshapeに接続する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        edges : list[tuple[str, str]]
            (始点ノードID, 終点ノードID) のタプルリスト。
        pos : dict[str, tuple[float, float]]
            ノードIDをキー、正規化座標を値とする辞書。
        node_shapes : dict[str, object]
            ノードIDをキー、Shapeオブジェクトを値とする辞書。
        left : int
            描画エリアの左端座標（EMU）。
        top : int
            描画エリアの上端座標（EMU）。
        width : int
            描画エリアの幅（EMU）。
        height : int
            描画エリアの高さ（EMU）。
        """
        for src, dst in edges:
            if src not in pos or dst not in pos:
                continue
            sx_norm, sy_norm = pos[src]
            dx_norm, dy_norm = pos[dst]
            sx, sy = self._pos_to_emu(sx_norm, sy_norm, left, top, width, height)
            dx, dy = self._pos_to_emu(dx_norm, dy_norm, left, top, width, height)

            # コネクターを追加する（直線コネクター: connector_type=1）
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                Emu(sx),
                Emu(sy),
                Emu(dx),
                Emu(dy),
            )

            # 接続ポイントインデックスを方向ベクトルから決定する
            src_cp, dst_cp = self._connection_indices(
                dx_norm - sx_norm, dy_norm - sy_norm
            )

            # shapeが存在する場合は begin_connect/end_connect でブロックに接続する
            src_shape = node_shapes.get(src)
            dst_shape = node_shapes.get(dst)
            if src_shape is not None:
                connector.begin_connect(src_shape, src_cp)
            if dst_shape is not None:
                connector.end_connect(dst_shape, dst_cp)

    def _render_state_diagram(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        stateDiagram-v2 を nodes リストと edges リストから描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"nodes"と"edges"キーを含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        raw_nodes = graph_data.get("nodes", [])
        raw_edges = graph_data.get("edges", [])
        if not raw_nodes:
            return

        # ノードIDリストとラベルマップを構築する
        nodes = [n["id"] for n in raw_nodes if "id" in n]
        label_map = {n["id"]: n.get("label", n["id"]) for n in raw_nodes}
        edges = [
            (e["start"], e["end"])
            for e in raw_edges
            if "start" in e and "end" in e
        ]

        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for src, dst in edges:
            G.add_edge(src, dst)
        pos: dict[str, tuple[float, float]] = nx.spring_layout(G, seed=42)

        node_shapes = self._draw_nodes(
            slide, nodes, pos, left, top, width, height, label_map=label_map
        )
        self._draw_edges(slide, edges, pos, node_shapes, left, top, width, height)

    def _render_class_diagram(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        classDiagram を UMLクラス図形式（ヘッダー・属性・メソッドの3段構成）で描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"classes"と"relations"キーを含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        classes = graph_data.get("classes", {})
        relations = graph_data.get("relations", [])
        if not classes:
            return

        nodes = list(classes.keys())
        edges = [
            (r["id1"], r["id2"])
            for r in relations
            if "id1" in r and "id2" in r
        ]

        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for src, dst in edges:
            G.add_edge(src, dst)
        pos: dict[str, tuple[float, float]] = nx.spring_layout(G, seed=42)

        # クラス図用レイアウト計算（クラス本体の高さに合わせます）
        max_class_h = max(
            (self._class_total_height(classes[nid]) for nid in nodes if nid in classes),
            default=_NODE_HEIGHT_EMU,
        )
        margin_x = _CLASS_WIDTH_EMU // 2
        margin_y = max_class_h // 2
        usable_w = max(width - _CLASS_WIDTH_EMU, _CLASS_WIDTH_EMU)
        usable_h = max(height - max_class_h, max_class_h)

        node_shapes: dict[str, object] = {}
        for node_id in nodes:
            if node_id not in pos or node_id not in classes:
                continue
            x_norm, y_norm = pos[node_id]
            x_ratio = (x_norm + 1.0) / 2.0
            y_ratio = (y_norm + 1.0) / 2.0
            cx = left + margin_x + int(x_ratio * usable_w)
            cy = top + margin_y + int(y_ratio * usable_h)
            anchor = self._draw_class_box(slide, node_id, classes[node_id], cx, cy)
            node_shapes[node_id] = anchor

        self._draw_edges(slide, edges, pos, node_shapes, left, top, width, height)

    def _class_total_height(self, class_data: dict) -> int:
        """
        クラスボックスの合計高さ（EMU）を返す。

        Parameters
        ----------
        class_data : dict
            mermaid-parser-pyが返すクラス情報辞書。

        Returns
        -------
        int
            ヘッダー・属性セクション・メソッドセクションの合計高さ（EMU）。
        """
        n_members = len(class_data.get("members", []))
        n_methods = len(class_data.get("methods", []))
        attrs_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, n_members * _CLASS_ROW_HEIGHT_EMU)
        methods_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, n_methods * _CLASS_ROW_HEIGHT_EMU)
        return _CLASS_HEADER_HEIGHT_EMU + attrs_h + methods_h

    def _draw_class_box(
        self,
        slide: Slide,
        node_id: str,
        class_data: dict,
        cx: int,
        cy: int,
    ) -> object:
        """
        UMLクラス図の1クラス分のボックスを描画してヘッダーShapeを返す。

        ヘッダー（クラス名）・属性セクション・メソッドセクションの3段構成の矩形を縦に重ねて配置する。
        コネクターのアンカーとしてヘッダーShapeを返す。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        node_id : str
            クラスID（クラス名）。
        class_data : dict
            mermaid-parser-pyが返すクラス情報辞書。
        cx : int
            クラスボックス中心のX座標（EMU）。
        cy : int
            クラスボックス中心のY座標（EMU）。

        Returns
        -------
        object
            ヘッダー部のShapeオブジェクト（コネクターのアンカーとして使用）。
        """
        members = class_data.get("members", [])
        methods = class_data.get("methods", [])
        attrs_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, len(members) * _CLASS_ROW_HEIGHT_EMU)
        methods_h = max(_CLASS_MIN_SECTION_HEIGHT_EMU, len(methods) * _CLASS_ROW_HEIGHT_EMU)
        total_h = _CLASS_HEADER_HEIGHT_EMU + attrs_h + methods_h

        box_left = cx - _CLASS_WIDTH_EMU // 2
        header_top = cy - total_h // 2

        # ---- ヘッダーセクション（クラス名・太字・中山寄せ）----
        header_shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(box_left), Emu(header_top),
            Emu(_CLASS_WIDTH_EMU), Emu(_CLASS_HEADER_HEIGHT_EMU),
        )
        tf = header_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(50000)
        tf.margin_right = Emu(50000)
        tf.margin_top = Emu(60000)
        tf.margin_bottom = Emu(60000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = class_data.get("label", node_id)
        run.font.bold = True
        run.font.size = Pt(11)

        # ---- 属性セクション ----
        attrs_top = header_top + _CLASS_HEADER_HEIGHT_EMU
        attrs_shape = slide.shapes.add_shape(
            1,
            Emu(box_left), Emu(attrs_top),
            Emu(_CLASS_WIDTH_EMU), Emu(attrs_h),
        )
        tf_a = attrs_shape.text_frame
        tf_a.word_wrap = True
        tf_a.margin_left = Emu(80000)
        tf_a.margin_right = Emu(50000)
        tf_a.margin_top = Emu(50000)
        tf_a.margin_bottom = Emu(50000)
        if members:
            first = True
            for member in members:
                # バックスラッシュエスケープを除去する
                raw = member.get("text", "").lstrip("\\")
                if first:
                    p = tf_a.paragraphs[0]
                    first = False
                else:
                    p = tf_a.add_paragraph()
                run = p.add_run()
                run.text = raw
                run.font.size = Pt(9)

        # ---- メソッドセクション ----
        methods_top = attrs_top + attrs_h
        methods_shape = slide.shapes.add_shape(
            1,
            Emu(box_left), Emu(methods_top),
            Emu(_CLASS_WIDTH_EMU), Emu(methods_h),
        )
        tf_m = methods_shape.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = Emu(80000)
        tf_m.margin_right = Emu(50000)
        tf_m.margin_top = Emu(50000)
        tf_m.margin_bottom = Emu(50000)
        if methods:
            first = True
            for method in methods:
                raw = method.get("text", "").lstrip("\\")
                if first:
                    p = tf_m.paragraphs[0]
                    first = False
                else:
                    p = tf_m.add_paragraph()
                run = p.add_run()
                run.text = raw
                run.font.size = Pt(9)

        return header_shape

    def _render_er_diagram(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        erDiagram を entities 辞書と relationships リストから描画する。

        entities のキーはラベル（エンティティ名）、値の id フィールドが内部ID。
        relationships は内部IDで参照するため逆引きマップを使用する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"entities"と"relationships"キーを含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        entities = graph_data.get("entities", {})
        relationships = graph_data.get("relationships", [])
        if not entities:
            return

        # 内部IDからエンティティ名への逆引きマップを構築する
        id_to_label: dict[str, str] = {
            ent.get("id", ""): label
            for label, ent in entities.items()
            if ent.get("id")
        }

        nodes = list(entities.keys())
        edges: list[tuple[str, str]] = [
            (id_to_label[rel["entityA"]], id_to_label[rel["entityB"]])
            for rel in relationships
            if rel.get("entityA") in id_to_label and rel.get("entityB") in id_to_label
        ]

        G = nx.DiGraph()
        G.add_nodes_from(nodes)
        for src, dst in edges:
            G.add_edge(src, dst)
        pos: dict[str, tuple[float, float]] = nx.spring_layout(G, seed=42)

        node_shapes = self._draw_nodes(slide, nodes, pos, left, top, width, height)
        self._draw_edges(slide, edges, pos, node_shapes, left, top, width, height)

    def _mindmap_collect_layout(
        self,
        node: dict,
        depth_map: dict[str, int],
        order_map: dict[str, float],
        counter: list[int],
    ) -> None:
        """
        マインドマップのツリーを後順走査して各ノードの深さと縦位置を収集する。

        Parameters
        ----------
        node : dict
            現在のノード（nodeId, level, children を持つ）。
        depth_map : dict[str, int]
            ノードIDをキーにした深さのマップ（ルートが1）。
        order_map : dict[str, float]
            ノードIDをキーにした縦位置（葉ノードは整数、内部ノードは子の中央）。
        counter : list[int]
            葉ノードの通し番号（リストで参照渡し）。
        """
        node_id = node.get("nodeId", "")
        if not node_id:
            return

        # レベル値は4の倍数（4=depth1, 8=depth2, 12=depth3）
        depth_map[node_id] = node.get("level", 4) // 4
        children = node.get("children", [])

        if not children:
            # 葉ノード: 縦位置としてカウンター値を割り当てる
            order_map[node_id] = float(counter[0])
            counter[0] += 1
        else:
            # 内部ノード: 先に子を走査してから子の縦位置の平均を取る
            for child in children:
                self._mindmap_collect_layout(child, depth_map, order_map, counter)
            child_orders = [
                order_map[c["nodeId"]]
                for c in children
                if c.get("nodeId") in order_map
            ]
            order_map[node_id] = (
                sum(child_orders) / len(child_orders) if child_orders else 0.0
            )

    def _mindmap_collect_edges(
        self,
        node: dict,
        edges: list[tuple[str, str]],
    ) -> None:
        """
        マインドマップのツリーを走査して親子エッジを収集する。

        Parameters
        ----------
        node : dict
            現在のノード。
        edges : list[tuple[str, str]]
            (親nodeId, 子nodeId) タプルを追加するリスト。
        """
        node_id = node.get("nodeId", "")
        for child in node.get("children", []):
            child_id = child.get("nodeId", "")
            if node_id and child_id:
                edges.append((node_id, child_id))
            self._mindmap_collect_edges(child, edges)

    def _render_mindmap(
        self,
        slide: Slide,
        graph_data: dict,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        mindmap をツリーレイアウトで左→右方向に描画する。

        ルートノードを左端に配置し、子ノードを深さに応じて右へ展開する。
        縦位置は葉ノードの数を基準に均等分配する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        graph_data : dict
            graph_data辞書（"nodes"キーにツリー構造を含む）。
        left, top, width, height : int
            描画エリアのEMU座標。
        """
        nodes_data = graph_data.get("nodes", [])
        if not nodes_data:
            return

        # 最初の要素がルートノード（children を含む完全ツリー）
        root = nodes_data[0]

        # ツリーを後順走査して深さと縦位置を収集する
        depth_map: dict[str, int] = {}
        order_map: dict[str, float] = {}
        counter: list[int] = [0]
        self._mindmap_collect_layout(root, depth_map, order_map, counter)

        if not depth_map:
            return

        # descr フィールドを表示ラベルとして使う
        label_map: dict[str, str] = {
            n.get("nodeId", ""): n.get("descr", n.get("nodeId", ""))
            for n in nodes_data
            if n.get("nodeId")
        }

        # 正規化座標（-1.0〜1.0）を計算する（深さ→X、縦位置→Y）
        max_depth = max(depth_map.values()) if depth_map else 1
        leaf_count = counter[0]  # 葉ノードの総数
        total = max(leaf_count - 1, 1)  # ゼロ除算を防ぐ

        pos: dict[str, tuple[float, float]] = {}
        for node_id, depth in depth_map.items():
            x_norm = (depth / max_depth) * 2.0 - 1.0
            y_val = order_map.get(node_id, 0.0)
            y_norm = (y_val / total) * 2.0 - 1.0
            pos[node_id] = (x_norm, y_norm)

        # 親子エッジを収集する
        edges: list[tuple[str, str]] = []
        self._mindmap_collect_edges(root, edges)

        all_node_ids = list(depth_map.keys())
        node_shapes = self._draw_nodes(
            slide, all_node_ids, pos, left, top, width, height, label_map=label_map
        )
        self._draw_edges(slide, edges, pos, node_shapes, left, top, width, height)

    def _render_fallback(
        self,
        slide: Slide,
        mermaid_text: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        Mermaidパース失敗時のフォールバック処理。
        テキストボックスにMermaidテキストをそのまま表示する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        mermaid_text : str
            Mermaidテキスト。
        left : int
            左端座標（EMU）。
        top : int
            上端座標（EMU）。
        width : int
            幅（EMU）。
        height : int
            高さ（EMU）。
        """
        shape = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = mermaid_text
        run.font.name = "Courier New"
        run.font.size = Pt(10)
