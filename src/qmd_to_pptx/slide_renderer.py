"""
スライドレンダラーモジュール。

YAMLパーサーから受け取ったメタデータとDOMトラバーサーからの各ノード情報を統合し、
テンプレートPPTXを基にスライドを生成・管理するオーケストレーター。
"""

from __future__ import annotations

import importlib.resources
import json
import xml.etree.ElementTree as ET
from pathlib import Path

from lxml import etree as lxml_etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from .dom_traverser import DOMTraverser
from .formula_renderer import FormulaRenderer
from .markdown_parser import MarkdownParser
from .mermaid_renderer import MermaidRenderer
from .models import (
    DOMNodeInfo,
    DOMNodeType,
    LayoutDef,
    LayoutJSON,
    PlaceholderInfo,
    SlideContent,
    SlideMetadata,
    SeparatorType,
)
from .slide_splitter import SlideSplitter
from .text_renderer import TextRenderer
from .yaml_parser import YAMLParser


class _NodeAdapter:
    """
    内部用: DOMNodeInfo と同じインタフェースを持つ軽量アダプター。

    スライドレンダラー内でタイトル・サブタイトル等の
    仮ノードを生成するために使用する。
    """

    def __init__(self, node_type: DOMNodeType, element: ET.Element) -> None:
        self.node_type = node_type
        self.element = element


class SlideRenderer:
    """
    スライドレンダラークラス。

    各コンポーネントを呼び出す中心的なオーケストレーター。
    YAMLパーサー・スライド分割器・Markdownパーサー・DOMトラバーサーを
    順に呼び出して全スライドを生成し、指定パスに保存する。
    """

    def __init__(self) -> None:
        """スライドレンダラーの初期化。各サブレンダラーを生成する。"""
        # テキストレンダラー
        self._text_renderer = TextRenderer()
        # Mermaidレンダラー
        self._mermaid_renderer = MermaidRenderer()
        # 数式レンダラー
        self._formula_renderer = FormulaRenderer()
        # レイアウトJSONを読み込む
        self._layout_json: LayoutJSON = self._load_layout_json()

    def render_all(
        self,
        normalized_text: str,
        output: str,
        reference_doc: str | None = None,
    ) -> None:
        """
        正規化済みQMDテキストから全スライドを生成して指定パスに保存する。

        Parameters
        ----------
        normalized_text : str
            前処理器で正規化済みのQMDテキスト。
        output : str
            出力先PPTXファイルのパス。
        reference_doc : str | None
            ベースとなるPPTXテンプレートファイルのパス（省略可）。
        """
        # YAMLパーサーでメタデータを取得する
        yaml_parser = YAMLParser()
        metadata: SlideMetadata = yaml_parser.parse(normalized_text)

        # reference_doc の優先順位: 引数 > YAMLフロントマター
        effective_ref_doc = reference_doc or metadata.reference_doc

        # スライド分割器でスライドリストを生成する
        splitter = SlideSplitter()
        slide_contents: list[SlideContent] = splitter.split(
            normalized_text, metadata.slide_level
        )

        # プレゼンテーションオブジェクトを生成する
        if effective_ref_doc and Path(effective_ref_doc).exists():
            prs = Presentation(effective_ref_doc)
        else:
            prs = Presentation()
            effective_ref_doc = None

        # タイトルスライドを追加する
        self._add_title_slide(prs, metadata, effective_ref_doc)

        # Markdownパーサーとトラバーサーを準備する
        md_parser = MarkdownParser()
        dom_traverser = DOMTraverser()

        # 各スライドコンテンツを処理する
        for content in slide_contents:
            dom_root: ET.Element = md_parser.parse(content.body_text)
            nodes: list[DOMNodeInfo] = dom_traverser.traverse(dom_root)

            layout_name = self._select_layout(content, nodes, metadata.slide_level)
            slide = self._add_slide(prs, layout_name, effective_ref_doc)

            # 背景画像を設定する
            if content.background_image:
                self._set_background_image(slide, content.background_image)

            # スライドタイトルを書き込む
            if content.title:
                self._write_title(
                    slide, content.title, layout_name, effective_ref_doc
                )

            # 各ノードを処理する
            self._render_nodes(
                slide, nodes, layout_name, effective_ref_doc, metadata
            )

        # presentation.xml の notesMasterIdLst を補完する（python-pptxのバグ回避）
        self._fix_notes_master_id_lst(prs)

        # PPTXファイルを保存する
        prs.save(output)

    def _fix_notes_master_id_lst(self, prs: Presentation) -> None:
        """
        presentation.xml に notesMasterIdLst 要素がなければ追加する。

        python-pptx は notes_slide アクセス時に notesMaster を生成して
        presentation.xml.rels には登録するが、presentation.xml 本体の
        <p:notesMasterIdLst> 要素を追加しない。PowerPoint はその不整合を
        「フォーマットエラー」として報告するため、ここで補完する。

        Parameters
        ----------
        prs : Presentation
            保存前の Presentation オブジェクト。
        """
        from lxml import etree

        P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
        R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        prs_part = prs.part
        prs_elem = prs_part._element

        # すでに notesMasterIdLst が存在すれば何もしない
        if prs_elem.find(f"{{{P_NS}}}notesMasterIdLst") is not None:
            return

        # presentation.xml.rels から notesMaster の rId を探す
        notes_master_rid = None
        for rId, rel in prs_part.rels.items():
            if "notesMaster" in rel.reltype:
                notes_master_rid = rId
                break

        # notesMaster が登録されていなければ補完不要
        if notes_master_rid is None:
            return

        # notesMasterIdLst 要素と子要素を生成する
        notes_master_id_lst = etree.Element(f"{{{P_NS}}}notesMasterIdLst")
        etree.SubElement(
            notes_master_id_lst,
            f"{{{P_NS}}}notesMasterId",
            {f"{{{R_NS}}}id": notes_master_rid},
        )

        # sldMasterIdLst の直後に挿入する（OOXMLスキーマの要素順序に従う）
        sld_master_id_lst = prs_elem.find(f"{{{P_NS}}}sldMasterIdLst")
        if sld_master_id_lst is not None:
            idx = list(prs_elem).index(sld_master_id_lst)
            prs_elem.insert(idx + 1, notes_master_id_lst)
        else:
            prs_elem.insert(0, notes_master_id_lst)

    def _load_layout_json(self) -> LayoutJSON:
        """
        パッケージ同梱の resources/default_layout.json を読み込み
        LayoutJSONオブジェクトを生成して返す。

        Returns
        -------
        LayoutJSON
            レイアウトJSON定義オブジェクト。
        """
        try:
            # Python 3.9+ の importlib.resources でパッケージリソースを読み込む
            ref = importlib.resources.files("qmd_to_pptx") / "resources" / "default_layout.json"
            json_text = ref.read_text(encoding="utf-8")
        except Exception:
            # フォールバック: ファイルシステムから直接読み込む
            json_path = Path(__file__).parent / "resources" / "default_layout.json"
            json_text = json_path.read_text(encoding="utf-8")

        data = json.loads(json_text)
        layouts: dict[str, LayoutDef] = {}
        for layout_name, layout_data in data.get("layouts", {}).items():
            placeholders = [
                PlaceholderInfo(
                    idx=ph["idx"],
                    role=ph["role"],
                    left=ph["left"],
                    top=ph["top"],
                    width=ph["width"],
                    height=ph["height"],
                )
                for ph in layout_data.get("placeholders", [])
            ]
            layouts[layout_name] = LayoutDef(placeholders=placeholders)

        return LayoutJSON(
            slide_width_emu=data.get("slide_width_emu", 9144000),
            slide_height_emu=data.get("slide_height_emu", 5143500),
            layouts=layouts,
        )

    def _select_layout(
        self,
        content: SlideContent,
        nodes: list[DOMNodeInfo],
        slide_level: int = 2,
    ) -> str:
        """
        SlideContentの区切り種別とノード構成を元にレイアウト名を決定して返す。

        設計書 QMD_TO_PPTX_DESIGN.md 4.8節のレイアウト自動選択ルールを適用する。

        Parameters
        ----------
        content : SlideContent
            スライド内容。
        nodes : list[DOMNodeInfo]
            DOMノードリスト。
        slide_level : int
            スライド区切りとして扱う見出しレベル（1または2）。デフォルトは2。

        Returns
        -------
        str
            PowerPointレイアウト名。
        """
        # Section Header: HEADING1による分割 かつ slide-level: 2 の場合のみ適用
        # slide-level: 1 の場合は HEADING1 もコンテンツスライドとして扱う
        if content.separator_type == SeparatorType.HEADING1 and slide_level == 2:
            return "Section Header"

        # コンテンツが空の場合（スピーカーノートのみ含む場合を含む）
        content_nodes = [
            n for n in nodes
            if n.node_type not in (DOMNodeType.NOTES,)
        ]
        if not content_nodes:
            return "Blank"

        # .columns divの有無を確認する
        columns_nodes = [
            n for n in nodes if n.node_type == DOMNodeType.COLUMNS
        ]
        if columns_nodes:
            # columnsが存在する場合は Two Content または Comparison を選択する
            col_element = columns_nodes[0].element
            column_divs = [
                child for child in col_element
                if child.get("class", "") == "column"
            ]
            # 各カラムのコンテンツ種別を確認する
            has_non_text = any(
                self._has_non_text_content(col)
                for col in column_divs
            )
            if has_non_text:
                return "Comparison"
            return "Two Content"

        # テキストの後に非テキスト要素が続くかを確認する（Content with Caption）
        if self._is_content_with_caption(nodes):
            return "Content with Caption"

        # 上記以外はすべて Title and Content
        return "Title and Content"

    def _has_non_text_content(self, element: ET.Element) -> bool:
        """
        要素内に非テキスト要素（図・表など）が含まれるかどうかを返す。

        Parameters
        ----------
        element : ET.Element
            確認対象の要素。

        Returns
        -------
        bool
            非テキスト要素が含まれる場合はTrue。
        """
        for child in element.iter():
            if child.tag in ("table", "img"):
                return True
            if child.tag == "code" and "language-mermaid" in child.get("class", ""):
                return True
        return False

    def _is_content_with_caption(self, nodes: list[DOMNodeInfo]) -> bool:
        """
        1カラム構成でテキスト系ノードと図系ノードの両方が存在するかを確認する。

        順序に依存せず、両種別のノードが混在していれば True を返す。

        Parameters
        ----------
        nodes : list[DOMNodeInfo]
            DOMノードリスト。

        Returns
        -------
        bool
            Content with Captionレイアウトを適用する場合はTrue。
        """
        # テキスト系ノードが存在するか確認する
        has_text = any(
            n.node_type in (
                DOMNodeType.PARAGRAPH,
                DOMNodeType.UL,
                DOMNodeType.OL,
                DOMNodeType.CODE,
            )
            for n in nodes
        )
        # 図系ノードが存在するか確認する
        has_diagram = any(
            n.node_type in (
                DOMNodeType.TABLE,
                DOMNodeType.MERMAID,
                DOMNodeType.FORMULA_BLOCK,
            )
            for n in nodes
        )
        return has_text and has_diagram

    def _resolve_placeholder(self, slide: Slide, idx: int) -> bool:
        """
        slide.placeholders に指定 idx が存在する場合は True を返す。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        idx : int
            確認するプレースホルダーのidx。

        Returns
        -------
        bool
            指定 idx が存在する場合はTrue。
        """
        return any(
            ph.placeholder_format.idx == idx
            for ph in slide.placeholders
        )

    def _get_placeholder(self, slide: Slide, idx: int) -> object | None:
        """
        slide.placeholders から指定 idx のプレースホルダーを取得する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        idx : int
            取得するプレースホルダーのidx。

        Returns
        -------
        object | None
            プレースホルダーオブジェクト、存在しない場合はNone。
        """
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _write_via_placeholder(
        self,
        slide: Slide,
        idx: int,
        node: DOMNodeInfo,
        layout_name: str,
        incremental: bool = False,
    ) -> None:
        """
        slide.placeholders[idx] を取得してコンテンツを書き込む。

        テキスト系ノードはそのshapeをTextRendererの対応メソッドに渡す。
        テーブルノードはプレースホルダーの座標を取得してrender_table()を呼び出す。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        idx : int
            プレースホルダーのidx。
        node : DOMNodeInfo
            処理対象のDOMノード情報。
        layout_name : str
            スライドレイアウト名（ログ用）。
        incremental : bool
            リストを逐次表示するかどうか。
        """
        ph = self._get_placeholder(slide, idx)
        if ph is None:
            return

        ntype = node.node_type
        elem = node.element

        if ntype == DOMNodeType.H1:
            self._text_renderer.render_heading(ph, elem, level=1)
        elif ntype == DOMNodeType.H2:
            self._text_renderer.render_heading(ph, elem, level=2)
        elif ntype == DOMNodeType.PARAGRAPH:
            # インライン数式を含む場合は専用メソッドで処理する
            if self._has_inline_formula(elem):
                self._render_paragraph_with_inline(ph, elem)
            else:
                self._text_renderer.render_paragraph(ph, elem)
        elif ntype in (DOMNodeType.UL, DOMNodeType.OL):
            self._text_renderer.render_list(ph, elem, incremental=incremental)
        elif ntype == DOMNodeType.CODE:
            self._text_renderer.render_code(ph, elem)
        elif ntype == DOMNodeType.TABLE:
            left = ph.left
            top = ph.top
            width = ph.width
            height = ph.height
            self._text_renderer.render_table(slide, elem, left, top, width, height)
        elif ntype == DOMNodeType.MERMAID:
            self._mermaid_renderer.render(
                slide, elem, ph.left, ph.top, ph.width, ph.height
            )
        elif ntype in (DOMNodeType.FORMULA_BLOCK,):
            # 既存textframeに追記することで前テキストとの重なりを回避する
            self._formula_renderer.render_block_into_frame(ph, elem)

    def _write_via_textbox(
        self,
        slide: Slide,
        role: str,
        layout_def: LayoutDef,
        node: DOMNodeInfo,
        incremental: bool = False,
    ) -> None:
        """
        LayoutDef.placeholders を role で検索して座標情報を取得し、
        add_textbox()でshapeを作成してコンテンツを書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        role : str
            コンテンツの役割（title / body / subtitle など）。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        node : DOMNodeInfo
            処理対象のDOMノード情報。
        incremental : bool
            リストを逐次表示するかどうか。
        """
        # roleに対応するPlaceholderInfoを線形探索する
        ph_info: PlaceholderInfo | None = None
        for ph in layout_def.placeholders:
            if ph.role == role:
                ph_info = ph
                break

        if ph_info is None:
            # role に対応するプレースホルダー定義がない場合はスキップする
            return

        ntype = node.node_type
        elem = node.element

        if ntype == DOMNodeType.TABLE:
            self._text_renderer.render_table(
                slide, elem,
                ph_info.left, ph_info.top,
                ph_info.width, ph_info.height,
            )
            return

        if ntype == DOMNodeType.MERMAID:
            self._mermaid_renderer.render(
                slide, elem,
                ph_info.left, ph_info.top,
                ph_info.width, ph_info.height,
            )
            return

        if ntype == DOMNodeType.FORMULA_BLOCK:
            self._formula_renderer.render_block(
                slide, elem,
                ph_info.left, ph_info.top,
                ph_info.width, ph_info.height,
            )
            return

        # テキスト系ノードはtextboxを追加してshapeを渡す
        shape = slide.shapes.add_textbox(
            Emu(ph_info.left), Emu(ph_info.top),
            Emu(ph_info.width), Emu(ph_info.height),
        )

        if ntype == DOMNodeType.H1:
            self._text_renderer.render_heading(shape, elem, level=1)
        elif ntype == DOMNodeType.H2:
            self._text_renderer.render_heading(shape, elem, level=2)
        elif ntype == DOMNodeType.PARAGRAPH:
            # インライン数式を含む場合は専用メソッドで処理する
            if self._has_inline_formula(elem):
                self._render_paragraph_with_inline(shape, elem)
            else:
                self._text_renderer.render_paragraph(shape, elem)
        elif ntype in (DOMNodeType.UL, DOMNodeType.OL):
            self._text_renderer.render_list(shape, elem, incremental=incremental)
        elif ntype == DOMNodeType.CODE:
            self._text_renderer.render_code(shape, elem)

    def _add_title_slide(
        self,
        prs: Presentation,
        metadata: SlideMetadata,
        reference_doc: str | None,
    ) -> None:
        """
        メタデータからタイトルスライドをプレゼンテーションに追加する。

        Parameters
        ----------
        prs : Presentation
            python-pptxのPresentationオブジェクト。
        metadata : SlideMetadata
            スライドメタデータ。
        reference_doc : str | None
            テンプレートファイルのパス。
        """
        layout_name = "Title Slide"
        slide = self._add_slide(prs, layout_name, reference_doc)
        layout_def = self._layout_json.layouts.get(layout_name, LayoutDef())

        # タイトルを書き込む（idx=0）
        self._write_title_to_slide(slide, metadata.title, 0, "title", layout_def)

        # サブタイトルを書き込む（idx=1、author と date を組み合わせる。空文字でも常に処理する）
        subtitle_parts = [metadata.author, metadata.date]
        subtitle_text = "\n".join(p for p in subtitle_parts if p)
        self._write_subtitle_to_slide(
            slide, subtitle_text, 1, "subtitle", layout_def
        )

    def _write_title_to_slide(
        self,
        slide: Slide,
        title_text: str,
        idx: int,
        role: str,
        layout_def: LayoutDef,
    ) -> None:
        """
        スライドにタイトルテキストを書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        title_text : str
            タイトルテキスト。
        idx : int
            プレースホルダーのidx。
        role : str
            コンテンツの役割。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        """
        # 仮のDOMノードを生成する
        title_elem = ET.Element("h1")
        title_elem.text = title_text

        node = _NodeAdapter(DOMNodeType.H1, title_elem)

        if self._resolve_placeholder(slide, idx):
            self._write_via_placeholder(slide, idx, node, "Title Slide")
        else:
            ph_info = next(
                (p for p in layout_def.placeholders if p.role == role), None
            )
            if ph_info:
                shape = slide.shapes.add_textbox(
                    Emu(ph_info.left), Emu(ph_info.top),
                    Emu(ph_info.width), Emu(ph_info.height),
                )
                self._text_renderer.render_heading(shape, title_elem, level=1)

    def _write_subtitle_to_slide(
        self,
        slide: Slide,
        subtitle_text: str,
        idx: int,
        role: str,
        layout_def: LayoutDef,
    ) -> None:
        """
        スライドにサブタイトルテキストを書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        subtitle_text : str
            サブタイトルテキスト（author + date の組み合わせ）。
        idx : int
            プレースホルダーのidx。
        role : str
            コンテンツの役割。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        """
        sub_elem = ET.Element("p")
        sub_elem.text = subtitle_text

        node = _NodeAdapter(DOMNodeType.PARAGRAPH, sub_elem)

        if self._resolve_placeholder(slide, idx):
            ph = self._get_placeholder(slide, idx)
            if ph:
                self._text_renderer.render_paragraph(ph, sub_elem)
        else:
            ph_info = next(
                (p for p in layout_def.placeholders if p.role == role), None
            )
            if ph_info:
                shape = slide.shapes.add_textbox(
                    Emu(ph_info.left), Emu(ph_info.top),
                    Emu(ph_info.width), Emu(ph_info.height),
                )
                self._text_renderer.render_paragraph(shape, sub_elem)

    def _add_slide(
        self,
        prs: Presentation,
        layout_name: str,
        reference_doc: str | None,
    ) -> Slide:
        """
        指定したレイアウト名でスライドをプレゼンテーションに追加して返す。

        レイアウト名が見つからない場合はパターンCとしてデフォルトレイアウトを使用する。

        Parameters
        ----------
        prs : Presentation
            python-pptxのPresentationオブジェクト。
        layout_name : str
            使用するスライドレイアウト名。
        reference_doc : str | None
            テンプレートファイルのパス（None の場合は新規生成）。

        Returns
        -------
        Slide
            追加されたSlideオブジェクト。
        """
        # レイアウト名でスライドレイアウトを検索する
        layout = self._find_layout(prs, layout_name)
        if layout is None:
            # コンテンツ系レイアウトが見つからない場合は Title and Content を試みる
            _CONTENT_FALLBACK_LAYOUTS = {"Two Content", "Comparison", "Content with Caption"}
            if layout_name in _CONTENT_FALLBACK_LAYOUTS:
                layout = self._find_layout(prs, "Title and Content")
            # それでも見つからない場合は最初のレイアウトを使用する
            if layout is None:
                layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        self._fix_slide_xml(slide)
        return slide

    def _fix_slide_xml(self, slide: Slide) -> None:
        """
        python-pptx が生成したスライド XML の OOXML スキーマ不備を修正する。

        p:spTree の p:grpSpPr が空（a:xfrm なし）の場合、
        PowerPoint が修復ダイアログを表示するため、ゼロ座標の a:xfrm を補完する。
        """
        spTree = slide.shapes._spTree
        grpSpPr = spTree.find(qn("p:grpSpPr"))
        if grpSpPr is not None and grpSpPr.find(qn("a:xfrm")) is None:
            xfrm = lxml_etree.SubElement(grpSpPr, qn("a:xfrm"))
            off = lxml_etree.SubElement(xfrm, qn("a:off"))
            off.set("x", "0")
            off.set("y", "0")
            ext = lxml_etree.SubElement(xfrm, qn("a:ext"))
            ext.set("cx", "0")
            ext.set("cy", "0")
            chOff = lxml_etree.SubElement(xfrm, qn("a:chOff"))
            chOff.set("x", "0")
            chOff.set("y", "0")
            chExt = lxml_etree.SubElement(xfrm, qn("a:chExt"))
            chExt.set("cx", "0")
            chExt.set("cy", "0")

    def _find_layout(
        self, prs: Presentation, layout_name: str
    ) -> object | None:
        """
        プレゼンテーションからレイアウト名でスライドレイアウトを検索する。

        Parameters
        ----------
        prs : Presentation
            python-pptxのPresentationオブジェクト。
        layout_name : str
            検索するレイアウト名。

        Returns
        -------
        object | None
            スライドレイアウトオブジェクト、見つからない場合はNone。
        """
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None

    def _write_title(
        self,
        slide: Slide,
        title: str,
        layout_name: str,
        reference_doc: str | None,
    ) -> None:
        """
        スライドにタイトルを書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        title : str
            タイトルテキスト。
        layout_name : str
            スライドレイアウト名。
        reference_doc : str | None
            テンプレートファイルのパス。
        """
        title_elem = ET.Element("h2")
        title_elem.text = title

        layout_def = self._layout_json.layouts.get(layout_name, LayoutDef())

        # タイトルプレースホルダー（idx=0）への書き込みを試みる
        if self._resolve_placeholder(slide, 0):
            ph = self._get_placeholder(slide, 0)
            if ph:
                # プレースホルダーのテキストを直接設定する
                try:
                    ph.text = title
                except Exception:
                    self._text_renderer.render_heading(ph, title_elem, level=2)
        else:
            # フォールバック: textboxを使用する
            ph_info = next(
                (p for p in layout_def.placeholders if p.role == "title"), None
            )
            if ph_info:
                shape = slide.shapes.add_textbox(
                    Emu(ph_info.left), Emu(ph_info.top),
                    Emu(ph_info.width), Emu(ph_info.height),
                )
                self._text_renderer.render_heading(shape, title_elem, level=2)

    def _render_nodes(
        self,
        slide: Slide,
        nodes: list[DOMNodeInfo],
        layout_name: str,
        reference_doc: str | None,
        metadata: SlideMetadata,
    ) -> None:
        """
        DOMノードリストをイテレートして各ノードをスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        nodes : list[DOMNodeInfo]
            DOMノードリスト。
        layout_name : str
            スライドレイアウト名。
        reference_doc : str | None
            テンプレートファイルのパス。
        metadata : SlideMetadata
            スライドメタデータ（incrementalフラグなど）。
        """
        layout_def = self._layout_json.layouts.get(layout_name, LayoutDef())

        # "Title and Content" かつテキスト系・図系の混在コンテンツの場合は縦分割で描画する
        if layout_name == "Title and Content":
            body_nodes = self._collect_body_nodes(nodes, metadata)
            _DIAGRAM_TYPES = (DOMNodeType.MERMAID, DOMNodeType.TABLE, DOMNodeType.FORMULA_BLOCK)
            _TEXT_TYPES = (DOMNodeType.PARAGRAPH, DOMNodeType.UL, DOMNodeType.OL, DOMNodeType.CODE)
            has_text_body = any(n.node_type in _TEXT_TYPES for n in body_nodes)
            has_diagram_body = any(n.node_type in _DIAGRAM_TYPES for n in body_nodes)
            if has_text_body and has_diagram_body:
                # スピーカーノートだけ先に処理してから縦分割描画する
                for node in nodes:
                    if node.node_type == DOMNodeType.NOTES:
                        self._text_renderer.render_notes(slide, node.element)
                self._render_body_nodes_vertical_split(
                    slide, body_nodes, layout_def, reference_doc, metadata.incremental
                )
                return

        for node in nodes:
            ntype = node.node_type
            elem = node.element

            # スピーカーノートを処理する
            if ntype == DOMNodeType.NOTES:
                self._text_renderer.render_notes(slide, elem)
                continue

            # 2カラムコンテナを処理する
            if ntype == DOMNodeType.COLUMNS:
                self._render_columns(
                    slide, elem, layout_name, layout_def, reference_doc
                )
                continue

            # インクリメンタルリストを処理する
            if ntype == DOMNodeType.INCREMENTAL:
                for child in elem:
                    if child.tag in ("ul", "ol"):
                        self._render_body_node(
                            slide, DOMNodeInfo(DOMNodeType.UL, child),
                            layout_def, reference_doc, incremental=True,
                            layout_name=layout_name,
                        )
                continue

            # 非インクリメンタルリストを処理する
            if ntype == DOMNodeType.NON_INCREMENTAL:
                for child in elem:
                    if child.tag in ("ul", "ol"):
                        self._render_body_node(
                            slide, DOMNodeInfo(DOMNodeType.UL, child),
                            layout_def, reference_doc, incremental=False,
                            layout_name=layout_name,
                        )
                continue

            # H1/H2はタイトルとして処理済みのためスキップする（bodyには含めない）
            # ただし slide-level: 1 の場合、H2は本文見出しとして処理するため
            # スキップせずにbodyとして描画する
            if ntype == DOMNodeType.H1:
                continue
            if ntype == DOMNodeType.H2 and metadata.slide_level != 1:
                continue

            # その他のノードをbodyとして処理する
            self._render_body_node(
                slide, node, layout_def, reference_doc,
                incremental=metadata.incremental,
                layout_name=layout_name,
            )

    def _render_body_node(
        self,
        slide: Slide,
        node: DOMNodeInfo,
        layout_def: LayoutDef,
        reference_doc: str | None,
        incremental: bool = False,
        layout_name: str = "",
    ) -> None:
        """
        ボディコンテンツのDOMノードをスライドに描画する。

        layout_name が "Content with Caption" の場合、図系ノードを body エリア（idx=1）に、
        テキスト系ノードを caption エリア（idx=2）にルーティングする。
        それ以外のレイアウトでは従来通り idx=1 / role="body" に描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        node : DOMNodeInfo
            処理対象のDOMノード情報。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        reference_doc : str | None
            テンプレートファイルのパス。
        incremental : bool
            リストを逐次表示するかどうか。
        layout_name : str
            スライドレイアウト名。"Content with Caption" 時にルーティング分岐に使用する。
        """
        # 図系ノード種別の定義
        _DIAGRAM_TYPES = (DOMNodeType.MERMAID, DOMNodeType.TABLE, DOMNodeType.FORMULA_BLOCK)

        if layout_name == "Content with Caption":
            if node.node_type in _DIAGRAM_TYPES:
                # 図系ノード → body エリア (idx=1 / role="body") に描画する
                if self._resolve_placeholder(slide, 1):
                    self._write_via_placeholder(
                        slide, 1, node, "", incremental=incremental
                    )
                else:
                    self._write_via_textbox(
                        slide, "body", layout_def, node, incremental=incremental
                    )
            else:
                # テキスト系ノード → caption エリア (idx=2 / role="caption") に描画する
                if self._resolve_placeholder(slide, 2):
                    self._write_via_placeholder(
                        slide, 2, node, "", incremental=incremental
                    )
                else:
                    self._write_via_textbox(
                        slide, "caption", layout_def, node, incremental=incremental
                    )
        else:
            # 従来動作: body プレースホルダー (idx=1 / role="body") に描画する
            if self._resolve_placeholder(slide, 1):
                self._write_via_placeholder(
                    slide, 1, node, "", incremental=incremental
                )
            else:
                self._write_via_textbox(
                    slide, "body", layout_def, node, incremental=incremental
                )

    def _collect_body_nodes(
        self,
        nodes: list[DOMNodeInfo],
        metadata: SlideMetadata,
    ) -> list[DOMNodeInfo]:
        """
        ノードリストから body 描画対象のノードを収集して返す。

        NOTES / COLUMNS / INCREMENTAL / NON_INCREMENTAL / H1 /
        H2（slide_level != 1 の場合）を除いた描画対象ノードを返す。

        Parameters
        ----------
        nodes : list[DOMNodeInfo]
            DOMノードリスト。
        metadata : SlideMetadata
            スライドメタデータ（slide_level など）。

        Returns
        -------
        list[DOMNodeInfo]
            body 描画対象のノードリスト。
        """
        _SKIP_TYPES = {
            DOMNodeType.NOTES,
            DOMNodeType.COLUMNS,
            DOMNodeType.INCREMENTAL,
            DOMNodeType.NON_INCREMENTAL,
        }
        result: list[DOMNodeInfo] = []
        for node in nodes:
            ntype = node.node_type
            if ntype in _SKIP_TYPES:
                continue
            if ntype == DOMNodeType.H1:
                continue
            if ntype == DOMNodeType.H2 and metadata.slide_level != 1:
                continue
            result.append(node)
        return result

    def _get_body_area_coords(
        self,
        slide: Slide,
        layout_def: LayoutDef,
    ) -> tuple[int, int, int, int] | None:
        """
        body エリアの座標 (left, top, width, height) を EMU 単位で返す。

        idx=1 のプレースホルダーが存在すればその座標を使用する。
        存在しない場合は layout_def の role="body" のプレースホルダー定義を使用する。
        どちらも存在しない場合は None を返す。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。

        Returns
        -------
        tuple[int, int, int, int] | None
            (left, top, width, height) の EMU タプル、取得できない場合は None。
        """
        # idx=1 プレースホルダーがあればその座標を使用する
        ph = self._get_placeholder(slide, 1)
        if ph is not None:
            return (int(ph.left), int(ph.top), int(ph.width), int(ph.height))

        # layout_def から role="body" の座標を使用する
        for ph_info in layout_def.placeholders:
            if ph_info.role == "body":
                return (ph_info.left, ph_info.top, ph_info.width, ph_info.height)

        return None

    def _render_body_nodes_vertical_split(
        self,
        slide: Slide,
        body_nodes: list[DOMNodeInfo],
        layout_def: LayoutDef,
        reference_doc: str | None,
        incremental: bool,
    ) -> None:
        """
        body エリアをテキスト系（上半分）と図系（下半分）に縦分割して描画する。

        テキスト系ノード (PARAGRAPH / UL / OL / CODE) を body エリア上半分に、
        図系ノード (MERMAID / TABLE / FORMULA_BLOCK) を body エリア下半分に描画する。
        プレースホルダー (idx=1) は使用せず、すべて textbox / 図の直接配置で描画するため
        コンテンツが重なる問題を回避する。

        Parameters
        ----------
        slide : Slide
            python-pptx の Slide オブジェクト。
        body_nodes : list[DOMNodeInfo]
            body 描画対象のノードリスト（_collect_body_nodes の返り値）。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        reference_doc : str | None
            テンプレートファイルのパス。
        incremental : bool
            リストを逐次表示するかどうか。
        """
        coords = self._get_body_area_coords(slide, layout_def)
        if coords is None:
            # body エリア座標が取得できない場合は縦分割せずにスキップする
            return

        left, top, width, height = coords
        half_height = height // 2

        _DIAGRAM_TYPES = (DOMNodeType.MERMAID, DOMNodeType.TABLE, DOMNodeType.FORMULA_BLOCK)
        _TEXT_TYPES = (DOMNodeType.PARAGRAPH, DOMNodeType.UL, DOMNodeType.OL, DOMNodeType.CODE)

        # テキスト系ノードと図系ノードに分類する
        text_nodes = [n for n in body_nodes if n.node_type in _TEXT_TYPES]
        diagram_nodes = [n for n in body_nodes if n.node_type in _DIAGRAM_TYPES]

        # --- 上半分: テキスト系ノードを等分割して textbox に描画する ---
        if text_nodes:
            slot_height = half_height // len(text_nodes)
            for i, node in enumerate(text_nodes):
                t = top + i * slot_height
                shape = slide.shapes.add_textbox(
                    Emu(left), Emu(t), Emu(width), Emu(slot_height)
                )
                ntype = node.node_type
                elem = node.element
                if ntype == DOMNodeType.PARAGRAPH:
                    if self._has_inline_formula(elem):
                        self._render_paragraph_with_inline(shape, elem)
                    else:
                        self._text_renderer.render_paragraph(shape, elem)
                elif ntype in (DOMNodeType.UL, DOMNodeType.OL):
                    self._text_renderer.render_list(shape, elem, incremental=incremental)
                elif ntype == DOMNodeType.CODE:
                    self._text_renderer.render_code(shape, elem)

        # --- 下半分: 図系ノードを等分割して直接描画する ---
        if diagram_nodes:
            diagram_top = top + half_height
            slot_height = half_height // len(diagram_nodes)
            for i, node in enumerate(diagram_nodes):
                t = diagram_top + i * slot_height
                ntype = node.node_type
                elem = node.element
                if ntype == DOMNodeType.TABLE:
                    self._text_renderer.render_table(
                        slide, elem, left, t, width, slot_height
                    )
                elif ntype == DOMNodeType.MERMAID:
                    self._mermaid_renderer.render(
                        slide, elem, left, t, width, slot_height
                    )
                elif ntype == DOMNodeType.FORMULA_BLOCK:
                    self._formula_renderer.render_block(
                        slide, elem, left, t, width, slot_height
                    )

    def _render_columns(
        self,
        slide: Slide,
        columns_elem: ET.Element,
        layout_name: str,
        layout_def: LayoutDef,
        reference_doc: str | None,
    ) -> None:
        """
        2カラムコンテナを処理してスライドに描画する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        columns_elem : ET.Element
            columns div要素。
        layout_name : str
            スライドレイアウト名。
        layout_def : LayoutDef
            レイアウト定義オブジェクト。
        reference_doc : str | None
            テンプレートファイルのパス。
        """
        column_divs = [
            child for child in columns_elem
            if child.get("class", "") == "column"
        ]

        # カラム別のroleを定義する
        column_roles = ["left_content", "right_content"]
        column_idx_map = {
            "left_content": 1,
            "right_content": 2,
        }

        for i, col_elem in enumerate(column_divs):
            if i >= len(column_roles):
                break
            role = column_roles[i]
            idx = column_idx_map.get(role, i + 1)

            # カラム内の各ノードを処理する
            for child in col_elem:
                child_class = child.get("class", "")
                # DOMトラバーサーで判定したノード種別を模倣する
                child_node = self._classify_child(child)
                if child_node is None:
                    continue

                if self._resolve_placeholder(slide, idx):
                    self._write_via_placeholder(slide, idx, child_node, layout_name)
                else:
                    self._write_via_textbox(slide, role, layout_def, child_node)

    def _has_inline_formula(self, element: ET.Element) -> bool:
        """
        段落要素の子にインライン数式（span.arithmatex）が含まれるかを確認する。

        Parameters
        ----------
        element : ET.Element
            確認対象の段落要素。

        Returns
        -------
        bool
            インライン数式が含まれる場合はTrue。
        """
        for child in element:
            if child.tag == "span" and "arithmatex" in child.get("class", ""):
                return True
        return False

    def _render_paragraph_with_inline(
        self, shape: object, element: ET.Element
    ) -> None:
        """
        インライン数式（span.arithmatex）を含む段落をshapeのテキストフレームに書き込む。

        段落の各子要素を走査し、arithmatex spanはOMMLとして挿入し、
        それ以外はrunのテキストとして追加する。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト。
        element : ET.Element
            段落要素（p）。インライン数式spanを子に持つ。
        """
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        para = tf.paragraphs[0]

        def add_text_run(text: str) -> None:
            """テキストを追加するrunを生成する。空文字列は無視する。"""
            if text:
                r = para.add_run()
                r.text = text
                r.font.size = Pt(18)

        # 段落先頭のテキストノードを追加する（element.text は最初の子要素の前のテキスト）
        add_text_run(element.text or "")

        for child in element:
            if child.tag == "span" and "arithmatex" in child.get("class", ""):
                # インライン数式: 空のrunを起点としてOMMLを挿入する
                r = para.add_run()
                r.text = ""
                self._formula_renderer.render_inline(r, child)
            else:
                # その他のインライン要素（em, strong 等）はテキストとして追加する
                add_text_run("".join(child.itertext()))
            # 各子要素の後続テキスト（tail）を追加する
            add_text_run(child.tail or "")

    def _classify_child(self, element: ET.Element) -> DOMNodeInfo | None:
        """
        要素を DOMNodeInfo に変換する。

        Parameters
        ----------
        element : ET.Element
            分類する要素。

        Returns
        -------
        DOMNodeInfo | None
            変換したDOMNodeInfo、分類できない場合はNone。
        """
        from .dom_traverser import DOMTraverser
        traverser = DOMTraverser()
        tag = element.tag
        css_class = element.get("class", "")
        node_type = traverser._classify_node(tag, css_class)
        if node_type is None:
            return None
        return DOMNodeInfo(node_type=node_type, element=element)

    def _set_background_image(
        self, slide: Slide, image_path: str
    ) -> None:
        """
        スライドの背景に指定された画像を設定する。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        image_path : str
            背景画像ファイルのパス。
        """
        if not Path(image_path).exists():
            return

        # スライド幅・高さを取得する
        width = slide.slide_layout.slide_master.slide_width
        height = slide.slide_layout.slide_master.slide_height

        # 画像を背景としてスライド全体に配置する
        try:
            pic = slide.shapes.add_picture(image_path, 0, 0, width, height)
            # 画像をスライドの最背面に移動する（z-orderを最小にする）
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        except Exception:
            pass
