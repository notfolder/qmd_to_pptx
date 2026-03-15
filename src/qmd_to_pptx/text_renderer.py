"""
テキストレンダラーモジュール。

DOMトラバーサーから受け取ったテキスト系ノード（見出し・段落・リスト・表・
コードブロック）をpython-pptxのShapeとして現在のスライドに追加する。
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt


class TextRenderer:
    """
    テキストレンダラークラス。

    SlideRendererからDOMNodeInfo.elementを受け取り、
    対応するpython-pptxのShapeをスライドに追加する。
    """

    # コードブロック用等幅フォント名
    _MONOSPACE_FONT: str = "Courier New"
    # コードブロック背景色（薄いグレー）
    _CODE_BG_COLOR: RGBColor = RGBColor(0xF0, 0xF0, 0xF0)
    # デフォルトフォントサイズ
    _DEFAULT_FONT_SIZE: Pt = Pt(18)
    # 見出しフォントサイズ（h1）
    _H1_FONT_SIZE: Pt = Pt(36)
    # 見出しフォントサイズ（h2）
    _H2_FONT_SIZE: Pt = Pt(28)

    def render_heading(
        self,
        shape: object,
        element: ET.Element,
        level: int,
    ) -> None:
        """
        shapeのテキストフレームにlevelに応じたスタイルで見出しテキストを書き込む。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト（プレースホルダーまたはtextbox）。
        element : ET.Element
            見出し要素（h1またはh2）。
        level : int
            見出しレベル（1または2）。
        """
        text = "".join(element.itertext()).strip()
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.bold = True
        if level == 1:
            run.font.size = self._H1_FONT_SIZE
        else:
            run.font.size = self._H2_FONT_SIZE

    def render_paragraph(
        self,
        shape: object,
        element: ET.Element,
    ) -> None:
        """
        shapeのテキストフレームに段落テキストを書き込む。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト。
        element : ET.Element
            段落要素（p）。
        """
        text = "".join(element.itertext()).strip()
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = self._DEFAULT_FONT_SIZE

    def render_list(
        self,
        shape: object,
        element: ET.Element,
        incremental: bool = False,
    ) -> None:
        """
        shapeのテキストフレームにインデント付きのリストを書き込む。

        ulタグの場合は箇条書き（bullet）、olタグの場合は番号付きリストとして処理する。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト。
        element : ET.Element
            リスト要素（ulまたはol）。
        incremental : bool
            Trueの場合、アニメーション逐次表示を設定する（現時点では将来拡張用）。
        """
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        first_para = True
        # 最上位のリストタグが ol かどうかを判定する
        ordered = element.tag == "ol"

        self._render_list_items(
            tf, element, level=0, first_para_ref=[first_para], ordered=ordered
        )

    def _render_list_items(
        self,
        tf: object,
        element: ET.Element,
        level: int,
        first_para_ref: list[bool],
        ordered: bool = False,
    ) -> None:
        """
        リストアイテムを再帰的にテキストフレームに書き込む。

        Parameters
        ----------
        tf : object
            python-pptxのTextFrameオブジェクト。
        element : ET.Element
            リストまたはリストアイテム要素。
        level : int
            ネストレベル（インデントに使用）。
        first_para_ref : list[bool]
            最初の段落かどうかのフラグ（リスト経由で参照共有）。
        ordered : bool
            Trueの場合、番号付きリストとして処理する。
        """
        from lxml import etree

        # DrawingML の名前空間
        _A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

        for item in element:
            if item.tag != "li":
                continue

            # li直下のテキストを取得する（子リストは除く）
            item_text = self._get_direct_text(item)

            if first_para_ref[0]:
                para = tf.paragraphs[0]
                first_para_ref[0] = False
            else:
                para = tf.add_paragraph()

            para.level = level
            run = para.add_run()
            run.text = item_text
            run.font.size = self._DEFAULT_FONT_SIZE

            if ordered:
                # 番号付きリスト: <a:buAutoNum> を pPr に設定する
                # python-pptx は buAutoNum を直接サポートしないため lxml で操作する
                p_elem = para._p
                # 既存の pPr を取得または生成する
                pPr = p_elem.find(f"{{{_A_NS}}}pPr")
                if pPr is None:
                    pPr = etree.SubElement(p_elem, f"{{{_A_NS}}}pPr")
                    p_elem.insert(0, pPr)
                # 既存の buChar/buNone/buAutoNum を除去する
                for tag in ("buNone", "buChar", "buAutoNum"):
                    old = pPr.find(f"{{{_A_NS}}}{tag}")
                    if old is not None:
                        pPr.remove(old)
                # 番号付きリストマーカーを設定する（arabicPeriod = 「1.」形式）
                etree.SubElement(
                    pPr,
                    f"{{{_A_NS}}}buAutoNum",
                    attrib={"type": "arabicPeriod"},
                )

            # ネストリスト（ul/ol）を再帰処理する
            for child in item:
                if child.tag in ("ul", "ol"):
                    # ネストされたリストは子タグのol/ulに従って順序を決定する
                    self._render_list_items(
                        tf, child, level + 1, first_para_ref,
                        ordered=(child.tag == "ol"),
                    )

    def _get_direct_text(self, element: ET.Element) -> str:
        """
        要素直下のテキストのみを取得する（子リスト要素のテキストは含まない）。

        Parameters
        ----------
        element : ET.Element
            対象要素。

        Returns
        -------
        str
            直下テキスト。
        """
        texts: list[str] = []
        if element.text:
            texts.append(element.text)
        for child in element:
            if child.tag not in ("ul", "ol"):
                texts.append("".join(child.itertext()))
            if child.tail:
                texts.append(child.tail)
        return "".join(texts).strip()

    def render_table(
        self,
        slide: Slide,
        element: ET.Element,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """
        slide.shapes.add_table()で指定座標にテーブルShapeを生成し、
        <table>要素の内容を書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        element : ET.Element
            テーブル要素（table）。
        left : int
            左端座標（EMU）。
        top : int
            上端座標（EMU）。
        width : int
            幅（EMU）。
        height : int
            高さ（EMU）。
        """
        # 行・列数を計算する
        rows_data = self._extract_table_data(element)
        if not rows_data:
            return

        num_rows = len(rows_data)
        num_cols = max(len(row) for row in rows_data)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, left, top, width, height
        )
        table = table_shape.table

        for r_idx, row_data in enumerate(rows_data):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text
                    # ヘッダー行は太字にする
                    if r_idx == 0:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.bold = True

    def _extract_table_data(
        self, element: ET.Element
    ) -> list[list[str]]:
        """
        <table>要素から行・列データを抽出する。

        Parameters
        ----------
        element : ET.Element
            テーブル要素。

        Returns
        -------
        list[list[str]]
            行ごとのセルテキストのリスト。
        """
        rows: list[list[str]] = []
        # thead / tbody を探してtrを収集する
        for section in element:
            if section.tag in ("thead", "tbody"):
                for tr in section:
                    if tr.tag == "tr":
                        row_data = [
                            "".join(cell.itertext()).strip()
                            for cell in tr
                            if cell.tag in ("th", "td")
                        ]
                        rows.append(row_data)
            elif section.tag == "tr":
                row_data = [
                    "".join(cell.itertext()).strip()
                    for cell in section
                    if cell.tag in ("th", "td")
                ]
                rows.append(row_data)
        return rows

    def render_code(
        self,
        shape: object,
        element: ET.Element,
    ) -> None:
        """
        shapeのテキストフレームに等幅フォントでコードテキストを書き込む。

        Parameters
        ----------
        shape : object
            python-pptxのShapeオブジェクト。
        element : ET.Element
            コード要素（code）。
        """
        text = "".join(element.itertext())
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = False
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.name = self._MONOSPACE_FONT
        run.font.size = Pt(14)

    def render_notes(
        self,
        slide: Slide,
        element: ET.Element,
    ) -> None:
        """
        elementのテキストをスライドのノートテキストフレームに書き込む。

        Parameters
        ----------
        slide : Slide
            python-pptxのSlideオブジェクト。
        element : ET.Element
            ノート要素（div class="notes"）。
        """
        text = "".join(element.itertext()).strip()
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        if text:
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = text
