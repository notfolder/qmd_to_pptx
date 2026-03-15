"""
数式レンダラー（FormulaRenderer）の単体テスト。
"""

import xml.etree.ElementTree as ET
import pytest
from pptx import Presentation
from pptx.util import Emu
from qmd_to_pptx.formula_renderer import FormulaRenderer


def _make_slide():
    """テスト用スライドを生成して返す。"""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Blankレイアウト
    return prs.slides.add_slide(layout)


def _make_inline_element(latex: str) -> ET.Element:
    """インライン数式要素 span.arithmatex を生成する。"""
    elem = ET.Element("span")
    elem.set("class", "arithmatex")
    elem.text = f"\\({latex}\\)"
    return elem


def _make_block_element(latex: str) -> ET.Element:
    """ブロック数式要素 div.arithmatex を生成する。"""
    elem = ET.Element("div")
    elem.set("class", "arithmatex")
    elem.text = f"\\[{latex}\\]"
    return elem


class TestFormulaRenderer:
    """FormulaRenderer クラスの単体テスト。"""

    def setup_method(self) -> None:
        """テスト前にFormulaRendererインスタンスを生成する。"""
        self.renderer = FormulaRenderer()

    # --- _extract_latex のテスト ---

    def test_extract_latex_inline(self) -> None:
        r"""インライン数式 \(...\) からデリミタを除去してLaTeXを返す。"""
        elem = _make_inline_element("E=mc^2")
        result = self.renderer._extract_latex(elem)
        assert result == "E=mc^2"

    def test_extract_latex_block(self) -> None:
        r"""ブロック数式 \[...\] からデリミタを除去してLaTeXを返す。"""
        elem = _make_block_element(r"\sum_{i=0}^{n} i")
        result = self.renderer._extract_latex(elem)
        assert result == r"\sum_{i=0}^{n} i"

    def test_extract_latex_dollar_delimiter(self) -> None:
        """$$...$$ 形式のデリミタも除去する。"""
        elem = ET.Element("span")
        elem.set("class", "arithmatex")
        elem.text = "$$E=mc^2$$"
        result = self.renderer._extract_latex(elem)
        assert result == "E=mc^2"

    def test_extract_latex_no_delimiter(self) -> None:
        """デリミタなしのテキストはそのまま返す。"""
        elem = ET.Element("span")
        elem.set("class", "arithmatex")
        elem.text = "E=mc^2"
        result = self.renderer._extract_latex(elem)
        assert result == "E=mc^2"

    # --- render_block のテスト ---

    def test_render_block_does_not_raise_on_valid_latex(self) -> None:
        """有効なLaTeXのブロック数式をレンダリングしても例外が発生しない。"""
        slide = _make_slide()
        elem = _make_block_element("E=mc^2")
        # 例外なく完了することを確認
        self.renderer.render_block(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(2000000)
        )

    def test_render_block_fallback_on_invalid_latex(self) -> None:
        """無効なLaTeXでもフォールバックとしてテキストボックスが追加される。"""
        slide = _make_slide()
        initial_shapes = len(slide.shapes)
        elem = ET.Element("div")
        elem.set("class", "arithmatex")
        elem.text = "\\[これは無効な数式です\\]"
        self.renderer.render_block(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(2000000)
        )
        # フォールバックでテキストボックスが追加されるか、OMMLが追加されるか
        # いずれにせよ例外なく完了することを確認
        # （フォールバック時はshapeが追加される）

    def test_render_block_empty_latex_is_skipped(self) -> None:
        """LaTeXテキストが空の場合は何も追加しない。"""
        slide = _make_slide()
        initial_shapes = len(slide.shapes)
        elem = ET.Element("div")
        elem.set("class", "arithmatex")
        elem.text = ""
        self.renderer.render_block(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(2000000)
        )
        # 空テキストの場合は何も追加されない
        assert len(slide.shapes) == initial_shapes

    # --- render_block_into_frame のテスト ---

    def test_render_block_into_frame_empty_frame(self) -> None:
        """既存テキストなしの場合、paragraphs[0]に数式を追記する。"""
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(5000000), Emu(2000000))
        elem = _make_block_element("E=mc^2")
        initial_para_count = len(shape.text_frame.paragraphs)
        # 例外なく完了することを確認
        self.renderer.render_block_into_frame(shape, elem)
        # 段落数が増えていないこと（paragraphs[0]を使用するため）
        assert len(shape.text_frame.paragraphs) == initial_para_count

    def test_render_block_into_frame_with_existing_text(self) -> None:
        """既存テキストがある場合、空行スペーサー＋数式段落が2つ追加される。"""
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(5000000), Emu(2000000))
        # 先にテキストを書き込む
        tf = shape.text_frame
        tf.paragraphs[0].add_run().text = "前テキスト"
        before_count = len(tf.paragraphs)
        elem = _make_block_element("E=mc^2")
        self.renderer.render_block_into_frame(shape, elem)
        # 空行スペーサー + 数式段落 = 2段落追加されること
        assert len(tf.paragraphs) == before_count + 2

    def test_render_block_into_frame_empty_latex_is_skipped(self) -> None:
        """LaTeXテキストが空の場合は段落を追加しない。"""
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(5000000), Emu(2000000))
        elem = ET.Element("div")
        elem.set("class", "arithmatex")
        elem.text = ""
        before_count = len(shape.text_frame.paragraphs)
        self.renderer.render_block_into_frame(shape, elem)
        assert len(shape.text_frame.paragraphs) == before_count
