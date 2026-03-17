"""
GanttParser / GanttRenderer のユニットテスト。

gantt_parser.py のカスタムパーサーが各種Mermaid Gantt構文を正しく解析することと、
gantt_renderer.py がスライドに正しい図形・テーブルを追加することを検証する。
"""

from __future__ import annotations

from datetime import date, timedelta

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.gantt_parser import (
    GanttChart,
    GanttSection,
    GanttTask,
    _parse_duration,
    _resolve_end,
    _resolve_start,
    parse_gantt,
)
from qmd_to_pptx.mermaid.gantt_renderer import GanttRenderer


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------

def _make_slide():
    """テスト用スライドオブジェクトを生成する。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _all_shape_texts(slide) -> list[str]:
    """スライド上の全シェイプからテキストを収集する（テーブルセルも含む）。"""
    texts = []
    for sp in slide.shapes:
        # テーブルシェイプの場合はセルからテキストを収集する
        if sp.has_table:
            try:
                for row in sp.table.rows:
                    for cell in row.cells:
                        t = cell.text_frame.text.strip()
                        if t:
                            texts.append(t)
            except Exception:
                pass
            continue
        # 通常のシェイプ（text_frame を持つもの）
        try:
            t = sp.text_frame.text.strip()
            if t:
                texts.append(t)
        except AttributeError:
            pass
    return texts


_LEFT = 457200
_TOP = 457200
_WIDTH = 8229600
_HEIGHT = 4800000


# ===========================================================================
# TestParseDuration: 期間文字列の変換テスト
# ===========================================================================

class TestParseDuration:
    """_parse_duration() のユニットテスト。"""

    def test_days(self):
        """Nd 形式の期間を正しく変換する。"""
        assert _parse_duration("7d") == timedelta(days=7)

    def test_weeks(self):
        """Nw 形式の期間を正しく変換する（週 = 7日）。"""
        assert _parse_duration("2w") == timedelta(weeks=2)

    def test_hours_rounded_up(self):
        """Nh 形式の期間を日単位に切り上げる。"""
        assert _parse_duration("24h") == timedelta(days=1)
        assert _parse_duration("25h") == timedelta(days=2)
        assert _parse_duration("1h") == timedelta(days=1)

    def test_invalid_raises(self):
        """不正な期間文字列は ValueError を送出する。"""
        with pytest.raises(ValueError):
            _parse_duration("invalid")
        with pytest.raises(ValueError):
            _parse_duration("")


# ===========================================================================
# TestResolveStart: start日付解決テスト
# ===========================================================================

class TestResolveStart:
    """_resolve_start() のユニットテスト。"""

    def test_date_string(self):
        """日付文字列を正しく解決する。"""
        result = _resolve_start("2024-01-15", {})
        assert result == date(2024, 1, 15)

    def test_after_single(self):
        """after taskId 参照を正しく解決する。"""
        task_end_map = {"t1": date(2024, 1, 20)}
        result = _resolve_start("after t1", task_end_map)
        assert result == date(2024, 1, 20)

    def test_after_multiple_takes_max(self):
        """after タスクが複数の場合、最も遅い終了日を返す。"""
        task_end_map = {
            "t1": date(2024, 1, 10),
            "t2": date(2024, 1, 20),
        }
        result = _resolve_start("after t1 t2", task_end_map)
        assert result == date(2024, 1, 20)

    def test_after_unknown_raises(self):
        """unknown タスクへの after 参照は ValueError を送出する。"""
        with pytest.raises(ValueError):
            _resolve_start("after unknown_task", {})

    def test_date_string_yyyy_mm(self):
        """YYYY-MM 形式の日付文字列を正しく解決する（月初を補完）。"""
        result = _resolve_start("2024-03", {}, date_format="YYYY-MM")
        assert result == date(2024, 3, 1)

    def test_date_string_yyyy_mm_does_not_match_yyyymmdd_format(self):
        """YYYY-MM-DD 形式指定時に YYYY-MM 文字列は ValueError になる。"""
        with pytest.raises(ValueError):
            _resolve_start("2024-03", {}, date_format="YYYY-MM-DD")


# ===========================================================================
# TestResolveEnd: end日付解決テスト
# ===========================================================================

class TestResolveEnd:
    """_resolve_end() のユニットテスト。"""

    def test_date_string(self):
        """日付文字列を正しく解決する。"""
        result = _resolve_end("2024-02-01", {}, date(2024, 1, 15))
        assert result == date(2024, 2, 1)

    def test_duration(self):
        """期間文字列を start_date からのオフセットとして解決する。"""
        result = _resolve_end("14d", {}, date(2024, 1, 15))
        assert result == date(2024, 1, 29)

    def test_until_reference(self):
        """until taskId 参照を正しく解決する。"""
        task_end_map = {"t1": date(2024, 2, 10)}
        result = _resolve_end("until t1", task_end_map, date(2024, 1, 1))
        assert result == date(2024, 2, 10)

    def test_date_string_yyyy_mm(self):
        """YYYY-MM 形式の日付文字列を正しく解決する（月初を補完）。"""
        result = _resolve_end("2024-06", {}, date(2024, 3, 1), date_format="YYYY-MM")
        assert result == date(2024, 6, 1)


# ===========================================================================
# TestParseGantt: parse_gantt() の統合テスト
# ===========================================================================

class TestParseGantt:
    """parse_gantt() のテスト群。"""

    def test_basic_parse(self):
        """基本的なガントチャートのパースが正しく動作する。"""
        text = """\
gantt
    title プロジェクト計画
    dateFormat YYYY-MM-DD
    section 設計
        要件定義  :done, des1, 2024-01-01, 2024-01-07
        基本設計  :done, des2, 2024-01-08, 7d
"""
        chart = parse_gantt(text)
        assert chart.title == "プロジェクト計画"
        assert chart.date_format == "YYYY-MM-DD"
        assert len(chart.sections) == 1
        section = chart.sections[0]
        assert section.name == "設計"
        assert len(section.tasks) == 2

        t1 = section.tasks[0]
        assert t1.title == "要件定義"
        assert t1.task_id == "des1"
        assert t1.start_date == date(2024, 1, 1)
        assert t1.end_date == date(2024, 1, 7)
        assert t1.is_done is True

        t2 = section.tasks[1]
        assert t2.title == "基本設計"
        assert t2.end_date == date(2024, 1, 15)  # 2024-01-08 + 7d

    def test_active_crit_tags(self):
        """active/crit タグが正しく設定される。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section 実装
        タスクA  :active, crit, a1, 2024-02-01, 10d
"""
        chart = parse_gantt(text)
        task = chart.sections[0].tasks[0]
        assert task.is_active is True
        assert task.is_crit is True
        assert task.is_done is False

    def test_milestone_normalized(self):
        """マイルストーンは終了日が開始日+1日に正規化される。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section リリース
        デプロイ :milestone, 2024-03-01, 7d
"""
        chart = parse_gantt(text)
        task = chart.sections[0].tasks[0]
        assert task.is_milestone is True
        assert task.start_date == date(2024, 3, 1)
        assert task.end_date == date(2024, 3, 2)  # 期間に関係なく +1日

    def test_after_reference(self):
        """after taskId 参照が正しく解決される。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section 実装
        タスクA  :done, ta, 2024-01-15, 14d
        タスクB  :tb, after ta, 7d
"""
        chart = parse_gantt(text)
        tasks = chart.sections[0].tasks
        ta = tasks[0]
        tb = tasks[1]
        assert ta.end_date == date(2024, 1, 29)          # 2024-01-15 + 14d
        assert tb.start_date == date(2024, 1, 29)        # after ta
        assert tb.end_date == date(2024, 2, 5)           # + 7d

    def test_multiple_sections(self):
        """複数セクションが正しく解析される。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section 設計
        設計タスク :d1, 2024-01-01, 7d
    section 実装
        実装タスク :i1, 2024-01-08, 14d
    section テスト
        テストタスク :t1, 2024-01-22, 7d
"""
        chart = parse_gantt(text)
        assert len(chart.sections) == 3
        assert chart.sections[0].name == "設計"
        assert chart.sections[1].name == "実装"
        assert chart.sections[2].name == "テスト"

    def test_excludes_directive(self):
        """excludes ディレクティブが正しく解析される。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    excludes weekends, 2024-01-01
    section S
        T :t1, 2024-01-01, 7d
"""
        chart = parse_gantt(text)
        assert "weekends" in chart.excludes
        assert "2024-01-01" in chart.excludes

    def test_auto_task_id(self):
        """task_id が指定されない場合は自動生成される (task0, task1, ...)。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section S
        タスクA :2024-01-01, 7d
        タスクB :2024-01-08, 7d
"""
        chart = parse_gantt(text)
        tasks = chart.sections[0].tasks
        assert tasks[0].task_id == "task0"
        assert tasks[1].task_id == "task1"

    def test_all_tasks_flat(self):
        """all_tasks() が全セクションのタスクをフラットに返す。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section A
        T1 :t1, 2024-01-01, 3d
    section B
        T2 :t2, 2024-01-04, 3d
        T3 :t3, 2024-01-07, 3d
"""
        chart = parse_gantt(text)
        all_tasks = chart.all_tasks()
        assert len(all_tasks) == 3
        assert all_tasks[0].task_id == "t1"
        assert all_tasks[2].task_id == "t3"

    def test_no_title(self):
        """title が省略された場合は空文字列になる。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section S
        T :t1, 2024-01-01, 1d
"""
        chart = parse_gantt(text)
        assert chart.title == ""

    def test_comments_ignored(self):
        """%% コメント行はスキップされる。"""
        text = """\
gantt
    %% これはコメントです
    dateFormat YYYY-MM-DD
    section S
        %% タスクのコメント
        T :t1, 2024-01-01, 7d
"""
        chart = parse_gantt(text)
        assert len(chart.sections[0].tasks) == 1


# ===========================================================================
# TestParseGanttYyyyMm: YYYY-MM dateFormat の統合テスト
# ===========================================================================

class TestParseGanttYyyyMm:
    """ガントチャートの dateFormat YYYY-MM を対象とした parse_gantt() のテスト群。"""

    def test_yyyy_mm_date_string(self):
        """YYYY-MM 形式の日付が正しくパースされる（月初を補完）。"""
        text = """\
gantt
    title ロードマップ
    dateFormat YYYY-MM
    section Q1
        設計  :des, 2024-01, 2024-03
        実装  :impl, 2024-03, 2024-06
"""
        chart = parse_gantt(text)
        assert chart.date_format == "YYYY-MM"
        tasks = chart.sections[0].tasks
        assert len(tasks) == 2

        # 設計: 2024-01-01 〜 2024-03-01
        assert tasks[0].start_date == date(2024, 1, 1)
        assert tasks[0].end_date == date(2024, 3, 1)

        # 実装: 2024-03-01 〜 2024-06-01
        assert tasks[1].start_date == date(2024, 3, 1)
        assert tasks[1].end_date == date(2024, 6, 1)

    def test_yyyy_mm_with_duration(self):
        """YYYY-MM 開始日と期間指定の組み合わせが正しく動作する。"""
        text = """\
gantt
    dateFormat YYYY-MM
    section S
        タスクA :a1, 2024-02, 4w
"""
        chart = parse_gantt(text)
        t = chart.sections[0].tasks[0]
        assert t.start_date == date(2024, 2, 1)
        assert t.end_date == date(2024, 2, 1) + 4 * __import__('datetime').timedelta(weeks=1)

    def test_yyyy_mm_with_after_reference(self):
        """YYYY-MM 形式で after 参照が正しく解決される。"""
        text = """\
gantt
    dateFormat YYYY-MM
    section S
        タスクA :a1, 2024-01, 2024-03
        タスクB :b1, after a1, 2024-06
"""
        chart = parse_gantt(text)
        tasks = chart.sections[0].tasks
        # タスクB は タスクA の終了日 (2024-03-01) から開始する
        assert tasks[1].start_date == date(2024, 3, 1)
        assert tasks[1].end_date == date(2024, 6, 1)

    def test_yyyy_mm_yyyymmdd_string_not_matched(self):
        """dateFormat YYYY-MM 指定時に YYYY-MM-DD 文字列はタスクとして認識されない。"""
        text = """\
gantt
    dateFormat YYYY-MM
    section S
        タスクA :a1, 2024-01-15, 2024-03-31
"""
        chart = parse_gantt(text)
        # 日付が認識できないためタスクは 0 件
        assert len(chart.sections[0].tasks) == 0

class TestGanttRenderer:
    """GanttRenderer.render() のスライド出力テスト。"""

    def _make_chart(self) -> GanttChart:
        """テスト用GanttChartオブジェクトを生成する。"""
        t1 = GanttTask(
            title="要件定義", task_id="t1", section="設計",
            start_date=date(2024, 1, 1), end_date=date(2024, 1, 8),
            is_done=True,
        )
        t2 = GanttTask(
            title="コーディング", task_id="t2", section="実装",
            start_date=date(2024, 1, 8), end_date=date(2024, 1, 22),
            is_active=True,
        )
        t3 = GanttTask(
            title="デプロイ", task_id="t3", section="リリース",
            start_date=date(2024, 1, 22), end_date=date(2024, 1, 23),
            is_milestone=True,
        )
        t_crit = GanttTask(
            title="クリティカル", task_id="tc", section="実装",
            start_date=date(2024, 1, 10), end_date=date(2024, 1, 20),
            is_crit=True,
        )
        sections = [
            GanttSection("設計", [t1]),
            GanttSection("実装", [t2, t_crit]),
            GanttSection("リリース", [t3]),
        ]
        return GanttChart(
            title="テストプロジェクト",
            date_format="YYYY-MM-DD",
            excludes=[],
            sections=sections,
        )

    def test_render_adds_table(self):
        """render() でスライドにテーブルが1つ追加される。"""
        slide = _make_slide()
        chart = self._make_chart()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        # テーブルシェイプ（shape_type = TABLE = 19）が存在することを確認する
        tables = [sp for sp in slide.shapes if sp.has_table]
        assert len(tables) == 1

    def test_render_adds_task_bar_shapes(self):
        """render() でタスクバー（rightArrow）図形が追加される。"""
        slide = _make_slide()
        chart = self._make_chart()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        # テーブル以外の自由図形（タスクバー）が存在する
        non_table = [sp for sp in slide.shapes if not sp.has_table]
        # 4タスク（t1=done, t2=active, t_crit=crit, t3=milestone）なので4個以上の図形
        assert len(non_table) >= 4

    def test_render_milestone_diamond(self):
        """マイルストーンタスクは DIAMOND 図形 (auto_shape_type=4) で描画される。"""
        slide = _make_slide()
        t_ms = GanttTask(
            title="MS", task_id="ms1", section="S",
            start_date=date(2024, 1, 15), end_date=date(2024, 1, 16),
            is_milestone=True,
        )
        chart = GanttChart("", "YYYY-MM-DD", [], [GanttSection("S", [t_ms])])
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        # DIAMOND形状が存在することをOXMLで確認する
        from pptx.oxml.ns import qn as _qn
        diamond_found = False
        for sp in slide.shapes:
            sp_elem = sp._element
            prstGeom = sp_elem.find(".//" + _qn("a:prstGeom"))
            if prstGeom is not None and prstGeom.get("prst") == "diamond":
                diamond_found = True
                break
        assert diamond_found, "DIAMOND(diamond)図形が見つかりません"

    def test_render_right_arrow(self):
        """通常タスクは rightArrow 図形で描画される。"""
        slide = _make_slide()
        t_normal = GanttTask(
            title="タスク", task_id="n1", section="S",
            start_date=date(2024, 1, 1), end_date=date(2024, 1, 15),
        )
        chart = GanttChart("", "YYYY-MM-DD", [], [GanttSection("S", [t_normal])])
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        from pptx.oxml.ns import qn as _qn
        arrow_found = False
        for sp in slide.shapes:
            sp_elem = sp._element
            prstGeom = sp_elem.find(".//" + _qn("a:prstGeom"))
            if prstGeom is not None and prstGeom.get("prst") == "rightArrow":
                arrow_found = True
                break
        assert arrow_found, "rightArrow図形が見つかりません"

    def test_render_empty_chart_no_crash(self):
        """タスクが空のチャートは何も描画せずにクラッシュしない。"""
        slide = _make_slide()
        chart = GanttChart("空", "YYYY-MM-DD", [], [])
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # テーブルなしで終わること
        tables = [sp for sp in slide.shapes if sp.has_table]
        assert len(tables) == 0

    def test_render_header_text_exists(self):
        """テーブルのヘッダー行に「タスク」テキストが存在する。"""
        slide = _make_slide()
        chart = self._make_chart()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_shape_texts(slide)
        assert any("タスク" in t for t in texts), f"「タスク」が見つかりません: {texts}"

    def test_render_section_name_in_table(self):
        """テーブルのセクション行にセクション名が記入されている。"""
        slide = _make_slide()
        chart = self._make_chart()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_shape_texts(slide)
        assert any("設計" in t for t in texts), f"「設計」が見つかりません: {texts}"

    def test_render_task_name_in_table(self):
        """テーブルのタスク行にタスク名が記入されている。"""
        slide = _make_slide()
        chart = self._make_chart()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_shape_texts(slide)
        assert any("要件定義" in t for t in texts), f"「要件定義」が見つかりません: {texts}"

    def test_render_from_parse_gantt(self):
        """parse_gantt() + GanttRenderer でクラッシュなく描画できる。"""
        text = """\
gantt
    title 統合テスト
    dateFormat YYYY-MM-DD
    section 設計
        要件定義  :done, des1, 2024-01-01, 2024-01-07
        基本設計  :done, des2, 2024-01-08, 7d
    section 実装
        コーディング :active, dev1, 2024-01-15, 14d
        テスト       :        dev2, after dev1, 7d
    section リリース
        デプロイ :milestone, 2024-02-05, 1d
"""
        chart = parse_gantt(text)
        slide = _make_slide()
        renderer = GanttRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        tables = [sp for sp in slide.shapes if sp.has_table]
        assert len(tables) == 1

    def test_render_via_mermaid_renderer(self):
        """MermaidRenderer 経由でガントチャートが描画される（統合テスト）。"""
        import xml.etree.ElementTree as ET

        from qmd_to_pptx.mermaid.renderer import MermaidRenderer

        gantt_xml = """\
gantt
    title テスト
    dateFormat YYYY-MM-DD
    section 設計
        要件定義 :done, d1, 2024-01-01, 7d
"""
        # Mermaidコードブロック要素を模擬する
        code_elem = ET.Element("code")
        code_elem.set("class", "language-mermaid")
        code_elem.text = gantt_xml

        slide = _make_slide()
        renderer = MermaidRenderer()
        renderer.render(slide, code_elem, _LEFT, _TOP, _WIDTH, _HEIGHT)

        tables = [sp for sp in slide.shapes if sp.has_table]
        assert len(tables) == 1

    def test_granularity_week(self):
        """30〜90日の期間は週単位の時間軸になる（列数が適切）。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section S
        T1 :t1, 2024-01-01, 45d
"""
        chart = parse_gantt(text)
        from qmd_to_pptx.mermaid.gantt_renderer import _GRANULARITY_WEEK
        renderer = GanttRenderer()
        _axis_start, _axis_end, granularity, col_dates = renderer._compute_time_axis(chart.all_tasks())
        assert granularity == _GRANULARITY_WEEK
        assert 4 <= len(col_dates) <= 12

    def test_granularity_month(self):
        """90日以上の期間は月単位の時間軸になる。"""
        text = """\
gantt
    dateFormat YYYY-MM-DD
    section S
        T1 :t1, 2024-01-01, 2024-12-31
"""
        chart = parse_gantt(text)
        from qmd_to_pptx.mermaid.gantt_renderer import _GRANULARITY_MONTH
        renderer = GanttRenderer()
        _axis_start, _axis_end, granularity, col_dates = renderer._compute_time_axis(chart.all_tasks())
        assert granularity == _GRANULARITY_MONTH
        assert 1 <= len(col_dates) <= 24
