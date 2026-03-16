"""
gitgraph_parser.py および GitGraphRenderer のテストモジュール。

テスト対象:
  - parse_gitgraph(): Mermaid gitGraph テキストの解析
  - GitGraphRenderer.render(): PPTX スライドへの描画（シェープ生成の検証）

テスト観点:
  - direction の解析（LR/TB/BT/デフォルト）
  - commit 各属性（id/msg/type/tag）
  - branch 作成・order 指定
  - checkout / switch によるカレントブランチ切り替え
  - merge（属性なし・id/type/tag 付き）
  - cherry-pick（parent 指定あり・なし）
  - コメント行のスキップ
  - 自動生成コミット ID
  - mainBranchName の検出
  - 複数ブランチのコミット帰属
  - ブランチ order によるソート
  - レンダラー: シェープ数の確認（コミット円・ラベル等）
  - レンダラー: LR / TB / BT の各方向
  - レンダラー: タグ矩形の生成
  - レンダラー: マージ接続線の生成
  - レンダラー: cherry-pick コミットタイプ
  - レンダラー: REVERSE / HIGHLIGHT コミットタイプ
  - レンダラー: 空グラフ（コミット 0 件）
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.gitgraph_parser import (
    GitBranch,
    GitCommit,
    GitGraph,
    parse_gitgraph,
)
from qmd_to_pptx.mermaid.gitgraph_renderer import (
    GitGraphRenderer,
    _sort_branches,
    _calc_positions_1d,
)


# ---------------------------------------------------------------------------
# ヘルパー
# ---------------------------------------------------------------------------

def _slide():
    """テスト用の空スライドを返す。"""
    prs = Presentation()
    layout = prs.slide_layouts[6]  # ブランク
    return prs.slides.add_slide(layout)


def _default_area():
    """テスト描画エリアの標準座標（EMU）を返す。"""
    return {
        "left": Emu(457200),
        "top": Emu(685800),
        "width": Emu(8229600),
        "height": Emu(5181600),
    }


# ---------------------------------------------------------------------------
# parse_gitgraph テスト群
# ---------------------------------------------------------------------------

class TestParseGitGraphDirection:
    """direction 解析のテスト。"""

    def test_default_direction_is_lr(self):
        """direction 未指定のとき LR になること。"""
        g = parse_gitgraph("gitGraph\n  commit")
        assert g.direction == "LR"

    def test_explicit_lr(self):
        """LR: を指定すると direction が LR になること。"""
        g = parse_gitgraph("gitGraph LR:\n  commit")
        assert g.direction == "LR"

    def test_tb_direction(self):
        """TB: を指定すると direction が TB になること。"""
        g = parse_gitgraph("gitGraph TB:\n  commit")
        assert g.direction == "TB"

    def test_bt_direction(self):
        """BT: を指定すると direction が BT になること。"""
        g = parse_gitgraph("gitGraph BT:\n  commit")
        assert g.direction == "BT"

    def test_case_insensitive_direction(self):
        """direction キーワードが大文字小文字を区別しないこと。"""
        g = parse_gitgraph("gitgraph tb:\n  commit")
        assert g.direction == "TB"


class TestParseGitGraphCommit:
    """commit 行の解析テスト。"""

    def test_simple_commit(self):
        """引数なし commit が NORMAL タイプで追加されること。"""
        g = parse_gitgraph("gitGraph\n  commit")
        assert len(g.commits) == 1
        assert g.commits[0].commit_type == "NORMAL"
        assert g.commits[0].branch == "main"

    def test_commit_with_custom_id(self):
        """id 属性が正しく取得されること。"""
        g = parse_gitgraph('gitGraph\n  commit id: "Alpha"')
        assert g.commits[0].commit_id == "Alpha"

    def test_commit_with_single_quote_id(self):
        """シングルクォートの id も取得できること。"""
        g = parse_gitgraph("gitGraph\n  commit id: 'Beta'")
        assert g.commits[0].commit_id == "Beta"

    def test_commit_type_reverse(self):
        """REVERSE タイプが正しく解析されること。"""
        g = parse_gitgraph('gitGraph\n  commit type: REVERSE')
        assert g.commits[0].commit_type == "REVERSE"

    def test_commit_type_highlight(self):
        """HIGHLIGHT タイプが正しく解析されること。"""
        g = parse_gitgraph('gitGraph\n  commit type: HIGHLIGHT')
        assert g.commits[0].commit_type == "HIGHLIGHT"

    def test_commit_type_normal_explicit(self):
        """NORMAL タイプを明示指定しても NORMAL になること。"""
        g = parse_gitgraph('gitGraph\n  commit type: NORMAL')
        assert g.commits[0].commit_type == "NORMAL"

    def test_commit_with_tag(self):
        """tag 属性が正しく取得されること。"""
        g = parse_gitgraph('gitGraph\n  commit tag: "v1.0.0"')
        assert g.commits[0].tag == "v1.0.0"

    def test_commit_with_msg(self):
        """msg 属性が正しく取得されること。"""
        g = parse_gitgraph('gitGraph\n  commit msg: "初期コミット"')
        assert g.commits[0].msg == "初期コミット"

    def test_commit_all_attributes(self):
        """id/type/tag/msg の全属性を同時に指定できること。"""
        text = 'gitGraph\n  commit id: "C1" type: HIGHLIGHT tag: "v1" msg: "init"'
        g = parse_gitgraph(text)
        c = g.commits[0]
        assert c.commit_id == "C1"
        assert c.commit_type == "HIGHLIGHT"
        assert c.tag == "v1"
        assert c.msg == "init"

    def test_auto_generated_id(self):
        """id を指定しないとき自動生成IDが付与されること（空文字列でないこと）。"""
        g = parse_gitgraph("gitGraph\n  commit")
        assert g.commits[0].commit_id != ""

    def test_multiple_commits(self):
        """3 件のコミットが正しく追加されること。"""
        g = parse_gitgraph("gitGraph\n  commit\n  commit\n  commit")
        assert len(g.commits) == 3

    def test_commit_parent_chain(self):
        """2番目以降のコミットは前のコミットを親として持つこと。"""
        text = 'gitGraph\n  commit id: "A"\n  commit id: "B"'
        g = parse_gitgraph(text)
        assert g.commits[1].parents == ["A"]

    def test_first_commit_has_no_parent(self):
        """最初のコミットは親を持たないこと。"""
        g = parse_gitgraph('gitGraph\n  commit id: "A"')
        assert g.commits[0].parents == []


class TestParseGitGraphBranch:
    """branch 行の解析テスト。"""

    def test_branch_creates_new_branch(self):
        """branch コマンドで新しいブランチが作成されること。"""
        g = parse_gitgraph("gitGraph\n  commit\n  branch develop")
        assert any(b.name == "develop" for b in g.branches)

    def test_branch_switches_current(self):
        """branch 作成後はそのブランチがカレントブランチになること。"""
        text = "gitGraph\n  commit\n  branch develop\n  commit id: \"X\""
        g = parse_gitgraph(text)
        # "X" は develop ブランチに属するはず
        x_commit = next(c for c in g.commits if c.commit_id == "X")
        assert x_commit.branch == "develop"

    def test_branch_with_order(self):
        """branch の order 属性が正しく解析されること。"""
        g = parse_gitgraph("gitGraph\n  branch feature order: 3")
        feat = next(b for b in g.branches if b.name == "feature")
        assert feat.order == 3

    def test_default_branch_always_present(self):
        """デフォルトブランチ（main）が常に存在すること。"""
        g = parse_gitgraph("gitGraph\n  commit")
        assert any(b.name == "main" for b in g.branches)

    def test_branch_with_slash_in_name(self):
        """スラッシュを含むブランチ名が解析できること。"""
        g = parse_gitgraph("gitGraph\n  branch feature/login")
        assert any(b.name == "feature/login" for b in g.branches)

    def test_branch_with_quoted_name(self):
        """ダブルクォートで囲まれたブランチ名が解析できること。"""
        g = parse_gitgraph('gitGraph\n  commit\n  branch "cherry-pick"')
        assert any(b.name == "cherry-pick" for b in g.branches)


class TestParseGitGraphCheckout:
    """checkout / switch 行の解析テスト。"""

    def test_checkout_switches_branch(self):
        """checkout でカレントブランチが切り替わること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\""
        )
        g = parse_gitgraph(text)
        c_commit = next(c for c in g.commits if c.commit_id == "C")
        assert c_commit.branch == "main"

    def test_switch_alias_works(self):
        """switch コマンドが checkout の別名として機能すること。"""
        text = (
            "gitGraph\n"
            "  branch dev\n"
            "  switch main\n"
            "  commit id: \"M\""
        )
        g = parse_gitgraph(text)
        m_commit = next(c for c in g.commits if c.commit_id == "M")
        assert m_commit.branch == "main"

    def test_checkout_unknown_branch_ignored(self):
        """存在しないブランチへの checkout はエラーにならず無視されること。"""
        text = "gitGraph\n  commit\n  checkout nonexistent\n  commit"
        # 例外が出ないことを確認する
        g = parse_gitgraph(text)
        assert len(g.commits) == 2


class TestParseGitGraphMerge:
    """merge 行の解析テスト。"""

    def test_simple_merge_creates_merge_commit(self):
        """merge コマンドで MERGE タイプのコミットが作成されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  merge dev"
        )
        g = parse_gitgraph(text)
        merge_commits = [c for c in g.commits if c.commit_type == "MERGE"]
        assert len(merge_commits) == 1

    def test_merge_has_two_parents(self):
        """マージコミットは 2 つの親を持つこと。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            "  merge dev"
        )
        g = parse_gitgraph(text)
        merge_commit = next(c for c in g.commits if c.commit_type == "MERGE")
        assert len(merge_commit.parents) == 2
        assert "C" in merge_commit.parents
        assert "B" in merge_commit.parents

    def test_merge_with_custom_id(self):
        """merge の id 属性が正しく取得されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            '  merge dev id: "merge-1"'
        )
        g = parse_gitgraph(text)
        merge_commit = next(c for c in g.commits if c.commit_type == "MERGE")
        assert merge_commit.commit_id == "merge-1"

    def test_merge_with_tag(self):
        """merge の tag 属性が正しく取得されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            '  merge dev tag: "release"'
        )
        g = parse_gitgraph(text)
        merge_commit = next(c for c in g.commits if c.commit_type in ("MERGE",))
        assert merge_commit.tag == "release"

    def test_merge_with_reverse_type(self):
        """merge の type: REVERSE が反映されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  merge dev type: REVERSE"
        )
        g = parse_gitgraph(text)
        merge_commit = next(c for c in g.commits if "B" in c.parents or len(c.parents) == 2)
        assert merge_commit.commit_type == "REVERSE"

    def test_merge_belongs_to_current_branch(self):
        """マージコミットはカレントブランチに属すること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  merge dev"
        )
        g = parse_gitgraph(text)
        merge_commit = next(c for c in g.commits if c.commit_type == "MERGE")
        assert merge_commit.branch == "main"


class TestParseGitGraphCherryPick:
    """cherry-pick 行の解析テスト。"""

    def test_cherry_pick_creates_cherry_pick_commit(self):
        """cherry-pick コマンドで CHERRY_PICK タイプのコミットが生成されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B"'
        )
        g = parse_gitgraph(text)
        cp = next((c for c in g.commits if c.commit_type == "CHERRY_PICK"), None)
        assert cp is not None

    def test_cherry_pick_tag_is_source_id(self):
        """cherry-pick コミットのタグはコピー元コミット ID になること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B"'
        )
        g = parse_gitgraph(text)
        cp = next(c for c in g.commits if c.commit_type == "CHERRY_PICK")
        assert cp.tag == "B"

    def test_cherry_pick_cherry_from(self):
        """cherry-pick の cherry_from フィールドにコピー元 ID が設定されること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B"'
        )
        g = parse_gitgraph(text)
        cp = next(c for c in g.commits if c.commit_type == "CHERRY_PICK")
        assert cp.cherry_from == "B"

    def test_cherry_pick_belongs_to_current_branch(self):
        """cherry-pick コミットはカレントブランチに属すること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B"'
        )
        g = parse_gitgraph(text)
        cp = next(c for c in g.commits if c.commit_type == "CHERRY_PICK")
        assert cp.branch == "main"

    def test_cherry_pick_without_id_skipped(self):
        """id を指定しない cherry-pick はスキップされること。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            "  cherry-pick"
        )
        g = parse_gitgraph(text)
        cp_commits = [c for c in g.commits if c.commit_type == "CHERRY_PICK"]
        assert len(cp_commits) == 0

    def test_cherry_pick_with_parent(self):
        """parent 属性が指定された cherry-pick もエラーにならないこと。"""
        text = (
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B" parent: "A"'
        )
        g = parse_gitgraph(text)
        cp_commits = [c for c in g.commits if c.commit_type == "CHERRY_PICK"]
        assert len(cp_commits) == 1


class TestParseGitGraphComment:
    """コメント行処理のテスト。"""

    def test_comment_lines_skipped(self):
        """'%%' で始まる行はスキップされること。"""
        text = (
            "gitGraph\n"
            "  %% これはコメント\n"
            "  commit"
        )
        g = parse_gitgraph(text)
        assert len(g.commits) == 1

    def test_multiple_comment_lines(self):
        """複数のコメント行がすべてスキップされること。"""
        text = (
            "gitGraph\n"
            "  %% コメント1\n"
            "  %% コメント2\n"
            "  commit\n"
            "  %% コメント3\n"
            "  commit"
        )
        g = parse_gitgraph(text)
        assert len(g.commits) == 2


class TestParseGitGraphMainBranch:
    """mainBranchName の解析テスト。"""

    def test_default_main_branch_name(self):
        """デフォルトブランチ名が 'main' であること。"""
        g = parse_gitgraph("gitGraph\n  commit")
        assert g.main_branch == "main"

    def test_custom_main_branch_name_from_text(self):
        """gitGraph テキスト内の mainBranchName が正しく検出されること。"""
        text = (
            "mainBranchName: master\n"
            "gitGraph\n"
            "  commit"
        )
        g = parse_gitgraph(text)
        assert g.main_branch == "master"

    def test_custom_main_branch_from_frontmatter(self):
        """YAML front-matter の mainBranchName が正しく検出されること。"""
        g = parse_gitgraph("gitGraph\n  commit", yaml_frontmatter="mainBranchName: develop")
        assert g.main_branch == "develop"


class TestSortBranches:
    """_sort_branches のテスト。"""

    def test_no_order_preserves_definition_order(self):
        """order 未指定のブランチは定義順を維持すること。"""
        branches = [
            GitBranch("main", order=None),
            GitBranch("dev", order=None),
            GitBranch("feat", order=None),
        ]
        result = _sort_branches(branches)
        assert [b.name for b in result] == ["main", "dev", "feat"]

    def test_branches_with_order_sorted_by_order(self):
        """order 指定のブランチは order 値の昇順でソートされること。"""
        branches = [
            GitBranch("main", order=None),
            GitBranch("test1", order=3),
            GitBranch("test2", order=1),
        ]
        result = _sort_branches(branches)
        assert result[0].name == "main"
        assert result[1].name == "test2"
        assert result[2].name == "test1"

    def test_mixed_order_and_no_order(self):
        """order なしが先、order ありが後になること。"""
        branches = [
            GitBranch("a", order=2),
            GitBranch("b", order=None),
            GitBranch("c", order=None),
            GitBranch("d", order=1),
        ]
        result = _sort_branches(branches)
        names = [b.name for b in result]
        assert names.index("b") < names.index("d")
        assert names.index("c") < names.index("d")
        assert names.index("d") < names.index("a")


class TestCalcPositions1D:
    """_calc_positions_1d のテスト。"""

    def test_single_element_centered(self):
        """要素1個のとき中心位置が返ること。"""
        result = _calc_positions_1d(1, 1000, 200, 0)
        assert result == [500]

    def test_two_elements(self):
        """要素2個のとき等間隔に配置されること。"""
        result = _calc_positions_1d(2, 1000, 500, 0)
        assert len(result) == 2
        assert result[0] < result[1]

    def test_empty(self):
        """要素0個のとき空リストが返ること。"""
        result = _calc_positions_1d(0, 1000, 200, 0)
        assert result == []


# ---------------------------------------------------------------------------
# GitGraphRenderer テスト群
# ---------------------------------------------------------------------------

class TestGitGraphRendererShapes:
    """GitGraphRenderer のシェープ生成テスト。"""

    def test_empty_graph_no_shapes(self):
        """コミット 0 件では何もシェープが描画されないこと。"""
        slide = _slide()
        initial_count = len(slide.shapes)
        renderer = GitGraphRenderer()
        empty_graph = GitGraph()
        area = _default_area()
        renderer.render(slide, empty_graph, **area)
        assert len(slide.shapes) == initial_count

    def test_single_commit_adds_shapes(self):
        """コミット 1 件で複数のシェープが追加されること。"""
        slide = _slide()
        initial_count = len(slide.shapes)
        renderer = GitGraphRenderer()
        g = parse_gitgraph("gitGraph\n  commit")
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > initial_count

    def test_multiple_commits_lr(self):
        """LR 方向で複数コミットが正常に描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph LR:\n"
            "  commit id: \"A\"\n"
            "  commit id: \"B\"\n"
            "  commit id: \"C\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        # 最低でも: コミット円3 + レーン線1 + ブランチラベル1 + IDラベル3 = 8以上
        assert len(slide.shapes) >= 8

    def test_tb_direction_renders(self):
        """TB 方向で描画がエラーなく完了すること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph TB:\n"
            "  commit id: \"A\"\n"
            "  commit id: \"B\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_bt_direction_renders(self):
        """BT 方向で描画がエラーなく完了すること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph BT:\n"
            "  commit id: \"A\"\n"
            "  commit id: \"B\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_tag_creates_extra_shape(self):
        """タグが指定されているとき追加のシェープが生成されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()

        # タグなしの場合
        g_no_tag = parse_gitgraph('gitGraph\n  commit id: "A"')
        renderer.render(slide, g_no_tag, **_default_area())
        count_no_tag = len(slide.shapes)

        # タグありの場合
        slide2 = _slide()
        g_tag = parse_gitgraph('gitGraph\n  commit id: "A" tag: "v1.0"')
        renderer.render(slide2, g_tag, **_default_area())
        count_tag = len(slide2.shapes)

        # タグありは多くのシェープが生成されること
        assert count_tag > count_no_tag

    def test_merge_renders_without_error(self):
        """マージ操作が含まれる場合に描画がエラーなく完了すること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  merge dev id: \"M\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_cherry_pick_renders_without_error(self):
        """cherry-pick 操作を含む場合に描画がエラーなく完了すること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph\n"
            "  commit id: \"A\"\n"
            "  branch dev\n"
            "  commit id: \"B\"\n"
            "  checkout main\n"
            "  commit id: \"C\"\n"
            '  cherry-pick id: "B"'
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_reverse_commit_renders(self):
        """REVERSE タイプのコミットが描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph("gitGraph\n  commit type: REVERSE")
        renderer.render(slide, g, **_default_area())
        assert len(slide.shapes) > 0

    def test_highlight_commit_renders(self):
        """HIGHLIGHT タイプのコミットが描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph("gitGraph\n  commit type: HIGHLIGHT")
        renderer.render(slide, g, **_default_area())
        assert len(slide.shapes) > 0

    def test_multiple_branches_render(self):
        """複数ブランチを含む gitGraph が正常に描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph\n"
            "  commit id: \"1\"\n"
            "  branch develop\n"
            "  commit id: \"2\"\n"
            "  branch feature\n"
            "  commit id: \"3\"\n"
            "  checkout develop\n"
            "  commit id: \"4\"\n"
            "  checkout main\n"
            "  merge develop id: \"5\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_japanese_branch_name(self):
        """日本語ブランチ名でも描画がエラーにならないこと。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        # ブランチ名に日本語を使う場合はクォート必須（パーサーの制約）
        g = parse_gitgraph(
            "gitGraph\n"
            "  commit id: \"A\"\n"
            '  branch "開発ブランチ"\n'
            "  commit id: \"B\""
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_branch_order_affects_lane_assignment(self):
        """branch の order 指定がレーン割り当てに影響すること（エラーなく描画）。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        g = parse_gitgraph(
            "gitGraph\n"
            "  commit\n"
            "  branch test1 order: 3\n"
            "  branch test2 order: 1\n"
            "  checkout test2\n"
            "  commit\n"
            "  checkout test1\n"
            "  commit"
        )
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0

    def test_full_workflow_lr(self):
        """完全な Git ワークフローが LR 方向で描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        text = (
            "gitGraph\n"
            "    commit id: \"initial\"\n"
            "    commit id: \"feat: add login\"\n"
            "    branch develop\n"
            "    checkout develop\n"
            "    commit id: \"fix: typo\"\n"
            "    commit id: \"feat: add API\"\n"
            "    checkout main\n"
            '    merge develop id: "merge develop"\n'
            "    commit id: \"release: v1.0\""
        )
        g = parse_gitgraph(text)
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(g.commits) == 6
        assert len(slide.shapes) > 0

    def test_complex_graph_tb(self):
        """複数ブランチ・マージを含む複雑な gitGraph が TB 方向で描画されること。"""
        slide = _slide()
        renderer = GitGraphRenderer()
        text = (
            "gitGraph TB:\n"
            "    commit id: \"ZERO\"\n"
            "    branch develop\n"
            "    branch release\n"
            "    commit id: \"A\"\n"
            "    checkout main\n"
            "    commit id: \"ONE\"\n"
            "    checkout develop\n"
            "    commit id: \"B\"\n"
            "    checkout main\n"
            '    merge develop id: "MERGE"\n'
            "    commit id: \"TWO\""
        )
        g = parse_gitgraph(text)
        area = _default_area()
        renderer.render(slide, g, **area)
        assert len(slide.shapes) > 0
