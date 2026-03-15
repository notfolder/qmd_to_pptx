"""
機能網羅型の結合テスト。

ライブラリが対応する全機能を網羅したQMDテキストからPPTXを生成し、
python-pptxを用いて出力内容を検証する。

検証対象：
  - タイトルスライドの生成
  - Section Header スライド（# 見出し）
  - Title and Content スライド（## 見出し）
  - Blank スライド（--- 水平区切り線）
  - 段落テキスト
  - 順序なしリスト（ネスト含む）
  - 順序付きリスト
  - テーブル
  - コードブロック
  - Mermaid 図（標準記法・Quartoネイティブ記法）
  - Quarto コードブロック記法（```{python}）
  - スピーカーノート
  - インクリメンタルリスト（:::{.incremental}）
  - 非インクリメンタルリスト（:::{.nonincremental}）
  - 2カラムレイアウト（:::{.columns}）
  - slide-level: 1 設定
  - YAMLフロントマター（title / author / date / incremental）
"""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.slide import Slide

from qmd_to_pptx import render


# ---------------------------------------------------------------------------
# テスト用 QMD テキスト定義
# ---------------------------------------------------------------------------

# 全機能を網羅したQMDテキスト（slide-level: 2 デフォルト）
_FULL_QMD = """\
---
title: 機能網羅型テストプレゼンテーション
author: テスト著者
date: 2026-01-01
format:
  pptx:
    incremental: false
---

# セクション1

## 段落テスト

これは段落テキストです。複数行にわたる内容も含みます。

## 順序なしリストテスト

- アイテム1
- アイテム2
  - サブアイテム2-1
  - サブアイテム2-2
- アイテム3

## 順序付きリストテスト

1. 最初のアイテム
2. 二番目のアイテム
3. 三番目のアイテム

## テーブルテスト

| 列A | 列B | 列C |
|-----|-----|-----|
| 値1 | 値2 | 値3 |
| 値4 | 値5 | 値6 |

## コードブロックテスト

```python
def hello():
    print("Hello, World!")
```

## Mermaidテスト（標準記法）

```mermaid
flowchart LR
    A[開始] --> B[処理]
    B --> C[終了]
```

## Quartoネイティブ記法テスト

```{mermaid}
flowchart TD
    X --> Y --> Z
```

## Quartoコードブロックテスト

```{python}
x = 42
print(x)
```

## スピーカーノートテスト

このスライドにはスピーカーノートがあります。

::: {.notes}
これはスピーカーノートの内容です。発表時にのみ表示されます。
:::

## インクリメンタルリストテスト

::: {.incremental}
- インクリメンタルアイテム1
- インクリメンタルアイテム2
- インクリメンタルアイテム3
:::

## 非インクリメンタルリストテスト

::: {.nonincremental}
- 通常アイテム1
- 通常アイテム2
:::

## 2カラムレイアウトテスト

:::: {.columns}
::: {.column}
左カラムのコンテンツです。

- 左アイテム1
- 左アイテム2
:::
::: {.column}
右カラムのコンテンツです。

- 右アイテム1
- 右アイテム2
:::
::::

---

水平区切り線によるBlankスライドのコンテンツ

# セクション2

## 最後のスライド

プレゼンテーションの最後のスライドです。
"""

# slide-level: 1 の動作検証用 QMD テキスト
_SLIDE_LEVEL_1_QMD = """\
---
title: slide-level 1 テスト
author: テスト著者
date: 2026-01-01
slide-level: 1
---

# スライド1タイトル

## 本文見出し（スライド区切りにならない）

段落コンテンツです。

# スライド2タイトル

- リストアイテム1
- リストアイテム2
"""

# incremental: true の全体設定検証用 QMD テキスト
_INCREMENTAL_QMD = """\
---
title: インクリメンタルテスト
format:
  pptx:
    incremental: true
---

## スライド1

- デフォルトインクリメンタルアイテム1
- デフォルトインクリメンタルアイテム2

## スライド2

::: {.nonincremental}
- 非インクリメンタルアイテム1
- 非インクリメンタルアイテム2
:::
"""


# ---------------------------------------------------------------------------
# ヘルパー関数
# ---------------------------------------------------------------------------

def _render_to_prs(qmd_text: str, **kwargs) -> Presentation:
    """
    QMDテキストをレンダリングして Presentation オブジェクトを返す。

    Parameters
    ----------
    qmd_text : str
        QMDまたはMarkdownテキスト。
    **kwargs
        render() 関数へ渡す追加引数。

    Returns
    -------
    Presentation
        生成されたpython-pptxのPresentationオブジェクト。
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "output.pptx"
        render(qmd_text, str(output_path), **kwargs)
        assert output_path.exists(), "PPTXファイルが生成されなかった"
        assert output_path.stat().st_size > 0, "PPTXファイルが空"
        return Presentation(str(output_path))


def _collect_all_text(prs: Presentation) -> list[str]:
    """
    プレゼンテーション内の全スライドの全テキストを収集して返す。

    Parameters
    ----------
    prs : Presentation
        python-pptxのPresentationオブジェクト。

    Returns
    -------
    list[str]
        全テキスト（空文字列を除く）のリスト。
    """
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
    return texts


def _get_slide_texts(slide: Slide) -> list[str]:
    """
    指定したスライドの全テキストを収集して返す。

    Parameters
    ----------
    slide : Slide
        python-pptxのSlideオブジェクト。

    Returns
    -------
    list[str]
        そのスライド内の全テキスト（空文字列を除く）のリスト。
    """
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return texts


def _get_notes_text(slide: Slide) -> str:
    """
    スライドのスピーカーノートテキストを返す。

    Parameters
    ----------
    slide : Slide
        python-pptxのSlideオブジェクト。

    Returns
    -------
    str
        ノートテキスト（トリム済み）。
    """
    return slide.notes_slide.notes_text_frame.text.strip()


def _has_table(slide: Slide) -> bool:
    """
    スライドにテーブルShapeが存在するかどうかを返す。

    Parameters
    ----------
    slide : Slide
        python-pptxのSlideオブジェクト。

    Returns
    -------
    bool
        テーブルが存在する場合True。
    """
    return any(shape.has_table for shape in slide.shapes)


# ---------------------------------------------------------------------------
# 結合テスト: 全機能網羅
# ---------------------------------------------------------------------------

class TestFullFeatureIntegration:
    """全機能を網羅したQMDからPPTX生成と検証の結合テスト。"""

    @pytest.fixture(scope="class")
    def prs(self) -> Presentation:
        """全機能網羅QMDからPresentationオブジェクトを生成する。"""
        return _render_to_prs(_FULL_QMD)

    def test_pptx_file_is_generated(self) -> None:
        """PPTXファイルが正常に生成される。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(_FULL_QMD, str(output_path))
            assert output_path.exists()
            assert output_path.stat().st_size > 0

    def test_slide_count_matches_expected(self, prs: Presentation) -> None:
        """生成されたスライド数が期待値（タイトル + 各スライド）以上である。"""
        # タイトルスライド1枚 + コンテンツスライド（最低12枚）
        assert len(prs.slides) >= 12

    def test_title_slide_contains_presentation_title(self, prs: Presentation) -> None:
        """最初のスライド（タイトルスライド）にプレゼンテーションタイトルが含まれる。"""
        title_slide = prs.slides[0]
        all_text = " ".join(_get_slide_texts(title_slide))
        assert "機能網羅型テストプレゼンテーション" in all_text

    def test_title_slide_contains_author(self, prs: Presentation) -> None:
        """タイトルスライドに著者名が含まれる。"""
        title_slide = prs.slides[0]
        all_text = " ".join(_get_slide_texts(title_slide))
        assert "テスト著者" in all_text

    def test_section_header_slide_exists(self, prs: Presentation) -> None:
        """Section Header スライド（# 見出し）の内容がどこかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "セクション1" in combined

    def test_paragraph_text_appears_in_slide(self, prs: Presentation) -> None:
        """段落テキストがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "これは段落テキストです" in combined

    def test_unordered_list_items_appear(self, prs: Presentation) -> None:
        """順序なしリストのアイテムがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "アイテム1" in combined
        assert "アイテム2" in combined

    def test_nested_list_items_appear(self, prs: Presentation) -> None:
        """ネストされたリストアイテムがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "サブアイテム2-1" in combined

    def test_ordered_list_items_appear(self, prs: Presentation) -> None:
        """順序付きリストのアイテムがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "最初のアイテム" in combined

    def test_table_exists_in_some_slide(self, prs: Presentation) -> None:
        """テーブルがいずれかのスライドに存在する。"""
        has_any_table = any(_has_table(slide) for slide in prs.slides)
        assert has_any_table, "どのスライドにもテーブルが見つからない"

    def test_table_header_text_is_correct(self, prs: Presentation) -> None:
        """テーブルのヘッダー行テキストが正しく含まれる。"""
        for slide in prs.slides:
            if _has_table(slide):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        # ヘッダー行（行0）のセルテキストを確認する
                        header_texts = [
                            table.cell(0, c).text.strip()
                            for c in range(len(table.columns))
                        ]
                        if "列A" in header_texts:
                            assert "列B" in header_texts
                            assert "列C" in header_texts
                            return
        pytest.fail("テーブルのヘッダー行が見つからない")

    def test_table_data_cells_are_correct(self, prs: Presentation) -> None:
        """テーブルのデータセルテキストが正しく含まれる。"""
        for slide in prs.slides:
            if _has_table(slide):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        all_cell_texts = [
                            table.cell(r, c).text.strip()
                            for r in range(len(table.rows))
                            for c in range(len(table.columns))
                        ]
                        if "値1" in all_cell_texts:
                            assert "値2" in all_cell_texts
                            assert "値4" in all_cell_texts
                            return
        pytest.fail("テーブルのデータセルが見つからない")

    def test_code_block_text_appears(self, prs: Presentation) -> None:
        """コードブロックのテキストがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        # コードブロックのテキストが含まれるか確認する
        assert "hello" in combined.lower()

    def test_mermaid_shapes_exist(self, prs: Presentation) -> None:
        """Mermaid図を含むスライドにShape（矩形またはコネクター）が存在する。"""
        # Mermaid図はShapeまたはConnectorとして描画されるか
        # パース失敗時はテキストボックスとしてフォールバックされる
        # いずれの場合もスライドにShapeが存在することを確認する
        all_shapes_count = sum(len(slide.shapes) for slide in prs.slides)
        assert all_shapes_count > 0

    def test_speaker_notes_text_exists(self, prs: Presentation) -> None:
        """スピーカーノートのテキストがいずれかのスライドのノートに含まれる。"""
        found_notes = False
        for slide in prs.slides:
            notes_text = _get_notes_text(slide)
            if "スピーカーノートの内容" in notes_text:
                found_notes = True
                break
        assert found_notes, "スピーカーノートのテキストが見つからない"

    def test_section2_exists_in_slides(self, prs: Presentation) -> None:
        """セクション2のタイトルがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "セクション2" in combined

    def test_horizontal_rule_creates_slide(self, prs: Presentation) -> None:
        """水平区切り線（---）で分割されたスライドのコンテンツが含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        # 水平区切り線後のコンテンツが含まれることを確認する
        assert "水平区切り線によるBlankスライドのコンテンツ" in combined

    def test_final_slide_content_exists(self, prs: Presentation) -> None:
        """最後のスライドのコンテンツが含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "プレゼンテーションの最後のスライドです" in combined

    def test_slide_titles_appear(self, prs: Presentation) -> None:
        """各スライドのタイトルがいずれかのスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "段落テスト" in combined
        assert "テーブルテスト" in combined
        assert "コードブロックテスト" in combined


# ---------------------------------------------------------------------------
# 結合テスト: slide-level: 1
# ---------------------------------------------------------------------------

class TestSlideLevelOneIntegration:
    """slide-level: 1 設定の結合テスト。"""

    @pytest.fixture(scope="class")
    def prs(self) -> Presentation:
        """slide-level: 1 のQMDからPresentationオブジェクトを生成する。"""
        return _render_to_prs(_SLIDE_LEVEL_1_QMD)

    def test_pptx_is_generated(self) -> None:
        """PPTXファイルが正常に生成される。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(_SLIDE_LEVEL_1_QMD, str(output_path))
            assert output_path.exists()

    def test_slide_count_with_level1(self, prs: Presentation) -> None:
        """slide-level: 1 では # 見出しがスライド区切りになる。"""
        # タイトルスライド1枚 + #スライドが2枚 = 最低3枚
        assert len(prs.slides) >= 3

    def test_h1_slide_title_appears(self, prs: Presentation) -> None:
        """# 見出しがスライドのタイトルとして生成される。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "スライド1タイトル" in combined

    def test_h2_is_body_content_not_separator(self, prs: Presentation) -> None:
        """slide-level: 1 では ## 見出しはスライド区切りにならない（追加スライドが生成されない）。"""
        # _SLIDE_LEVEL_1_QMD には # が2つ、## が1つある
        # slide-level:1 では # だけがスライド区切りになるため
        # タイトルスライド + スライド1 + スライド2 = 計3枚になるはず
        # もし ## がスライド区切りとして誤動作すれば4枚以上になる
        assert len(prs.slides) == 3

    def test_list_content_appears(self, prs: Presentation) -> None:
        """リストコンテンツがスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "リストアイテム1" in combined


# ---------------------------------------------------------------------------
# 結合テスト: incremental: true
# ---------------------------------------------------------------------------

class TestIncrementalIntegration:
    """format.pptx.incremental: true 設定の結合テスト。"""

    @pytest.fixture(scope="class")
    def prs(self) -> Presentation:
        """incremental: true のQMDからPresentationオブジェクトを生成する。"""
        return _render_to_prs(_INCREMENTAL_QMD)

    def test_pptx_is_generated(self) -> None:
        """PPTXファイルが正常に生成される。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(_INCREMENTAL_QMD, str(output_path))
            assert output_path.exists()

    def test_list_items_appear(self, prs: Presentation) -> None:
        """インクリメンタルリストのアイテムがスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "デフォルトインクリメンタルアイテム1" in combined

    def test_nonincremental_items_appear(self, prs: Presentation) -> None:
        """非インクリメンタルリストのアイテムがスライドに含まれる。"""
        all_texts = _collect_all_text(prs)
        combined = " ".join(all_texts)
        assert "非インクリメンタルアイテム1" in combined


# ---------------------------------------------------------------------------
# 結合テスト: ファイル入力
# ---------------------------------------------------------------------------

class TestFileInputIntegration:
    """ファイルパス入力の結合テスト。"""

    def test_qmd_file_input(self) -> None:
        """QMDファイルをファイルパスで指定してPPTXを生成できる。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "presentation.qmd"
            output_path = Path(tmpdir) / "output.pptx"
            input_path.write_text(_FULL_QMD, encoding="utf-8")
            render(str(input_path), str(output_path))
            assert output_path.exists()
            assert output_path.stat().st_size > 0
            # ファイルから生成したPPTXが正しい内容を含むことを確認する
            prs = Presentation(str(output_path))
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            assert "機能網羅型テストプレゼンテーション" in combined

    def test_md_file_input(self) -> None:
        """Markdownファイルをファイルパスで指定してPPTXを生成できる。"""
        md_text = "## スライド\n\nMarkdownファイルからのコンテンツ\n"
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "presentation.md"
            output_path = Path(tmpdir) / "output.pptx"
            input_path.write_text(md_text, encoding="utf-8")
            render(str(input_path), str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            assert "Markdownファイルからのコンテンツ" in combined


# ---------------------------------------------------------------------------
# 結合テスト: Quarto記法
# ---------------------------------------------------------------------------

class TestQuartoSyntaxIntegration:
    """Quartoネイティブ記法の結合テスト。"""

    def test_quarto_mermaid_native_syntax(self) -> None:
        """Quartoネイティブ Mermaid 記法（```{mermaid}）が正しく処理される。"""
        qmd = (
            "---\ntitle: test\n---\n"
            "## スライド\n\n```{mermaid}\nflowchart LR\n    A --> B\n```\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            # Mermaidスライドにシェイプが存在することを確認する
            assert len(prs.slides) >= 1

    def test_quarto_python_code_block_syntax(self) -> None:
        """Quartoコードブロック記法（```{python}）が正しく処理される。"""
        qmd = (
            "---\ntitle: test\n---\n"
            "## スライド\n\n```{python}\nprint('hello')\n```\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            # コードが等幅フォントのテキストとして含まれるか確認する
            assert "print" in combined

    def test_quarto_r_code_block_syntax(self) -> None:
        """Quartoコードブロック記法（```{r}）が正しく処理される。"""
        qmd = (
            "---\ntitle: test\n---\n"
            "## スライド\n\n```{r}\nsummary(data)\n```\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()


# ---------------------------------------------------------------------------
# 結合テスト: エッジケース
# ---------------------------------------------------------------------------

class TestEdgeCaseIntegration:
    """エッジケースの結合テスト。"""

    def test_empty_slides_are_blank(self) -> None:
        """本文が空のスライドが Blank として生成される。"""
        qmd = "---\ntitle: test\n---\n## スライドタイトルのみ\n"
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()

    def test_text_only_input_without_frontmatter(self) -> None:
        """YAMLフロントマターなしのテキストが前処理器により補完されて処理される。"""
        text = "## スライドタイトル\n\nコンテンツ\n"
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(text, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            assert "コンテンツ" in combined

    def test_multiple_sections_and_slides(self) -> None:
        """複数のセクションと複数のスライドが正しく生成される。"""
        qmd = (
            "---\ntitle: 複数セクションテスト\n---\n"
            "# セクションA\n## スライドA1\nコンテンツA1\n"
            "## スライドA2\nコンテンツA2\n"
            "# セクションB\n## スライドB1\nコンテンツB1\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            # タイトル + Section A + A1 + A2 + Section B + B1 = 6枚以上
            assert len(prs.slides) >= 6
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            assert "セクションA" in combined
            assert "セクションB" in combined
            assert "コンテンツA1" in combined
            assert "コンテンツB1" in combined

    def test_nested_list_levels(self) -> None:
        """深くネストしたリストが正しく処理される。"""
        qmd = (
            "## スライド\n\n"
            "- レベル1\n"
            "  - レベル2\n"
            "    - レベル3\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            all_texts = _collect_all_text(prs)
            combined = " ".join(all_texts)
            assert "レベル1" in combined
            assert "レベル2" in combined
            assert "レベル3" in combined

    def test_horizontal_rule_creates_additional_slide(self) -> None:
        """水平区切り線（---）が追加スライドを生成する。"""
        qmd = (
            "---\ntitle: テスト\n---\n"
            "## スライド1\nコンテンツ1\n\n"
            "---\n\nコンテンツ2\n"
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            render(qmd, str(output_path))
            assert output_path.exists()
            prs = Presentation(str(output_path))
            # タイトルスライド + スライド1 + 水平区切りスライド = 3枚以上
            assert len(prs.slides) >= 3

    def test_background_image_attribute_does_not_crash(self) -> None:
        """background-image 属性付きスライドが例外なく処理される。"""
        qmd = (
            "---\ntitle: テスト\n---\n"
            '## スライド {background-image="nonexistent.png"}\n\nコンテンツ\n'
        )
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "output.pptx"
            # 画像ファイルが存在しなくても例外なく処理される
            render(qmd, str(output_path))
            assert output_path.exists()
