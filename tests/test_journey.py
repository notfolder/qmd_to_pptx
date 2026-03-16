"""
JourneyParser / JourneyRenderer のユニットテスト。

journey_parser.py のカスタムパーサーが各種 Mermaid journey 構文を正しく解析すること、
journey_renderer.py がスライドにエラーなくジャーニー図を追加することを検証する。
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.journey_parser import (
    JourneyChart,
    JourneyTask,
    parse_journey,
)
from qmd_to_pptx.mermaid.journey_renderer import JourneyRenderer, _score_to_emoji


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------

def _make_slide():
    """テスト用スライドオブジェクトを生成する。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _all_text(slide) -> list[str]:
    """スライド上の全シェープからテキストを収集する（空文字は除外）。"""
    texts = []
    for sp in slide.shapes:
        try:
            t = sp.text_frame.text.strip()
            if t:
                texts.append(t)
        except AttributeError:
            pass
    return texts


_LEFT = 457200
_TOP = 457200
_WIDTH = 8229600
_HEIGHT = 4800000


# ===========================================================================
# TestJourneyParser: parse_journey() のテスト
# ===========================================================================

class TestJourneyParser:
    """parse_journey() の正常系・異常系テスト。"""

    def test_タイトルのパース(self):
        """title 行が正しくパースされること。"""
        text = """\
journey
    title My working day
    section Go to work
      Make tea: 5: Me
"""
        result = parse_journey(text)
        assert result.title == "My working day"

    def test_セクション名リストのパース(self):
        """section 行が出現順に sections リストへ格納されること。"""
        text = """\
journey
    title Test
    section Alpha
      Task A: 4: Me
    section Beta
      Task B: 2: Me
    section Gamma
      Task C: 3
"""
        result = parse_journey(text)
        assert result.sections == ["Alpha", "Beta", "Gamma"]

    def test_タスクの基本パース(self):
        """タスク名・スコア・アクターが正しくパースされること。"""
        text = """\
journey
    title Test
    section Work
      Make tea: 5: Me
      Go upstairs: 3: Me
      Do work: 1: Me, Cat
"""
        result = parse_journey(text)
        assert len(result.tasks) == 3
        assert result.tasks[0].task == "Make tea"
        assert result.tasks[0].score == 5
        assert result.tasks[0].people == ["Me"]
        assert result.tasks[0].section == "Work"

        assert result.tasks[2].task == "Do work"
        assert result.tasks[2].score == 1
        assert result.tasks[2].people == ["Me", "Cat"]

    def test_アクターなしのタスクパース(self):
        """アクターなし（"タスク名 : スコア" のみ）で people=[] になること。"""
        text = """\
journey
    section Alone
      Solo task: 3
"""
        result = parse_journey(text)
        assert result.tasks[0].people == []

    def test_複数アクターのパース(self):
        """カンマ区切りの複数アクターが正しくリスト化されること。"""
        text = """\
journey
    section Team
      Big meeting: 4: Alice, Bob, Carol
"""
        result = parse_journey(text)
        assert result.tasks[0].people == ["Alice", "Bob", "Carol"]

    def test_アクターリストの重複排除と順序保持(self):
        """全アクターが出現順で重複排除されること。"""
        text = """\
journey
    section S1
      T1: 5: Alice, Bob
      T2: 3: Bob, Carol
      T3: 4: Alice
"""
        result = parse_journey(text)
        # Alice -> Bob -> Carol の順（初出順、重複なし）
        assert result.actors == ["Alice", "Bob", "Carol"]

    def test_スコアのクランプ_上限(self):
        """スコアが 5 を超えた場合に 5 にクランプされること。"""
        text = """\
journey
    section S
      TooHigh: 9: Me
"""
        result = parse_journey(text)
        assert result.tasks[0].score == 5

    def test_スコアのクランプ_下限(self):
        """スコアが 1 未満（0 等）の場合に 1 にクランプされること。"""
        text = """\
journey
    section S
      TooLow: 0: Me
"""
        result = parse_journey(text)
        assert result.tasks[0].score == 1

    def test_タイトルなし(self):
        """title がない場合は title が空文字列になること。"""
        text = """\
journey
    section Go
      Run: 4: Me
"""
        result = parse_journey(text)
        assert result.title == ""

    def test_セクションなしタスク(self):
        """section 行がなくてもタスクをパースできること、sections=[""] になること。"""
        text = """\
journey
    No section task: 3: Me
"""
        result = parse_journey(text)
        assert result.tasks[0].task == "No section task"
        assert result.sections == [""]

    def test_コメント行をスキップ(self):
        """"%%" ではじまるコメント行がスキップされること。"""
        text = """\
journey
    %% This is a comment
    title Commented
    section Go
      %% ignore this
      Valid task: 4: Me
"""
        result = parse_journey(text)
        assert result.title == "Commented"
        assert len(result.tasks) == 1
        assert result.tasks[0].task == "Valid task"

    def test_複数セクションにまたがるタスク所属(self):
        """各タスクに正しいセクション名が付与されること。"""
        text = """\
journey
    section First
      T1: 5: Me
      T2: 4: Me
    section Second
      T3: 2: Me
"""
        result = parse_journey(text)
        assert result.tasks[0].section == "First"
        assert result.tasks[1].section == "First"
        assert result.tasks[2].section == "Second"

    def test_空テキストは空チャートを返す(self):
        """空文字列でも例外を発生させず空チャートを返すこと。"""
        result = parse_journey("")
        assert result.tasks == []
        assert result.title == ""

    def test_アクターのトリミング(self):
        """アクター名の前後の空白がトリミングされること。"""
        text = """\
journey
    section S
      T: 3:  Alice ,  Bob 
"""
        result = parse_journey(text)
        assert result.tasks[0].people == ["Alice", "Bob"]


# ===========================================================================
# TestScoreToEmoji: スコア → 絵文字変換テスト
# ===========================================================================

class TestScoreToEmoji:
    """_score_to_emoji() の境界値テスト。"""

    def test_score_5_is_happy(self):
        """score=5 は 😊 (笑顔) になること。"""
        assert _score_to_emoji(5) == "😊"

    def test_score_4_is_happy(self):
        """score=4 は 😊 (笑顔) になること（>3）。"""
        assert _score_to_emoji(4) == "😊"

    def test_score_3_is_neutral(self):
        """score=3 は 😐 (中立) になること（==3）。"""
        assert _score_to_emoji(3) == "😐"

    def test_score_2_is_sad(self):
        """score=2 は 😢 (悲しい) になること（<3）。"""
        assert _score_to_emoji(2) == "😢"

    def test_score_1_is_sad(self):
        """score=1 は 😢 (悲しい) になること。"""
        assert _score_to_emoji(1) == "😢"


# ===========================================================================
# TestJourneyRenderer: JourneyRenderer.render() のテスト
# ===========================================================================

class TestJourneyRenderer:
    """JourneyRenderer.render() の動作検証テスト。"""

    def test_基本レンダリングがクラッシュしない(self):
        """標準的な journey 構文でレンダリングが例外なく完了すること。"""
        text = """\
journey
    title My working day
    section Go to work
      Make tea: 5: Me
      Go upstairs: 3: Me
      Do work: 1: Me, Cat
    section Go home
      Go downstairs: 5: Me
      Sit down: 5: Me
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # 例外なく完了したことを確認する

    def test_シェープ数が最低要件を満たす(self):
        """タスク・感情アイコン・タスクカード・凡例のシェープが生成されること。"""
        text = """\
journey
    title Test
    section S1
      Task A: 5: Alice
      Task B: 2: Bob
    section S2
      Task C: 3: Alice
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        shapes = list(slide.shapes)
        # タスク数 3 × (感情+カード) + セクション×2 + 凡例 + タイトル≧ 8
        assert len(shapes) >= 8

    def test_タイトルがテキストとして描画される(self):
        """タイトルがスライドのシェープテキストに含まれること。"""
        text = """\
journey
    title My Journey Title
    section S
      T: 4: Me
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_text(slide)
        assert "My Journey Title" in texts

    def test_タスク名がテキストとして描画される(self):
        """タスク名がスライドのシェープテキストに含まれること。"""
        text = """\
journey
    title Test
    section Work
      Unique Task XYZ: 3: Me
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_text(slide)
        assert "Unique Task XYZ" in texts

    def test_感情絵文字がテキストとして描画される(self):
        """感情絵文字（😊 / 😐 / 😢）がスライドのシェープテキストに存在すること。"""
        text = """\
journey
    section S
      Happy: 5: Me
      Neutral: 3: Me
      Sad: 1: Me
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        all_texts = " ".join(_all_text(slide))
        assert "😊" in all_texts
        assert "😐" in all_texts
        assert "😢" in all_texts

    def test_アクター名が凡例に描画される(self):
        """アクター名がスライドのシェープテキストに含まれること。"""
        text = """\
journey
    section S
      Task: 4: Alice, Bob
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_text(slide)
        assert "Alice" in texts
        assert "Bob" in texts

    def test_タスク0件は何も描画しない(self):
        """タスクが 0 件の場合、スライドにシェープが追加されないこと。"""
        chart = JourneyChart(title="Empty", sections=[], tasks=[], actors=[])
        slide = _make_slide()
        initial_count = len(list(slide.shapes))
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # シェープ数が増えていないことを確認する
        assert len(list(slide.shapes)) == initial_count

    def test_アクターなしでもクラッシュしない(self):
        """アクター指定なし（people=[]）でもレンダリングが正常に完了すること。"""
        text = """\
journey
    section Work
      Solo: 4
      Pair: 2
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # 例外なく完了したことを確認する

    def test_セクション名がテキストとして描画される(self):
        """セクション名がスライドのシェープテキストに含まれること。"""
        text = """\
journey
    section かいしゃに行く
      電車に乗る: 3: 社員
    section うちに帰る
      電車に乗る: 4: 社員
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)

        texts = _all_text(slide)
        assert "かいしゃに行く" in texts
        assert "うちに帰る" in texts

    def test_多数タスクでクラッシュしない(self):
        """タスク数が多い場合（10件）でもレンダリングが正常完了すること。"""
        lines = ["journey", "    title Many Tasks"]
        for i in range(10):
            sec = i // 3
            lines.append(f"    section Section{sec}")
            lines.append(f"      Task{i}: {(i % 5) + 1}: User")
        text = "\n".join(lines)

        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # 例外なく完了したことを確認する

    def test_日本語タスク名でクラッシュしない(self):
        """日本語タスク名・セクション名でもレンダリングが正常完了すること。"""
        text = """\
journey
    title ユーザージャーニー
    section 朝の準備
      起床する: 2: ユーザー
      朝食を食べる: 4: ユーザー
    section 出社
      電車に乗る: 3: ユーザー
      会社に到着: 5: ユーザー
"""
        slide = _make_slide()
        chart = parse_journey(text)
        renderer = JourneyRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "ユーザージャーニー" in texts
