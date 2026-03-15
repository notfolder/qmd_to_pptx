"""
FlowchartRenderer の単体テスト。

ノード形状14種類・エッジ矢印種別・線種・ラベル描画の動作を確認する。
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.flowchart import FlowchartRenderer, _SHAPE_MAP, _EDGE_ARROW_MAP
from qmd_to_pptx.mermaid_renderer import MermaidRenderer


def _make_slide():
    """テスト用スライドを生成して返す。"""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Blankレイアウト
    return prs.slides.add_slide(layout)


def _make_mermaid_element(mermaid_text: str) -> ET.Element:
    """テスト用Mermaid code要素を生成する。"""
    elem = ET.Element("code")
    elem.set("class", "language-mermaid")
    elem.text = mermaid_text
    return elem


# ------------------------------------------------------------------
# _SHAPE_MAP の定義テスト
# ------------------------------------------------------------------

class TestShapeMap:
    """_SHAPE_MAPが全14種類のMermaidノード形状を網羅することを確認する。"""

    _EXPECTED_TYPES = [
        "square", "round", "stadium", "subroutine",
        "cylinder", "circle", "odd", "diamond",
        "hexagon", "lean_right", "lean_left",
        "trapezoid", "inv_trapezoid", "doublecircle",
    ]

    def test_all_14_types_defined(self) -> None:
        """_SHAPE_MAPがMermaid標準14種のノード形状すべてを含む。"""
        for node_type in self._EXPECTED_TYPES:
            assert node_type in _SHAPE_MAP, f"{node_type} が _SHAPE_MAP に定義されていない"

    def test_shape_type_values_are_positive_int(self) -> None:
        """_SHAPE_MAPの値は正の整数（MSO_AUTO_SHAPE_TYPE）であること。"""
        for node_type, shape_val in _SHAPE_MAP.items():
            assert isinstance(shape_val, int) and shape_val > 0, (
                f"{node_type} のシェイプ値 {shape_val} が正の整数でない"
            )


# ------------------------------------------------------------------
# _EDGE_ARROW_MAP の定義テスト
# ------------------------------------------------------------------

class TestEdgeArrowMap:
    """_EDGE_ARROW_MAPが全7種類のMermaidエッジ矢印種別を網羅することを確認する。"""

    _EXPECTED_TYPES = [
        "arrow_open", "arrow_point", "double_arrow_point",
        "arrow_circle", "double_arrow_circle",
        "arrow_cross", "double_arrow_cross",
    ]

    def test_all_7_types_defined(self) -> None:
        """_EDGE_ARROW_MAPがMermaid標準7種のエッジタイプすべてを含む。"""
        for edge_type in self._EXPECTED_TYPES:
            assert edge_type in _EDGE_ARROW_MAP, (
                f"{edge_type} が _EDGE_ARROW_MAP に定義されていない"
            )

    def test_arrow_point_has_headend(self) -> None:
        """通常矢印（arrow_point）はheadEndを持つ。"""
        assert "headEnd" in _EDGE_ARROW_MAP["arrow_point"]

    def test_double_arrow_point_has_both_ends(self) -> None:
        """両方向矢印（double_arrow_point）はheadEndとtailEndを持つ。"""
        conf = _EDGE_ARROW_MAP["double_arrow_point"]
        assert "headEnd" in conf
        assert "tailEnd" in conf

    def test_arrow_open_has_no_ends(self) -> None:
        """開放矢印（arrow_open）は矢印指定なし。"""
        assert _EDGE_ARROW_MAP["arrow_open"] == {}


# ------------------------------------------------------------------
# FlowchartRenderer クラスのテスト
# ------------------------------------------------------------------

class TestFlowchartRenderer:
    """FlowchartRenderer の描画メソッドを確認する。"""

    def setup_method(self) -> None:
        """テスト前にFlowchartRendererインスタンスを生成する。"""
        self.renderer = FlowchartRenderer()

    def _make_graph_data(
        self,
        vertices: dict,
        edges: list[dict] | None = None,
    ) -> dict:
        """テスト用のgraph_data辞書を生成する。"""
        return {
            "vertices": vertices,
            "edges": edges or [],
        }

    def test_render_single_node_does_not_raise(self) -> None:
        """単一ノードをレンダリングしても例外が発生しない。"""
        slide = _make_slide()
        graph_data = self._make_graph_data({"A": {"text": "開始", "type": "circle"}})
        self.renderer.render(slide, graph_data, "fallback", Emu(0), Emu(0), Emu(9000000), Emu(5000000))

    def test_render_adds_node_shapes(self) -> None:
        """ノードがスライドに追加されることを確認する。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = self._make_graph_data({
            "A": {"text": "A", "type": "square"},
            "B": {"text": "B", "type": "diamond"},
        })
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        # 2ノード分のShapeが追加される（コネクターは含まない）
        assert len(slide.shapes) >= initial + 2

    def test_render_edge_adds_connector(self) -> None:
        """エッジがコネクターとして追加されることを確認する。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = self._make_graph_data(
            vertices={"A": {"text": "A", "type": "square"}, "B": {"text": "B", "type": "square"}},
            edges=[{"start": "A", "end": "B", "stroke": "normal", "type": "arrow_point", "text": ""}],
        )
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        # 2ノード + 1コネクター以上のShapeが追加される
        assert len(slide.shapes) >= initial + 3

    def test_render_invisible_edge_skips_connector(self) -> None:
        """invisible エッジはコネクターを生成しない。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = self._make_graph_data(
            vertices={"A": {"text": "A", "type": "square"}, "B": {"text": "B", "type": "square"}},
            edges=[{"start": "A", "end": "B", "stroke": "invisible", "type": "arrow_point", "text": ""}],
        )
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        # コネクターなしのため追加されるのは2ノード分のみ（コネクター0個）
        added = len(slide.shapes) - initial
        assert added == 2

    def test_render_edge_with_label_adds_textbox(self) -> None:
        """ラベル付きエッジがテキストボックスを追加することを確認する。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = self._make_graph_data(
            vertices={"A": {"text": "A", "type": "square"}, "B": {"text": "B", "type": "square"}},
            edges=[{"start": "A", "end": "B", "stroke": "normal", "type": "arrow_point", "text": "条件"}],
        )
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        # 2ノード + 1コネクター + 1テキストボックス以上のShapeが追加される
        assert len(slide.shapes) >= initial + 4

    def test_render_fallback_on_empty_vertices(self) -> None:
        """vertices が空の場合はフォールバック（テキストボックス）が追加される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = self._make_graph_data(vertices={})
        self.renderer.render(slide, graph_data, "fallback text", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        assert len(slide.shapes) >= initial + 1

    @pytest.mark.parametrize("node_type", [
        "square", "round", "stadium", "subroutine",
        "cylinder", "circle", "odd", "diamond",
        "hexagon", "lean_right", "lean_left",
        "trapezoid", "inv_trapezoid", "doublecircle",
    ])
    def test_render_all_node_shapes(self, node_type: str) -> None:
        """全14種類のノード形状を持つフローチャートがエラーなく描画できる。"""
        slide = _make_slide()
        graph_data = self._make_graph_data({"X": {"text": "テスト", "type": node_type}})
        # 例外なく完了することを確認する
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))
        assert len(slide.shapes) >= 1

    @pytest.mark.parametrize("stroke", ["normal", "dotted", "thick"])
    def test_render_edge_stroke_types(self, stroke: str) -> None:
        """normal/dotted/thick の各線種でエラーなく描画できる。"""
        slide = _make_slide()
        graph_data = self._make_graph_data(
            vertices={"A": {"text": "A", "type": "square"}, "B": {"text": "B", "type": "square"}},
            edges=[{"start": "A", "end": "B", "stroke": stroke, "type": "arrow_point", "text": ""}],
        )
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))

    @pytest.mark.parametrize("edge_type", [
        "arrow_open", "arrow_point", "double_arrow_point",
        "arrow_circle", "double_arrow_circle",
        "arrow_cross", "double_arrow_cross",
    ])
    def test_render_all_edge_arrow_types(self, edge_type: str) -> None:
        """全7種類のエッジ矢印種別でエラーなく描画できる。"""
        slide = _make_slide()
        graph_data = self._make_graph_data(
            vertices={"A": {"text": "A", "type": "square"}, "B": {"text": "B", "type": "square"}},
            edges=[{"start": "A", "end": "B", "stroke": "normal", "type": edge_type, "text": ""}],
        )
        self.renderer.render(slide, graph_data, "", Emu(0), Emu(0), Emu(9000000), Emu(5000000))


# ------------------------------------------------------------------
# MermaidRenderer からの統合テスト（flowchart経由）
# ------------------------------------------------------------------

class TestMermaidRendererFlowchart:
    """MermaidRenderer 経由でのフローチャート描画を確認する。"""

    def setup_method(self) -> None:
        """テスト前にMermaidRendererインスタンスを生成する。"""
        self.renderer = MermaidRenderer()

    def test_flowchart_with_diamond_node(self) -> None:
        """ひし形ノードを含むフローチャートが例外なく描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element("flowchart LR\n    A[開始] --> B{判定}\n    B --> C[終了]")
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )
        assert len(slide.shapes) > 0

    def test_flowchart_with_edge_label(self) -> None:
        """ラベル付きエッジを含むフローチャートが例外なく描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "flowchart LR\n    A --> |条件| B"
        )
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )
        assert len(slide.shapes) > 0

    def test_flowchart_dotted_edge(self) -> None:
        """点線エッジを含むフローチャートが例外なく描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "flowchart LR\n    A -.-> B"
        )
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )
        assert len(slide.shapes) > 0
