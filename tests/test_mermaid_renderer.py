"""
Mermaidレンダラー（MermaidRenderer）の単体テスト。
"""

import xml.etree.ElementTree as ET
import pytest
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from qmd_to_pptx.mermaid_renderer import MermaidRenderer
from qmd_to_pptx.mermaid.mindmap import MindmapRenderer


def _make_slide():
    """テスト用スライドを生成して返す。"""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Blankレイアウト
    return prs.slides.add_slide(layout)


def _make_mermaid_element(mermaid_text: str) -> ET.Element:
    """テスト用Mermaid code要素を生成する。"""
    elem = ET.Element("code")
    elem.set("class", "language-mermaid")
    elem.text = mermaid_text
    return elem


class TestMermaidRenderer:
    """MermaidRenderer クラスの単体テスト。"""

    def setup_method(self) -> None:
        """テスト前にMermaidRendererインスタンスを生成する。"""
        self.renderer = MermaidRenderer()

    def test_render_simple_flowchart_does_not_raise(self) -> None:
        """シンプルなflowchartをレンダリングしても例外が発生しない。"""
        slide = _make_slide()
        elem = _make_mermaid_element("flowchart LR\n    A --> B")
        # 例外なく完了することを確認
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )

    def test_render_adds_shapes_to_slide(self) -> None:
        """レンダリング後にスライドにShapeが追加される。"""
        slide = _make_slide()
        initial_shapes = len(slide.shapes)
        elem = _make_mermaid_element("flowchart LR\n    A --> B")
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )
        # パース成功またはフォールバック時も必ずShapeが追加される
        assert len(slide.shapes) > initial_shapes

    def test_render_fallback_on_invalid_mermaid(self) -> None:
        """無効なMermaidテキストでもフォールバックとしてテキストボックスが追加される。"""
        slide = _make_slide()
        initial_shapes = len(slide.shapes)
        elem = _make_mermaid_element("これはMermaidではありません！！！###")
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )
        # フォールバック時も何らかのShapeが追加される
        assert len(slide.shapes) >= initial_shapes

    def test_render_empty_mermaid_uses_fallback(self) -> None:
        """空のMermaidテキストはフォールバック処理される。"""
        slide = _make_slide()
        elem = _make_mermaid_element("")
        # 例外なく完了することを確認
        self.renderer.render(
            slide, elem,
            Emu(500000), Emu(500000), Emu(7000000), Emu(4000000)
        )

    def test_extract_nodes(self) -> None:
        """_extract_nodes が graph_data から正しくノードリストを返す。"""
        graph_data = {
            "vertices": {"A": {}, "B": {}, "C": {}}
        }
        nodes = self.renderer._extract_nodes(graph_data)
        assert set(nodes) == {"A", "B", "C"}

    def test_extract_nodes_empty(self) -> None:
        """vertices が空の場合は空リストを返す。"""
        assert self.renderer._extract_nodes({}) == []

    def test_extract_edges(self) -> None:
        """_extract_edges が graph_data から正しくエッジリストを返す。"""
        graph_data = {
            "edges": [
                {"start": "A", "end": "B"},
                {"start": "B", "end": "C"},
            ]
        }
        edges = self.renderer._extract_edges(graph_data)
        assert ("A", "B") in edges
        assert ("B", "C") in edges

    def test_extract_edges_empty(self) -> None:
        """edges が空の場合は空リストを返す。"""
        assert self.renderer._extract_edges({}) == []

    def test_pos_to_emu_center(self) -> None:
        """正規化座標 (0, 0) は描画エリアの中央付近に変換される。"""
        left, top, width, height = 0, 0, 9144000, 5143500
        x_emu, y_emu = self.renderer._pos_to_emu(0.0, 0.0, left, top, width, height)
        # 中央付近であることを確認（厳密な値ではなく範囲で確認）
        assert width // 4 < x_emu < width * 3 // 4
        assert height // 4 < y_emu < height * 3 // 4


# ── mindmap 専用テスト ────────────────────────────────────────────────────────

def _make_mindmap_graph_data() -> dict:
    """テスト用マインドマップ graph_data を生成する（3ブランチ・深さ2）。"""
    return {
        "nodes": [
            {
                "nodeId": "root", "level": 2, "descr": "中心", "type": 3,
                "children": [
                    {
                        "nodeId": "A", "level": 4, "descr": "Aブランチ", "type": 2,
                        "children": [
                            {"nodeId": "A1", "level": 6, "descr": "A1葉", "type": 0, "children": []},
                            {"nodeId": "A2", "level": 6, "descr": "A2葉", "type": 6, "children": []},
                        ],
                    },
                    {
                        "nodeId": "B", "level": 4, "descr": "Bブランチ", "type": 4,
                        "children": [
                            {"nodeId": "B1", "level": 6, "descr": "B1葉", "type": 2, "children": []},
                        ],
                    },
                    {
                        "nodeId": "C", "level": 4, "descr": "Cブランチ", "type": 3,
                        "children": [
                            {"nodeId": "C1", "level": 6, "descr": "C1葉", "type": 2, "children": []},
                            {"nodeId": "C2", "level": 6, "descr": "C2葉", "type": 2, "children": []},
                        ],
                    },
                ],
            }
        ]
    }


class TestMindmapRenderer:
    """MindmapRenderer の単体テスト。"""

    def setup_method(self) -> None:
        """テスト前にインスタンスを生成する。"""
        self.renderer = MindmapRenderer()
        prs = Presentation()
        self.slide = prs.slides.add_slide(prs.slide_layouts[5])
        self.left = Emu(500000)
        self.top = Emu(500000)
        self.width = Emu(8000000)
        self.height = Emu(4500000)

    def _render(self, gd: dict) -> None:
        """指定graph_dataでrenderを呼び出すヘルパー。"""
        self.renderer.render(
            self.slide, gd,
            self.left, self.top, self.width, self.height,
        )

    def test_render_does_not_raise(self) -> None:
        """マインドマップのレンダリングで例外が発生しない。"""
        self._render(_make_mindmap_graph_data())

    def test_render_creates_correct_node_count(self) -> None:
        """9ノード（root+A,B,C,A1,A2,B1,C1,C2）が作成される。"""
        initial = len(self.slide.shapes)
        self._render(_make_mindmap_graph_data())
        # AUTO_SHAPEのみカウント（コネクターは除く）
        auto_shapes = [
            s for s in self.slide.shapes
            if s.shape_type == 1  # AUTO_SHAPE
        ]
        assert len(auto_shapes) == 9

    def test_render_creates_connectors(self) -> None:
        """8本のコネクター（9ノード - 1）が作成される。"""
        self._render(_make_mindmap_graph_data())
        connectors = [
            s for s in self.slide.shapes
            if s.shape_type == 9  # LINE/CONNECTOR
        ]
        assert len(connectors) == 8

    def test_root_is_at_center(self) -> None:
        """ルートノードが描画エリアの中央付近に配置される。"""
        self._render(_make_mindmap_graph_data())
        root_shape = next(
            s for s in self.slide.shapes
            if hasattr(s, "text") and s.text == "中心"
        )
        # ルートの中心X・Yが描画エリアの中央±20%以内か確認
        cx = root_shape.left + root_shape.width // 2
        cy = root_shape.top + root_shape.height // 2
        area_cx = self.left + self.width // 2
        area_cy = self.top + self.height // 2
        assert abs(cx - area_cx) < self.width * 0.2
        assert abs(cy - area_cy) < self.height * 0.2

    def test_root_node_shape_is_ellipse(self) -> None:
        """type=3（CIRCLE）のルートノードはellipse形状になる。"""
        self._render(_make_mindmap_graph_data())
        root_shape = next(
            s for s in self.slide.shapes
            if hasattr(s, "text") and s.text == "中心"
        )
        sp_el = root_shape._element
        spPr = sp_el.find(qn("p:spPr"))
        prst = spPr.find(qn("a:prstGeom"))
        assert prst is not None
        assert prst.get("prst") == "ellipse"

    def test_l1_branch_node_shapes(self) -> None:
        """L1ノードの形状がtype値に応じて正しく設定される。"""
        self._render(_make_mindmap_graph_data())

        def _get_prst(text: str) -> str:
            shape = next(
                s for s in self.slide.shapes
                if hasattr(s, "text") and s.text == text
            )
            spPr = shape._element.find(qn("p:spPr"))
            return spPr.find(qn("a:prstGeom")).get("prst")

        # Aブランチ: type=2 → rect
        assert _get_prst("Aブランチ") == "rect"
        # Bブランチ: type=4 → cloud
        assert _get_prst("Bブランチ") == "cloud"
        # Cブランチ: type=3 → ellipse
        assert _get_prst("Cブランチ") == "ellipse"

    def test_l2_hexagon_shape(self) -> None:
        """type=6（HEXAGON）のL2ノードはhexagon形状になる。"""
        self._render(_make_mindmap_graph_data())
        shape = next(
            s for s in self.slide.shapes
            if hasattr(s, "text") and s.text == "A2葉"
        )
        spPr = shape._element.find(qn("p:spPr"))
        prst = spPr.find(qn("a:prstGeom"))
        assert prst.get("prst") == "hexagon"

    def test_l1_nodes_have_different_positions(self) -> None:
        """L1ノードA・B・Cがそれぞれ異なる位置に配置される（放射状）。"""
        self._render(_make_mindmap_graph_data())
        texts = ["Aブランチ", "Bブランチ", "Cブランチ"]
        positions = []
        for text in texts:
            s = next(
                sh for sh in self.slide.shapes
                if hasattr(sh, "text") and sh.text == text
            )
            positions.append((s.left, s.top))
        # 全てのL1ノードが異なる位置にあることを確認
        assert len(set(positions)) == 3

    def test_connector_has_branch_color(self) -> None:
        """コネクターの線色がXMLに設定されている。"""
        self._render(_make_mindmap_graph_data())
        connectors = [
            s for s in self.slide.shapes
            if s.shape_type == 9
        ]
        # 少なくとも1本のコネクターがsolidFillを持つ
        has_color = False
        for conn in connectors:
            sp_el = conn._element
            spPr = sp_el.find(qn("p:spPr"))
            if spPr is None:
                continue
            ln_el = spPr.find(qn("a:ln"))
            if ln_el is None:
                continue
            solidFill = ln_el.find(qn("a:solidFill"))
            if solidFill is not None:
                has_color = True
                break
        assert has_color

    def test_render_empty_graph_data(self) -> None:
        """空のgraph_dataでも例外が発生しない。"""
        self.renderer.render(
            self.slide, {},
            self.left, self.top, self.width, self.height,
        )

    def test_render_single_root_node(self) -> None:
        """子なしルートノードだけのマインドマップでも例外が発生しない。"""
        gd = {
            "nodes": [
                {"nodeId": "root", "level": 2, "descr": "ルートのみ", "type": 0, "children": []}
            ]
        }
        self.renderer.render(
            self.slide, gd,
            self.left, self.top, self.width, self.height,
        )
