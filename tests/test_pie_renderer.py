"""
PieParser / PieChartRenderer のユニットテスト。

pie_parser.py のカスタムパーサーが各種 Mermaid pie 構文を正しく解析すること、
pie_renderer.py がスライドにエラーなく円グラフを追加することを検証する。
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

from qmd_to_pptx.mermaid.pie_parser import PieChart, PieSection, parse_pie
from qmd_to_pptx.mermaid.pie_renderer import PieChartRenderer


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------

def _make_slide():
    """テスト用スライドオブジェクトを生成する。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _chart_on_slide(slide):
    """スライド上の最初のチャートフレームを返す（なければ None）。"""
    for shape in slide.shapes:
        if shape.has_chart:
            return shape.chart
    return None


# ---------------------------------------------------------------------------
# PieParser テスト
# ---------------------------------------------------------------------------

class TestPieParser:
    """parse_pie() の正常系・異常系テスト。"""

    def test_基本構文のパース(self):
        """最低限の pie 構文（title なし showData なし）を正しく解析できること。"""
        text = '''pie
    "犬" : 386
    "猫" : 85
    "ネズミ" : 15'''
        result = parse_pie(text)
        assert result.title == ""
        assert result.show_data is False
        assert result.text_position == 0.75
        assert len(result.sections) == 3
        assert result.sections[0].label == "犬"
        assert result.sections[0].value == 386.0
        assert result.sections[1].label == "猫"
        assert result.sections[1].value == 85.0
        assert result.sections[2].label == "ネズミ"
        assert result.sections[2].value == 15.0

    def test_title_インライン(self):
        """ヘッダー行に title がインラインで記述された場合を正しく解析できること。"""
        text = '''pie title ペットの内訳
    "犬" : 386
    "猫" : 85'''
        result = parse_pie(text)
        assert result.title == "ペットの内訳"

    def test_title_単独行(self):
        """title が単独行で記述された場合を正しく解析できること。"""
        text = '''pie
    title テスト円グラフ
    "A" : 50
    "B" : 50'''
        result = parse_pie(text)
        assert result.title == "テスト円グラフ"

    def test_showData_キーワード検出(self):
        """showData キーワードが検出されたとき show_data が True になること。"""
        text = '''pie showData title 売上
    "製品A" : 100
    "製品B" : 200'''
        result = parse_pie(text)
        assert result.show_data is True
        assert result.title == "売上"

    def test_showData_のみ_title_なし(self):
        """showData のみでタイトルなしの場合を正しく解析できること。"""
        text = '''pie showData
    "X" : 40
    "Y" : 60'''
        result = parse_pie(text)
        assert result.show_data is True
        assert result.title == ""

    def test_小数点数値のパース(self):
        """小数点を含む数値が float として正しくパースされること。"""
        text = '''pie title テスト
    "A" : 42.96
    "B" : 50.05
    "C" : 10.01
    "D" : 5'''
        result = parse_pie(text)
        assert result.sections[0].value == pytest.approx(42.96)
        assert result.sections[1].value == pytest.approx(50.05)
        assert result.sections[2].value == pytest.approx(10.01)
        assert result.sections[3].value == pytest.approx(5.0)

    def test_コメント行スキップ(self):
        """%% コメント行がスキップされること。"""
        text = '''pie title テスト
    %% これはコメント
    "A" : 50
    %% もう一つのコメント
    "B" : 50'''
        result = parse_pie(text)
        assert len(result.sections) == 2

    def test_ゼロ以下の値はスキップ(self):
        """値が 0 以下のセクションはスキップされること。"""
        text = '''pie title テスト
    "A" : 50
    "B" : 0
    "C" : 30'''
        result = parse_pie(text)
        # B は値が 0 なのでスキップされる
        assert len(result.sections) == 2
        assert result.sections[0].label == "A"
        assert result.sections[1].label == "C"

    def test_YAML_front_matter_付きテキスト(self):
        """YAML front-matter を含むテキストから textPosition を取得できること。"""
        text = '''---
config:
  pie:
    textPosition: 0.5
---
pie title テスト
    "X" : 70
    "Y" : 30'''
        result = parse_pie(text)
        assert result.text_position == pytest.approx(0.5)
        assert result.title == "テスト"
        assert len(result.sections) == 2

    def test_YAML_なしはデフォルトtextPosition(self):
        """YAML front-matter がない場合は textPosition がデフォルト 0.75 であること。"""
        text = '''pie
    "A" : 100'''
        result = parse_pie(text)
        assert result.text_position == pytest.approx(0.75)

    def test_ヘッダーなしはValueError(self):
        """pie キーワードで始まらないテキストは ValueError を送出すること。"""
        text = '''flowchart TD
    A --> B'''
        with pytest.raises(ValueError):
            parse_pie(text)

    def test_大文字小文字混在のpieキーワード(self):
        """PIE / Pie など大文字小文字混在でも解析できること。"""
        text = '''PIE title 大文字テスト
    "A" : 60
    "B" : 40'''
        result = parse_pie(text)
        assert result.title == "大文字テスト"
        assert len(result.sections) == 2

    def test_セクション順序保持(self):
        """セクションの順序が Mermaid テキストの記述順と一致すること。"""
        text = '''pie
    "Z" : 10
    "A" : 60
    "M" : 30'''
        result = parse_pie(text)
        assert [s.label for s in result.sections] == ["Z", "A", "M"]


# ---------------------------------------------------------------------------
# PieChartRenderer テスト
# ---------------------------------------------------------------------------

class TestPieChartRenderer:
    """PieChartRenderer.render() の正常系テスト。"""

    def test_基本描画でチャートが追加される(self):
        """セクションありの PieChart をレンダリングするとスライドにチャートが追加される。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(
            title="テスト",
            sections=[PieSection("A", 60.0), PieSection("B", 40.0)],
        )
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart is not None
        assert chart.chart_type == XL_CHART_TYPE.PIE

    def test_タイトルあり(self):
        """title を設定した場合、チャートに has_title=True が設定される。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(
            title="売上構成",
            sections=[PieSection("A", 70.0), PieSection("B", 30.0)],
        )
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.has_title is True
        assert chart.chart_title.text_frame.text == "売上構成"

    def test_タイトルなし(self):
        """title が空文字列の場合、has_title=False が設定される。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(
            title="",
            sections=[PieSection("A", 50.0), PieSection("B", 50.0)],
        )
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.has_title is False

    def test_凡例が表示される(self):
        """常に凡例（has_legend=True）が設定される。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(sections=[PieSection("X", 100.0)])
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.has_legend is True

    def test_データラベルが有効(self):
        """has_data_labels が True になること。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(sections=[PieSection("A", 60.0), PieSection("B", 40.0)])
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.plots[0].has_data_labels is True

    def test_セクションなしは何も描画しない(self):
        """sections が空のとき スライドに何も追加されない。"""
        slide = _make_slide()
        shape_count_before = len(slide.shapes)
        renderer = PieChartRenderer()
        pie = PieChart(title="空グラフ", sections=[])
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        assert len(slide.shapes) == shape_count_before

    def test_showData_True_でshow_value設定(self):
        """show_data=True のとき data_labels.show_value が True になること。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(
            show_data=True,
            sections=[PieSection("A", 80.0), PieSection("B", 20.0)],
        )
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.plots[0].data_labels.show_value is True

    def test_showData_False_でshow_percentageのみ(self):
        """show_data=False のとき show_percentage=True かつ show_value が False。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(
            show_data=False,
            sections=[PieSection("A", 80.0), PieSection("B", 20.0)],
        )
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        dl = chart.plots[0].data_labels
        assert dl.show_percentage is True
        assert dl.show_value is False

    def test_vary_by_categories_有効(self):
        """スライスごとに色が変わる (vary_by_categories=True) こと。"""
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = PieChart(sections=[PieSection("A", 50.0), PieSection("B", 50.0)])
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart.plots[0].vary_by_categories is True

    def test_parse_pieと連携した統合描画(self):
        """parse_pie() の結果を PieChartRenderer に渡してエラーなく描画できること。"""
        mermaid_text = '''pie showData title 製品別売上
    "製品A" : 120.5
    "製品B" : 87.3
    "製品C" : 54.0'''
        slide = _make_slide()
        renderer = PieChartRenderer()
        pie = parse_pie(mermaid_text)
        renderer.render(slide, pie, 914400, 914400, 6858000, 4572000)

        chart = _chart_on_slide(slide)
        assert chart is not None
        assert chart.has_title is True
        assert chart.chart_title.text_frame.text == "製品別売上"


# ---------------------------------------------------------------------------
# textPosition → dLblPos 変換テスト
# ---------------------------------------------------------------------------

class TestTextPositionMapping:
    """_text_position_to_dLblPos() の変換ロジックのテスト。"""

    def setup_method(self):
        self.renderer = PieChartRenderer()

    def test_0_0は_ctr(self):
        assert self.renderer._text_position_to_dLblPos(0.0) == "ctr"

    def test_0_39は_ctr(self):
        assert self.renderer._text_position_to_dLblPos(0.39) == "ctr"

    def test_0_4は_inEnd(self):
        assert self.renderer._text_position_to_dLblPos(0.4) == "inEnd"

    def test_0_69は_inEnd(self):
        assert self.renderer._text_position_to_dLblPos(0.69) == "inEnd"

    def test_0_75はbestFit(self):
        assert self.renderer._text_position_to_dLblPos(0.75) == "bestFit"

    def test_0_99はbestFit(self):
        assert self.renderer._text_position_to_dLblPos(0.99) == "bestFit"

    def test_1_0は_outEnd(self):
        assert self.renderer._text_position_to_dLblPos(1.0) == "outEnd"

    def test_1_0超は_outEnd(self):
        assert self.renderer._text_position_to_dLblPos(1.5) == "outEnd"
