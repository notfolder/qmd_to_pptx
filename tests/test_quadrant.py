"""
QuadrantParser / QuadrantRenderer のユニットテスト。

quadrant_parser.py のカスタムパーサーが各種 Mermaid quadrantChart 構文を
正しく解析すること、quadrant_renderer.py がスライドにエラーなく
象限チャートを追加することを検証する。
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.quadrant_parser import (
    PointStyle,
    QuadrantChart,
    QuadrantPoint,
    parse_quadrant,
    _parse_color,
    _parse_styles,
    _clamp_coord,
)
from qmd_to_pptx.mermaid.quadrant_renderer import QuadrantRenderer, _resolve_style


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------


def _make_slide():
    """テスト用スライドオブジェクトを生成する。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _all_text(slide) -> list[str]:
    """スライド上の全シェープからテキストを収集する（空文字は除外）。"""
    texts = []
    for sp in slide.shapes:
        try:
            t = sp.text_frame.text.strip()
            if t:
                texts.append(t)
        except AttributeError:
            pass
    return texts


_LEFT = 457200
_TOP = 457200
_WIDTH = 8229600
_HEIGHT = 4800000


# ---------------------------------------------------------------------------
# フルサンプルテキスト（Mermaid 公式サンプル準拠）
# ---------------------------------------------------------------------------

_FULL_SAMPLE = """\
quadrantChart
    title Reach and engagement of campaigns
    x-axis Low Reach --> High Reach
    y-axis Low Engagement --> High Engagement
    quadrant-1 We should expand
    quadrant-2 Need to promote
    quadrant-3 Re-evaluate
    quadrant-4 May be improved
    Campaign A: [0.3, 0.6]
    Campaign B: [0.45, 0.23]
    Campaign C: [0.57, 0.69]
    Campaign D: [0.78, 0.34]
    Campaign E: [0.40, 0.34]
    Campaign F: [0.35, 0.78]
"""

_JAPANESE_SAMPLE = """\
quadrantChart
    title 製品ポートフォリオ分析
    x-axis 低成長 --> 高成長
    y-axis 低シェア --> 高シェア
    quadrant-1 スター
    quadrant-2 問題児
    quadrant-3 負け犬
    quadrant-4 金のなる木
    製品A: [0.7, 0.8]
    製品B: [0.3, 0.7]
    製品C: [0.2, 0.3]
    製品D: [0.8, 0.3]
"""

_STYLE_SAMPLE = """\
quadrantChart
    title Styling demo
    x-axis Low --> High
    y-axis Low --> High
    quadrant-1 Q1
    quadrant-2 Q2
    quadrant-3 Q3
    quadrant-4 Q4
    Point A:::class1: [0.9, 0.0]
    Point B:::class2: [0.8, 0.1] color: #ff3300, radius: 10
    Point C: [0.7, 0.2] radius: 25, color: #00ff33, stroke-color: #10f0f0
    Point D: [0.6, 0.3] stroke-color: #00ff0f, stroke-width: 5px, color: #ff33f0
    classDef class1 color: #109060
    classDef class2 color: #908342, radius: 10, stroke-color: #310085, stroke-width: 10px
"""


# ===========================================================================
# TestParseColor: _parse_color() のテスト
# ===========================================================================

class TestParseColor:
    """_parse_color() ユーティリティ関数のテスト。"""

    def test_正常な色文字列にシャープあり(self):
        """#rrggbb 形式の色文字列を正しく正規化すること。"""
        assert _parse_color("#ff3300") == "#ff3300"

    def test_正常な色文字列にシャープなし(self):
        """rrggbb 形式の色文字列に # を付けて返すこと。"""
        assert _parse_color("109060") == "#109060"

    def test_大文字でも小文字に正規化されること(self):
        """大文字の16進数を小文字に変換すること。"""
        assert _parse_color("#FF3300") == "#ff3300"

    def test_無効な色文字列はNoneを返すこと(self):
        """無効な形式は None を返すこと。"""
        assert _parse_color("invalid") is None

    def test_空文字はNoneを返すこと(self):
        """空文字列は None を返すこと。"""
        assert _parse_color("") is None


# ===========================================================================
# TestParseStyles: _parse_styles() のテスト
# ===========================================================================

class TestParseStyles:
    """_parse_styles() ユーティリティ関数のテスト。"""

    def test_colorの解析(self):
        """color属性が正しく解析されること。"""
        style = _parse_styles("color: #ff3300")
        assert style.color == "#ff3300"

    def test_radiusの解析(self):
        """radius 属性が整数として解析されること。"""
        style = _parse_styles("radius: 10")
        assert style.radius == 10

    def test_stroke_widthのpx付き解析(self):
        """stroke-width: 5px の px を除去して整数変換すること。"""
        style = _parse_styles("stroke-width: 5px")
        assert style.stroke_width == 5

    def test_stroke_widthのpxなし解析(self):
        """stroke-width: 3 の整数変換が正しいこと。"""
        style = _parse_styles("stroke-width: 3")
        assert style.stroke_width == 3

    def test_stroke_colorの解析(self):
        """stroke-color 属性が正しく解析されること。"""
        style = _parse_styles("stroke-color: #10f0f0")
        assert style.stroke_color == "#10f0f0"

    def test_複数属性のカンマ区切り解析(self):
        """カンマ区切りの複数属性を一度に解析できること。"""
        style = _parse_styles("color: #ff3300, radius: 10, stroke-width: 5px")
        assert style.color == "#ff3300"
        assert style.radius == 10
        assert style.stroke_width == 5

    def test_未知の属性は無視されること(self):
        """認識できない属性名は無視して他の属性を正しく解析すること。"""
        style = _parse_styles("unknown: xyz, radius: 7")
        assert style.radius == 7
        assert style.color is None

    def test_空文字列はデフォルトスタイルを返すこと(self):
        """空のスタイル文字列はすべて None のスタイルを返すこと。"""
        style = _parse_styles("")
        assert style.color is None
        assert style.radius is None
        assert style.stroke_width is None
        assert style.stroke_color is None


# ===========================================================================
# TestClampCoord: _clamp_coord() のテスト
# ===========================================================================

class TestClampCoord:
    """_clamp_coord() ユーティリティ関数のテスト。"""

    def test_0から1の範囲内はそのまま返ること(self):
        """範囲内の値はそのまま返ること。"""
        assert _clamp_coord(0.5) == 0.5

    def test_0の境界値はそのまま返ること(self):
        """最小値 0.0 はそのまま返ること。"""
        assert _clamp_coord(0.0) == 0.0

    def test_1の境界値はそのまま返ること(self):
        """最大値 1.0 はそのまま返ること。"""
        assert _clamp_coord(1.0) == 1.0

    def test_負の値は0にクランプされること(self):
        """0 未満の値は 0.0 にクランプされること。"""
        assert _clamp_coord(-0.5) == 0.0

    def test_1より大きい値は1にクランプされること(self):
        """1 より大きい値は 1.0 にクランプされること。"""
        assert _clamp_coord(1.5) == 1.0


# ===========================================================================
# TestQuadrantParser: parse_quadrant() のテスト
# ===========================================================================

class TestQuadrantParser:
    """parse_quadrant() の正常系・異常系テスト。"""

    def test_タイトルのパース(self):
        """title 行が正しくパースされること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        assert result.title == "Reach and engagement of campaigns"

    def test_X軸ラベルの両端パース(self):
        """x-axis の左ラベルと右ラベルが両方パースされること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        assert result.x_label_left == "Low Reach"
        assert result.x_label_right == "High Reach"

    def test_Y軸ラベルの両端パース(self):
        """y-axis の下ラベルと上ラベルが両方パースされること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        assert result.y_label_bottom == "Low Engagement"
        assert result.y_label_top == "High Engagement"

    def test_X軸ラベルの片方のみのパース(self):
        """x-axis に右ラベルがない場合、右ラベルが空文字列であること。"""
        text = "quadrantChart\n    x-axis Left only\n"
        result = parse_quadrant(text)
        assert result.x_label_left == "Left only"
        assert result.x_label_right == ""

    def test_象限ラベル1から4のパース(self):
        """quadrant-1〜4 のラベルが辞書に格納されること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        assert result.quadrant_labels[1] == "We should expand"
        assert result.quadrant_labels[2] == "Need to promote"
        assert result.quadrant_labels[3] == "Re-evaluate"
        assert result.quadrant_labels[4] == "May be improved"

    def test_象限ラベル一部省略時(self):
        """一部の象限ラベルが省略された場合、省略した番号はキーに存在しないこと。"""
        text = "quadrantChart\n    quadrant-1 Only Q1\n"
        result = parse_quadrant(text)
        assert 1 in result.quadrant_labels
        assert 2 not in result.quadrant_labels

    def test_ポイントの件数(self):
        """ポイントが正しい件数でパースされること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        assert len(result.points) == 6

    def test_ポイントの座標パース(self):
        """ポイントの x, y 座標が正しくパースされること。"""
        result = parse_quadrant(_FULL_SAMPLE)
        campaign_a = next(p for p in result.points if p.name == "Campaign A")
        assert campaign_a.x == pytest.approx(0.3)
        assert campaign_a.y == pytest.approx(0.6)

    def test_ポイント座標クランプ(self):
        """座標が 0.0〜1.0 の範囲にクランプされること。"""
        text = "quadrantChart\n    OverPoint: [1.5, -0.3]\n"
        result = parse_quadrant(text)
        assert len(result.points) == 1
        assert result.points[0].x == pytest.approx(1.0)
        assert result.points[0].y == pytest.approx(0.0)

    def test_日本語タイトルとラベルのパース(self):
        """日本語のタイトル・ラベル・ポイント名が正しくパースされること。"""
        result = parse_quadrant(_JAPANESE_SAMPLE)
        assert result.title == "製品ポートフォリオ分析"
        assert result.x_label_left == "低成長"
        assert result.x_label_right == "高成長"
        assert result.quadrant_labels[1] == "スター"
        assert any(p.name == "製品A" for p in result.points)

    def test_classNameの参照付きポイントのパース(self):
        """:::ClassName 付きポイントがクラス名を保持していること。"""
        result = parse_quadrant(_STYLE_SAMPLE)
        point_a = next(p for p in result.points if p.name == "Point A")
        assert point_a.class_name == "class1"

    def test_インラインスタイルの解析(self):
        """インラインスタイル（color, radius 等）が正しく解析されること。"""
        result = parse_quadrant(_STYLE_SAMPLE)
        point_c = next(p for p in result.points if p.name == "Point C")
        assert point_c.inline_style.color == "#00ff33"
        assert point_c.inline_style.radius == 25
        assert point_c.inline_style.stroke_color == "#10f0f0"

    def test_classDef定義のパース(self):
        """classDef 行がクラス名をキーとして格納されること。"""
        result = parse_quadrant(_STYLE_SAMPLE)
        assert "class1" in result.class_defs
        assert result.class_defs["class1"].color == "#109060"

    def test_classDef複数属性のパース(self):
        """classDef で複数属性が定義された場合にすべてパースされること。"""
        result = parse_quadrant(_STYLE_SAMPLE)
        cls2 = result.class_defs["class2"]
        assert cls2.color == "#908342"
        assert cls2.radius == 10
        assert cls2.stroke_color == "#310085"
        assert cls2.stroke_width == 10

    def test_コメント行はスキップされること(self):
        """"%%" で始まる行がスキップされること。"""
        text = "quadrantChart\n%% これはコメントです\n    title コメントなし\n"
        result = parse_quadrant(text)
        assert result.title == "コメントなし"
        assert len(result.points) == 0

    def test_空入力で空のチャートを返すこと(self):
        """空文字列の場合にデフォルト値の QuadrantChart が返ること。"""
        result = parse_quadrant("")
        assert result.title == ""
        assert len(result.points) == 0
        assert len(result.quadrant_labels) == 0

    def test_ヘッダーのみの最小構成(self):
        """ヘッダー行のみの最小構成でエラーなく解析できること。"""
        result = parse_quadrant("quadrantChart")
        assert isinstance(result, QuadrantChart)
        assert len(result.points) == 0

    def test_クラス名と戻り値による同時インラインスタイルのパース(self):
        """:::class + インラインスタイルが両方保持されること。"""
        result = parse_quadrant(_STYLE_SAMPLE)
        point_b = next(p for p in result.points if p.name == "Point B")
        # クラス名が保持されること
        assert point_b.class_name == "class2"
        # インラインスタイルが保持されること（classDef より優先されること）
        assert point_b.inline_style.color == "#ff3300"
        assert point_b.inline_style.radius == 10


# ===========================================================================
# TestResolveStyle: _resolve_style() のテスト
# ===========================================================================

class TestResolveStyle:
    """_resolve_style() スタイル解決関数のテスト。"""

    def test_インラインスタイルがclassDefより優先されること(self):
        """インラインの color が classDef の color より優先されること。"""
        point = QuadrantPoint(
            name="P",
            x=0.5,
            y=0.5,
            class_name="cls",
            inline_style=PointStyle(color="#ff0000"),
        )
        class_defs = {"cls": PointStyle(color="#00ff00")}
        fill_rgb, radius_px, stroke_w_px, stroke_rgb = _resolve_style(point, class_defs)
        assert fill_rgb == (255, 0, 0)

    def test_classDefスタイルがデフォルトより優先されること(self):
        """classDef の radius がデフォルトより優先されること。"""
        point = QuadrantPoint(
            name="P",
            x=0.5,
            y=0.5,
            class_name="cls",
            inline_style=PointStyle(),  # インラインは空
        )
        class_defs = {"cls": PointStyle(radius=15)}
        fill_rgb, radius_px, stroke_w_px, stroke_rgb = _resolve_style(point, class_defs)
        assert radius_px == 15

    def test_classNameが未定義の場合はデフォルト値になること(self):
        """classDef に対応するクラスがない場合はデフォルト値を使用すること。"""
        point = QuadrantPoint(
            name="P",
            x=0.5,
            y=0.5,
            class_name="nonexistent",
            inline_style=PointStyle(),
        )
        fill_rgb, radius_px, stroke_w_px, stroke_rgb = _resolve_style(point, {})
        # デフォルトの半径 5 が返ること
        assert radius_px == 5

    def test_スタイルなしの場合はデフォルト値になること(self):
        """スタイルがない場合にデフォルト値が適用されること。"""
        point = QuadrantPoint(name="P", x=0.5, y=0.5)
        fill_rgb, radius_px, stroke_w_px, stroke_rgb = _resolve_style(point, {})
        assert radius_px == 5
        assert stroke_w_px == 1


# ===========================================================================
# TestQuadrantRenderer: QuadrantRenderer のテスト
# ===========================================================================

class TestQuadrantRenderer:
    """QuadrantRenderer.render() のテスト。"""

    def test_ポイントなしでクラッシュしないこと(self):
        """ポイントが 0 件の場合もクラッシュなく描画できること。"""
        slide = _make_slide()
        chart = QuadrantChart(
            title="Empty Chart",
            x_label_left="Low",
            x_label_right="High",
            y_label_bottom="Bottom",
            y_label_top="Top",
            quadrant_labels={1: "Q1", 2: "Q2", 3: "Q3", 4: "Q4"},
        )
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # 少なくとも象限背景 + 外枠 + 軸線 + ラベル等が存在すること
        assert len(slide.shapes) > 0

    def test_タイトルが描画されること(self):
        """タイトルテキストがスライドに表示されること。"""
        slide = _make_slide()
        chart = QuadrantChart(title="テストタイトル")
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "テストタイトル" in texts

    def test_タイトルなしでクラッシュしないこと(self):
        """title が空文字列の場合もクラッシュなく描画できること。"""
        slide = _make_slide()
        chart = QuadrantChart(title="")
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # タイトルボックスは描画されず、エラーもないこと
        texts = _all_text(slide)
        assert "" not in texts  # 空テキストは除外されること

    def test_象限ラベルが描画されること(self):
        """quadrant-1〜4 のラベルがすべてスライドに描画されること。"""
        slide = _make_slide()
        chart = QuadrantChart(
            quadrant_labels={1: "Q1label", 2: "Q2label", 3: "Q3label", 4: "Q4label"}
        )
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "Q1label" in texts
        assert "Q2label" in texts
        assert "Q3label" in texts
        assert "Q4label" in texts

    def test_X軸ラベルが描画されること(self):
        """x-axis の左右ラベルがスライドに描画されること。"""
        slide = _make_slide()
        chart = QuadrantChart(
            x_label_left="左ラベル",
            x_label_right="右ラベル",
        )
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "左ラベル" in texts
        assert "右ラベル" in texts

    def test_Y軸ラベルが描画されること(self):
        """y-axis の下上ラベルがスライドに描画されること。"""
        slide = _make_slide()
        chart = QuadrantChart(
            y_label_bottom="下ラベル",
            y_label_top="上ラベル",
        )
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "下ラベル" in texts
        assert "上ラベル" in texts

    def test_ポイントありでシェープが生成されること(self):
        """ポイントがある場合にシェープ数が増加すること（ポイント + ラベル）。"""
        slide_no_points = _make_slide()
        slide_with_points = _make_slide()

        chart_base = QuadrantChart(
            title="Chart",
            quadrant_labels={1: "Q1", 2: "Q2", 3: "Q3", 4: "Q4"},
        )
        chart_with_points = QuadrantChart(
            title="Chart",
            quadrant_labels={1: "Q1", 2: "Q2", 3: "Q3", 4: "Q4"},
            points=[
                QuadrantPoint(name="A", x=0.3, y=0.7),
                QuadrantPoint(name="B", x=0.7, y=0.3),
            ],
        )
        renderer = QuadrantRenderer()
        renderer.render(slide_no_points, chart_base, _LEFT, _TOP, _WIDTH, _HEIGHT)
        renderer.render(slide_with_points, chart_with_points, _LEFT, _TOP, _WIDTH, _HEIGHT)

        # ポイントあり（2件）では 2個の OVAL + 2個のラベル = 4シェープ分多いこと
        diff = len(slide_with_points.shapes) - len(slide_no_points.shapes)
        assert diff == 4

    def test_ポイントラベルが描画されること(self):
        """ポイント名がスライドに描画されること。"""
        slide = _make_slide()
        chart = QuadrantChart(
            points=[QuadrantPoint(name="Campaign X", x=0.5, y=0.5)]
        )
        renderer = QuadrantRenderer()
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "Campaign X" in texts

    def test_日本語チャートがクラッシュなく描画されること(self):
        """日本語テキストを含む象限チャートが正常に描画されること。"""
        slide = _make_slide()
        result = parse_quadrant(_JAPANESE_SAMPLE)
        renderer = QuadrantRenderer()
        renderer.render(slide, result, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "製品ポートフォリオ分析" in texts
        assert "スター" in texts

    def test_フルサンプルがクラッシュなく描画されること(self):
        """Mermaid 公式フルサンプルが正常に描画されること。"""
        slide = _make_slide()
        result = parse_quadrant(_FULL_SAMPLE)
        renderer = QuadrantRenderer()
        renderer.render(slide, result, _LEFT, _TOP, _WIDTH, _HEIGHT)
        texts = _all_text(slide)
        assert "Reach and engagement of campaigns" in texts
        assert "Campaign A" in texts
        assert "Campaign F" in texts

    def test_スタイルサンプルがクラッシュなく描画されること(self):
        """classDef とインラインスタイルを含むチャートが正常に描画されること。"""
        slide = _make_slide()
        result = parse_quadrant(_STYLE_SAMPLE)
        renderer = QuadrantRenderer()
        renderer.render(slide, result, _LEFT, _TOP, _WIDTH, _HEIGHT)
        # すべてのポイントラベルが描画されること
        texts = _all_text(slide)
        assert "Point A" in texts
        assert "Point B" in texts
        assert "Point C" in texts
        assert "Point D" in texts

    def test_4象限背景矩形が4個描画されること(self):
        """4象限の背景矩形シェープが4個生成されること。"""
        slide = _make_slide()
        chart = QuadrantChart()
        renderer = QuadrantRenderer()
        n_before = len(slide.shapes)
        renderer.render(slide, chart, _LEFT, _TOP, _WIDTH, _HEIGHT)
        n_after = len(slide.shapes)
        # 最低でも 4 個（象限背景）+ 1 個（外枠）+ 2 本（軸線）が追加されること
        assert n_after - n_before >= 7
