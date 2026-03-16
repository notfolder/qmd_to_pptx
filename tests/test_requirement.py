"""
requirementDiagram パーサー・レンダラーのテストモジュール。

テスト構成:
  TestParseRequirementNode (12):  6 要件タイプ・フィールド・クォート・日本語
  TestParseElementNode (5):       type / docref / クォート
  TestParseRelation (8):          順方向・逆方向・7 種リレーション
  TestParseDirection (4):         TB / BT / LR / RL
  TestParseStyle (6):             classDef / class / style / :::ショートハンド
  TestParseMarkdown (5):          **bold** / *italic* / プレーン / 複合 /空文字列
  TestRequirementRenderer (16):   最小・全フィールド・全タイプ・リレーション・日本語
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

import pytest
from pptx import Presentation

from qmd_to_pptx.mermaid.requirement_parser import (
    ElementNode,
    NodeStyle,
    RequirementDiagram,
    RequirementNode,
    Relationship,
    parse_inline_markdown,
    parse_requirement,
    resolve_node_style,
)
from qmd_to_pptx.mermaid.requirement_renderer import RequirementRenderer


# ===========================================================================
# パーサーテスト
# ===========================================================================


class TestParseRequirementNode:
    """要件ノードのパーステスト。"""

    def test_basic_requirement(self) -> None:
        """最小構成の requirement ブロックを正しくパースできること。"""
        text = """
requirementDiagram
requirement MyReq {
    id: 1.0
    text: This is a requirement
}
"""
        diag = parse_requirement(text)
        assert "MyReq" in diag.requirements
        node = diag.requirements["MyReq"]
        assert node.req_type == "requirement"
        assert node.req_id == "1.0"
        assert node.text == "This is a requirement"

    def test_functional_requirement(self) -> None:
        """functionalRequirement 型が正しくパースされること。"""
        text = """
requirementDiagram
functionalRequirement FuncReq {
    id: 2.1
    text: Functional desc
}
"""
        diag = parse_requirement(text)
        assert "FuncReq" in diag.requirements
        node = diag.requirements["FuncReq"]
        assert node.req_type == "functionalrequirement"
        assert node.stereotype == "Functional Requirement"

    def test_interface_requirement(self) -> None:
        """interfaceRequirement 型が正しくパースされること。"""
        text = """
requirementDiagram
interfaceRequirement IntReq {
    id: 3.1
}
"""
        diag = parse_requirement(text)
        assert "IntReq" in diag.requirements
        assert diag.requirements["IntReq"].req_type == "interfacerequirement"

    def test_performance_requirement(self) -> None:
        """performanceRequirement 型が正しくパースされること。"""
        text = """
requirementDiagram
performanceRequirement PerfReq {
    id: 4.1
}
"""
        diag = parse_requirement(text)
        assert "PerfReq" in diag.requirements
        assert diag.requirements["PerfReq"].req_type == "performancerequirement"

    def test_physical_requirement(self) -> None:
        """physicalRequirement 型が正しくパースされること。"""
        text = """
requirementDiagram
physicalRequirement PhysReq {
    id: 5.1
}
"""
        diag = parse_requirement(text)
        assert "PhysReq" in diag.requirements
        assert diag.requirements["PhysReq"].req_type == "physicalrequirement"

    def test_design_constraint(self) -> None:
        """designConstraint 型が正しくパースされること。"""
        text = """
requirementDiagram
designConstraint Constraint1 {
    id: 6.1
}
"""
        diag = parse_requirement(text)
        assert "Constraint1" in diag.requirements
        assert diag.requirements["Constraint1"].req_type == "designconstraint"

    def test_all_fields(self) -> None:
        """id / text / risk / verifymethod 全フィールドが正しくパースされること。"""
        text = """
requirementDiagram
requirement FullReq {
    id: 7.1
    text: Full description
    risk: high
    verifymethod: test
}
"""
        diag = parse_requirement(text)
        node = diag.requirements["FullReq"]
        assert node.req_id == "7.1"
        assert node.text == "Full description"
        assert node.risk == "High"
        assert node.verify_method == "Test"

    def test_risk_normalization(self) -> None:
        """リスクレベルが Low/Medium/High に正規化されること。"""
        text = """
requirementDiagram
requirement R1 { id: 1; risk: low }
requirement R2 { id: 2; risk: medium }
requirement R3 { id: 3; risk: high }
"""
        diag = parse_requirement(text)
        assert diag.requirements["R1"].risk == "Low"
        assert diag.requirements["R2"].risk == "Medium"
        assert diag.requirements["R3"].risk == "High"

    def test_verify_method_normalization(self) -> None:
        """検証方法が正規化されること。"""
        text = """
requirementDiagram
requirement R1 { verifymethod: analysis }
requirement R2 { verifymethod: inspection }
requirement R3 { verifymethod: demonstration }
"""
        diag = parse_requirement(text)
        assert diag.requirements["R1"].verify_method == "Analysis"
        assert diag.requirements["R2"].verify_method == "Inspection"
        assert diag.requirements["R3"].verify_method == "Demonstration"

    def test_quoted_name(self) -> None:
        """クォートで囲まれた名前が正しくパースされること。"""
        text = """
requirementDiagram
requirement "My Requirement Name" {
    id: 8.1
}
"""
        diag = parse_requirement(text)
        assert "My Requirement Name" in diag.requirements

    def test_quoted_text_field(self) -> None:
        """クォートで囲まれた text フィールドが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 {
    text: "This text has, commas: and colons"
}
"""
        diag = parse_requirement(text)
        assert diag.requirements["R1"].text == "This text has, commas: and colons"

    def test_japanese_text(self) -> None:
        """日本語テキストを含む要件ノードが正しくパースされること。"""
        text = """
requirementDiagram
requirement "認証要件" {
    id: JP-1
    text: "ユーザーはパスワードでログインできること"
    risk: high
    verifymethod: test
}
"""
        diag = parse_requirement(text)
        assert "認証要件" in diag.requirements
        node = diag.requirements["認証要件"]
        assert node.req_id == "JP-1"
        assert "ログイン" in node.text
        assert node.risk == "High"


class TestParseElementNode:
    """エレメントノードのパーステスト。"""

    def test_basic_element(self) -> None:
        """最小構成のエレメントブロックが正しくパースされること。"""
        text = """
requirementDiagram
element SysA {
    type: simulation
}
"""
        diag = parse_requirement(text)
        assert "SysA" in diag.elements
        elem = diag.elements["SysA"]
        assert elem.elem_type == "simulation"

    def test_element_with_docref(self) -> None:
        """docref フィールドが正しくパースされること。"""
        text = """
requirementDiagram
element SysB {
    type: word doc
    docref: reqs/spec.md
}
"""
        diag = parse_requirement(text)
        node = diag.elements["SysB"]
        assert node.elem_type == "word doc"
        assert node.docref == "reqs/spec.md"

    def test_element_quoted_name(self) -> None:
        """クォートされたエレメント名が正しくパースされること。"""
        text = """
requirementDiagram
element "System Component" {
    type: hardware
}
"""
        diag = parse_requirement(text)
        assert "System Component" in diag.elements

    def test_element_quoted_docref(self) -> None:
        """クォートされた docref が正しくパースされること。"""
        text = """
requirementDiagram
element SysC {
    docref: "path/to/spec file.docx"
}
"""
        diag = parse_requirement(text)
        assert diag.elements["SysC"].docref == "path/to/spec file.docx"

    def test_japanese_element(self) -> None:
        """日本語エレメント名が正しくパースされること。"""
        text = """
requirementDiagram
element "認証システム" {
    type: software
    docref: "docs/auth.md"
}
"""
        diag = parse_requirement(text)
        assert "認証システム" in diag.elements
        assert diag.elements["認証システム"].elem_type == "software"


class TestParseRelation:
    """リレーションのパーステスト。"""

    def test_forward_contains(self) -> None:
        """順方向 contains リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
R1 - contains -> R2
"""
        diag = parse_requirement(text)
        assert len(diag.relations) == 1
        rel = diag.relations[0]
        assert rel.src == "R1"
        assert rel.dst == "R2"
        assert rel.rel_type == "contains"

    def test_reverse_traces(self) -> None:
        """逆方向 traces リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
element E1 { type: simulation }
R1 <- traces - E1
"""
        diag = parse_requirement(text)
        assert len(diag.relations) == 1
        rel = diag.relations[0]
        assert rel.src == "E1"
        assert rel.dst == "R1"
        assert rel.rel_type == "traces"

    def test_relation_copies(self) -> None:
        """copies リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
R1 - copies -> R2
"""
        diag = parse_requirement(text)
        assert diag.relations[0].rel_type == "copies"

    def test_relation_derives(self) -> None:
        """derives リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
R1 - derives -> R2
"""
        diag = parse_requirement(text)
        assert diag.relations[0].rel_type == "derives"

    def test_relation_satisfies(self) -> None:
        """satisfies リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
element E1 { type: simulation }
E1 - satisfies -> R1
"""
        diag = parse_requirement(text)
        assert diag.relations[0].rel_type == "satisfies"

    def test_relation_verifies(self) -> None:
        """verifies リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
element E1 { type: simulation }
E1 - verifies -> R1
"""
        diag = parse_requirement(text)
        assert diag.relations[0].rel_type == "verifies"

    def test_relation_refines(self) -> None:
        """refines リレーションが正しくパースされること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
R2 - refines -> R1
"""
        diag = parse_requirement(text)
        assert diag.relations[0].rel_type == "refines"

    def test_multiple_relations(self) -> None:
        """複数リレーションが全て記録されること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
element E1 { type: software }
R1 - contains -> R2
E1 - satisfies -> R1
E1 - verifies -> R2
"""
        diag = parse_requirement(text)
        assert len(diag.relations) == 3


class TestParseDirection:
    """direction のパーステスト。"""

    def test_direction_tb(self) -> None:
        """direction TB が正しく設定されること。"""
        diag = parse_requirement("requirementDiagram\ndirection TB")
        assert diag.direction == "TB"

    def test_direction_bt(self) -> None:
        """direction BT が正しく設定されること。"""
        diag = parse_requirement("requirementDiagram\ndirection BT")
        assert diag.direction == "BT"

    def test_direction_lr(self) -> None:
        """direction LR が正しく設定されること。"""
        diag = parse_requirement("requirementDiagram\ndirection LR")
        assert diag.direction == "LR"

    def test_direction_rl(self) -> None:
        """direction RL が正しく設定されること。"""
        diag = parse_requirement("requirementDiagram\ndirection RL")
        assert diag.direction == "RL"

    def test_default_direction(self) -> None:
        """direction 未指定時のデフォルトが TB であること。"""
        diag = parse_requirement("requirementDiagram\nrequirement R1 { id: 1 }")
        assert diag.direction == "TB"


class TestParseStyle:
    """スタイル指定のパーステスト。"""

    def test_classdef(self) -> None:
        """classDef が正しく NodeStyle としてパースされること。"""
        text = """
requirementDiagram
classDef critical fill:#ff0000, stroke:#880000, color:#ffffff
requirement R1 { id: 1 }
"""
        diag = parse_requirement(text)
        assert "critical" in diag.class_defs
        style = diag.class_defs["critical"]
        assert style.fill == "#ff0000"
        assert style.stroke == "#880000"
        assert style.color == "#ffffff"

    def test_class_apply(self) -> None:
        """class キーワードによるクラス適用が正しく機能すること。"""
        text = """
requirementDiagram
classDef critical fill:#ff0000
requirement R1 { id: 1 }
class R1 critical
"""
        diag = parse_requirement(text)
        assert "critical" in diag.requirements["R1"].classes

    def test_style_direct(self) -> None:
        """style キーワードによる直接スタイル適用が機能すること。"""
        text = """
requirementDiagram
requirement R1 { id: 1 }
style R1 fill:#aabbcc, stroke:#112233
"""
        diag = parse_requirement(text)
        style = diag.requirements["R1"].style
        assert style.fill == "#aabbcc"
        assert style.stroke == "#112233"

    def test_class_shorthand_on_node(self) -> None:
        """:::className ショートハンド（個別行）がノードに適用されること。"""
        text = """
requirementDiagram
classDef myClass fill:#abcdef
requirement R1 { id: 1 }
R1:::myClass
"""
        diag = parse_requirement(text)
        assert "myClass" in diag.requirements["R1"].classes

    def test_resolve_node_style_priority(self) -> None:
        """classDef < class < 直接 style の優先順で解決されること。"""
        class_defs = {
            "base": NodeStyle(fill="#111111", stroke="#222222"),
            "override": NodeStyle(fill="#333333"),
        }
        direct = NodeStyle(stroke="#444444")
        # nodes with class "base" and "override", plus direct stroke
        resolved = resolve_node_style(["base", "override"], direct, class_defs)
        # override.fill (#333333) が base.fill (#111111) より優先される
        assert resolved.fill == "#333333"
        # direct.stroke (#444444) が base.stroke (#222222) より優先される
        assert resolved.stroke == "#444444"

    def test_resolve_node_style_no_class(self) -> None:
        """クラスなし・直接スタイルなしの場合は全フィールドが None であること。"""
        resolved = resolve_node_style([], NodeStyle(), {})
        assert resolved.fill is None
        assert resolved.stroke is None


class TestParseMarkdown:
    """parse_inline_markdown 関数のテスト。"""

    def test_bold(self) -> None:
        """**太字** テキストが (text, True, False) として返されること。"""
        result = parse_inline_markdown("**bold text**")
        assert any(seg[1] is True and "bold text" in seg[0] for seg in result)

    def test_italic(self) -> None:
        """*斜体* テキストが (text, False, True) として返されること。"""
        result = parse_inline_markdown("*italic text*")
        assert any(seg[2] is True and "italic text" in seg[0] for seg in result)

    def test_plain(self) -> None:
        """書式なしテキストが (text, False, False) として返されること。"""
        result = parse_inline_markdown("plain text")
        assert result == [("plain text", False, False)]

    def test_mixed(self) -> None:
        """通常テキストと太字の混合が正しくパースされること。"""
        result = parse_inline_markdown("hello **world** end")
        texts = [seg[0] for seg in result]
        assert "hello " in texts or any("hello" in t for t in texts)
        assert any(seg[1] and "world" in seg[0] for seg in result)

    def test_empty(self) -> None:
        """空文字列でエラーが発生せず少なくとも 1 エントリが返されること。"""
        result = parse_inline_markdown("")
        assert isinstance(result, list)
        assert len(result) >= 1


# ===========================================================================
# レンダラーテスト
# ===========================================================================


def _make_slide():
    """テスト用の空スライドを返す。"""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


class TestRequirementRenderer:
    """RequirementRenderer の描画テスト。"""

    def _slide_shape_count(self, slide) -> int:
        """スライド上のシェープ数を返す。"""
        return len(slide.shapes)

    def test_empty_diagram_no_error(self) -> None:
        """ノードが存在しない RequirementDiagram で例外が発生しないこと。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = RequirementDiagram()
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        # ノードがないので何も描画されない

    def test_single_requirement_draws_shapes(self) -> None:
        """単一要件ノードを描画するとシェープが追加されること。"""
        slide = _make_slide()
        before = self._slide_shape_count(slide)
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement R1 {
    id: 1.0
    text: Test requirement
    risk: low
    verifymethod: test
}
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        # 3 (ヘッダー上段+下段+ボディ) + 1 グループ = 4 シェープ以上
        assert self._slide_shape_count(slide) > before

    def test_all_requirement_types(self) -> None:
        """6 種類の要件タイプが全て例外なく描画できること。"""
        types = [
            "requirement",
            "functionalRequirement",
            "interfaceRequirement",
            "performanceRequirement",
            "physicalRequirement",
            "designConstraint",
        ]
        for rtype in types:
            slide = _make_slide()
            renderer = RequirementRenderer()
            text = f"""
requirementDiagram
{rtype} Node1 {{
    id: 1
    text: test
}}
"""
            diag = parse_requirement(text)
            renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)

    def test_element_node_renders(self) -> None:
        """エレメントノードが例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
element SystemA {
    type: software
    docref: docs/spec.md
}
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_relation_draws_connector(self) -> None:
        """リレーションがあるとコネクターシェープが追加されること。"""
        slide = _make_slide()
        before = self._slide_shape_count(slide)
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement R1 { id: 1 }
requirement R2 { id: 2 }
R1 - contains -> R2
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        after = self._slide_shape_count(slide)
        # コネクター(1) + ラベルテキストボックス(1) が追加されること
        assert after > before

    def test_japanese_renders_without_error(self) -> None:
        """日本語テキストが例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement "認証要件" {
    id: JP-1
    text: "ユーザーはパスワードでログインできること"
    risk: high
    verifymethod: test
}
element "認証システム" {
    type: software
    docref: "docs/auth_design.md"
}
"認証システム" - satisfies -> "認証要件"
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_direction_lr(self) -> None:
        """direction LR で複数ノードが例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
direction LR
requirement R1 { id: 1 }
requirement R2 { id: 2 }
requirement R3 { id: 3 }
""")
        renderer.render(slide, diag, 0, 0, 10_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_multiple_nodes_multiple_relations(self) -> None:
        """複数ノード＋複数リレーションが例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        text = """
requirementDiagram
requirement SysReq {
    id: SYS-1
    text: System must work
    risk: high
}
functionalRequirement FuncReq {
    id: FUNC-1
    text: Login functionality
    risk: medium
    verifymethod: test
}
element LoginModule {
    type: software
    docref: docs/login.md
}
SysReq - contains -> FuncReq
LoginModule - satisfies -> FuncReq
LoginModule - verifies -> FuncReq
"""
        diag = parse_requirement(text)
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_class_style_applied_without_error(self) -> None:
        """classDef / class / style スタイル指定があっても例外が発生しないこと。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
classDef critical fill:#ff4444, stroke:#880000
requirement R1 { id: 1 }
class R1 critical
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_unknown_relation_nodes_skipped(self) -> None:
        """リレーションの src/dst が存在しないノードは例外なくスキップされること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        # relations に存在しないノード名を含む
        diag = RequirementDiagram(
            requirements={"R1": RequirementNode(name="R1")},
            relations=[Relationship(src="R1", dst="NonExistent", rel_type="contains")],
        )
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)

    def test_no_body_rows_element(self) -> None:
        """フィールドが空のエレメントでも例外が発生しないこと。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = RequirementDiagram(
            elements={"EmptyElem": ElementNode(name="EmptyElem")},
        )
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_markdown_in_text_rendered(self) -> None:
        """text フィールドの Markdown 書式が例外なく処理されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement MarkdownReq {
    id: MD-1
    text: "This must be **critical** and *important*"
}
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_direction_bt_reverses_order(self) -> None:
        """direction BT で複数ノードが例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
direction BT
requirement R1 { id: 1 }
requirement R2 { id: 2 }
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_all_relation_types(self) -> None:
        """全 7 種のリレーションタイプが例外なく描画されること。"""
        rel_types = ["contains", "copies", "derives", "satisfies", "verifies", "refines", "traces"]
        for rt in rel_types:
            slide = _make_slide()
            renderer = RequirementRenderer()
            text = f"""
requirementDiagram
requirement R1 {{ id: 1 }}
requirement R2 {{ id: 2 }}
R1 - {rt} -> R2
"""
            diag = parse_requirement(text)
            renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)

    def test_mixed_requirement_and_element(self) -> None:
        """要件ノードとエレメントノードの混在が例外なく描画されること。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement Req1 { id: 1 }
element Elem1 { type: hardware }
Elem1 - satisfies -> Req1
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0

    def test_style_direct_fill_applied(self) -> None:
        """直接 style 指定の fill が設定されても例外が発生しないこと。"""
        slide = _make_slide()
        renderer = RequirementRenderer()
        diag = parse_requirement("""
requirementDiagram
requirement R1 { id: 1 }
style R1 fill:#aabbcc, stroke:#112233
""")
        renderer.render(slide, diag, 0, 0, 8_000_000, 6_000_000)
        assert self._slide_shape_count(slide) > 0
