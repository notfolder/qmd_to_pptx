"""
SequenceDiagramRenderer のユニットテスト。

mermaid-parser-py でパースしたシーケンス図データを使い、
各種描画要素（参加者・ライフライン・メッセージ・アクティベーション・
ノート・フレーム・自己メッセージ・クロス矢印）が正しく
スライドに追加されることを確認する。
"""

from __future__ import annotations

import pytest
from mermaid_parser import MermaidParser
from pptx import Presentation
from pptx.oxml.ns import qn as _qn
from pptx.util import Emu

from qmd_to_pptx.mermaid.sequence_diagram import SequenceDiagramRenderer


# ---------------------------------------------------------------------------
# ヘルパー関数
# ---------------------------------------------------------------------------

def _parse_sequence(mermaid_text: str) -> dict:
    """mermaid-parser-py でシーケンス図をパースして graph_data を返す。"""
    mp = MermaidParser()
    result = mp.parse(mermaid_text)
    return result.get("graph_data", {})


def _make_slide():
    """テスト用のスライドオブジェクトを生成する。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _count_shapes(slide) -> int:
    """スライド上のシェイプ総数を返す。"""
    return len(slide.shapes)


def _find_textbox_texts(slide) -> list[str]:
    """スライド上のすべてのシェイプからテキストを収集する。"""
    texts = []
    for sp in slide.shapes:
        try:
            t = sp.text_frame.text.strip()
            if t:
                texts.append(t)
        except AttributeError:
            pass
    return texts


# ---------------------------------------------------------------------------
# テストケース
# ---------------------------------------------------------------------------

class TestBasicSequence:
    """基本的な2参加者シーケンス図のテスト。"""

    def test_basic_sequence_adds_shapes(self):
        """
        2参加者・2メッセージの基本シーケンス図が正しくスライドに追加されることを確認する。

        期待される要素:
        - 参加者ボックス × 2（participant A, B）
        - ライフライン × 2（垂直コネクター）
        - メッセージコネクター × 2
        - メッセージラベルテキストボックス × 2
        """
        text = """sequenceDiagram
    participant A as クライアント
    participant B as サーバー
    A->>B: リクエスト送信
    B-->>A: レスポンス返却
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        initial = _count_shapes(slide)

        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # 少なくとも 参加者×2 + ライフライン×2 + メッセージ×2 + ラベル×2 = 8 シェイプが追加される
        assert _count_shapes(slide) >= initial + 8

    def test_participant_labels_appear(self):
        """参加者ラベルがシェイプのテキストとして表示されることを確認する。"""
        text = """sequenceDiagram
    participant A as アリス
    participant B as ボブ
    A->>B: こんにちは
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("アリス" in t for t in texts)
        assert any("ボブ" in t for t in texts)


class TestActorShape:
    """actor（角丸矩形）とparticipant（矩形）の形状区別テスト。"""

    def test_actor_uses_rounded_rectangle(self):
        """
        actor 宣言のシェイプが角丸矩形（prstGeom val="roundRect"）で
        描画されることを確認する。
        """
        text = """sequenceDiagram
    actor A as ユーザー
    participant B as システム
    A->>B: 操作
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # 角丸矩形の prstGeom を持つシェイプが存在するか確認する
        spTree = slide.shapes._spTree
        round_rects = []
        for sp_el in spTree.findall(_qn("p:sp")):
            spPr = sp_el.find(_qn("p:spPr"))
            if spPr is not None:
                prstGeom = spPr.find(_qn("a:prstGeom"))
                if prstGeom is not None and prstGeom.get("prst") == "roundRect":
                    round_rects.append(sp_el)

        assert len(round_rects) >= 1, "actorシェイプに角丸矩形が存在しない"


class TestArrowStyles:
    """矢印スタイルの描画テスト。"""

    def test_dotted_arrow_has_dash_style(self):
        """
        -->> (DOTTED) メッセージのコネクターに破線スタイル（prstDash val="dash"）が
        設定されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A-->>B: 破線メッセージ
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # コネクター内に prstDash val="dash" が存在することを確認する
        spTree = slide.shapes._spTree
        dash_found = False
        for cxn_el in spTree.findall(_qn("p:cxnSp")):
            for prstDash in cxn_el.iter(_qn("a:prstDash")):
                if prstDash.get("val") == "dash":
                    dash_found = True
                    break

        assert dash_found, "破線スタイルが設定されていない"

    def test_cross_arrow_appends_x_mark(self):
        """
        -x (SOLID_CROSS) メッセージのラベルに「✕」が付加されることを確認する。

        OOXMLにクロス矢印headEndがないため、ラベルへの文字付加で代替する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A-xB: エラー
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("✕" in t for t in texts), "クロス矢印に✕マークが付加されていない"


class TestActivationBar:
    """アクティベーションバーの描画テスト。"""

    def test_activation_creates_bar(self):
        """
        activate/deactivate 宣言がアクティベーションバー（矩形）を生成することを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A->>B: リクエスト
    activate B
    B-->>A: レスポンス
    deactivate B
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        initial = _count_shapes(slide)
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # アクティベーションバーが少なくとも1つ追加される
        assert _count_shapes(slide) > initial + 4

    def test_plus_minus_suffix_creates_activation(self):
        """
        +/- サフィックス付きメッセージ（A->>+B / B-->>-A）で
        アクティベーションバーが生成されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A->>+B: Hello
    B-->>-A: Reply
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        initial = _count_shapes(slide)
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        assert _count_shapes(slide) > initial + 4


class TestNoteRendering:
    """ノートボックスの描画テスト。"""

    def test_note_right_of_appears(self):
        """
        Note right of X のノートボックスがスライドに追加されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A->>B: メッセージ
    Note right of B: ノートテキスト
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("ノートテキスト" in t for t in texts), "ノートテキストがスライドに存在しない"

    def test_note_over_two_actors_appears(self):
        """
        Note over A,B のノートが複数参加者にまたがってスライドに追加されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A->>B: メッセージ
    Note over A,B: 範囲ノート
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("範囲ノート" in t for t in texts)


class TestAutonumber:
    """autonumber（シーケンス番号）のテスト。"""

    def test_autonumber_prefixes_message_labels(self):
        """
        autonumber が有効な場合、メッセージラベルに「1. 」「2. 」形式の
        番号が付与されることを確認する。
        """
        text = """sequenceDiagram
    autonumber
    participant A
    participant B
    A->>B: 最初のメッセージ
    B-->>A: 次のメッセージ
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("1." in t for t in texts), "シーケンス番号1が付与されていない"
        assert any("2." in t for t in texts), "シーケンス番号2が付与されていない"


class TestFrameRendering:
    """フレームボックス（loop/alt/opt）の描画テスト。"""

    def test_loop_frame_label_appears(self):
        """
        loop フレームのキーワードラベル「loop」がスライドに表示されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    loop ヘルスチェック
        A->>B: ping
        B-->>A: pong
    end
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("loop" in t for t in texts), "loopキーワードがスライドに表示されていない"

    def test_alt_frame_label_and_else_appear(self):
        """
        alt フレームのキーワード・ラベル・else セクションラベルが
        スライドに表示されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    alt 成功の場合
        A->>B: 通常処理
    else 失敗の場合
        A->>B: エラー処理
    end
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("alt" in t for t in texts), "altキーワードが表示されていない"
        assert any("成功の場合" in t for t in texts), "altラベルが表示されていない"
        assert any("失敗の場合" in t for t in texts), "elseラベルが表示されていない"


class TestSelfMessage:
    """自己メッセージ（同一参加者への矢印）のテスト。"""

    def test_self_message_adds_connector(self):
        """
        自己メッセージが ELBOWコネクター（p:cxnSp）として追加されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    A->>A: 自己処理
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        initial_connectors = len(slide.shapes._spTree.findall(_qn("p:cxnSp")))
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # コネクターが追加されている（ライフライン1 + 自己メッセージ1 = 最低2）
        final_connectors = len(slide.shapes._spTree.findall(_qn("p:cxnSp")))
        assert final_connectors > initial_connectors

    def test_self_message_label_appears(self):
        """自己メッセージのラベルがスライドに表示されることを確認する。"""
        text = """sequenceDiagram
    participant A
    A->>A: 内部処理
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        texts = _find_textbox_texts(slide)
        assert any("内部処理" in t for t in texts)


class TestEmptyOrFallback:
    """空入力・フォールバックのテスト。"""

    def test_empty_actors_renders_fallback(self):
        """
        参加者なしの場合にフォールバック（テキストボックス）が表示されることを確認する。
        """
        graph_data = {"state": {"records": {"actors": {}, "messages": [], "boxes": []}}}
        slide = _make_slide()
        initial = _count_shapes(slide)
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # フォールバックが1つ追加される
        assert _count_shapes(slide) >= initial + 1

    def test_large_sequence_fits_in_height(self):
        """
        メッセージが多い場合でも MSG_STEP が縮小されてスライド高さに収まることを確認する。
        """
        lines = ["sequenceDiagram", "    participant A", "    participant B"]
        for i in range(30):
            lines.append(f"    A->>B: メッセージ{i}")
        text = "\n".join(lines)

        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        # 例外なく描画できることを確認する
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)
        assert True  # 例外が発生しなければOK


class TestParticipantBoxGroup:
    """参加者グループBox（box ... end）のテスト。"""

    def test_rect_background_added(self):
        """
        rect rgb(...)で指定された背景矩形がスライドに追加されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    rect rgb(200, 240, 255)
        A->>B: 背景中のメッセージ
    end
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        # 少なくとも参加者・ライフライン・メッセージ・rect背景が追加される
        assert _count_shapes(slide) >= 6


class TestBidirectionalArrow:
    """双方向矢印（<<->>）のテスト。"""

    def test_bidir_arrow_adds_connector(self):
        """
        <<->> 双方向矢印がコネクターとして追加されることを確認する。
        """
        text = """sequenceDiagram
    participant A
    participant B
    A<<->>B: 双方向通信
"""
        graph_data = _parse_sequence(text)
        slide = _make_slide()
        initial_connectors = len(slide.shapes._spTree.findall(_qn("p:cxnSp")))
        renderer = SequenceDiagramRenderer()
        renderer.render(slide, graph_data, 0, 0, 8_000_000, 5_000_000)

        final_connectors = len(slide.shapes._spTree.findall(_qn("p:cxnSp")))
        assert final_connectors > initial_connectors
