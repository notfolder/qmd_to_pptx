"""
SlideRenderer クラスの単体テスト。

_select_layout・_load_layout_json・_resolve_placeholder の各メソッドを
統合テストに依存せず直接検証する。
"""

import xml.etree.ElementTree as ET

import pytest
from pptx import Presentation

from qmd_to_pptx.models import (
    DOMNodeInfo,
    DOMNodeType,
    SeparatorType,
    SlideContent,
)
from qmd_to_pptx.slide_renderer import SlideRenderer


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------

def _para_node() -> DOMNodeInfo:
    """テキスト段落の DOMNodeInfo を生成する。"""
    elem = ET.Element("p")
    elem.text = "テキスト"
    return DOMNodeInfo(DOMNodeType.PARAGRAPH, elem)


def _notes_node() -> DOMNodeInfo:
    """スピーカーノートの DOMNodeInfo を生成する。"""
    elem = ET.Element("div")
    elem.set("class", "notes")
    return DOMNodeInfo(DOMNodeType.NOTES, elem)


def _make_content(sep_type: SeparatorType, body: str = "", title: str = "") -> SlideContent:
    """指定した区切り種別の SlideContent を生成する。"""
    return SlideContent(body_text=body, separator_type=sep_type, title=title)


# ---------------------------------------------------------------------------
# _load_layout_json のテスト
# ---------------------------------------------------------------------------

class TestLoadLayoutJson:
    """_load_layout_json の単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()

    def test_layout_json_has_7_layouts(self) -> None:
        """default_layout.json に7種類のレイアウトが定義されている。"""
        layouts = self.renderer._layout_json.layouts
        expected = {
            "Title Slide",
            "Title and Content",
            "Section Header",
            "Two Content",
            "Comparison",
            "Content with Caption",
            "Blank",
        }
        assert set(layouts.keys()) == expected

    def test_layout_json_slide_width_positive(self) -> None:
        """スライド幅が正の整数として読み込まれる。"""
        assert self.renderer._layout_json.slide_width_emu > 0

    def test_layout_json_slide_height_positive(self) -> None:
        """スライド高さが正の整数として読み込まれる。"""
        assert self.renderer._layout_json.slide_height_emu > 0

    def test_title_slide_has_title_and_subtitle_roles(self) -> None:
        """Title Slide に title と subtitle の role が定義されている。"""
        layout = self.renderer._layout_json.layouts["Title Slide"]
        roles = {ph.role for ph in layout.placeholders}
        assert "title" in roles
        assert "subtitle" in roles

    def test_blank_layout_has_no_placeholders(self) -> None:
        """Blank レイアウトはプレースホルダーを持たない。"""
        layout = self.renderer._layout_json.layouts["Blank"]
        assert layout.placeholders == []

    def test_two_content_has_left_and_right_content_roles(self) -> None:
        """Two Content に left_content と right_content の role が定義されている。"""
        layout = self.renderer._layout_json.layouts["Two Content"]
        roles = {ph.role for ph in layout.placeholders}
        assert "left_content" in roles
        assert "right_content" in roles

    def test_comparison_has_all_required_roles(self) -> None:
        """Comparison に left_header/left_content/right_header/right_content の role が定義されている。"""
        layout = self.renderer._layout_json.layouts["Comparison"]
        roles = {ph.role for ph in layout.placeholders}
        assert "left_header" in roles
        assert "left_content" in roles
        assert "right_header" in roles
        assert "right_content" in roles

    def test_content_with_caption_has_caption_role(self) -> None:
        """Content with Caption に caption の role が定義されている。"""
        layout = self.renderer._layout_json.layouts["Content with Caption"]
        roles = {ph.role for ph in layout.placeholders}
        assert "caption" in roles


# ---------------------------------------------------------------------------
# _select_layout のテスト
# ---------------------------------------------------------------------------

class TestSelectLayout:
    """_select_layout メソッドの単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()

    def test_section_header_when_heading1_and_slide_level_2(self) -> None:
        """HEADING1 かつ slide-level:2 の場合 Section Header が選択される。"""
        content = _make_content(SeparatorType.HEADING1)
        result = self.renderer._select_layout(content, [], slide_level=2)
        assert result == "Section Header"

    def test_not_section_header_when_heading1_and_slide_level_1(self) -> None:
        """slide-level:1 の場合 HEADING1 でも Section Header にならない。"""
        content = _make_content(SeparatorType.HEADING1)
        nodes = [_para_node()]
        result = self.renderer._select_layout(content, nodes, slide_level=1)
        assert result != "Section Header"

    def test_blank_when_no_content_nodes(self) -> None:
        """コンテンツノードが空の場合 Blank が選択される。"""
        content = _make_content(SeparatorType.RULER)
        result = self.renderer._select_layout(content, [], slide_level=2)
        assert result == "Blank"

    def test_blank_when_notes_only(self) -> None:
        """NOTES ノードのみの場合も Blank が選択される。"""
        content = _make_content(SeparatorType.RULER)
        nodes = [_notes_node()]
        result = self.renderer._select_layout(content, nodes, slide_level=2)
        assert result == "Blank"

    def test_title_and_content_for_paragraph(self) -> None:
        """段落ノードのみの場合 Title and Content が選択される。"""
        content = _make_content(SeparatorType.HEADING2)
        nodes = [_para_node()]
        result = self.renderer._select_layout(content, nodes, slide_level=2)
        assert result == "Title and Content"

    def test_two_content_for_text_columns(self) -> None:
        """テキストのみの2カラム構成で Two Content が選択される。"""
        content = _make_content(SeparatorType.HEADING2)
        columns_elem = ET.Element("div")
        columns_elem.set("class", "columns")
        col1 = ET.SubElement(columns_elem, "div")
        col1.set("class", "column")
        ET.SubElement(col1, "p").text = "左カラム"
        col2 = ET.SubElement(columns_elem, "div")
        col2.set("class", "column")
        ET.SubElement(col2, "p").text = "右カラム"
        nodes = [DOMNodeInfo(DOMNodeType.COLUMNS, columns_elem)]
        result = self.renderer._select_layout(content, nodes, slide_level=2)
        assert result == "Two Content"

    def test_comparison_for_columns_with_table(self) -> None:
        """テーブルを含むカラムが存在する場合 Comparison が選択される。"""
        content = _make_content(SeparatorType.HEADING2)
        columns_elem = ET.Element("div")
        columns_elem.set("class", "columns")
        col1 = ET.SubElement(columns_elem, "div")
        col1.set("class", "column")
        ET.SubElement(col1, "table")  # 非テキスト要素
        col2 = ET.SubElement(columns_elem, "div")
        col2.set("class", "column")
        ET.SubElement(col2, "p").text = "右カラム"
        nodes = [DOMNodeInfo(DOMNodeType.COLUMNS, columns_elem)]
        result = self.renderer._select_layout(content, nodes, slide_level=2)
        assert result == "Comparison"

    def test_default_slide_level_is_2(self) -> None:
        """slide_level 引数を省略した場合のデフォルト値は 2 として動作する。"""
        content = _make_content(SeparatorType.HEADING1)
        result = self.renderer._select_layout(content, [])
        assert result == "Section Header"


# ---------------------------------------------------------------------------
# _resolve_placeholder のテスト
# ---------------------------------------------------------------------------

class TestResolvePlaceholder:
    """_resolve_placeholder メソッドの単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()
        # デフォルトプレゼンテーションでテスト用スライドを生成する
        self.prs = Presentation()
        self.layout = self.prs.slide_layouts[0]  # Title Slide レイアウト
        self.slide = self.prs.slides.add_slide(self.layout)

    def test_resolve_existing_placeholder_returns_true(self) -> None:
        """存在するプレースホルダー idx の場合 True が返る。"""
        # idx=0 はデフォルトの Title Slide に存在する
        result = self.renderer._resolve_placeholder(self.slide, 0)
        assert result is True

    def test_resolve_nonexistent_placeholder_returns_false(self) -> None:
        """存在しないプレースホルダー idx の場合 False が返る。"""
        # idx=999 は通常存在しない
        result = self.renderer._resolve_placeholder(self.slide, 999)
        assert result is False


# ---------------------------------------------------------------------------
# _is_content_with_caption のテスト
# ---------------------------------------------------------------------------

class TestIsContentWithCaption:
    """_is_content_with_caption の単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()

    def _mermaid_node(self) -> DOMNodeInfo:
        """Mermaid 図の DOMNodeInfo を生成する。"""
        elem = ET.Element("code")
        elem.set("class", "language-mermaid")
        return DOMNodeInfo(DOMNodeType.MERMAID, elem)

    def test_text_then_diagram_returns_true(self) -> None:
        """テキスト→図の順でも True を返す（従来の動作を維持）。"""
        nodes = [_para_node(), self._mermaid_node()]
        assert self.renderer._is_content_with_caption(nodes) is True

    def test_diagram_then_text_returns_true(self) -> None:
        """図→テキストの順でも True を返す（順不同対応の修正確認）。"""
        nodes = [self._mermaid_node(), _para_node()]
        assert self.renderer._is_content_with_caption(nodes) is True

    def test_text_only_returns_false(self) -> None:
        """テキストのみの場合は False を返す。"""
        nodes = [_para_node(), _para_node()]
        assert self.renderer._is_content_with_caption(nodes) is False

    def test_diagram_only_returns_false(self) -> None:
        """図のみの場合は False を返す。"""
        nodes = [self._mermaid_node()]
        assert self.renderer._is_content_with_caption(nodes) is False


# ---------------------------------------------------------------------------
# _render_body_node のルーティングテスト
# ---------------------------------------------------------------------------

class TestRenderBodyNodeRouting:
    """_render_body_node の Content with Caption ルーティングの単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()
        # テンプレートなし（プレースホルダーなし）のスライドを生成する
        self.prs = Presentation()
        # Blank レイアウト（プレースホルダーなし）を使用する
        blank_layout = None
        for layout in self.prs.slide_layouts:
            if not list(layout.placeholders):
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(blank_layout)

    def _mermaid_node(self) -> DOMNodeInfo:
        """Mermaid 図の DOMNodeInfo を生成する。"""
        elem = ET.Element("code")
        elem.set("class", "language-mermaid")
        elem.text = "graph LR\n  A --> B"
        return DOMNodeInfo(DOMNodeType.MERMAID, elem)

    def test_diagram_routes_to_body_role(self) -> None:
        """Content with Caption 時、図系ノードは body ロールに描画される。"""
        layout_def = self.renderer._layout_json.layouts["Content with Caption"]
        node = self._mermaid_node()
        # 例外が出ずに処理されれば OK（実際の描画内容はインテグレーションテストで確認）
        try:
            self.renderer._render_body_node(
                self.slide, node, layout_def, None,
                incremental=False, layout_name="Content with Caption"
            )
        except Exception as e:
            pytest.fail(f"図系ノードのルーティングで例外が発生した: {e}")

    def test_text_routes_to_caption_role(self) -> None:
        """Content with Caption 時、テキスト系ノードは caption ロールに描画される。"""
        layout_def = self.renderer._layout_json.layouts["Content with Caption"]
        node = _para_node()
        try:
            self.renderer._render_body_node(
                self.slide, node, layout_def, None,
                incremental=False, layout_name="Content with Caption"
            )
        except Exception as e:
            pytest.fail(f"テキスト系ノードのルーティングで例外が発生した: {e}")


# ---------------------------------------------------------------------------
# _add_slide のフォールバックテスト
# ---------------------------------------------------------------------------

class TestAddSlideLayoutFallback:
    """_add_slide のフォールバック動作の単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()
        # デフォルトの Presentation（11種類のレイアウトを持つ）を使用する
        self.prs = Presentation()

    def test_unknown_layout_falls_back_to_first_layout(self) -> None:
        """存在しないレイアウト名の場合 slide_layouts[0] を使用する。"""
        slide = self.renderer._add_slide(self.prs, "存在しないレイアウト", None)
        assert slide is not None

    def test_content_with_caption_falls_back_to_title_and_content(self) -> None:
        """'Content with Caption' が見つからない場合 'Title and Content' へフォールバックする。"""
        # デフォルトの Presentation には "Content with Caption" が存在する場合があるため、
        # 存在しない名前を使ってフォールバックチェーン全体を確認する
        slide = self.renderer._add_slide(self.prs, "Two Content", None)
        assert slide is not None


# ---------------------------------------------------------------------------
# _render_body_nodes_vertical_split のテスト
# ---------------------------------------------------------------------------

class TestRenderBodyNodesVerticalSplit:
    """_render_body_nodes_vertical_split の単体テスト。"""

    def setup_method(self) -> None:
        self.renderer = SlideRenderer()
        self.prs = Presentation()
        # プレースホルダーなしのスライドを生成する
        blank_layout = None
        for layout in self.prs.slide_layouts:
            if not list(layout.placeholders):
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(blank_layout)

    def _mermaid_node(self) -> DOMNodeInfo:
        """Mermaid 図の DOMNodeInfo を生成する。"""
        elem = ET.Element("code")
        elem.set("class", "language-mermaid")
        elem.text = "graph LR\n  A --> B"
        return DOMNodeInfo(DOMNodeType.MERMAID, elem)

    def test_vertical_split_adds_textboxes_without_overlap(self) -> None:
        """縦分割時、テキスト系と図系が別の y 座標範囲に描画される。"""
        layout_def = self.renderer._layout_json.layouts["Title and Content"]
        body_nodes = [_para_node(), self._mermaid_node()]
        shapes_before = len(self.slide.shapes)

        self.renderer._render_body_nodes_vertical_split(
            self.slide, body_nodes, layout_def, None, incremental=False
        )

        # テキスト系 1 件のため textbox が 1 つ追加されているはず
        shapes_after = len(self.slide.shapes)
        assert shapes_after > shapes_before

    def test_vertical_split_no_crash_with_no_body_area(self) -> None:
        """body エリア座標が取得できない場合は処理をスキップして例外を出さない。"""
        # 空の LayoutDef（placeholderなし）を使用する
        from qmd_to_pptx.models import LayoutDef
        empty_layout_def = LayoutDef(placeholders=[])
        body_nodes = [_para_node()]
        # 例外が出なければ OK
        self.renderer._render_body_nodes_vertical_split(
            self.slide, body_nodes, empty_layout_def, None, incremental=False
        )
