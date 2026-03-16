"""
StateDiagramRenderer のユニットテスト（改版）。

state_diagram.py が各 shape フィールドを正しい OOXML 形状で描画すること、
複合状態・choice・fork/join・遷移ラベルを正しく処理することを確認する。
"""

from __future__ import annotations

import xml.etree.ElementTree as ET

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from qmd_to_pptx.mermaid.state_diagram import StateDiagramRenderer
from qmd_to_pptx.mermaid_renderer import MermaidRenderer


# ---------------------------------------------------------------------------
# テスト用ヘルパー
# ---------------------------------------------------------------------------

def _make_slide():
    """テスト用スライドを生成して返す。"""
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[5])


def _prst(shape) -> str | None:
    """
    python-pptx の Shape から OOXML の prstGeom prst 属性を取得する。

    Parameters
    ----------
    shape : Shape
        python-pptx の Shape オブジェクト。

    Returns
    -------
    str | None
        prstGeom の prst 属性値（例: "ellipse", "roundRect" など）。
        存在しない場合は None を返す。
    """
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return None
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is None:
        return None
    return prstGeom.get("prst")


def _make_graph_data(
    nodes: list[dict],
    edges: list[dict] | None = None,
) -> dict:
    """テスト用 graph_data 辞書を生成する。"""
    return {"nodes": nodes, "edges": edges or []}


def _collect_all_texts(slide) -> list[str]:
    """
    スライド内の全テキストを収集する。

    GroupShape 内の子シェイプも再帰的に探索し、
    ラベルテキストボックスがグループ内に移動されていなシェイプも見つける。
    """
    results: list[str] = []

    def _recurse(shapes) -> None:
        for sp in shapes:
            try:
                t = sp.text_frame.text.strip()
                if t:
                    results.append(t)
            except AttributeError:
                pass
            # GroupShape の場合は子シェイプを再帰的に探索する
            try:
                _recurse(sp.shapes)
            except AttributeError:
                pass

    _recurse(slide.shapes)
    return results


def _make_mermaid_element(text: str) -> ET.Element:
    """テスト用 Mermaid code 要素を生成する。"""
    elem = ET.Element("code")
    elem.set("class", "language-mermaid")
    elem.text = text
    return elem


_L, _T, _W, _H = Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)


# ---------------------------------------------------------------------------
# StateDiagramRenderer 単体テスト
# ---------------------------------------------------------------------------

class TestStateDiagramRendererBasic:
    """StateDiagramRenderer の基本的な描画動作テスト。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_render_empty_nodes_does_not_raise(self) -> None:
        """nodes が空のとき例外が発生しない。"""
        slide = _make_slide()
        self.renderer.render(slide, _make_graph_data([]), _L, _T, _W, _H)

    def test_render_single_normal_state_adds_shape(self) -> None:
        """通常状態ノード 1 件でスライドに Shape が追加される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "S1", "label": "待機", "shape": "rect"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        assert len(slide.shapes) > initial

    def test_render_does_not_raise_with_edge(self) -> None:
        """エッジありの状態図が例外なく描画できる。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {"id": "B", "label": "B", "shape": "rect"},
            ],
            edges=[{"start": "A", "end": "B", "label": ""}],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)


# ---------------------------------------------------------------------------
# ノード形状の OOXML prst 検証
# ---------------------------------------------------------------------------

class TestStateNodeShapes:
    """各 shape フィールドが正しい OOXML 形状に変換されることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_normal_state_uses_round_rect(self) -> None:
        """shape="rect" の通常状態は ROUNDED_RECTANGLE (roundRect) で描画される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "S1", "label": "通常", "shape": "rect"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "roundRect" in prst_list, (
            f"shape='rect' で roundRect が見つからない: {prst_list}"
        )

    def test_start_state_uses_ellipse(self) -> None:
        """shape="stateStart" の開始状態は OVAL (ellipse) で描画される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "start", "label": "", "shape": "stateStart"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "ellipse" in prst_list, (
            f"stateStart で ellipse が見つからない: {prst_list}"
        )

    def test_end_state_adds_two_ovals(self) -> None:
        """
        shape="stateEnd" の終了状態は bull's-eye として 2 枚の OVAL を追加する。

        OOXMLに bull's-eye プリセットがないため、
        外側黒OVAL + 内側白OVAL の 2 図形重ね合わせで表現する。
        """
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "end", "label": "", "shape": "stateEnd"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        ellipse_shapes = [s for s in added if _prst(s) == "ellipse"]
        assert len(ellipse_shapes) >= 2, (
            f"stateEnd で OVAL が 2 枚以上ない: 追加された ellipse = {len(ellipse_shapes)}"
        )

    def test_choice_state_uses_diamond(self) -> None:
        """shape="choice" の選択状態は DIAMOND で描画される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "ch", "label": "", "shape": "choice"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "diamond" in prst_list, (
            f"choice で diamond が見つからない: {prst_list}"
        )

    def test_fork_state_uses_flat_rectangle(self) -> None:
        """shape="fork" は扁平黒 RECTANGLE (rect) で描画される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "fk", "label": "", "shape": "fork"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "rect" in prst_list, (
            f"fork で rect が見つからない: {prst_list}"
        )

    def test_join_state_uses_flat_rectangle(self) -> None:
        """shape="join" は扁平黒 RECTANGLE (rect) で描画される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "jn", "label": "", "shape": "join"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "rect" in prst_list, (
            f"join で rect が見つからない: {prst_list}"
        )

    def test_unknown_shape_falls_back_to_round_rect(self) -> None:
        """未知の shape 値は通常状態（roundRect）にフォールバックされる。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "X", "label": "X", "shape": "unknownShapeType"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "roundRect" in prst_list, (
            f"未知 shape で roundRect にフォールバックされていない: {prst_list}"
        )


# ---------------------------------------------------------------------------
# ラベル確認テスト
# ---------------------------------------------------------------------------

class TestStateNodeLabels:
    """ノードのラベルがシェイプのテキストとして描画されることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def _all_texts(self, slide, initial: int) -> list[str]:
        """追加されたシェイプのテキストリストを返す。"""
        texts = []
        for shape in list(slide.shapes)[initial:]:
            try:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            except AttributeError:
                pass
        return texts

    def test_normal_state_label_is_set(self) -> None:
        """通常状態のラベルがシェイプのテキストに設定される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "S1", "label": "待機中", "shape": "rect"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        texts = self._all_texts(slide, initial)
        assert "待機中" in texts, f"ラベル '待機中' が見つからない: {texts}"

    def test_multiple_state_labels(self) -> None:
        """複数の通常状態がそれぞれのラベルを持つ。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {"id": "Idle", "label": "アイドル", "shape": "rect"},
            {"id": "Running", "label": "実行中", "shape": "rect"},
            {"id": "Stopped", "label": "停止", "shape": "rect"},
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        texts = self._all_texts(slide, initial)
        for label in ["アイドル", "実行中", "停止"]:
            assert label in texts, f"ラベル '{label}' が見つからない: {texts}"


# ---------------------------------------------------------------------------
# 遷移ラベルのテスト
# ---------------------------------------------------------------------------

class TestTransitionLabels:
    """エッジの label が遷移ラベルテキストボックスとして描画されることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_edge_with_label_adds_extra_textbox(self) -> None:
        """label 付きエッジは遷移ラベルテキストをスライドに追加する。"""
        slide = _make_slide()
        graph_data_no_label = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {"id": "B", "label": "B", "shape": "rect"},
            ],
            edges=[{"start": "A", "end": "B", "label": ""}],
        )
        self.renderer.render(slide, graph_data_no_label, _L, _T, _W, _H)
        texts_no_label = _collect_all_texts(slide)

        slide2 = _make_slide()
        graph_data_with_label = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {"id": "B", "label": "B", "shape": "rect"},
            ],
            edges=[{"start": "A", "end": "B", "label": "イベント発生"}],
        )
        self.renderer.render(slide2, graph_data_with_label, _L, _T, _W, _H)
        texts_with_label = _collect_all_texts(slide2)

        # ラベル付きの場合はラベルテキストがスライド内（グループ内含む）に現れる
        assert any("イベント発生" in t for t in texts_with_label), (
            f"ラベル 'イベント発生' が見つからない: {texts_with_label}"
        )
        assert not any("イベント発生" in t for t in texts_no_label), (
            f"ラベルなしの場合にラベルテキストが現れた: {texts_no_label}"
        )

    def test_transition_label_text_appears_in_shapes(self) -> None:
        """遷移ラベルのテキストがスライドのシェイプ（グループ内含む）に現れる。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "Idle", "label": "アイドル", "shape": "rect"},
                {"id": "Run", "label": "実行", "shape": "rect"},
            ],
            edges=[{"start": "Idle", "end": "Run", "label": "スタート"}],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = _collect_all_texts(slide)
        assert any("スタート" in t for t in all_texts), (
            f"遷移ラベル 'スタート' がシェイプに見つからない: {all_texts}"
        )

    def test_multiple_edges_with_labels(self) -> None:
        """複数のラベル付きエッジがすべて描画される。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {"id": "B", "label": "B", "shape": "rect"},
                {"id": "C", "label": "C", "shape": "rect"},
            ],
            edges=[
                {"start": "A", "end": "B", "label": "go"},
                {"start": "B", "end": "C", "label": "next"},
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = _collect_all_texts(slide)
        assert any("go" in t for t in all_texts), "ラベル 'go' が見つからない"
        assert any("next" in t for t in all_texts), "ラベル 'next' が見つからない"


# ---------------------------------------------------------------------------
# 複合状態のテスト
# ---------------------------------------------------------------------------

class TestCompositeState:
    """複合状態（isGroup=True の roundedWithTitle）の描画テスト。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_composite_state_adds_background_shape(self) -> None:
        """複合状態コンテナの背景 ROUNDED_RECTANGLE が追加される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data([
            {
                "id": "Composite1",
                "label": "複合状態",
                "shape": "roundedWithTitle",
                "isGroup": True,
            },
            {
                "id": "Inner1",
                "label": "内部1",
                "shape": "rect",
                "parentId": "Composite1",
            },
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        added = list(slide.shapes)[initial:]
        prst_list = [_prst(s) for s in added if _prst(s) is not None]
        assert "roundRect" in prst_list, (
            f"複合状態の外枠 roundRect が見つからない: {prst_list}"
        )

    def test_composite_state_title_text_appears(self) -> None:
        """複合状態のタイトルラベルがスライドに描画される。"""
        slide = _make_slide()
        graph_data = _make_graph_data([
            {
                "id": "VehicleState",
                "label": "VehicleState",
                "shape": "roundedWithTitle",
                "isGroup": True,
            },
            {
                "id": "Moving",
                "label": "Moving",
                "shape": "rect",
                "parentId": "VehicleState",
            },
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = []
        for shape in slide.shapes:
            try:
                t = shape.text_frame.text.strip()
                if t:
                    all_texts.append(t)
            except AttributeError:
                pass
        assert any("VehicleState" in t for t in all_texts), (
            f"複合状態タイトル 'VehicleState' が見つからない: {all_texts}"
        )

    def test_composite_child_node_is_drawn(self) -> None:
        """複合状態の子ノードが描画される。"""
        slide = _make_slide()
        graph_data = _make_graph_data([
            {
                "id": "Comp",
                "label": "Comp",
                "shape": "roundedWithTitle",
                "isGroup": True,
            },
            {
                "id": "Child",
                "label": "子ノード",
                "shape": "rect",
                "parentId": "Comp",
            },
        ])
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = []
        for shape in slide.shapes:
            try:
                t = shape.text_frame.text.strip()
                if t:
                    all_texts.append(t)
            except AttributeError:
                pass
        assert any("子ノード" in t for t in all_texts), (
            f"子ノードラベル '子ノード' が見つからない: {all_texts}"
        )

    def test_composite_state_with_multiple_children(self) -> None:
        """複数の子ノードをもつ複合状態が例外なく描画される。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {
                    "id": "Vehicle",
                    "label": "Vehicle",
                    "shape": "roundedWithTitle",
                    "isGroup": True,
                },
                {"id": "Idle", "label": "Idle", "shape": "rect", "parentId": "Vehicle"},
                {"id": "Moving", "label": "Moving", "shape": "rect", "parentId": "Vehicle"},
                {"id": "Stopped", "label": "Stopped", "shape": "rect", "parentId": "Vehicle"},
            ],
            edges=[
                {"start": "Idle", "end": "Moving", "label": "start"},
                {"start": "Moving", "end": "Stopped", "label": "stop"},
            ],
        )
        # 例外なく完了することを確認する
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# 完全な stateDiagram-v2 の統合テスト
# ---------------------------------------------------------------------------

class TestStateDiagramFullRender:
    """StateDiagramRenderer の完全な graph_data を使った統合テスト。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_full_diagram_with_all_node_types(self) -> None:
        """start / end / rect / choice / fork / join をすべて含む図が描画できる。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "[*]_start", "label": "", "shape": "stateStart"},
                {"id": "Idle", "label": "Idle", "shape": "rect"},
                {"id": "ch1", "label": "", "shape": "choice"},
                {"id": "fk1", "label": "", "shape": "fork"},
                {"id": "Moving", "label": "Moving", "shape": "rect"},
                {"id": "jn1", "label": "", "shape": "join"},
                {"id": "[*]_end", "label": "", "shape": "stateEnd"},
            ],
            edges=[
                {"start": "[*]_start", "end": "Idle", "label": ""},
                {"start": "Idle", "end": "ch1", "label": ""},
                {"start": "ch1", "end": "fk1", "label": "速い"},
                {"start": "fk1", "end": "Moving", "label": ""},
                {"start": "Moving", "end": "jn1", "label": ""},
                {"start": "jn1", "end": "[*]_end", "label": ""},
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        # 7 ノード以上のシェイプが追加される（コネクター・ラベルを含む）
        assert len(slide.shapes) >= 7

    def test_diagram_edge_count(self) -> None:
        """エッジの数だけコネクターが追加される（コネクター = 有向接続線）。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {"id": "B", "label": "B", "shape": "rect"},
                {"id": "C", "label": "C", "shape": "rect"},
            ],
            edges=[
                {"start": "A", "end": "B", "label": ""},
                {"start": "B", "end": "C", "label": ""},
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        # 3 ノード + 2 コネクター = 5 シェイプ以上
        assert len(slide.shapes) >= initial + 5


# ---------------------------------------------------------------------------
# MermaidRenderer 経由の統合テスト
# ---------------------------------------------------------------------------

class TestMermaidRendererStateIntegration:
    """MermaidRenderer 経由で stateDiagram-v2 が正しくディスパッチされることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に MermaidRenderer インスタンスを生成する。"""
        self.renderer = MermaidRenderer()

    def test_state_diagram_dispatches_without_error(self) -> None:
        """stateDiagram-v2 が MermaidRenderer から例外なく描画される。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> Idle\n"
            "    Idle --> Moving : start\n"
            "    Moving --> [*] : stop\n"
        )
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )

    def test_state_diagram_adds_shapes(self) -> None:
        """stateDiagram-v2 のレンダリングでスライドに Shape が追加される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> Idle\n"
            "    Idle --> Moving : go\n"
            "    Moving --> [*] : stop\n"
        )
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > initial

    def test_state_diagram_with_choice_via_mermaid_renderer(self) -> None:
        """choice を含む stateDiagram-v2 が MermaidRenderer 経由で描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> Check\n"
            "    Check --> OK : valid\n"
            "    Check --> Error : invalid\n"
            "    OK --> [*]\n"
            "    Error --> [*]\n"
        )
        # 例外なく完了することを確認する
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# 向き（direction）のテスト
# ---------------------------------------------------------------------------

class TestStateDirection:
    """direction キーワードで描画方向が切り替わることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に MermaidRenderer インスタンスを生成する。"""
        self.renderer = MermaidRenderer()

    def _center_y_of_shapes(self, slide) -> list[int]:
        """スライド上の全シェイプの中心 Y 座標（EMU）をリストで返す。"""
        result = []
        for sp in slide.shapes:
            try:
                result.append(int(sp.top) + int(sp.height) // 2)
            except Exception:
                pass
        return result

    def _center_x_of_shapes(self, slide) -> list[int]:
        """スライド上の全シェイプの中心 X 座標（EMU）をリストで返す。"""
        result = []
        for sp in slide.shapes:
            try:
                result.append(int(sp.left) + int(sp.width) // 2)
            except Exception:
                pass
        return result

    def test_direction_tb_does_not_raise(self) -> None:
        """direction TB（デフォルト）で例外なく描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    direction TB\n"
            "    [*] --> A\n"
            "    A --> B\n"
            "    B --> [*]\n"
        )
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0

    def test_direction_lr_does_not_raise(self) -> None:
        """direction LR で例外なく描画できる。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    direction LR\n"
            "    [*] --> Idle\n"
            "    Idle --> Running\n"
            "    Running --> [*]\n"
        )
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0

    def test_no_direction_keyword_defaults_to_tb(self) -> None:
        """direction キーワードなしの場合でも例外なく描画できる（デフォルト TB）。"""
        slide = _make_slide()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> A\n"
            "    A --> B\n"
            "    B --> [*]\n"
        )
        self.renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# ノート（note）のテスト
# ---------------------------------------------------------------------------

class TestStateNotes:
    """note ノードが黄色テキストボックスとして描画されることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_note_group_adds_shape(self) -> None:
        """noteGroup ノードが存在するときシェイプが追加される。"""
        slide = _make_slide()
        initial = len(slide.shapes)
        graph_data = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {
                    "id": "A----parent",
                    "label": "これはノートです",
                    "shape": "noteGroup",
                    "isGroup": True,
                    "position": "right of",
                },
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        # noteGroup の分だけシェイプが増える
        assert len(slide.shapes) > initial

    def test_note_text_appears_in_shapes(self) -> None:
        """ノートのラベルテキストがスライドのシェイプに現れる。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "A", "label": "A", "shape": "rect"},
                {
                    "id": "A----parent",
                    "label": "重要なメモ",
                    "shape": "noteGroup",
                    "isGroup": True,
                    "position": "right of",
                },
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = []
        for shape in slide.shapes:
            try:
                t = shape.text_frame.text.strip()
                if t:
                    all_texts.append(t)
            except AttributeError:
                pass
        assert any("重要なメモ" in t for t in all_texts), (
            f"ノートテキスト '重要なメモ' が見つからない: {all_texts}"
        )

    def test_note_via_mermaid_renderer(self) -> None:
        """note を含む stateDiagram-v2 が MermaidRenderer 経由で例外なく描画できる。"""
        slide = _make_slide()
        renderer = MermaidRenderer()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> Active\n"
            "    Active --> [*]\n"
            "    note right of Active : テストノート\n"
        )
        renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# 並行（concurrent / divider）のテスト
# ---------------------------------------------------------------------------

class TestStateConcurrent:
    """並行セクション（--）が divider として描画されることを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_divider_shape_does_not_raise(self) -> None:
        """shape="divider" ノードが存在しても例外が発生しない。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {
                    "id": "Active",
                    "label": "Active",
                    "shape": "roundedWithTitle",
                    "isGroup": True,
                },
                {
                    "id": "div1",
                    "label": "",
                    "shape": "divider",
                    "isGroup": True,
                    "parentId": "Active",
                },
                {
                    "id": "NumLockOff",
                    "label": "NumLockOff",
                    "shape": "rect",
                    "parentId": "div1",
                },
                {
                    "id": "div2",
                    "label": "",
                    "shape": "divider",
                    "isGroup": True,
                    "parentId": "Active",
                },
                {
                    "id": "CapsLockOff",
                    "label": "CapsLockOff",
                    "shape": "rect",
                    "parentId": "div2",
                },
            ],
            edges=[
                {"start": "NumLockOff", "end": "NumLockOff", "label": ""},
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

    def test_divider_children_labels_appear(self) -> None:
        """divider 内の子ノードラベルがスライドに描画される。"""
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {
                    "id": "Active",
                    "label": "Active",
                    "shape": "roundedWithTitle",
                    "isGroup": True,
                },
                {
                    "id": "div1",
                    "label": "",
                    "shape": "divider",
                    "isGroup": True,
                    "parentId": "Active",
                },
                {
                    "id": "NumLock",
                    "label": "NumLock状態",
                    "shape": "rect",
                    "parentId": "div1",
                },
                {
                    "id": "div2",
                    "label": "",
                    "shape": "divider",
                    "isGroup": True,
                    "parentId": "Active",
                },
                {
                    "id": "CapsLock",
                    "label": "CapsLock状態",
                    "shape": "rect",
                    "parentId": "div2",
                },
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)

        all_texts = []
        for shape in slide.shapes:
            try:
                t = shape.text_frame.text.strip()
                if t:
                    all_texts.append(t)
            except AttributeError:
                pass
        assert any("NumLock状態" in t for t in all_texts), (
            f"divider 内のラベルが見つからない: {all_texts}"
        )
        assert any("CapsLock状態" in t for t in all_texts), (
            f"divider 内のラベルが見つからない: {all_texts}"
        )

    def test_concurrent_state_via_mermaid_renderer(self) -> None:
        """-- 記法の並行状態が MermaidRenderer 経由で例外なく描画できる。"""
        slide = _make_slide()
        renderer = MermaidRenderer()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> Active\n"
            "    state Active {\n"
            "        [*] --> NumLockOff\n"
            "        NumLockOff --> NumLockOn : press\n"
            "        --\n"
            "        [*] --> CapsLockOff\n"
            "        CapsLockOff --> CapsLockOn : press\n"
            "    }\n"
            "    Active --> [*]\n"
        )
        renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# BFS ランクレイアウト（サイクルを含む状態図の縦配置）テスト
# ---------------------------------------------------------------------------

class TestBfsRankLayout:
    """サイクルのある状態図でも TB 方向で縦に並ぶことを確認する。"""

    def setup_method(self) -> None:
        """テスト前に StateDiagramRenderer インスタンスを生成する。"""
        self.renderer = StateDiagramRenderer()

    def test_cyclic_diagram_does_not_raise(self) -> None:
        """
        サイクル（エラー → 待機中 → 処理中 → エラー）を含む状態図が
        例外なく描画できる（demo_full.qmd 5a に対応）。
        """
        slide = _make_slide()
        graph_data = _make_graph_data(
            nodes=[
                {"id": "__start__", "label": "", "shape": "stateStart"},
                {"id": "待機中", "label": "待機中", "shape": "rect"},
                {"id": "処理中", "label": "処理中", "shape": "rect"},
                {"id": "完了", "label": "完了", "shape": "rect"},
                {"id": "エラー", "label": "エラー", "shape": "rect"},
                {"id": "__end__", "label": "", "shape": "stateEnd"},
            ],
            edges=[
                {"start": "__start__", "end": "待機中", "label": ""},
                {"start": "待機中", "end": "処理中", "label": "開始"},
                {"start": "処理中", "end": "完了", "label": "成功"},
                {"start": "処理中", "end": "エラー", "label": "失敗"},
                {"start": "エラー", "end": "待機中", "label": "リトライ"},  # サイクル辺
                {"start": "完了", "end": "__end__", "label": ""},
            ],
        )
        self.renderer.render(slide, graph_data, _L, _T, _W, _H)
        assert len(slide.shapes) > 0

    def test_cyclic_diagram_tb_arranges_vertically(self) -> None:
        """
        サイクルを含む状態図を TB 方向で描画したとき、
        ノードの Y 座標が複数の値を取る（縦に並んでいる）ことを確認する。
        """
        import networkx as nx

        # StateDiagramRenderer 内部の _bfs_rank_layout を直接テストする
        G = nx.DiGraph()
        nodes = ["__start__", "待機中", "処理中", "完了", "エラー", "__end__"]
        G.add_nodes_from(nodes)
        # サイクル辺: エラー → 待機中
        edges = [
            ("__start__", "待機中"),
            ("待機中", "処理中"),
            ("処理中", "完了"),
            ("処理中", "エラー"),
            ("エラー", "待機中"),  # サイクル
            ("完了", "__end__"),
        ]
        G.add_edges_from(edges)

        pos = self.renderer._bfs_rank_layout(G, "TB")

        # 全ノードがレイアウトに含まれる
        assert set(pos.keys()) == set(nodes)

        # TB 方向なので Y 座標に複数の異なる値がある（=縦に並んでいる）
        y_values = [v[1] for v in pos.values()]
        assert len(set(y_values)) > 1, (
            f"TB 方向のはずが Y 座標が単一値: {y_values}"
        )

    def test_cyclic_diagram_via_mermaid_renderer(self) -> None:
        """
        サイクルを含む stateDiagram-v2 が MermaidRenderer 経由で
        例外なく描画でき、シェイプが追加される。
        """
        slide = _make_slide()
        renderer = MermaidRenderer()
        elem = _make_mermaid_element(
            "stateDiagram-v2\n"
            "    [*] --> 待機中\n"
            "    待機中 --> 処理中 : 開始\n"
            "    処理中 --> 完了 : 成功\n"
            "    処理中 --> エラー : 失敗\n"
            "    エラー --> 待機中 : リトライ\n"
            "    完了 --> [*]\n"
        )
        renderer.render(
            slide, elem,
            Emu(457200), Emu(457200), Emu(8229600), Emu(4800000)
        )
        assert len(slide.shapes) > 0
