"""
テキストレンダラー（TextRenderer）の単体テスト。
"""

import xml.etree.ElementTree as ET
import pytest
from pptx import Presentation
from pptx.util import Emu
from qmd_to_pptx.text_renderer import TextRenderer


def _make_slide():
    """テスト用スライドを生成して返す。"""
    prs = Presentation()
    layout = prs.slide_layouts[5]  # Blankレイアウト
    return prs.slides.add_slide(layout), prs


def _make_textbox(slide):
    """テスト用テキストボックスを生成して返す。"""
    return slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(3000000), Emu(1000000))


class TestTextRenderer:
    """TextRenderer クラスの単体テスト。"""

    def setup_method(self) -> None:
        """テスト前にTextRendererインスタンスを生成する。"""
        self.renderer = TextRenderer()

    # --- render_heading のテスト ---

    def test_render_heading_h1(self) -> None:
        """h1 見出しを太字・大フォントサイズでテキストボックスに書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.Element("h1")
        elem.text = "見出し1テスト"
        self.renderer.render_heading(shape, elem, level=1)
        tf = shape.text_frame
        runs = [run for para in tf.paragraphs for run in para.runs]
        assert len(runs) == 1
        assert runs[0].text == "見出し1テスト"
        assert runs[0].font.bold is True

    def test_render_heading_h2(self) -> None:
        """h2 見出しを太字でテキストボックスに書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.Element("h2")
        elem.text = "見出し2テスト"
        self.renderer.render_heading(shape, elem, level=2)
        tf = shape.text_frame
        runs = [run for para in tf.paragraphs for run in para.runs]
        assert runs[0].text == "見出し2テスト"
        assert runs[0].font.bold is True

    # --- render_paragraph のテスト ---

    def test_render_paragraph(self) -> None:
        """段落テキストをテキストボックスに書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.Element("p")
        elem.text = "段落テキストです"
        self.renderer.render_paragraph(shape, elem)
        tf = shape.text_frame
        runs = [run for para in tf.paragraphs for run in para.runs]
        assert runs[0].text == "段落テキストです"

    # --- render_list のテスト ---

    def test_render_list_ul(self) -> None:
        """箇条書きリストをインデント付きでテキストボックスに書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        # <ul><li>アイテム1</li><li>アイテム2</li></ul>
        elem = ET.fromstring("<ul><li>アイテム1</li><li>アイテム2</li></ul>")
        self.renderer.render_list(shape, elem)
        tf = shape.text_frame
        texts = [run.text for para in tf.paragraphs for run in para.runs]
        assert "アイテム1" in texts
        assert "アイテム2" in texts

    def test_render_list_nested(self) -> None:
        """ネストしたリストをインデントレベル付きで書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.fromstring(
            "<ul><li>親<ul><li>子</li></ul></li></ul>"
        )
        self.renderer.render_list(shape, elem)
        tf = shape.text_frame
        levels = [para.level for para in tf.paragraphs]
        assert 0 in levels  # 親
        assert 1 in levels  # 子

    def test_render_list_incremental_flag_accepted(self) -> None:
        """incremental=True フラグを渡しても例外なく動作する。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.fromstring("<ul><li>アイテム</li></ul>")
        # 例外が発生しないことを確認
        self.renderer.render_list(shape, elem, incremental=True)

    # --- render_code のテスト ---

    def test_render_code(self) -> None:
        """コードテキストを等幅フォントでテキストボックスに書き込む。"""
        slide, _ = _make_slide()
        shape = _make_textbox(slide)
        elem = ET.Element("code")
        elem.text = "print('hello')"
        self.renderer.render_code(shape, elem)
        tf = shape.text_frame
        runs = [run for para in tf.paragraphs for run in para.runs]
        assert runs[0].text == "print('hello')"
        assert runs[0].font.name == "Courier New"

    # --- render_table のテスト ---

    def test_render_table(self) -> None:
        """テーブル要素をPowerPointテーブルShapeとして追加する。"""
        slide, _ = _make_slide()
        elem = ET.fromstring(
            "<table><thead><tr><th>列1</th><th>列2</th></tr></thead>"
            "<tbody><tr><td>A</td><td>B</td></tr></tbody></table>"
        )
        initial_shapes = len(slide.shapes)
        self.renderer.render_table(
            slide, elem,
            Emu(100000), Emu(100000), Emu(5000000), Emu(2000000)
        )
        # テーブルShapeが追加されていることを確認
        assert len(slide.shapes) == initial_shapes + 1

    def test_render_table_header_bold(self) -> None:
        """テーブルのヘッダー行（1行目）は太字になる。"""
        slide, _ = _make_slide()
        elem = ET.fromstring(
            "<table><thead><tr><th>ヘッダー</th></tr></thead>"
            "<tbody><tr><td>データ</td></tr></tbody></table>"
        )
        self.renderer.render_table(
            slide, elem,
            Emu(100000), Emu(100000), Emu(5000000), Emu(2000000)
        )
        table_shape = slide.shapes[-1]
        header_cell = table_shape.table.cell(0, 0)
        # ヘッダーセルのrunが太字かどうかを確認
        for para in header_cell.text_frame.paragraphs:
            for run in para.runs:
                assert run.font.bold is True

    # --- render_notes のテスト ---

    def test_render_notes(self) -> None:
        """スピーカーノートをスライドのノートテキストフレームに書き込む。"""
        slide, _ = _make_slide()
        elem = ET.Element("div")
        elem.set("class", "notes")
        elem.text = "スピーカーノートのテキスト"
        self.renderer.render_notes(slide, elem)
        notes_tf = slide.notes_slide.notes_text_frame
        notes_text = "".join(
            run.text
            for para in notes_tf.paragraphs
            for run in para.runs
        )
        assert "スピーカーノートのテキスト" in notes_text
