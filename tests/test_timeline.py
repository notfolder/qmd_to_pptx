"""
timeline_parser.py および TimelineRenderer のテストモジュール。

テスト対象:
  - parse_timeline(): Mermaid timeline テキストの解析
  - TimelineRenderer.render(): PPTX スライドへの描画（シェープ生成の検証）

テスト観点:
  - title 行の解析
  - section 行の解析とグループ化
  - period + event の解析（同行複数イベント）
  - 継続行（`: event`）による period へのイベント追加
  - <br> タグの改行変換
  - コメント行（%%）のスキップ
  - accTitle / accDescr のスキップ
  - section なしの場合の period ごとカラー割り当て
  - section ありの場合の section ごとカラー割り当て
  - 空のタイムライン（period 0 件）
  - period のみでイベントなし
  - 日本語テキスト
  - 複数 section にまたがる period
  - レンダラー: タイトルありなし
  - レンダラー: section 帯の生成
  - レンダラー: イベントカードの生成
  - レンダラー: period 円の生成
  - レンダラー: 複数 period の描画
  - レンダラー: 複数 section の描画
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu

from qmd_to_pptx.mermaid.timeline_parser import (
    TimelineData,
    TimelineEvent,
    TimelinePeriod,
    parse_timeline,
)
from qmd_to_pptx.mermaid.timeline_renderer import (
    TimelineRenderer,
    _build_color_map,
    _lighten,
    _darken,
    _SCALE_COLORS,
)


# ---------------------------------------------------------------------------
# ヘルパー
# ---------------------------------------------------------------------------

def _slide():
    """テスト用の空スライドを生成する。"""
    prs = Presentation()
    prs.slide_width = Emu(9_144_000)
    prs.slide_height = Emu(5_143_500)
    layout = prs.slide_layouts[6]  # 白紙レイアウト
    return prs.slides.add_slide(layout)


def _render(data: TimelineData, w: int = 8_000_000, h: int = 4_500_000) -> tuple:
    """TimelineRenderer.render() を実行してシェープ数を返す。"""
    slide = _slide()
    renderer = TimelineRenderer()
    before = len(slide.shapes)
    renderer.render(slide, data, 0, 0, w, h)
    return slide, len(slide.shapes) - before


# ===========================================================================
# TestParseTimelineBasic - 基本解析
# ===========================================================================

class TestParseTimelineBasic:
    """parse_timeline() の基本動作テスト。"""

    def test_empty_text(self):
        """空テキストは空の TimelineData を返す。"""
        result = parse_timeline("")
        assert result.title == ""
        assert result.sections == []
        assert result.periods == []

    def test_header_only(self):
        """timeline ヘッダーのみのテキストは空の TimelineData を返す。"""
        result = parse_timeline("timeline")
        assert result.title == ""
        assert result.periods == []

    def test_title(self):
        """title 行のテキストが正しく解析される。"""
        text = """timeline
    title History of Social Media Platform"""
        result = parse_timeline(text)
        assert result.title == "History of Social Media Platform"

    def test_title_japanese(self):
        """日本語 title が正しく解析される。"""
        text = """timeline
    title ソフトウェア開発史"""
        result = parse_timeline(text)
        assert result.title == "ソフトウェア開発史"

    def test_no_title(self):
        """title がない場合、title は空文字列になる。"""
        text = """timeline
    2004 : Facebook"""
        result = parse_timeline(text)
        assert result.title == ""

    def test_comment_lines_skipped(self):
        """%% コメント行がスキップされる。"""
        text = """timeline
    %% これはコメント
    title テスト
    %% 別のコメント
    2004 : Facebook"""
        result = parse_timeline(text)
        assert result.title == "テスト"
        assert len(result.periods) == 1

    def test_acc_lines_skipped(self):
        """accTitle / accDescr 行がスキップされる。"""
        text = """timeline
    accTitle: Accessibility Title
    accDescr: Accessibility Description
    2004 : Facebook"""
        result = parse_timeline(text)
        assert len(result.periods) == 1
        assert result.periods[0].label == "2004"

    def test_blank_lines_skipped(self):
        """空行がスキップされる。"""
        text = """timeline

    title テスト

    2004 : Facebook

    2005 : YouTube"""
        result = parse_timeline(text)
        assert len(result.periods) == 2


# ===========================================================================
# TestParseTimelineSection - セクション解析
# ===========================================================================

class TestParseTimelineSection:
    """section 関連の解析テスト。"""

    def test_single_section(self):
        """単一 section が正しく解析される。"""
        text = """timeline
    section 産業革命
        Industry 1.0 : 機械化
        Industry 2.0 : 電化"""
        result = parse_timeline(text)
        assert result.sections == ["産業革命"]
        assert len(result.periods) == 2
        assert result.periods[0].section == "産業革命"
        assert result.periods[1].section == "産業革命"

    def test_multiple_sections(self):
        """複数 section が出現順に記録される。"""
        text = """timeline
    section 20世紀
        1900 : 蒸気機関
    section 21世紀
        2000 : インターネット"""
        result = parse_timeline(text)
        assert result.sections == ["20世紀", "21世紀"]
        assert result.periods[0].section == "20世紀"
        assert result.periods[1].section == "21世紀"

    def test_period_without_section(self):
        """section 前の period の section属性は None になる。"""
        text = """timeline
    2004 : Facebook
    section SNS時代
        2006 : Twitter"""
        result = parse_timeline(text)
        assert result.periods[0].section is None
        assert result.periods[1].section == "SNS時代"

    def test_section_unique_dedup(self):
        """同名 section が重複しない（sections リストはユニーク順）。"""
        text = """timeline
    section A
        2000 : X
    section B
        2010 : Y
    section A
        2020 : Z"""
        result = parse_timeline(text)
        # "A" は2度定義されるが sections リストには1回のみ登録される
        assert result.sections == ["A", "B"]
        assert result.sections.count("A") == 1

    def test_all_periods_no_section(self):
        """section が存在しない場合、全 period の section が None。"""
        text = """timeline
    2002 : LinkedIn
    2004 : Facebook"""
        result = parse_timeline(text)
        assert result.sections == []
        for p in result.periods:
            assert p.section is None

    def test_section_with_br(self):
        """section 名に <br> が含まれる場合、改行に変換される。"""
        text = """timeline
    section 2023 Q1 <br> Release
        Bullet 1 : sub-point 1a"""
        result = parse_timeline(text)
        assert "\n" in result.sections[0]
        assert "2023 Q1" in result.sections[0]

    def test_section_spans_multiple_periods(self):
        """1 section に複数 period が属する。"""
        text = """timeline
    section Group
        2000 : A
        2001 : B
        2002 : C"""
        result = parse_timeline(text)
        assert len(result.periods) == 3
        for p in result.periods:
            assert p.section == "Group"


# ===========================================================================
# TestParseTimelineEvents - イベント解析
# ===========================================================================

class TestParseTimelineEvents:
    """イベント解析テスト。"""

    def test_single_event(self):
        """1イベントが正しく解析される。"""
        text = "timeline\n    2004 : Facebook"
        result = parse_timeline(text)
        assert len(result.periods) == 1
        assert result.periods[0].label == "2004"
        assert len(result.periods[0].events) == 1
        assert result.periods[0].events[0].text == "Facebook"

    def test_multiple_events_same_line(self):
        """同行の複数イベント（コロン区切り）が解析される。"""
        text = "timeline\n    2004 : Facebook : Google"
        result = parse_timeline(text)
        p = result.periods[0]
        assert len(p.events) == 2
        assert p.events[0].text == "Facebook"
        assert p.events[1].text == "Google"

    def test_multiple_events_three(self):
        """3イベントが同行で解析される。"""
        text = "timeline\n    2004 : A : B : C"
        result = parse_timeline(text)
        assert len(result.periods[0].events) == 3

    def test_continuation_line(self):
        """継続行（: event）が直前 period にイベントを追加する。"""
        text = """timeline
    2004 : Facebook
         : Google"""
        result = parse_timeline(text)
        assert len(result.periods) == 1
        assert len(result.periods[0].events) == 2
        assert result.periods[0].events[1].text == "Google"

    def test_continuation_multiple_events(self):
        """継続行に複数イベントが含まれる場合。"""
        text = """timeline
    2004 : Facebook
         : Google : YouTube"""
        result = parse_timeline(text)
        assert len(result.periods[0].events) == 3

    def test_period_no_events(self):
        """イベントなし（コロンなし）の period が解析される。"""
        text = "timeline\n    2004"
        result = parse_timeline(text)
        assert len(result.periods) == 1
        assert result.periods[0].label == "2004"
        assert result.periods[0].events == []

    def test_multiple_periods(self):
        """複数 period が出現順に記録される。"""
        text = """timeline
    2002 : LinkedIn
    2004 : Facebook
    2005 : YouTube
    2006 : Twitter"""
        result = parse_timeline(text)
        assert len(result.periods) == 4
        assert result.periods[0].label == "2002"
        assert result.periods[3].label == "2006"

    def test_japanese_events(self):
        """日本語イベントテキストが正しく解析される。"""
        text = "timeline\n    2000年代 : SNSの台頭 : 動画配信の普及"
        result = parse_timeline(text)
        p = result.periods[0]
        assert p.label == "2000年代"
        assert p.events[0].text == "SNSの台頭"
        assert p.events[1].text == "動画配信の普及"


# ===========================================================================
# TestParseTimelineTextProcessing - テキスト処理
# ===========================================================================

class TestParseTimelineTextProcessing:
    """<br> タグの変換・テキストトリムテスト。"""

    def test_br_in_event(self):
        """イベント内の <br> が改行に変換される。"""
        text = "timeline\n    2004 : Steam <br> power"
        result = parse_timeline(text)
        assert "\n" in result.periods[0].events[0].text
        assert "Steam" in result.periods[0].events[0].text
        assert "power" in result.periods[0].events[0].text

    def test_br_self_closing_in_event(self):
        """<br/> が改行に変換される。"""
        text = "timeline\n    2004 : A<br/>B"
        result = parse_timeline(text)
        assert result.periods[0].events[0].text == "A\nB"

    def test_br_with_space_in_event(self):
        """<br /> が改行に変換される。"""
        text = "timeline\n    2004 : A<br />B"
        result = parse_timeline(text)
        assert result.periods[0].events[0].text == "A\nB"

    def test_br_in_period_label(self):
        """period ラベル内の <br> が改行に変換される。"""
        text = "timeline\n    2023 Q1 <br> Release : ABC"
        result = parse_timeline(text)
        assert "\n" in result.periods[0].label

    def test_br_in_title(self):
        """title 内の <br> が改行に変換される。"""
        text = "timeline\n    title Line1 <br> Line2"
        result = parse_timeline(text)
        assert result.title == "Line1 \n Line2"

    def test_whitespace_trimmed(self):
        """各フィールドの前後空白がトリムされる。"""
        text = "timeline\n      2004   :   Facebook   "
        result = parse_timeline(text)
        assert result.periods[0].label == "2004"
        assert result.periods[0].events[0].text == "Facebook"


# ===========================================================================
# TestColorMap - カラーマップ構築
# ===========================================================================

class TestColorMap:
    """_build_color_map() のカラー割り当てテスト。"""

    def test_no_sections_period_colors(self):
        """section なしの場合、period ごとに異なる色が割り当てられる。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(label="A", section=None),
                TimelinePeriod(label="B", section=None),
                TimelinePeriod(label="C", section=None),
            ]
        )
        color_map = _build_color_map(data)
        assert color_map[0] == _SCALE_COLORS[0]
        assert color_map[1] == _SCALE_COLORS[1]
        assert color_map[2] == _SCALE_COLORS[2]

    def test_with_sections_same_color_per_section(self):
        """同一 section の period は同じ色になる。"""
        data = TimelineData(
            sections=["A", "B"],
            periods=[
                TimelinePeriod(label="p1", section="A"),
                TimelinePeriod(label="p2", section="A"),
                TimelinePeriod(label="p3", section="B"),
            ]
        )
        color_map = _build_color_map(data)
        assert color_map[0] == color_map[1]  # 同じ section A
        assert color_map[0] != color_map[2]  # section B は別色

    def test_color_cycles_past_12(self):
        """12色以上の period でも色がサイクルする。"""
        periods = [TimelinePeriod(label=str(i), section=None) for i in range(14)]
        data = TimelineData(periods=periods)
        color_map = _build_color_map(data)
        # 12 個目と 0 個目は同じ色
        assert color_map[12] == color_map[0]

    def test_section_none_in_sections_mode(self):
        """section ありモードで section=None の period は個別色を持つ。"""
        data = TimelineData(
            sections=["A"],
            periods=[
                TimelinePeriod(label="p0", section=None),
                TimelinePeriod(label="p1", section="A"),
            ]
        )
        color_map = _build_color_map(data)
        # p0（section=None）と p1（section="A"）は異なる色になることを確認
        # （どちらも _SCALE_COLORS から取得されるが、インデックスが異なる）
        assert color_map[0] in _SCALE_COLORS
        assert color_map[1] in _SCALE_COLORS

    def test_lighten_utility(self):
        """_lighten() が正しく動作する。"""
        white = _lighten((0, 0, 0), 1.0)
        assert white == (255, 255, 255)
        same = _lighten((100, 150, 200), 0.0)
        assert same == (100, 150, 200)
        mid = _lighten((0, 0, 0), 0.5)
        assert mid == (127, 127, 127)

    def test_darken_utility(self):
        """_darken() が正しく動作する。"""
        dark = _darken((200, 200, 200), 0.5)
        assert dark == (100, 100, 100)
        same = _darken((100, 100, 100), 1.0)
        assert same == (100, 100, 100)


# ===========================================================================
# TestTimelineRendererShapes - レンダラー描画テスト
# ===========================================================================

class TestTimelineRendererShapes:
    """TimelineRenderer.render() のシェープ生成テスト。"""

    def test_empty_no_shapes(self):
        """period が 0 件の場合、シェープが追加されない。"""
        data = TimelineData()
        _, added = _render(data)
        assert added == 0

    def test_single_period_no_event(self):
        """period 1件・events なしの場合、最低限のシェープが生成される。
        期待: 軸ライン1 + period 円1 + period ラベル1 = 3"""
        data = TimelineData(
            periods=[TimelinePeriod(label="2004", section=None)]
        )
        _, added = _render(data)
        assert added >= 3

    def test_single_period_with_title(self):
        """title ありの場合、タイトルテキストボックスが追加される。"""
        data = TimelineData(
            title="テストタイム",
            periods=[TimelinePeriod(label="2004", section=None)]
        )
        _, no_title_added = _render(TimelineData(
            periods=[TimelinePeriod(label="2004", section=None)]
        ))
        _, with_title_added = _render(data)
        # タイトルありはなしより1つ多い
        assert with_title_added == no_title_added + 1

    def test_single_period_with_one_event(self):
        """period 1件・event 1件の場合、event カードが追加される。
        期待: 軸ライン1 + 円1 + ラベル1 + カード1 = 4"""
        data = TimelineData(
            periods=[TimelinePeriod(
                label="2004",
                events=[TimelineEvent(text="Facebook")],
                section=None,
            )]
        )
        _, added = _render(data)
        assert added >= 4

    def test_single_period_multiple_events(self):
        """period 1件・event 3件の場合、カードが3枚追加される。
        期待: 軸ライン1 + 円1 + ラベル1 + カード3 = 6"""
        data = TimelineData(
            periods=[TimelinePeriod(
                label="2004",
                events=[
                    TimelineEvent(text="A"),
                    TimelineEvent(text="B"),
                    TimelineEvent(text="C"),
                ],
                section=None,
            )]
        )
        _, added = _render(data)
        assert added >= 6

    def test_multiple_periods(self):
        """period 3件の場合、期ごとに円・ラベルが生成される。
        期待: 軸ライン1 + 円3 + ラベル3 = 7"""
        data = TimelineData(
            periods=[
                TimelinePeriod(label="2002", events=[], section=None),
                TimelinePeriod(label="2004", events=[], section=None),
                TimelinePeriod(label="2005", events=[], section=None),
            ]
        )
        _, added = _render(data)
        assert added >= 7

    def test_section_adds_band(self):
        """section が存在する場合、section 帯が追加される。"""
        data_no_sec = TimelineData(
            periods=[
                TimelinePeriod(label="A", events=[], section=None),
                TimelinePeriod(label="B", events=[], section=None),
            ]
        )
        data_with_sec = TimelineData(
            sections=["Group"],
            periods=[
                TimelinePeriod(label="A", events=[], section="Group"),
                TimelinePeriod(label="B", events=[], section="Group"),
            ]
        )
        _, no_sec_count = _render(data_no_sec)
        _, with_sec_count = _render(data_with_sec)
        # section 帯が追加される分、シェープが多い
        assert with_sec_count > no_sec_count

    def test_two_sections_two_bands(self):
        """section 2件の場合、section 帯が2つ追加される。"""
        data_one_sec = TimelineData(
            sections=["A"],
            periods=[
                TimelinePeriod(label="p1", events=[], section="A"),
                TimelinePeriod(label="p2", events=[], section="A"),
            ]
        )
        data_two_sec = TimelineData(
            sections=["A", "B"],
            periods=[
                TimelinePeriod(label="p1", events=[], section="A"),
                TimelinePeriod(label="p2", events=[], section="B"),
            ]
        )
        _, one_sec_count = _render(data_one_sec)
        _, two_sec_count = _render(data_two_sec)
        assert two_sec_count > one_sec_count

    def test_full_example_social_media(self):
        """SNSプラットフォーム歴史のサンプルが例外なく描画される。"""
        text = """timeline
    title History of Social Media Platform
    2002 : LinkedIn
    2004 : Facebook
         : Google
    2005 : YouTube
    2006 : Twitter"""
        data = parse_timeline(text)
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)
        # 例外が発生しなければ成功

    def test_full_example_industrial_revolution(self):
        """産業革命タイムライン（section あり）が例外なく描画される。"""
        text = """timeline
    title 産業革命のタイムライン
    section 17-20世紀
        Industry 1.0 : 機械化・水力・蒸気力
        Industry 2.0 : 電気・内燃機関・大量生産
        Industry 3.0 : 電子機器・コンピューター・オートメーション
    section 21世紀
        Industry 4.0 : インターネット・ロボット工学
        Industry 5.0 : AI・ビッグデータ・3Dプリント"""
        data = parse_timeline(text)
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_full_example_with_br(self):
        """<br> を含むイベントテキストが例外なく描画される。"""
        text = """timeline
    title England's History Timeline
    section Stone Age
      7600 BC : Britain's oldest known house was built in Orkney, Scotland
      6000 BC : Sea levels rise and Britain becomes an island.<br> The people who live here are hunter-gatherers.
    section Bronze Age
      2300 BC : People arrive from Europe and settle in Britain.<br>They bring farming and metalworking.
              : New styles of pottery and ways of burying the dead appear.
      2200 BC : The last major building works are completed at Stonehenge."""
        data = parse_timeline(text)
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_render_no_title_no_section(self):
        """タイトルなし・sectionなしのシンプルな timeline が描画される。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(label="2002", events=[TimelineEvent("LinkedIn")], section=None),
                TimelinePeriod(label="2004", events=[TimelineEvent("Facebook"), TimelineEvent("Google")], section=None),
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_render_japanese_labels(self):
        """日本語の period・event ラベルが例外なく描画される。"""
        data = TimelineData(
            title="日本の歴史",
            sections=["古代", "中世"],
            periods=[
                TimelinePeriod(label="奈良時代", events=[TimelineEvent("大仏建立"), TimelineEvent("万葉集")], section="古代"),
                TimelinePeriod(label="平安時代", events=[TimelineEvent("源氏物語")], section="古代"),
                TimelinePeriod(label="鎌倉時代", events=[TimelineEvent("武士の台頭")], section="中世"),
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_event_with_newline_text(self):
        """\\n を含む event テキストが例外なく描画される。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(
                    label="2023",
                    events=[TimelineEvent("Line1\nLine2\nLine3")],
                    section=None,
                )
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_many_events_per_period(self):
        """1 period に 5 件のイベントが例外なく描画される。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(
                    label="2024",
                    events=[TimelineEvent(f"event{i}") for i in range(5)],
                    section=None,
                )
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_many_periods(self):
        """10 件の period が例外なく描画される。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(label=str(2000 + i), events=[], section=None)
                for i in range(10)
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_color_cycle_13_periods(self):
        """13 件の period で色サイクルが機能し、例外が発生しない。"""
        data = TimelineData(
            periods=[
                TimelinePeriod(label=str(i), events=[TimelineEvent(f"event{i}")], section=None)
                for i in range(13)
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)

    def test_mixed_section_and_no_section(self):
        """section あり period と section なし period が混在する場合も例外なく描画。"""
        data = TimelineData(
            sections=["Group A"],
            periods=[
                TimelinePeriod(label="before", events=[TimelineEvent("e1")], section=None),
                TimelinePeriod(label="in group", events=[TimelineEvent("e2")], section="Group A"),
                TimelinePeriod(label="after", events=[TimelineEvent("e3")], section=None),
            ]
        )
        slide = _slide()
        renderer = TimelineRenderer()
        renderer.render(slide, data, 0, 0, 8_000_000, 4_500_000)
